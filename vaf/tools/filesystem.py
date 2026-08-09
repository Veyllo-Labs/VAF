# SPDX-FileCopyrightText: 2026 Veyllo GmbH
# SPDX-License-Identifier: AGPL-3.0-or-later
# Additional permissions and terms under AGPL Section 7: see LICENSING.md
import os
import shutil
from vaf.tools.base import BaseTool
from pathlib import Path
from typing import List, Tuple
import heapq

# Common Safety Logic
BLOCKED_DIRS = [
    "Windows", "Program Files", "Program Files (x86)", "System32",  # Windows
    "/etc", "/usr", "/sys", "/proc", "/var", "/boot",  # Linux/Mac
    ".git", ".ssh", "node_modules", ".env", "id_rsa"
]

# VAF program root - agent must NEVER access this (source code, config, secrets)
_VAF_PROJECT_ROOT = Path(__file__).resolve().parents[2]

# ── Per-user filesystem jail ─────────────────────────────────────────────────────────────────────
# While a jailed tool run executes, is_safe_path additionally enforces that the tool can access only
# the CALLER's OWN data - never another user's VAF_Projects/<uid8>/. This is a contextvar so it is
# scoped to that run only: when unset (the default) it has ZERO effect. Set/reset via
# set_librarian_scope/reset_librarian_scope; the jail info comes from compute_user_jail (single
# source - do not re-derive it per tool). Installed by MECHANISM, not by roster: every tool
# that declares `file_access` on BaseTool enters this jail around run() (base.py wraps it in
# user_jail - that declaration is the supported form), plus the librarian's env-carried lane,
# which cannot use the declaration and installs it by hand. An earlier version of this comment
# listed the declaring tools by name and was stale within two rounds - the guards are
# `tests/test_path_tools_ask_the_shared_rule.py` (a new path-taking tool cannot quietly skip
# asking) and `tests/test_file_access_declaration.py` (the declaration wires the jail).
import contextlib as _contextlib
import contextvars as _contextvars
_librarian_scope_ctx = _contextvars.ContextVar("vaf_librarian_scope", default=None)


def set_librarian_scope(scope_info):
    """scope_info: None (no jail) or dict {is_admin: bool, uid8: str, allowed_roots: list[Path]}.
    Returns a token to pass to reset_librarian_scope."""
    return _librarian_scope_ctx.set(scope_info)


def reset_librarian_scope(token):
    try:
        _librarian_scope_ctx.reset(token)
    except Exception:
        pass


def _visible_skill_roots(user_scope_id):
    """Folders of the skills this user may see, for the READ jail only.

    use_skill hands the model absolute paths to a skill's bundled files and tells it to
    open them with read_file, so a read jail that allowed only the user's own project
    tree would break every skill shipping reference material. Visibility is per user and
    already has an authority - the manifest's shared_with, via
    get_visible_skill_ids_for_user - but the FOLDERS all sit in one directory, so without
    this list "allow skills" would mean "allow everyone's skills", including a skill kept
    private to another user.

    Read-only on purpose: editing authority is a different question with a different
    answer (can_user_edit_skill), so these roots never enter the write jail.
    Fail-closed to an empty list - the caller keeps their own root and loses skill reads,
    rather than the lookup failure widening anything."""
    try:
        from vaf.core.skills_registry import get_skills_dir, get_visible_skill_ids_for_user
        base = get_skills_dir()
        return [base / sid for sid in get_visible_skill_ids_for_user(user_scope_id)]
    except Exception:
        return []


def compute_user_jail(user_scope_id, user_role=None, *, mode="write"):
    """Jail info for set_librarian_scope, shared by every tool that confines file
    access to the caller (librarian, write_file, edit_file, read tools, send_mail
    attachments).

    `mode` picks the allowed roots, because reading legitimately reaches further than
    writing: "read" adds the folders of the skills visible to this user (see
    _visible_skill_roots). "write" - the default, and what every existing caller gets -
    stays at the caller's own project tree. A tool installs the jail for the operation it
    performs; is_safe_path itself cannot tell a read from a write.

    Admin means the same thing here as everywhere else in VAF (is_admin_identity):
    the DB role, OR the configured local-admin scope. The scope half is not
    redundant - a direct consumer and the tokenless desktop have no role, and an
    owner session carries the admin's real UUID rather than None, so a plain
    truthiness check would jail the machine owner out of their own files (live
    regression class). The role half is not redundant either: a SECOND admin
    account carries its own scope UUID and used to be jailed here while every
    other gate in the app already treated it as a full admin.

    user_role must come from the caller's trusted context (the injected, JWT-derived
    role), never from an argument the model chose - the dispatcher overwrites it.
    Fail-closed: on any error the caller gets a jail with no allowed roots rather
    than no jail at all."""
    try:
        from vaf.core.config import is_admin_identity
        from vaf.core.session import get_user_projects_root
        scope = str(user_scope_id or "")
        if (not scope) or is_admin_identity(user_role, scope):
            return {"is_admin": True, "uid8": None, "allowed_roots": []}
        own_root = get_user_projects_root(scope)
        roots = [own_root] if own_root else []
        if mode == "read":
            roots = roots + _visible_skill_roots(scope)
        return {"is_admin": False, "uid8": scope.replace("-", "").lower()[:8],
                "allowed_roots": roots}
    except Exception:
        return {"is_admin": False, "uid8": "", "allowed_roots": []}


@_contextlib.contextmanager
def user_jail(user_scope_id, user_role=None, *, mode="write"):
    """Install the per-user jail for the duration of one tool body.

    Every jailed tool needs the same four steps (resolve, skip for an admin, set the
    contextvar, reset it in a finally) and each hand-written copy is a chance to forget
    the finally. It must be entered INSIDE the tool, and the reason is NOT that the
    dispatcher cannot reach it - the bounded run copies the caller's context into the worker
    thread on purpose, so a contextvar set outside does arrive (vaf/core/bounded_run.py).
    The reason is that the dispatcher is not always there: the coder, the workflow engine
    and automations call these tools directly, so run() is the only place that always runs.

    A falsy scope means a direct consumer (coder, workflow engine, automations) - no jail,
    exactly as before this existed.

    NESTING ONLY EVER NARROWS. An inner jail may shrink the roots it inherits and may never
    add one, and it can never lift an outer jail by resolving to an admin. Without that rule
    nesting WIDENS, and not hypothetically: the librarian installs the default `write` mode
    around its whole run while its read sub-tools enter `read`, which reaches further by the
    folders of the skills this user may see. That gap opened with the Phase 3 unification -
    before it the sub-tools received no scope at all, so the inner call was a no-op and there
    was nothing to widen.
    Stated at its true size, because an overstated example devalues the rule it justifies:
    those extra roots are the skills THIS user is allowed to see, decided by a separate
    authority (`get_visible_skill_ids_for_user`). It widened onto data the caller is
    separately entitled to, not onto someone else's. The MECHANISM is the forbidden one all
    the same, and the next widening will not be that harmless."""
    token = None
    if user_scope_id:
        info = compute_user_jail(user_scope_id, user_role, mode=mode)
        outer = _librarian_scope_ctx.get(None)
        if outer:
            if info.get("is_admin"):
                # DEFENSIVE, AND TODAY EQUIVALENT TO DOING NOTHING - deliberately kept, and
                # deliberately not covered by a test of its own. Replacing this line with
                # `pass` leaves the whole suite green, because an admin resolution also fails
                # the `if not info.get("is_admin")` below: no token is set and the outer
                # contextvar survives by itself. Both routes end in the same place.
                # It is load-bearing for one property that is NOT visible here: this holds
                # only as long as user_jail never CLEARS the contextvar, but merely sets it
                # or does nothing. Rewrite it to always install a token and this line becomes
                # the thing that stops an inner admin from lifting an outer jail. The
                # PROPERTY is tested (test_an_inner_admin_cannot_lift_an_outer_jail); this
                # line is not, and pinning an equivalent line would be the ballast the test
                # rules call a reason for deletion.
                info = outer
            else:
                # NOT an intersection, deliberately: an inner root is KEPT whole when it lies
                # inside an outer one and DROPPED whole otherwise. An inner root that is a
                # PARENT of an outer root is therefore discarded rather than trimmed to the
                # overlap, and if that was the only root the tool is left with none, i.e.
                # deny-all. That is the safe direction and it is the intended one - but it is
                # a silent total refusal rather than access to the shared part, so it is
                # written down here. Do not "fix" it by widening to the parent.
                outer_roots = [os.path.abspath(str(r)) for r in (outer.get("allowed_roots") or [])]
                info = {**info, "allowed_roots": [
                    r for r in (info.get("allowed_roots") or [])
                    if any(os.path.abspath(str(r)) == o
                           or os.path.abspath(str(r)).startswith(o.rstrip("/") + os.sep)
                           for o in outer_roots)]}
        if not info.get("is_admin"):
            token = set_librarian_scope(info)
    try:
        yield
    finally:
        if token is not None:
            reset_librarian_scope(token)


def _librarian_jail_ok(abs_path) -> bool:
    """True if abs_path is allowed under the active librarian jail. Fail-closed. No jail set => True."""
    info = _librarian_scope_ctx.get()
    if not info:
        return True
    try:
        import re as _re_j
        from vaf.core.platform import Platform
        target = Path(abs_path).resolve()
        if info.get("is_admin"):
            return True  # local admin / machine owner: full access (still bounded by VAF/BLOCKED checks)
        # HARD cross-user invariant: any VAF_Projects/<8-hex> that is not the caller's own is DENIED.
        projects_root = (Platform.documents_dir() / "VAF_Projects").resolve()
        if target == projects_root or target.is_relative_to(projects_root):
            rel = target.relative_to(projects_root)
            first = rel.parts[0] if rel.parts else ""
            if _re_j.fullmatch(r"[0-9a-f]{8}", first) and first != (info.get("uid8") or ""):
                return False
        # Positive allow-list: a remote (non-admin) user may only read inside their own allowed roots.
        for r in (info.get("allowed_roots") or []):
            try:
                rp = Path(r).resolve()
                if target == rp or target.is_relative_to(rp):
                    return True
            except Exception:
                continue
        return False
    except Exception:
        return False  # fail-closed

def _resolve_folder_alias(path_str: str) -> str:
    """Resolve folder aliases like 'Desktop', 'Documents' to actual paths.
    If Desktop is not accessible, automatically falls back to Documents."""
    from pathlib import Path
    from vaf.core.platform import Platform
    
    path_lower = path_str.lower()
    home = Path.home()

    # LLMs sometimes emit /home/user/ as the home directory.
    # Remap it to the actual home so writes land in the right place.
    import re as _re
    # Use a function replacement so the home path is inserted literally. On
    # Windows str(home) is e.g. "C:\\Users\\me"; passing it as a plain
    # replacement string makes re.sub parse "\U" as a regex escape and raise
    # "bad escape \U" (and would also mangle "\g"/backslashes elsewhere).
    path_str = _re.sub(
        r'^/home/user(?=/|$)',
        lambda _m: str(home),
        path_str,
    )

    # Check if path contains folder alias
    folder_aliases = {
        "desktop": home / "Desktop",
        "documents": home / "Documents",
        "downloads": home / "Downloads",
        "pictures": home / "Pictures",
        "videos": home / "Videos",
        "music": home / "Music",
    }
    
    # German folder names (Windows often uses these, but check on all platforms)
    # Some Linux/macOS users might also have German folder names
    german_mappings = {
        "desktop": ["Desktop", "Arbeitsplatz"],
        "documents": ["Documents", "Dokumente"],
        "pictures": ["Pictures", "Bilder"],
        "videos": ["Videos"],
        "music": ["Music", "Musik"],
        "downloads": ["Downloads", "Herunterladen"],
    }
    
    for key, variants in german_mappings.items():
        for variant in variants:
            variant_path = home / variant
            if variant_path.exists():
                folder_aliases[key] = variant_path
    
    # Helper to check if a path is writable
    def is_writable(path: Path) -> bool:
        """Check if a directory is writable."""
        if not path.exists():
            return False
        try:
            # Try to create a test file
            test_file = path / ".vaf_write_test"
            test_file.touch()
            test_file.unlink()
            return True
        except (PermissionError, OSError):
            return False
    
    # Check if path starts with or contains a folder alias
    for alias, alias_path in folder_aliases.items():
        # Check if path starts with the alias AT A PATH BOUNDARY (the alias alone,
        # or the alias followed by a separator). Without the boundary check a bare
        # filename like "Documentsfile.txt" would be misrouted into ~/Documents/.
        _matched = None
        for _cand in (alias, alias.capitalize()):
            if path_str == _cand or path_str.startswith(_cand + "/") or path_str.startswith(_cand + "\\"):
                _matched = _cand
                break
        if _matched:
            # Replace alias with actual path
            remaining = path_str[len(_matched):].lstrip('\\/')
            
            # Special handling for Desktop: check if writable, fallback to Documents
            if alias.lower() == "desktop":
                if not is_writable(alias_path):
                    # Desktop not writable, use Documents as fallback
                    documents_path = folder_aliases.get("documents", home / "Documents")
                    if documents_path.exists():
                        if remaining:
                            return str(documents_path / remaining)
                        else:
                            return str(documents_path)
            
            if remaining:
                return str(alias_path / remaining)
            else:
                return str(alias_path)
        # Also check for Windows-style paths (Desktop\file.txt)
        if '\\' in path_str or '/' in path_str:
            parts = path_str.replace('\\', '/').split('/')
            if parts[0].lower() == alias:
                remaining = '/'.join(parts[1:])
                
                # Special handling for Desktop: check if writable, fallback to Documents
                if alias.lower() == "desktop":
                    if not is_writable(alias_path):
                        # Desktop not writable, use Documents as fallback
                        documents_path = folder_aliases.get("documents", home / "Documents")
                        if documents_path.exists():
                            if remaining:
                                return str(documents_path / remaining)
                            else:
                                return str(documents_path)
                
                return str(alias_path / remaining) if remaining else str(alias_path)
    
    return path_str

def is_safe_path(path):
    try:
        # First resolve folder aliases
        resolved = _resolve_folder_alias(path)
        abs_path = os.path.abspath(os.path.expanduser(resolved))

        # CRITICAL: Never allow access to VAF program root (source, config, secrets)
        try:
            target = Path(abs_path).resolve()
            if target == _VAF_PROJECT_ROOT or target.is_relative_to(_VAF_PROJECT_ROOT):
                return False, "Access denied: VAF program directory is protected"
        except (ValueError, OSError):
            pass

        # Block system locations by PATH COMPONENT / PREFIX, not raw substring. A
        # substring match wrongly rejected legitimate paths that merely *contain* a
        # blocked token: the macOS temp dir "/private/var/folders/..." contains "/var",
        # "~/.environment" contains ".env", "C:\\Users\\Windows10Fan" contains "Windows".
        # Absolute roots (/etc, /var, ...) now match the dir itself or anything strictly
        # under it; name-based tokens (.ssh, node_modules, Windows, ...) match only when
        # they are a whole path component.
        norm = abs_path.replace("\\", "/")
        components = norm.split("/")
        for blocked in BLOCKED_DIRS:
            if blocked.startswith("/"):
                if norm == blocked or norm.startswith(blocked + "/"):
                    return False, f"Access denied: {blocked}"
            elif blocked in components:
                return False, f"Access denied: {blocked}"
        # VAF's own data directory is the credential store, and no file tool has a
        # legitimate reason to open it: settings are read through Config.get(), never
        # through read_file. It holds config.json (every api_key_*, the JWT secret) and
        # its .bak copies, secrets/, ssl/ private keys, browser_sessions/ (live logged-in
        # cookies), speaker_profiles/ (voice biometrics) and sessions/ (every
        # conversation). This is blocked for EVERYONE, the machine owner included: an API
        # key extracted through a chat prompt keeps working outside VAF and outlives any
        # access this instance could revoke, so it is not a per-user permission question.
        #
        # Deny-by-default, so a new folder under ~/.vaf is protected the day it is added
        # rather than the day someone remembers to list it. The exceptions are the folders
        # holding USER-CREATED CONTENT, which is ordinary material the agent works with and
        # which is governed by the per-user jail below like anything else - not by this
        # block, whose job is only the credential store:
        #   skills/    - use_skill hands the model absolute paths to a skill's bundled
        #                files and tells it to open them with read_file (use_skill.py:33,92)
        #   workflows/ - agent-created workflows; the same category as skills and as the
        #                custom tools in data_dir/custom_tools, which this block never
        #                touched because they live outside ~/.vaf
        # Fail-closed: if the data dir cannot be resolved, nothing under .vaf is allowed.
        if ".vaf" in components:
            _content_ok = False
            try:
                from vaf.core.platform import Platform as _PlatFS
                _vaf_root = Path(_PlatFS.vaf_dir()).resolve()
                for _sub in ("skills", "workflows"):
                    _root = str(_vaf_root / _sub).replace("\\", "/")
                    if norm == _root or norm.startswith(_root + "/"):
                        _content_ok = True
                        break
            except Exception:
                _content_ok = False
            if not _content_ok:
                return False, "Access denied: VAF's own data directory"
        # Librarian per-user jail (no-op unless a librarian run set the scope contextvar).
        if not _librarian_jail_ok(abs_path):
            return False, "Access denied: outside your own data"
        # Every check above (except the VAF-root check and the jail, which resolve for
        # themselves) ran against the UNRESOLVED path - os.path.abspath follows no
        # symlinks - so a link inside an allowed folder could smuggle in a protected
        # target: the downstream open() follows the link. Resolve, and when the real
        # path differs, ask this same rule about the real one. Termination is
        # structural: a fully resolved path resolves to itself, so the nested call
        # cannot recurse a second time. The RETURN VALUE stays the unresolved
        # abs_path - callers pin that.
        real = str(Path(abs_path).resolve())
        if real != abs_path:
            ok, verdict = is_safe_path(real)
            if not ok:
                return False, verdict
        return True, abs_path
    except Exception:
        return False, "Invalid path"

class ListFilesTool(BaseTool):
    name = "list_files"
    identity_kwargs = ("user_role", "user_scope_id")
    file_access = "read"
    permission_level = "read"
    side_effect_class = "none"
    description = "Lists files in a directory."

    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Absolute path"},
            "sort_by": {"type": "string", "enum": ["name", "date", "size"], "description": "Sort order (default: name)"},
            "limit": {"type": "integer", "description": "Max files to return (default: 100)"}
        },
        "required": ["path"]
    }


    def run(self, **kwargs) -> str:
        path = kwargs.get('path', '.')
        sort_by = kwargs.get('sort_by', 'name')
        limit = kwargs.get('limit', 100)
        
        safe, res = is_safe_path(path)
        if not safe: return res
        
        try:
            entries = []
            with os.scandir(res) as it:
                for entry in it:
                    try:
                        stat = entry.stat()
                        entries.append({
                            "name": entry.name,
                            "is_dir": entry.is_dir(),
                            "size": stat.st_size,
                            "mtime": stat.st_mtime
                        })
                    except: pass
            
            # Sorting
            if sort_by == 'date':
                entries.sort(key=lambda x: x['mtime'], reverse=True) # Newest first
            elif sort_by == 'size':
                entries.sort(key=lambda x: x['size'], reverse=True) # Largest first
            else:
                entries.sort(key=lambda x: x['name'].lower())

            # Formatting
            import datetime
            output = f"Directory Listing: {res} (Total files found: {len(entries)}) (Sorted by: {sort_by})\n"
            output += f"{'Type':<6} {'Date':<18} {'Size':<10} {'Name'}\n"
            output += "-"*60 + "\n"
            
            count = 0
            for e in entries:
                if count >= limit:
                    output += f"... (and {len(entries) - limit} more)\n"
                    break
                    
                dt = datetime.datetime.fromtimestamp(e['mtime']).strftime('%Y-%m-%d %H:%M')
                
                if e['is_dir']:
                    size_str = "<DIR>"
                    type_str = "[DIR]"
                else:
                    type_str = "[FILE]"
                    # Size formatting
                    if e['size'] < 1024: size_str = f"{e['size']} B"
                    elif e['size'] < 1024*1024: size_str = f"{e['size']/1024:.1f} KB"
                    else: size_str = f"{e['size']/(1024*1024):.1f} MB"
                
                output += f"{type_str:<6} {dt:<18} {size_str:<10} {e['name']}\n"
                count += 1
                
            return output if output else "Empty Directory"
        except Exception as e: return str(e)


class FolderSizeTool(BaseTool):
    name = "folder_size"
    identity_kwargs = ("user_role", "user_scope_id")
    file_access = "read"
    permission_level = "read"
    side_effect_class = "none"
    description = "Calculates the total size of a folder (recursive), with optional largest-files preview."

    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Folder path (absolute or relative)"},
            "top_n": {"type": "integer", "description": "Show N largest files (default: 10)"},
            "max_files": {"type": "integer", "description": "Safety cap for scanned files (default: 200000)"},
        },
        "required": ["path"],
    }


    def run(self, **kwargs) -> str:
        path = kwargs.get("path", "")
        top_n = int(kwargs.get("top_n", 10) or 10)
        max_files = int(kwargs.get("max_files", 200000) or 200000)

        if not path:
            return "Error: No path provided."

        safe, abs_path = is_safe_path(path)
        if not safe:
            return abs_path

        root = Path(abs_path)
        if not root.exists():
            return f"Error: Path not found: {abs_path}"
        if not root.is_dir():
            return f"Error: Not a directory: {abs_path}"

        def fmt_size(num_bytes: int) -> str:
            if num_bytes < 1024:
                return f"{num_bytes} B"
            if num_bytes < 1024 * 1024:
                return f"{num_bytes / 1024:.1f} KB"
            if num_bytes < 1024 * 1024 * 1024:
                return f"{num_bytes / (1024 * 1024):.1f} MB"
            return f"{num_bytes / (1024 * 1024 * 1024):.2f} GB"

        total_bytes = 0
        file_count = 0
        dir_count = 0
        largest: List[Tuple[int, str]] = []  # min-heap of (size, path)
        top_n = max(0, min(top_n, 50))

        try:
            for dirpath, dirnames, filenames in os.walk(root, topdown=True, onerror=None, followlinks=False):
                dir_count += 1

                for name in filenames:
                    file_count += 1
                    if file_count > max_files:
                        raise RuntimeError(f"Scan aborted: too many files (> {max_files}).")

                    fpath = Path(dirpath) / name
                    try:
                        size = fpath.stat().st_size
                    except Exception:
                        continue

                    total_bytes += int(size)

                    if top_n > 0:
                        pstr = str(fpath)
                        if len(largest) < top_n:
                            heapq.heappush(largest, (int(size), pstr))
                        else:
                            if int(size) > largest[0][0]:
                                heapq.heapreplace(largest, (int(size), pstr))
        except RuntimeError as e:
            # Return partial results with explicit warning
            warning = str(e)
        else:
            warning = ""

        largest_sorted = sorted(largest, key=lambda x: x[0], reverse=True)

        out = []
        out.append("### Folder Size")
        out.append(f"Path: {abs_path}")
        out.append(f"Total size: **{fmt_size(total_bytes)}**")
        out.append(f"Files: {file_count}")
        out.append(f"Directories: {dir_count}")
        if warning:
            out.append(f"\n⚠️ {warning}")
        if largest_sorted:
            out.append("\n**Largest files:**")
            for size, pstr in largest_sorted:
                out.append(f"- {fmt_size(size)} - {pstr}")

        return "\n".join(out)

class ReadFileTool(BaseTool):
    name = "read_file"
    identity_kwargs = ("user_role", "user_scope_id")
    file_access = "read"
    permission_level = "read"
    side_effect_class = "none"
    description = """Reads the content of a file. Supports text, PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx).

Use this when:
- User asks to read a specific file
- Sub-agent created a file (look for "🔗 EXTRACTED FILE PATHS" in sub-agent result!)
- Need to analyze document content quickly

**IMPORTANT:** If a sub-agent just created a file, the path is already in the conversation context!
Example: research_agent says "Saved to: C:\\...\\report.html" → use read_file("C:\\...\\report.html")

For detailed analysis of large files, consider using librarian_agent instead."""

    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string"},
            "start_line": {"type": "integer", "description": "First line to return (1-indexed). Use to read a file in sections."},
            "end_line": {"type": "integer", "description": "Last line to return (inclusive). Omit to read to end of file."},
            "first_page": {"type": "integer", "description": "PDF only: first page to read (1-indexed). Default 1."},
            "last_page": {"type": "integer", "description": "PDF only: last page to read (inclusive). Omit for the default window from first_page."},
        },
        "required": ["path"]
    }

    # PDF pages returned when no explicit range is asked for. The limit is
    # NAMED to the caller: the honest header says which pages of how many were
    # read and how to continue (first_page=N) - never a bare "(truncated)".
    _PDF_READ_DEFAULT_PAGES = 50


    def run(self, **kwargs) -> str:
        path = kwargs.get('path', '')
        start_line = kwargs.get('start_line')
        end_line = kwargs.get('end_line')
        safe, res = is_safe_path(path)
        if not safe: return res
        try:
            if not os.path.exists(res): return "Error: File not found."
            if not os.path.isfile(res): return "Error: Not a file."
            
            # Get file extension
            from pathlib import Path
            file_path = Path(res)
            ext = file_path.suffix.lower()
            
            # ═══════════════════════════════════════════════════════════
            # PDF Files
            # ═══════════════════════════════════════════════════════════
            if ext == '.pdf':
                try:
                    from vaf.core.pdf_extract import extract_pdf_markdown, format_pdf_read_result
                    # Shared extractor: pdfplumber markdown (headings + tables), PyPDF2 + OCR
                    # fallbacks. The shared formatter renders the result honestly (page range
                    # covered, continuation hint, named OCR reason) - the hand-rolled 15k cut
                    # with its guessed "no embedded text?" message lived here twice (this tool
                    # + the librarian) and destroyed the facts the extractor already knew.
                    first = max(1, int(kwargs.get('first_page') or 1))
                    last = kwargs.get('last_page')
                    window = (max(0, int(last) - first + 1) if last
                              else self._PDF_READ_DEFAULT_PAGES)
                    res_md = extract_pdf_markdown(
                        file_path, max_pages=window, ocr_fallback=True, first_page=first)
                    return format_pdf_read_result(res_md, file_name=file_path.name)

                except ImportError:
                    return "Error: PDF support not installed. Run: pip install pdfplumber PyPDF2"
                except Exception as e:
                    err_str = str(e)
                    hint = " For AES-encrypted PDFs run: pip install pycryptodome" if ("PyCryptodome" in err_str or "AES" in err_str) else ""
                    return f"Error reading PDF: {e}{hint}"
            
            # ═══════════════════════════════════════════════════════════
            # Word Documents (.docx)
            # ═══════════════════════════════════════════════════════════
            elif ext == '.docx':
                try:
                    from docx import Document
                    doc = Document(res)
                    
                    content = []
                    
                    # Extract paragraphs
                    for para in doc.paragraphs:
                        text = para.text.strip()
                        if text:
                            content.append(text)
                    
                    # Extract tables
                    if doc.tables:
                        content.append("\n--- Tables ---")
                        for i, table in enumerate(doc.tables[:5], 1):
                            content.append(f"\nTable {i}:")
                            for row in table.rows[:10]:
                                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                                if row_text:
                                    content.append(row_text)
                    
                    full_text = "\n".join(content)
                    
                    # Truncate if too long
                    if len(full_text) > 15000:
                        full_text = full_text[:15000] + "\n\n... (truncated)"
                    
                    return f"### Word Document: {file_path.name}\n**Paragraphs:** {len(doc.paragraphs)}\n**Tables:** {len(doc.tables)}\n\n{full_text}"
                    
                except ImportError:
                    return "Error: Word document support not installed. Run: pip install python-docx"
                except Exception as e:
                    return f"Error reading Word document: {e}"
            
            # ═══════════════════════════════════════════════════════════
            # Excel Files (.xlsx, .xls)
            # ═══════════════════════════════════════════════════════════
            elif ext in ['.xlsx', '.xls']:
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(res, read_only=True, data_only=True)
                    
                    content = []
                    content.append(f"**Sheets:** {', '.join(wb.sheetnames)}\n")
                    
                    # Read first 3 sheets
                    for sheet_name in wb.sheetnames[:3]:
                        sheet = wb[sheet_name]
                        content.append(f"\n--- Sheet: {sheet_name} ---")
                        
                        # Get dimensions
                        max_row = min(sheet.max_row, 50)
                        max_col = min(sheet.max_column, 20)
                        
                        # Read data
                        for row in sheet.iter_rows(min_row=1, max_row=max_row, max_col=max_col, values_only=True):
                            row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                            if row_text.strip():
                                content.append(row_text)
                    
                    if len(wb.sheetnames) > 3:
                        content.append(f"\n... ({len(wb.sheetnames) - 3} more sheets not shown)")
                    
                    full_text = "\n".join(content)
                    
                    # Truncate if too long
                    if len(full_text) > 15000:
                        full_text = full_text[:15000] + "\n\n... (truncated)"
                    
                    return f"### Excel File: {file_path.name}\n{full_text}"
                    
                except ImportError:
                    return "Error: Excel support not installed. Run: pip install openpyxl"
                except Exception as e:
                    return f"Error reading Excel file: {e}"
            
            # ═══════════════════════════════════════════════════════════
            # PowerPoint Files (.pptx)
            # ═══════════════════════════════════════════════════════════
            elif ext == '.pptx':
                try:
                    from pptx import Presentation
                    prs = Presentation(res)
                    
                    content = []
                    content.append(f"**Slides:** {len(prs.slides)}\n")
                    
                    # Read first 20 slides
                    for i, slide in enumerate(prs.slides[:20], 1):
                        content.append(f"\n--- Slide {i} ---")
                        
                        # Extract text from shapes
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                content.append(shape.text)
                    
                    if len(prs.slides) > 20:
                        content.append(f"\n... ({len(prs.slides) - 20} more slides not shown)")
                    
                    full_text = "\n".join(content)
                    
                    # Truncate if too long
                    if len(full_text) > 15000:
                        full_text = full_text[:15000] + "\n\n... (truncated)"
                    
                    return f"### PowerPoint: {file_path.name}\n{full_text}"
                    
                except ImportError:
                    return "Error: PowerPoint support not installed. Run: pip install python-pptx"
                except Exception as e:
                    return f"Error reading PowerPoint file: {e}"
            
            # ═══════════════════════════════════════════════════════════
            # Text Files (default)
            # ═══════════════════════════════════════════════════════════
            else:
                with open(res, 'r', encoding='utf-8', errors='ignore') as f:
                    lines = f.readlines()
                total_lines = len(lines)
                if start_line is not None or end_line is not None:
                    s = max(1, int(start_line or 1)) - 1  # convert to 0-indexed
                    e = int(end_line) if end_line is not None else total_lines
                    selected = lines[s:e]
                    header = f"[Lines {s+1}–{min(e, total_lines)} of {total_lines} total]\n"
                    return header + "".join(selected)
                content = "".join(lines)
                return content
                
        except Exception as e: return str(e)

class WriteFileTool(BaseTool):
    name = "write_file"
    identity_kwargs = ("user_role", "user_scope_id")
    file_access = "write"
    permission_level = "write"
    side_effect_class = "reversible"
    description = (
        "Writes content to a file (creates or overwrites). Use for saving any single-file "
        "artifact the content of which you already have (html, svg, txt, code, ...). "
        "A relative path lands in the current chat workspace; an absolute or ~ path is "
        "honored (VAF's own directory and system locations are protected). "
        "For files PRODUCED BY CODE in python_sandbox (images, PDFs), do NOT route them "
        "through here as base64 - use the sandbox's export_files parameter instead (large "
        "base64 gets truncated in context). content_base64 exists only for SMALL binary "
        "data you already hold. For multi-file code projects use coding_agent instead."
    )
    input_examples = [
        {"path": "chart.svg", "content": "<svg xmlns=\"http://www.w3.org/2000/svg\">...</svg>"},
        {"path": "report/summary.md", "content": "# Summary\n..."},
        {"path": "diagram.png", "content_base64": "iVBORw0KGgoAAAANSUhEUgAA..."},
    ]

    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string"},
            "content": {"type": "string", "description": "Text content (UTF-8). Use content_base64 for binary files."},
            "content_base64": {"type": "string", "description": "Base64-encoded binary content (png/jpg/pdf/...). Mutually exclusive with content."}
        },
        "required": ["path"]
    }

    # Weak models often send synonyms; repair_tool_input (R0) remaps them to
    # the canonical names before dispatch. Kept OFF the model-facing schema
    # above (live incident: a local model sent file_path/message and the
    # write was silently lost).
    input_aliases = {
        "path": ["file_path", "filepath", "filename", "file"],
        # "file_content" observed live (live incident: a 4B model burned
        # four calls on it before finding a mapped name).
        "content": ["message", "text", "body", "data", "file_content", "contents"],
    }


    def run(self, **kwargs) -> str:
        path = kwargs.get('path', '')
        content = kwargs.get('content', '')
        # Binary lane: content_base64 decodes to raw bytes (png/jpg/pdf/... - the
        # sandbox can render and print base64, but a str content cannot carry
        # binary). Exactly one of content / content_base64 must be provided;
        # explicit content="" (create an empty file) stays valid.
        _b64 = kwargs.get('content_base64') or ''
        if _b64:
            if (content or '') != '':
                return "Tool Error: pass either content OR content_base64, not both."
            try:
                import base64 as _b64mod
                import binascii as _binascii
                _data = _b64mod.b64decode(_b64, validate=True)
            except (_binascii.Error, ValueError):
                return "Tool Error: content_base64 is not valid base64 data."
        elif 'content' not in kwargs or kwargs.get('content') is None:
            return "Tool Error: provide content (text) or content_base64 (binary)."
        else:
            _data = (content or "").encode("utf-8")
        # Main-agent calls inject the chat workspace: a relative path ("chart.svg")
        # resolves there instead of against the process cwd (which is the protected
        # VAF root when the backend runs from the repo, so it would be denied).
        _ws = kwargs.get('_session_workspace') or ""
        if _ws and path and not os.path.isabs(os.path.expanduser(str(path))):
            path = os.path.join(str(_ws), str(path))
        safe, res = is_safe_path(path)
        if not safe: return res

        # Reroute writes that would land loose in the home directory into
        # Documents/VAF, so the agent neither creates arbitrary home-level dirs
        # (e.g. ~/my_project/) nor leaves a cwd-relative bare filename at the home
        # root (e.g. a workflow write_file step with path="draft" run from a
        # process whose cwd is the home dir). An EXPLICIT absolute or ~-anchored
        # target is the user's deliberate choice (e.g. read-modify-write of
        # ~/.bashrc) and is written in place.
        _explicit_target = os.path.isabs(os.path.expanduser(path or ""))
        try:
            from vaf.core.platform import Platform as _Plat
            _home = Path.home()
            # Resolve BOTH sides: is_safe_path's os.path.abspath symlink-resolves the
            # cwd, but $HOME may itself be an unresolved symlink - comparing raw would
            # miss a home write on such setups and re-open the litter-the-home bug.
            _rel = Path(res).resolve().relative_to(_home.resolve())   # ValueError if not under home
            _first = _rel.parts[0] if _rel.parts else None
            _multi_segment = len(_rel.parts) > 1
            _known_home_dirs = {
                "Documents", "Downloads", "Desktop", "Pictures", "Videos", "Music",
                "Dokumente", "Bilder", "Musik", "Herunterladen",  # German
                "VAF", ".vaf", ".config", ".local", ".cache", ".ssh",
            }
            _out = _Plat.get_vaf_output_dir()
            if _first and _first not in _known_home_dirs and not (_home / _first).exists():
                # New home-level target. An unrecognized multi-segment home dir is always
                # rerouted; an explicit single new file (e.g. "create ~/.vimrc") is the
                # user's choice and written in place.
                if _multi_segment or not _explicit_target:
                    res = str(_out / _rel)
            elif len(_rel.parts) == 1 and _first not in _known_home_dirs and not _explicit_target:
                # Existing bare file at the home root reached via a cwd-relative name
                # - reroute so it lands in a servable root.
                # An explicit ~/x or /home/x overwrite happens in place.
                res = str(_out / _rel)
        except ValueError:
            pass  # not under home - keep the resolved path
        except Exception as _guard_err:
            # Log if the home-reroute guard fails instead of swallowing the error,
            # so unexpected failures stay visible.
            import logging
            logging.getLogger(__name__).warning("write_file home-reroute guard skipped: %s", _guard_err)

        import time

        # Retry mechanism for file locking issues (especially on Windows)
        max_retries = 3
        retry_delay = 0.1  # 100ms between retries
        
        for attempt in range(max_retries):
            try:
                # Ensure parent directory exists
                parent_dir = os.path.dirname(res)
                if parent_dir and not os.path.exists(parent_dir):
                    try:
                        os.makedirs(parent_dir, exist_ok=True)
                    except PermissionError as pe:
                        # Desktop might have permission issues, try Documents as fallback (cross-platform)
                        from vaf.core.platform import Platform
                        if "desktop" in res.lower():
                            # NOTE: no local pathlib import here - a function-local
                            # 'from pathlib import Path' shadows the module-level Path
                            # for the WHOLE function scope and killed the home-reroute
                            # guard above with an UnboundLocalError on every call.
                            fallback = Path.home() / "Documents"
                            if fallback.exists():
                                # Replace Desktop with Documents in path
                                filename = os.path.basename(res)
                                res = str(fallback / filename)
                                parent_dir = str(fallback)
                                try:
                                    os.makedirs(parent_dir, exist_ok=True)
                                except:
                                    return f"❌ Permission error: Could not write to Desktop or Documents. Path: {res}"
                        else:
                            # Generic permission error message (OS-independent)
                            return f"❌ Permission error: Access denied to '{parent_dir}'"
                
                # On Windows, files can be locked even after closing
                # Use atomic write pattern: write to temp file, then rename
                import tempfile
                import shutil
                
                # Create temp file in same directory (for atomic rename)
                temp_dir = os.path.dirname(res) or '.'
                temp_fd, temp_path = tempfile.mkstemp(
                    prefix='vaf_',
                    suffix='.tmp',
                    dir=temp_dir
                )
                
                try:
                    # Write to temp file
                    # Bytes were prepared above (text encoded / base64 decoded);
                    # writing bytes avoids platform newline translation breaking
                    # the size verification.
                    data = _data
                    with os.fdopen(temp_fd, "wb") as f:
                        f.write(data)
                        f.flush()  # Ensure data is written
                        os.fsync(f.fileno())  # Force write to disk (OS-independent)
                    
                    # Atomic rename (works on all OS)
                    # On Windows, this will fail if file is locked
                    if os.path.exists(res):
                        # Remove existing file first (on Windows, rename fails if target exists)
                        try:
                            os.remove(res)
                        except PermissionError:
                            # File is locked - wait and retry
                            os.remove(temp_path)  # Clean up temp file
                            if attempt < max_retries - 1:
                                time.sleep(retry_delay * (attempt + 1))  # Exponential backoff
                                continue
                            raise
                    
                    # Atomic rename
                    try:
                        shutil.move(temp_path, res)
                    except PermissionError:
                        # Desktop might not be writable, try Documents as fallback (cross-platform)
                        from vaf.core.platform import Platform
                        if "desktop" in res.lower():
                            fallback = Path.home() / "Documents"
                            if fallback.exists():
                                # Replace Desktop with Documents in path
                                filename = os.path.basename(res)
                                fallback_path = str(fallback / filename)
                                # Move temp file to Documents instead
                                try:
                                    shutil.move(temp_path, fallback_path)
                                    res = fallback_path  # Update res for verification
                                except Exception as e:
                                    os.remove(temp_path)  # Clean up temp file
                                    if attempt < max_retries - 1:
                                        time.sleep(retry_delay * (attempt + 1))
                                        continue
                                    raise
                        else:
                            # Not a Desktop path, re-raise
                            os.remove(temp_path)  # Clean up temp file
                            raise
                    
                    # Verify file was written correctly
                    # Check if we used Documents fallback (original path had Desktop but res now has Documents)
                    used_fallback = "desktop" in path.lower() and "documents" in res.lower()
                    success = False
                    if os.path.exists(res) and os.path.getsize(res) == len(data):
                        success = True
                    else:
                        try:
                            if os.path.exists(res) and os.path.getsize(res) > 0:
                                success = True
                        except Exception:
                            pass
                    if success:
                        # Emit file_created for all written files so Web UI shows a download/open link.
                        # Session resolution order matters (emit-site scoping invariant): the
                        # session injected for THIS call wins; the env var / process-global
                        # fallback is only for legacy sub-agent processes, where it is
                        # per-process anyway. On a multi-user server the global fallback could
                        # attribute the file to another user's session.
                        try:
                            # One resolver: context first, process boundary second.
                            from vaf.core.subagent_ipc import get_current_session_id
                            _sid = kwargs.get('_session_id') or get_current_session_id()
                            if _sid:
                                from vaf.core.web_interface import notify_file_created
                                notify_file_created(_sid, res, title=os.path.basename(res))
                        except Exception:
                            pass
                        # When a document file is saved, also open it in the Document Editor.
                        # HTML files are NOT opened in the doc editor — they get the blue download chip instead.
                        _doc_extensions = (".md", ".txt", ".docx")
                        if res.lower().endswith(_doc_extensions):
                            try:
                                from vaf.core.subagent_ipc import get_current_session_id
                                _sid2 = kwargs.get('_session_id') or get_current_session_id()
                                if _sid2:
                                    from vaf.core.web_interface import notify_document_created
                                    notify_document_created(
                                        _sid2, res,
                                        title=os.path.basename(res)
                                    )
                            except Exception:
                                pass
                        # Main-agent calls (workspace injected): flag writes that land
                        # OUTSIDE the chat workspace in the same turn - the UI file
                        # browser only shows the workspace, so the model can correct
                        # itself before telling the user "the file is ready" (live
                        # case: deliverable copied to the VAF_Projects root, invisible
                        # to LAN clients). Informational only - explicit absolute
                        # targets stay allowed for the admin.
                        _ws_note = ""
                        try:
                            if _ws:
                                _res_r = Path(res).resolve()
                                _ws_r = Path(str(_ws)).resolve()
                                if not (_res_r == _ws_r or _res_r.is_relative_to(_ws_r)):
                                    _ws_note = (
                                        " (WARNING: this file is OUTSIDE the chat workspace and "
                                        "will NOT appear in the user's UI file browser. If it is "
                                        "the user's deliverable, save it AGAIN now with just the "
                                        "filename as a relative path so it lands in the workspace, "
                                        "and report THAT path to the user)"
                                    )
                        except Exception:
                            _ws_note = ""
                        if used_fallback:
                            return f"File written successfully to {res} (Desktop not writable, saved to Documents instead){_ws_note}"
                        return f"File written successfully to {res}{_ws_note}"
                    return f"⚠️ File written but size verification failed: {res}"
                        
                except Exception as e:
                    # Clean up temp file on error
                    try:
                        if os.path.exists(temp_path):
                            os.remove(temp_path)
                    except:
                        pass
                    raise
                    
            except PermissionError as e:
                if attempt < max_retries - 1:
                    time.sleep(retry_delay * (attempt + 1))
                    continue
                # Final attempt failed - return detailed error
                error_msg = str(e).lower()
                if 'denied' in error_msg or 'permission' in error_msg:
                    return (
                        f"❌ Permission denied: Cannot write to {res}\n"
                        f"Possible causes:\n"
                        f"- File is open in another program (editor, browser, etc.)\n"
                        f"- Insufficient file permissions\n"
                        f"- File is read-only\n"
                        f"Solution: Close any programs using this file and try again."
                    )
                return f"❌ Permission error: {e}"
            except OSError as e:
                if attempt < max_retries - 1:
                    time.sleep(retry_delay * (attempt + 1))
                    continue
                # Final attempt failed - return detailed error
                error_msg = str(e).lower()
                if 'locked' in error_msg or 'in use' in error_msg or 'being used' in error_msg:
                    return (
                        f"❌ File is locked: {res}\n"
                        f"The file is currently being used by another program.\n"
                        f"Please close any programs that have this file open and try again."
                    )
                elif 'no space' in error_msg or 'disk full' in error_msg:
                    return f"❌ Disk full: Cannot write to {res}. Free up disk space and try again."
                elif 'path too long' in error_msg:
                    return f"❌ Path too long: {res}. Use a shorter path."
                return f"❌ OS error: {e}"
            except Exception as e:
                # Generic error with clear message
                error_type = type(e).__name__
                return f"❌ Error writing file ({error_type}): {e}\nFile: {res}"


class EditFileTool(BaseTool):
    """Surgical search/replace edit of an EXISTING file - the alternative to rewriting the
    whole file with write_file. Modeled on Claude Code's exact old->new string replacement:
    each edit's `search` must match the file EXACTLY and UNIQUELY. All edits are applied
    all-or-nothing against the original buffer, then written once via WriteFileTool (so the
    home-reroute and atomic write are inherited). The per-user jail is installed by run()
    around the whole body - it cannot be left to the inner WriteFileTool call, which is
    handed no scope and would run unconfined."""


    identity_kwargs = ("user_role", "user_scope_id")
    file_access = "write"
    name = "edit_file"
    permission_level = "write"
    side_effect_class = "reversible"
    description = (
        "Change ONLY specific parts of an EXISTING file via search/replace, instead of "
        "rewriting the whole file. For each edit give the exact `search` text (with enough "
        "surrounding context that it is UNIQUE in the file) and its `replace`. PREFERRED over "
        "write_file for editing an existing file - it preserves everything you do not touch."
    )
    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Path to the EXISTING file to edit."},
            "edits": {
                "type": "array",
                "description": "One or more search/replace edits, applied together (all-or-nothing).",
                "items": {
                    "type": "object",
                    "properties": {
                        "search": {"type": "string", "description": "Exact text to find - must be unique in the file."},
                        "replace": {"type": "string", "description": "Text to replace it with."},
                    },
                    "required": ["search", "replace"],
                },
            },
        },
        "required": ["path", "edits"],
    }

    @staticmethod
    def _nearby(text: str, needle_first_line: str) -> str:
        """A few lines around the closest line to help the model retarget on a miss."""
        lines = text.splitlines()
        key = (needle_first_line or "").strip()
        for i, ln in enumerate(lines):
            if key and key in ln:
                lo, hi = max(0, i - 3), min(len(lines), i + 4)
                return "\n".join(lines[lo:hi])
        return "\n".join(lines[:6])

    @staticmethod
    def _locate(original: str, search: str):
        """Byte span (start, end) of `search` in `original`, or 'not_found' / 'ambiguous'.

        Exact match first; then a line-based fallback tolerant of carriage returns and
        trailing whitespace. The fallback only LOCATES the span - the original bytes are what
        get replaced, so nothing is silently normalized."""
        count = original.count(search)
        if count == 1:
            s = original.index(search)
            return (s, s + len(search))
        if count > 1:
            return "ambiguous"
        # Fallback: match line-by-line ignoring CR and trailing whitespace.
        def _n(line: str) -> str:
            return line.replace("\r", "").rstrip()
        o_lines = original.splitlines(keepends=True)
        s_norm = [_n(l) for l in search.splitlines()]
        if not s_norm:
            return "not_found"
        hits = []
        for i in range(len(o_lines) - len(s_norm) + 1):
            if [_n(o_lines[i + j]) for j in range(len(s_norm))] == s_norm:
                start = sum(len(o_lines[k]) for k in range(i))
                end = start + sum(len(o_lines[i + j]) for j in range(len(s_norm)))
                hits.append((start, end))
        if len(hits) == 1:
            return hits[0]
        return "ambiguous" if len(hits) > 1 else "not_found"


    def run(self, **kwargs) -> str:
        path = kwargs.get("path", "")
        edits = kwargs.get("edits")
        # Weak models sometimes pass a single {search,replace} or flat search/replace.
        if isinstance(edits, dict):
            edits = [edits]
        if not edits and ("search" in kwargs and "replace" in kwargs):
            edits = [{"search": kwargs.get("search", ""), "replace": kwargs.get("replace", "")}]
        if not path:
            return "Error: edit_file needs a 'path'."
        if not edits or not isinstance(edits, list):
            return "Error: edit_file needs 'edits' - a list of {search, replace} objects."

        safe, res = is_safe_path(path)
        if not safe:
            return res
        if not os.path.isfile(res):
            return (f"Error: edit_file target does not exist: {path}. "
                    "edit_file changes an EXISTING file; use write_file to create a new one.")
        try:
            with open(res, "r", encoding="utf-8") as f:
                original = f.read()
        except Exception as e:
            return f"Error: could not read {path}: {e}"

        # Phase 1 - locate every edit as a UNIQUE byte span in the ORIGINAL (all-or-nothing).
        spans = []  # (start, end, replace, edit_number)
        for i, ed in enumerate(edits, 1):
            if not isinstance(ed, dict) or "search" not in ed or "replace" not in ed:
                return f"Error: edit {i} must be an object with 'search' and 'replace'."
            search, replace = ed.get("search", ""), ed.get("replace", "")
            if not search:
                return f"Error: edit {i} has an empty 'search'."
            if search == replace:
                continue  # no-op edit
            loc = self._locate(original, search)
            if loc == "not_found":
                # Idempotency: old anchor gone AND the replacement already present -> already
                # applied (e.g. a retry). Skip rather than fail the whole call.
                _sfirst = (search.splitlines()[0].strip() if search.splitlines() else search.strip())
                if replace and replace in original and _sfirst and _sfirst not in original:
                    continue
                # Whole-file rewrite disguised as an edit: a weaker model "edits" by pasting the
                # entire old file as `search` and the entire new file as `replace`; a few chars of
                # drift then break the exact match. When the search covers essentially the WHOLE
                # file, `replace` IS the intended new file, so rescue the work as a write_file
                # rather than trashing it. Gated to the sole edit covering ~the whole file — a
                # partial huge chunk's `replace` is only a fragment and writing it would destroy
                # the rest of the file (there is nothing safe to rescue there).
                if len(edits) == 1 and replace.strip() and len(search) >= 0.9 * max(1, len(original)):
                    write_res = WriteFileTool().run(path=res, content=replace)
                    _wr = str(write_res)
                    if _wr.startswith("❌") or _wr.startswith("Error"):
                        return write_res
                    return (f"Edited {os.path.basename(res)}: the search covered essentially the whole "
                            f"file, so this was a full rewrite - applied via write_file. For a rewrite, "
                            f"call write_file directly; keep edit_file for small, surgical changes.")
                _hint = ""
                if len(search) > 2000:
                    _hint = ("\nThis search block is very large. To rewrite the whole file use write_file "
                             "with the new content; for a targeted change use a smaller, exact search.")
                return (f"EDIT FAILED: edit {i}'s search block was not found - read_file the region "
                        f"and copy the EXACT text (including indentation). Nothing was written.{_hint}\n"
                        f"--- nearest region ---\n{self._nearby(original, _sfirst)}")
            if loc == "ambiguous":
                return (f"EDIT FAILED: edit {i}'s search block is not unique - add more surrounding "
                        f"context so it matches exactly once. Nothing was written.")
            spans.append((loc[0], loc[1], replace, i))

        if not spans:
            return "Error: edit_file made no change (nothing to apply, or every edit was already applied)."

        # Reject overlapping edits - a sequential replace would corrupt them.
        spans.sort()
        for a, b in zip(spans, spans[1:]):
            if a[1] > b[0]:
                return (f"EDIT FAILED: edits {a[3]} and {b[3]} target overlapping regions - "
                        f"combine them into one edit. Nothing was written.")

        # Apply span-based (order-independent, atomic), then ONE write via WriteFileTool.
        parts, last = [], 0
        for start, end, replace, _ in spans:
            parts.append(original[last:start])
            parts.append(replace)
            last = end
        parts.append(original[last:])
        new_content = "".join(parts)
        if new_content == original:
            return "Error: edit_file produced no change."

        write_res = WriteFileTool().run(path=res, content=new_content)
        _wr = str(write_res)
        if _wr.startswith("❌") or _wr.startswith("Error"):
            return write_res
        return f"Edited {os.path.basename(res)}: applied {len(spans)} change(s), everything else preserved."


class MoveFileTool(BaseTool):
    name = "move_file"
    permission_level = "write"
    side_effect_class = "reversible"
    description = "Moves or renames a file. Use for renaming: src=current_path, dst=same_folder/new_name.ext"

    parameters = {
        "type": "object",
        "properties": {
            "src": {"type": "string"},
            "dst": {"type": "string"}
        },
        "required": ["src", "dst"]
    }

    def run(self, **kwargs) -> str:
        src = kwargs.get('src', '')
        dst = kwargs.get('dst', '')
        
        safe_src, res_src = is_safe_path(src)
        if not safe_src: return f"Source Error: {res_src}"
        
        safe_dst, res_dst = is_safe_path(dst)
        if not safe_dst: return f"Dest Error: {res_dst}"
        
        try:
            shutil.move(res_src, res_dst)
            return f"Moved {res_src} to {res_dst}"
        except Exception as e: return str(e)

class TreeTool(BaseTool):
    name = "tree"
    identity_kwargs = ("user_role", "user_scope_id")
    file_access = "read"
    permission_level = "read"
    side_effect_class = "none"
    description = "Generates an ASCII tree view of a directory structure."

    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Root path (default: current dir)"},
            "depth": {"type": "integer", "description": "Max depth (default: 2)"}
        },
        "required": []
    }


    def run(self, **kwargs) -> str:
        path = kwargs.get('path', '.')
        try:
            depth = int(kwargs.get('depth', 2))
        except: depth = 2
        
        safe, res = is_safe_path(path)
        if not safe: return res
        
        if not os.path.exists(res): return "Error: Path not found."
        if not os.path.isdir(res): return "Error: Path is not a directory."
        
        tree_str = f"{os.path.basename(res)}/\n"
        
        for root, dirs, files in os.walk(res):
            level = root.replace(res, '').count(os.sep)
            if level >= depth: 
                del dirs[:] # Stop recursing
                continue
            
            indent = "│   " * (level)
            subindent = "├── "
            
            # Limited output for sanity
            if len(files) + len(dirs) > 50:
                 files = files[:40]
                 files.append(f"... (+{len(files)-40} more)")
            
            for d in dirs:
                tree_str += f"{indent}{subindent}{d}/\n"
            for f in files:
                tree_str += f"{indent}{subindent}{f}\n"
                
        return tree_str

class FinderTool(BaseTool):
    name = "find_files"
    identity_kwargs = ("user_role", "user_scope_id")
    file_access = "read"
    permission_level = "read"
    side_effect_class = "none"
    description = "Finds files matching a glob pattern (e.g. *.py) recursively."
    
    parameters = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Root path to search in"},
            "pattern": {"type": "string", "description": "Glob pattern (e.g. *.json, app.py, *test*)"}
        },
        "required": ["pattern"]
    }


    def run(self, **kwargs) -> str:
        import fnmatch
        path = kwargs.get('path', '.')
        pattern = kwargs.get('pattern', '*')
        
        safe, res = is_safe_path(path)
        if not safe: return res
        
        matches = []
        try:
            for root, dirs, files in os.walk(res):
                for name in files:
                    if fnmatch.fnmatch(name, pattern):
                        full_path = os.path.join(root, name)
                        # Return relative to search root for readability, or absolute?
                        # Using absolute is clearer for tools.
                        matches.append(full_path)
                        
                if len(matches) > 100:
                    break # Safety limit
            
            if not matches: return "No files found matching pattern."
            
            # Limit output
            count = len(matches)
            output = f"Found {count} files:\n" + "\n".join(matches[:50])
            if count > 50: output += f"\n... ({count-50} more)"
            return output
            
        except Exception as e: return str(e)
