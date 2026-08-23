# SPDX-FileCopyrightText: 2026 Veyllo GmbH
# SPDX-License-Identifier: AGPL-3.0-or-later
# Additional permissions and terms under AGPL Section 7: see LICENSING.md
from vaf.startup_logger import log
log("WebServer", "Module load started")

from fastapi import FastAPI, WebSocket, WebSocketDisconnect, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse, FileResponse
from fastapi import HTTPException, Query
from starlette.requests import Request
import asyncio
import uvicorn
import threading
import inspect
import html
import os
import time
import queue
log("WebServer", "Basic imports done")

from vaf.core.web_interface import get_web_interface
from vaf.core.session import SessionManager, Session
from vaf.cli.autosuggest import SmartAutoSuggest
import json
from vaf.core.config import Config
from vaf.version import __version__
from vaf.core.log_helper import append_domain_log, get_dated_log_path, is_debug_logging_enabled
from pathlib import Path
from vaf.core.path_jail import PathEscape, contained_path, safe_entry_name
from typing import Dict, Optional, List
import logging
from vaf.core.tray_context import TrayContext
log("WebServer", "VAF imports done")

log_uvicorn = logging.getLogger("uvicorn")

app = FastAPI(title="VAF Local Server", version=__version__)


@app.get("/api/version")
async def api_version():
    """The running VAF version (single source of truth: vaf/version.py)."""
    return {"version": __version__}

# Active model download cancel events keyed by websocket id (for cancel_model_download)
_active_model_download_cancels: dict[int, threading.Event] = {}

# Active attachment-indexing tasks keyed by session id, so the stop button can cancel
# in-flight document indexing (the LLM RAG-indexes attachments in the background).
_active_index_tasks: "dict[str, set[asyncio.Task]]" = {}
_headless_agent_thread: Optional[threading.Thread] = None
_headless_agent_lock = threading.Lock()


def _ensure_headless_agent_runner(origin: str = "web_server") -> None:
    """
    Ensure exactly one in-process headless runner thread is alive.
    Queue producers (WebSocket/API) run in this process, so the consumer must also
    live here to avoid QUEUE_ADD without QUEUE_GET.
    """
    global _headless_agent_thread
    with _headless_agent_lock:
        if _headless_agent_thread is not None and _headless_agent_thread.is_alive():
            return
        from vaf.core.headless_runner import run_headless_agent

        _headless_agent_thread = threading.Thread(
            target=run_headless_agent,
            daemon=True,
            name="HeadlessAgent",
        )
        _headless_agent_thread.start()
        log("WebServer", f"Headless agent runner started ({origin})")


@app.exception_handler(Exception)
async def json_exception_handler(request, exc):
    """Ensure unhandled exceptions return JSON (for API clients) instead of HTML."""
    from fastapi.responses import JSONResponse
    # HTTPException: return JSON; other exceptions: 500 with error message
    if isinstance(exc, HTTPException):
        return JSONResponse(status_code=exc.status_code, content={"detail": exc.detail, "error": str(exc.detail)})
    log_uvicorn.exception("Unhandled exception: %s", exc)
    return JSONResponse(status_code=500, content={"ok": False, "error": str(exc), "detail": str(exc)})


# CORS: Use regex to allow localhost AND all RFC 1918 private network origins.
# This is safe because Layer 2 (IPValidationMiddleware) already blocks non-local IPs.
_CORS_ORIGIN_REGEX = r"^https?://(localhost|127\.0\.0\.1|10\.\d+\.\d+\.\d+|172\.(1[6-9]|2\d|3[0-1])\.\d+\.\d+|192\.168\.\d+\.\d+)(:\d+)?$"

app.add_middleware(
    CORSMiddleware,
    allow_origin_regex=_CORS_ORIGIN_REGEX,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# The sidebar list's limit and its ONE projection both live with the store that
# answers them (vaf/core/session.py), so a surface that is not the web server -
# the classic CLI pushes this list too - can build the same rows without
# importing the web application.
from vaf.core.session import SESSION_LIST_LIMIT, session_list_payload  # noqa: E402

# Max room frames painted into the browser at once. A room is a finite conversation
# by design (the named limit in the protocol doc), and every frame is a separate
# encrypted file, so the tail is what a reader wants and the whole is what
# `vaf a2a export` is for.
ROOM_TRANSCRIPT_LIMIT = 500

# How long a keypress in the room's input box counts as composing. The client
# refreshes it every ~2 seconds while keys are actually pressed and the 3s room
# poll carries it to the other viewers, so the dots die 3-6 seconds after the
# last key - a bubble that outlives the typing is the lie this replaces.
ROOM_KEYS_TYPING_WINDOW_S = 6.0

# Who is PRESSING KEYS in which room's input box, {room_id: {peer_id: last_ts}}.
# In-memory and ephemeral BY DESIGN: typing is presence, not conversation. It
# never touches the store or the wire (the protocol deliberately has no typing
# frame kind), and a restart forgetting it is correct behaviour. Written by the
# `room_typing` websocket message, read and expired by the room projection.
_ROOM_KEYS_TYPING: Dict[str, Dict[str, float]] = {}


def derive_room_presence(*, activity: Dict[str, Dict[str, float]],
                         members: Dict[str, dict], labels: Dict[str, str],
                         acting: str, keys: Dict[str, float],
                         busy_agent: str, now: float) -> "tuple[list, list]":
    """(typing, read_positions) for the room projection. Pure, so it is testable.

    Typing keeps exactly two honest signals: the host agent's TURN MARKER (it is
    composing an answer right now) and a human's KEYPRESSES in the input box.
    The old third signal - "took the newest message recently and has not
    answered" - painted every reader as composing for two minutes, which made an
    agent that merely monitors its room look permanently busy. That reading is
    not discarded, it is PROMOTED: every member's read position goes out as a
    read receipt, which says what the derivation actually knew ("has read up to
    here") instead of what it guessed ("is about to speak").
    """
    def label_of(peer: str) -> str:
        return labels.get(peer) or members[peer]["display"]

    typing: list = []
    read_positions: list = []
    for peer, facts in activity.items():
        if peer == acting or peer not in members:
            # The viewer knows what they have read; their own receipt is noise.
            continue
        if int(facts.get("read_to") or 0) > 0:
            read_positions.append({"peer": peer, "label": label_of(peer),
                                   "readTo": int(facts["read_to"]),
                                   "readAt": float(facts.get("read_at") or 0.0)})
    for peer, pressed_at in list(keys.items()):
        if (now - float(pressed_at)) > ROOM_KEYS_TYPING_WINDOW_S:
            keys.pop(peer, None)
            continue
        if peer == acting or peer not in members:
            continue
        typing.append({"peer": peer, "label": label_of(peer), "kind": "keys"})
    if busy_agent and busy_agent in members:
        typing.append({"peer": busy_agent, "label": label_of(busy_agent),
                       "kind": "turn"})
    return typing, read_positions

# The sidebar filter (channel + thinking sessions out) lives on the engine as
# SessionManager.list_ui - one rule for every surface, the terminal app's
# panel included. The channel prefixes come from the dispatch registry there;
# this file carried a private copy of both until they diverged from nothing
# but luck.


#: How long finished work stays on the LIVE board. Half an hour, which is long enough
#: to see that the thing you just asked for is done and short enough that a room does
#: not turn into its own changelog. What leaves it is not gone: it moves to the room's
#: history, which is fetched on demand rather than riding along in a payload that is
#: polled every three seconds.
#:
#: The pattern is the ordinary one (Google Tasks folds, Linear and Gmail archive,
#: monday auto-archives after N days) with one VAF-specific simplification: the archive
#: already exists and is authoritative - the transcript is write-once and
#: `vaf a2a tasks` prints every task there has ever been - so the panel needs a WINDOW,
#: never an archive of its own.
ROOM_DONE_WINDOW_S = 1800.0


def _room_board(room, *, open_cap: int = 40, done_cap: int = 8,
                done_window_s: float = ROOM_DONE_WINDOW_S) -> list:
    """The task board a browser gets: every piece of OPEN work, and what was finished
    recently.

    The cap used to be a flat slice of the whole board, which cuts by recency and
    therefore cuts the wrong end: an agent that finishes twenty things in an afternoon
    pushes the one task still running off a list sorted with finished work in it. Open
    work is what a surface is for, so it is never dropped for something that is over.

    Finished work is bounded twice, by age and by count, and neither is a deletion: the
    record lives in the transcript and in `vaf a2a tasks`, which is where a long history
    belongs. A panel is not an archive.
    """
    import time as _t
    done_states = ("completed", "failed", "rejected", "canceled")
    board = room.tasks()
    cutoff = _t.time() - float(done_window_s or 0.0)
    open_work = [t for t in board if t["status"] not in done_states]
    finished = [t for t in board if t["status"] in done_states
                and (done_window_s <= 0 or float(t.get("updated_ts") or 0.0) >= cutoff)]
    return open_work[:open_cap] + finished[:done_cap]


def _worker_card_entry(task) -> dict:
    """One worker card, the shape the reference design draws: name, status, the
    task line, progress, and WHEN it started (the meta line reads
    'Coder · 1 Worker · gestartet HH:MM')."""
    started = ""
    try:
        raw = str(getattr(task, "created_at", "") or "")
        if raw:
            from datetime import datetime as _dt2
            started = _dt2.fromisoformat(raw).strftime("%H:%M")
    except Exception:
        started = ""
    return {
        "type": str(getattr(task, "agent_type", "") or ""),
        "status": str(getattr(task, "status", "") or ""),
        "task": str(getattr(task, "task_description", "") or "")[:80],
        "done": getattr(task, "progress_done", None),
        "total": getattr(task, "progress_total", None),
        "started": started,
    }


def _viewer_agent_workers(user_scope_id) -> list:
    """The sub-agents the VIEWER'S own agent is running right now, for the room view.

    Mirrors what the chat's TasksLine already shows - the IPC active list - rather
    than inventing a second status lane. The filter is the point, and it FAILS
    CLOSED: the IPC file is global to every user and its records carry a session
    id, not a scope, so a task is shown only when its session PROVABLY belongs to
    the viewer (same ownership rule the workspace resolver applies: a scopeless
    legacy session belongs to the local admin, never to everyone). A task with
    neither a session nor a room is never shown - its description is user
    content, and "probably mine" is not an ownership answer. A SESSIONLESS task
    that a room ordered is proven by that room's manifest instead: room turns
    may legitimately run without a session, and the room's tenant is exactly who
    the feed belongs to.

    Foreign room members never get worker cards from here: their workers run on
    their machines, which this store cannot see. A named boundary, not a gap -
    the protocol's card is where a peer may SAY what it is running.
    """
    out = []
    try:
        from vaf.core.config import get_local_admin_scope_id
        from vaf.core.subagent_ipc import get_ipc
        for task in get_ipc().get_active_tasks(None):
            sid = getattr(task, "session_id", None)
            if not sid:
                # A room turn may spawn with NO session at all (the runner's
                # room frame binds no chat). For those, the ordering room's own
                # manifest is the ownership proof - stronger than the session
                # proxy, and still fail-closed: no room, or a room whose tenant
                # is not the viewer, shows nothing.
                room_id = str(getattr(task, "room_id", "") or "").strip()
                if not room_id:
                    continue
                owner = get_web_interface().room_owner_scope(room_id)
                if owner is None or str(owner) != str(user_scope_id):
                    continue
                out.append(_worker_card_entry(task))
                continue
            try:
                session_scope = (getattr(session_mgr.load(sid), "metadata", None)
                                 or {}).get("user_scope_id")
            except Exception:
                continue
            owned = (str(session_scope) == str(user_scope_id)
                     if session_scope is not None
                     else str(user_scope_id) == str(get_local_admin_scope_id()))
            if not owned:
                continue
            out.append(_worker_card_entry(task))
    except Exception:
        return []
    return out[:8]


def _derive_or_empty(lane: str, user_scope_id, room_id: str) -> str:
    """A derived room handle, or "" when it cannot be derived. Never raises: the
    transcript must paint even for a viewer whose scope resolves nothing."""
    try:
        from vaf.core.a2a.room import derive_peer_id, participant_key
        return derive_peer_id(participant_key(lane, user_scope_id), room_id)
    except Exception:
        return ""


async def _send_room_transcript(websocket, room, user_scope_id: Optional[str]) -> None:
    """The room as the browser reads it, from the store, every time.

    One place, because three commands answer with it - opening a room, writing into
    one, closing one - and a browser that repainted from what it ASSUMED had happened
    would drift from the transcript the moment anything else wrote at the same time.
    The room has N writers; only the store knows what is in it.
    """
    from vaf.core.a2a.room import (CAPABILITIES, derive_peer_id, describe,
                                   owner_tenant, participant_key)
    from vaf.core.session import _room_rows

    row = next((r for r in _room_rows(user_scope_id) if r["room_id"] == room.room_id), {})
    entries = room.transcript()[-ROOM_TRANSCRIPT_LIMIT:]
    members, labels, hosts = room.members(), room.labels(), room.host_peers()
    # Who belongs to whom, once for the whole payload rather than per member: it is a
    # derivation over the same accounts every time, and asking it per row would walk
    # the room's own member list once per member.
    try:
        pairs = room.pairs()
    except Exception:
        pairs = {}
    # The handle the browser ACTS as, which is the CLI lane - not whichever lane the
    # sidebar row happened to resolve first. Asking about the wrong one is how the
    # remove buttons stayed hidden for the host of their own room: the row answered for
    # the agent while the browser was acting as the person.
    try:
        acting = derive_peer_id(participant_key("cli", user_scope_id), room.room_id)
    except Exception:
        acting = ""
    # The name the room shows for this person is corrected the moment they LOOK at the
    # room, not only when they act in it. The heal used to live in the command branch
    # (say/close/kick/rename), which meant a member record written as the lane literal
    # by an older CLI stayed "terminal" for anyone who only ever read their room - and
    # reading is what a person does most. Every room command ends here, so this is the
    # one place that covers all of them. Admin-only by the same measure that makes it
    # safe: today the room's own tenant is the only account that can act on it at all
    # (Room._check_tenant), so the account name IS the local admin's; a per-user name
    # lookup becomes necessary exactly when cross-account rooms do.
    try:
        if acting and user_scope_id is not None:
            from vaf.core.config import (get_local_admin_scope_id,
                                         get_local_admin_username)
            record = room.store.member(acting) or {}
            if (record and str(record.get("display") or "") in ("terminal", "guest", "")
                    and str(user_scope_id) == str(get_local_admin_scope_id())):
                wanted_name = str(get_local_admin_username() or "").strip()
                healed = room.identity_for(participant_key("cli", user_scope_id))
                if wanted_name and healed is not None:
                    room.introduce(healed, display=wanted_name)
                    members, labels = room.members(), room.labels()
    except Exception:
        pass
    # Who is composing, and who has read how far. Composing keeps two honest
    # signals - the host agent's turn marker (precise) and a human's keypresses
    # in the input box (reported by the browser, expired by the clock). What used
    # to be the third signal, "read the newest message recently and answered
    # nothing", is promoted into READ RECEIPTS instead of being painted as
    # typing: it made a merely-monitoring agent look busy for two minutes at a
    # time. None of this is a wire concept - see Room.activity() for why there
    # is no "typing" frame kind.
    typing = []
    read_positions = []
    try:
        import time as _time
        latest = room.store.highest_lamport()
        if latest and not room.closed:
            # The host agent's handle, needed to resolve the turn marker to a
            # peer. Only the marker may paint this agent as composing: its
            # cursor advances right after every turn, including one in which it
            # deliberately said nothing, so any cursor-based composing rule
            # would show it busy precisely when it chose to stay quiet.
            host_agent = ""
            try:
                owner = owner_tenant(room.manifest.get("owner_scope"))
                if owner:
                    host_agent = derive_peer_id(participant_key("agent", owner),
                                                room.room_id)
            except Exception:
                host_agent = ""
            busy_agent = ""
            try:
                turn = getattr(getattr(manager, "agent_instance", None),
                               "_room_turn", None)
                if isinstance(turn, dict) and str(turn.get("room_id")) == room.room_id:
                    # The marker names the room, not the peer; the answering agent
                    # on this machine is the host account's agent lane.
                    busy_agent = host_agent
            except Exception:
                busy_agent = ""
            typing, read_positions = derive_room_presence(
                activity=room.activity(), members=members, labels=labels,
                acting=acting,
                keys=_ROOM_KEYS_TYPING.setdefault(room.room_id, {}),
                busy_agent=busy_agent, now=_time.time())
    except Exception:
        typing = []
        read_positions = []
    await websocket.send_json({
        "type": "room_transcript",
        "room": {
            "id": row.get("id") or f"room:{room.room_id}",
            "roomId": room.room_id,
            "title": row.get("name") or room.manifest.get("topic") or room.room_id,
            "roomKind": room.kind,
            "role": row.get("role"),
            "closed": bool(room.closed),
            "members": len(members),
            # Who is in it, by the name the ROOM resolved. A header that listed join
            # names would show two agents called "Codex" and no way to tell which is
            # which, which is the same reason the tag exists at all.
            # Whether the viewer may remove anybody here at all. Answered by the room
            # rather than guessed from the role, because the host may act beyond its
            # role and a browser that decided for itself would be a second rule.
            "createdAt": room.manifest.get("created_at"),
            "topic": room.manifest.get("topic") or "",
            "me": acting or row.get("peer"),
            # The viewer's own AGENT in this room, so the view can draw it as the
            # agent it is (the living avatar) rather than as initials. Derived the
            # same way every identity here is; empty when their agent is not a
            # member. A foreign agent never gets this treatment: an agent that is
            # not ours is a full agent of its own, never a second face of ours.
            "agentPeer": next((p for p in [
                _derive_or_empty("agent", user_scope_id, room.room_id)] if p in members), ""),
            # How far the viewer has allowed THEIR agent to act on this room's
            # messages (observe/assist/autonomous). Shown and set from the member
            # panel; empty when their agent is not in the room.
            "agentMode": next((room.mode_of(p) for p in [
                _derive_or_empty("agent", user_scope_id, room.room_id)] if p in members), ""),
            # The workers the viewer's agent is running RIGHT NOW - the same list
            # the chat's TasksLine reads, scoped hard to the viewer. Refreshed by
            # the 3s poll like everything else here.
            "agentWorkers": (_viewer_agent_workers(user_scope_id)
                             if _derive_or_empty("agent", user_scope_id, room.room_id)
                             in members else []),
            "typing": typing,
            # Read receipts: every other member's read position as facts, so the
            # view can stack each reader's face under the last message they took.
            # Remote wire peers keep their position in their own seat file and
            # never appear here - the documented protocol boundary, not a bug.
            "readPositions": read_positions,
            # The task board, derived exactly like typing is (Room.tasks): a
            # directive, or anything a report chain answers via reply_to. The 3s
            # poll keeps it fresh; nothing travels on the wire for it. Capped the
            # way the transcript is - a board is a glance, not an archive.
            "tasks": [
                {"id": t["id"], "title": t["title"], "status": t["status"],
                 "requester": t["requester_label"], "assignee": t["assignee_label"],
                 "reports": t["reports"], "ts": t["updated_ts"],
                 # Nothing said about it for hours. Not "finished" - nobody said that -
                 # but not work in progress either, and a surface that counts it as
                 # such fills up with entries nobody is doing.
                 "quiet": bool(t.get("quiet")),
                 "silentFor": float(t.get("silent_for_s") or 0.0),
                 "result": str(t.get("result") or "")[:400],
                 # Rebuilt field by field here, so every new one has to be named
                 # or it vanishes without a word - which has shipped twice. This
                 # one is how far the work has come, from the last report that
                 # said anything about it.
                 "progress": t.get("progress")}
                for t in _room_board(room)
            ],
            # Open votes, folded like the board is. Capped for the same reason:
            # a glance, not an archive.
            "votes": [
                {"id": v["id"], "question": v["question"], "options": v["options"],
                 "askedBy": v["asked_by_label"], "tally": v["tally"],
                 "voted": v["voted"], "waitingFor": v["waiting_for"],
                 "closed": v["closed"],
                 # The countdown runs in the browser off this one number, rather than
                 # off a seconds-left the server recomputes every poll: a clock that
                 # only moves when a poll lands stutters, and this way the card is
                 # exact between polls and cannot drift apart from the deadline the
                 # room will actually act on.
                 "deadline": v["deadline"],
                 "everyoneVoted": v["everyone_voted"],
                 "mine": next((b["choice"] for b in v["ballots"]
                               if b["peer"] == acting), ""),
                 "ballots": [{"label": b["label"], "choice": b["choice"]}
                             for b in v["ballots"]][:20]}
                # A vote leaves this list when the ROOM has said how it ended, not
                # when the clock ran out: the card and the result message change
                # places, so there is never a moment with neither on screen. The
                # sweep writes that result within its own tick.
                for v in room.votes() if not v["concluded"]
            ][:6],
            "mission": str(room.manifest.get("mission") or ""),
            "canManage": bool(acting) and (
                acting in hosts or room.role_of(acting) == "leader"),
            # Built from members(), which already resolves the role, the card and the
            # lease. Reading those out of the store again here would be a second answer
            # to "who is in this room and are they awake".
            "members_list": [
                {
                    "peer": peer,
                    "label": labels.get(peer) or record["display"],
                    "role": record["role"],
                    "card": record.get("card") or {},
                    "stale": bool(record.get("stale")),
                    # The raw lease, so a reader can say WHEN rather than only that it
                    # lapsed. Nobody heartbeats a file-only room, so "not responding"
                    # on its own reads as a fault where none exists.
                    "lease": float((room.store.member(peer) or {}).get("lease") or 0.0),
                    # What this role may send, from the table that enforces it. An
                    # abilities list written anywhere else would be a second answer.
                    "can": sorted(CAPABILITIES.get(record["role"], ())),
                    # The room's own host handles. Sent so the view can leave the remove
                    # button off entirely rather than offer one that is refused: an
                    # action offered and then denied reads as a fault, and this is a
                    # deliberate rule with a different answer (close the room).
                    "protected": peer in hosts,
                    # WHO this member is, and whose. Derived by the room from the
                    # account each handle was built from, never claimed by the member -
                    # a `speaks_for` in a member's own file would be a peer naming its
                    # own partner, and in a room of strangers that is the one claim
                    # nobody may make. A guest that arrived on an invitation named no
                    # account and stays "unknown" rather than being guessed at.
                    "kind": (pairs.get(peer) or {}).get("kind") or "unknown",
                    "partner": (pairs.get(peer) or {}).get("partner") or "",
                    "partnerLabel": (pairs.get(peer) or {}).get("partner_label") or "",
                }
                for peer, record in sorted(
                    members.items(), key=lambda kv: labels.get(kv[0]) or kv[1]["display"])
            ],
        },
        "messages": [
            {"id": e["id"], "peer": e["peer"], "label": e["label"],
             "role": e["role"], "kind": e["kind"],
             "text": describe(e), "ts": e["ts"],
             # Forwarded EXPLICITLY, like every field this projection carries:
             # the handler rebuilds the row field by field, so anything not
             # named here is silently dropped on its way to the browser - the
             # shape that has cost this file two features already.
             "files": e.get("files") or [],
             "lamport": e["lamport"], "to": e.get("to") or {}}
            for e in entries
            # A check-in is the room talking to ONE agent about its own attention,
            # not something anybody said, and drawing it as a message would put
            # words in a member's mouth that no member wrote. It stays in the log
            # (and in `vaf a2a log`) where an audit can find it.
            if e["kind"] != "ping"
        ],
    })

    # Looking at the room IS reading it. The person's cursor (the cli lane the
    # browser shares with the terminal) moves to the newest frame - AFTER the
    # transcript went out, never before, so an interruption costs a repeat and not
    # a swallowed line - and the sidebar badge that counts from this cursor goes
    # out with it. Needs no membership: a reading position is the reader's own
    # file. The notify fires only on actual movement, or the 3-second room poll
    # would broadcast a sidebar refresh forever.
    try:
        if acting and not room.closed:
            latest = room.store.highest_lamport()
            if latest and room.store.cursor(acting) < latest:
                room.store.set_cursor(acting, latest)
                try:
                    from vaf.core.web_interface import notify_rooms_changed
                    notify_rooms_changed(user_scope_id)
                except Exception:
                    pass
    except Exception:
        pass


log("WebServer", "Getting WebInterfaceManager...")
manager = get_web_interface()

# ═══════════════════════════════════════════════════════════════════════════════
# WHISPER MODEL SINGLETON - Prevents memory leak from reloading model per request
# ═══════════════════════════════════════════════════════════════════════════════
_whisper_model = None
_whisper_model_lock = threading.Lock()

# Live-call state per WebSocket connection (voice-agent first layer):
# {id(websocket): {"history": [...], "lang": str, "scope": str}}. Entries are removed on
# voice_call_end AND on socket teardown.
#
# A stale entry is NOT harmless any more (the comment here used to say it was, back when
# nothing read this dict as a signal). The tray's idle watchdog now treats an entry as "a
# call is live, keep the local model loaded", so a single orphan would pin the model for the
# life of the process and silently disable idle unloading. voice_call_end is sent by exactly
# one frontend site, inside endCall, so any abrupt teardown - closed tab, refresh, a wifi
# blip that makes the client reconnect on a NEW id(websocket) - would leave one behind.
# Hence the teardown pop, plus a liveness intersection at the reading end.
_VOICE_CALLS: dict = {}

def _rss_mb() -> float:
    """Resident memory in MB for the load log, 0.0 when psutil is unavailable.

    Diagnostics must never decide whether transcription works.
    """
    try:
        import psutil
        return psutil.Process().memory_info().rss / (1024 * 1024)
    except Exception:
        return 0.0


def get_whisper_model():
    """Get or lazily load the Whisper STT model (singleton)."""
    global _whisper_model
    if _whisper_model is None:
        with _whisper_model_lock:
            if _whisper_model is None:  # Double-check
                model_size = Config.get("speech_stt_whisper_model", "base")
                # ONLY the faster-whisper import may raise ImportError out of
                # here: the caller turns that into "local STT is not usable",
                # and the memory bookkeeping below used to sit in the same try,
                # so a missing psutil was reported as a missing faster-whisper.
                # The original message travels with it - an installed package
                # whose native dependency fails to load (arch mismatch, a
                # frequent macOS case) says something different from an absent
                # one, and the user needs to see which.
                try:
                    import importlib

                    whisper_module = importlib.import_module("faster_whisper")
                    WhisperModel = getattr(whisper_module, "WhisperModel")
                except ImportError as exc:
                    raise ImportError(
                        f"faster-whisper is not usable ({exc}). "
                        'Install it with: pip install "vaf[speech]"'
                    ) from exc

                mem_before = _rss_mb()
                log("WebServer", f"Loading WhisperModel ({model_size}, CPU, int8) - Memory before: {mem_before:.0f}MB")
                _whisper_model = WhisperModel(model_size, device="cpu", compute_type="int8")
                mem_after = _rss_mb()
                log("WebServer", f"WhisperModel ({model_size}) loaded - Memory after: {mem_after:.0f}MB (delta: {mem_after-mem_before:.0f}MB)")

                # Log to file (consolidated in memory.log)
                append_domain_log("memory", f"[WHISPER] WhisperModel({model_size}): {mem_before:.0f}MB -> {mem_after:.0f}MB (delta: {mem_after-mem_before:.0f}MB)")
    return _whisper_model

def unload_whisper_model():
    """Unload Whisper model to free memory."""
    global _whisper_model
    with _whisper_model_lock:
        if _whisper_model is not None:
            del _whisper_model
            _whisper_model = None
            import gc
            gc.collect()
            log("WebServer", "WhisperModel unloaded")
log("WebServer", "Getting SessionManager...")
session_mgr = SessionManager()
log("WebServer", "SmartAutoSuggest will be lazy loaded...")
autosuggest = None
tray_context = TrayContext()


def _ws_session_owner_ok(websocket, session_id, *, loaded=None, allow_missing=False):
    """Ownership gate for WebSocket session commands (chat / load / delete / rename / hide / artifact_edit).

    Returns (allowed: bool, loaded_session_or_None). This resolves WHO is asking from the
    connection; the rule itself is `SessionManager.may_access`, so the HTTP endpoints that
    name a session id answer identically. Admin is detected role-aware (connection role ==
    'admin' OR connection scope == the local-admin scope) so the desktop owner is never
    locked out even when its scope is None. `allow_missing=True` passes a not-yet-created
    id, for the chat's first-message-into-a-new-session flow.
    """
    from vaf.core.config import get_local_admin_scope_id
    user_scope_id = manager.get_connection_user(websocket)
    role = manager.get_connection_user_role(websocket)
    is_admin = (str(role or "").lower() == "admin") or (
        user_scope_id is not None and str(user_scope_id) == str(get_local_admin_scope_id())
    )
    if loaded is not None:
        if is_admin:
            return (True, loaded)
        session_scope = (getattr(loaded, "metadata", None) or {}).get("user_scope_id")
        allowed = session_scope is not None and str(session_scope) == str(user_scope_id)
        return (allowed, loaded)
    from vaf.core.session import session_access_allowed
    return session_access_allowed(
        session_mgr,
        session_id,
        user_scope_id=user_scope_id,
        is_admin=is_admin,
        allow_missing=allow_missing,
    )

def _session_record_exists(session_id) -> bool:
    """Is there a file for this chat at all?

    The ownership gate cannot answer it: a record that is missing and one that
    cannot be READ both come back as `(allowed, None)`, and the two need
    opposite treatment when a delete arrives without a confirmation. An
    unreadable record must be asked about; an absent one has nothing to ask
    about, and a dialog for a chat that is already gone (a double-clicked trash
    icon, a second tab) would be a question with no subject.
    """
    try:
        return any((session_mgr.storage_dir / f"{session_id}{ext}").exists()
                   for ext in (".json", ".json.gz"))
    except Exception:
        return True          # cannot tell -> treat as present -> ask


def _clean_history_text(text: str) -> str:
    """Strip legacy rich-markup from a stored assistant message.

    Old sessions carry terminal markup ([dim] ... [/dim]) where the web UI
    expects think tags, so a reloaded chat would otherwise show the markup as
    literal text.
    """
    import re as _re
    if not text:
        return ""
    text = _re.sub(r'\[dim\]', '<think>', text, flags=_re.IGNORECASE)
    text = _re.sub(r'\[white dim\]', '<think>', text, flags=_re.IGNORECASE)
    text = _re.sub(r'\[.*?dim.*?\]', '<think>', text, flags=_re.IGNORECASE)
    text = text.replace('[/dim]', '</think>')
    if '<think>' in text and '</think>' not in text:
        text = text.replace('[/]', '</think>')
    text = _re.sub(r'\[\/?[^\]]+\]', '', text)
    return text


def _history_projection(loaded, sid) -> list:
    """The stored session as the chat transcript the web UI renders.

    One function because there were two hand copies of this loop, and they had
    already drifted: only one of them turned a user message's images back into
    URLs, so the same chat showed its own attachments on one path and not on the
    other. Every surface that repaints a transcript goes through here.

    Not pure: for a thinking session it POPS the stored user reply, which is
    what makes that reply appear exactly once. Both copies did this, and moving
    it here keeps that property rather than duplicating the side effect.
    """
    out = []
    for msg in getattr(loaded, "messages", None) or []:
        role = msg.get("role") if isinstance(msg, dict) else getattr(msg, "role", "user")
        content = msg.get("content") if isinstance(msg, dict) else getattr(msg, "content", "")
        timestamp = msg.get("timestamp") if isinstance(msg, dict) else getattr(msg, "timestamp", None)
        meta = (msg.get("metadata") if isinstance(msg, dict) else getattr(msg, "metadata", None)) or {}
        if role == "assistant":
            content = _clean_history_text(content)
        entry = {"role": role, "content": content, "timestamp": timestamp}
        _kind = msg.get("kind") if isinstance(msg, dict) else getattr(msg, "kind", None)
        if _kind:
            entry["kind"] = _kind   # proactive bubble tag drives the avatar animation on reload
        if role == "user" and meta and meta.get("images"):
            from urllib.parse import quote as _urlquote
            _img_entries = []
            for img in meta["images"]:
                if img.get("path"):
                    # File on disk -> serve via /api/file (per-user isolation enforced there).
                    _img_entries.append({
                        "url": f"/api/file?path={_urlquote(str(img['path']))}",
                        "name": img.get("name", "image"),
                    })
                elif img.get("data"):
                    # Legacy inline base64 -> data: URL.
                    _img_entries.append({
                        "url": f"data:{img.get('mime_type', 'image/jpeg')};base64,{img['data']}",
                        "name": img.get("name", "image"),
                    })
            entry["images"] = _img_entries
        if role == "tool":
            # Metadata dict first, then the top-level keys the backend actually stores.
            tool_name = (meta.get("toolName") if meta else None) or (
                msg.get("name") if isinstance(msg, dict) else getattr(msg, "name", None))
            tool_id = (meta.get("toolId") if meta else None) or (
                msg.get("tool_call_id") if isinstance(msg, dict) else getattr(msg, "tool_call_id", None))
            tool_status = (meta.get("toolStatus") if meta else None)
            if tool_name is not None:
                entry["toolName"] = tool_name
            if tool_id is not None:
                entry["toolId"] = tool_id
            # Content present means the tool answered.
            entry["toolStatus"] = tool_status or ("completed" if content else "running")
        out.append(entry)

    if sid and str(sid).startswith("thinking_"):
        try:
            from vaf.core.thinking_mode import pop_user_reply_for_session
            reply_data = pop_user_reply_for_session(sid)
            if reply_data:
                preview = (reply_data.get("reply") or "").strip()
                if preview:
                    out.append({
                        "role": "user",
                        "content": f"User replied: {preview}",
                        "timestamp": reply_data.get("at"),
                    })
        except Exception:
            pass
    return out


def _session_is_active(session_id) -> bool:
    """Is a turn running or queued in THIS chat - the queue is the authority.

    The former construction (registered agent's _session_id plus the global
    latest_state status) could not answer this: only worker 1 is registered,
    latest_state.status is one field for the whole process and the headless
    product path never writes it, so parallel workers made the answer wrong in
    both directions. TaskQueue.is_busy_for_session reads the same in-flight
    set the workers maintain, so it is right for every worker and every user.
    """
    if not session_id:
        return False
    try:
        from vaf.core.task_queue import TaskQueue
        return TaskQueue().is_busy_for_session(str(session_id))
    except Exception:
        return False


def _session_status_label(session_id) -> str:
    """Status text for an ACTIVE session's history_update.

    latest_state.status is process-global, so it may only be shown when the one
    registered agent is provably on this session; any other busy session gets
    the neutral label rather than another chat's status.
    """
    try:
        if (manager.agent_instance and
                getattr(manager.agent_instance, "_session_id", None) == str(session_id)):
            return manager.latest_state.get("status", "working")
    except Exception:
        pass
    return "working"


# Mount Memory System routes if enabled
if Config.get("memory_enabled", True):
    try:
        from vaf.memory.routes import memory_router
        app.include_router(memory_router, prefix="/api/memory", tags=["memory"])
        log("WebServer", "Memory system routes mounted at /api/memory")
    except ImportError as e:
        log("WebServer", f"Memory system not available: {e}")
    except Exception as e:
        log("WebServer", f"Failed to mount memory routes: {e}")

# Mount Discord Integration routes
try:
    from vaf.api.discord_routes import router as discord_router
    app.include_router(discord_router)
    log("WebServer", "Discord integration routes mounted at /api/discord")
except ImportError as e:
    log("WebServer", f"Discord integration not available: {e}")
except Exception as e:
    log("WebServer", f"Failed to mount Discord routes: {e}")

# Mount Telegram Integration routes
try:
    from vaf.api.telegram_routes import router as telegram_router
    app.include_router(telegram_router)
    log("WebServer", "Telegram integration routes mounted at /api/telegram")
except ImportError as e:
    log("WebServer", f"Telegram integration not available: {e}")
except Exception as e:
    log("WebServer", f"Failed to mount Telegram routes: {e}")

# Mount WhatsApp Integration routes
try:
    from vaf.api.whatsapp_routes import router as whatsapp_router
    app.include_router(whatsapp_router)
    log("WebServer", "WhatsApp integration routes mounted at /api/whatsapp")
except ImportError as e:
    log("WebServer", f"WhatsApp integration not available: {e}")
except Exception as e:
    log("WebServer", f"Failed to mount WhatsApp routes: {e}")

# Mount Contacts routes (central contact list with personal file)
try:
    from vaf.api.contact_routes import router as contact_router
    app.include_router(contact_router)
    log("WebServer", "Contacts routes mounted at /api/contacts")
except ImportError as e:
    log("WebServer", f"Contacts routes not available: {e}")
except Exception as e:
    log("WebServer", f"Failed to mount Contacts routes: {e}")

# Mount Auth routes (Local Network Authentication)
try:
    from vaf.api.auth_routes import router as auth_router
    app.include_router(auth_router)
    log("WebServer", "Auth routes mounted at /api/auth")
except ImportError as e:
    log("WebServer", f"Auth routes not available: {e}")
except Exception as e:
    log("WebServer", f"Failed to mount auth routes: {e}")

# Mount User Management routes (Admin only)
try:
    from vaf.api.user_routes import router as user_router
    app.include_router(user_router)
    log("WebServer", "User management routes mounted at /api/users")
except ImportError as e:
    log("WebServer", f"User routes not available: {e}")
except Exception as e:
    log("WebServer", f"Failed to mount user routes: {e}")

# Mount Network routes (topology, status)
try:
    from vaf.api.network_routes import router as network_router
    app.include_router(network_router)
    log("WebServer", "Network routes mounted at /api/network")
except ImportError as e:
    log("WebServer", f"Network routes not available: {e}")
except Exception as e:
    log("WebServer", f"Failed to mount network routes: {e}")

# Mount User Persona routes
try:
    from vaf.api.user_persona_routes import router as persona_router
    app.include_router(persona_router)
    log("WebServer", "User persona routes mounted at /api/user")
except Exception as e:
    log("WebServer", f"Failed to mount persona routes: {e}")

# Mount Config REST API (for onboarding connections step; config also available via WebSocket)
try:
    from vaf.api.config_routes import router as config_router
    app.include_router(config_router)
    log("WebServer", "Config REST API mounted at /api/config")
except Exception as e:
    log("WebServer", f"Failed to mount config routes: {e}")

# Mount TTS (Text-to-Speech) routes
try:
    from vaf.api.tts_routes import router as tts_router
    app.include_router(tts_router)
    log("WebServer", "TTS routes mounted at /api/tts")
except Exception as e:
    log("WebServer", f"Failed to mount TTS routes: {e}")

# Mount voice provider catalog routes (admin-only; ElevenLabs model/voice pickers)
try:
    from vaf.api.voice_routes import router as voice_router
    app.include_router(voice_router)
    log("WebServer", "Voice catalog routes mounted at /api/voice")
except Exception as e:
    log("WebServer", f"Failed to mount voice catalog routes: {e}")

# Mount Email connection routes (OAuth2 PKCE + accounts CRUD)
try:
    from vaf.api.email_routes import router as email_router
    app.include_router(email_router)
    log("WebServer", "Email integration routes mounted at /api/email")
except Exception as e:
    log("WebServer", f"Failed to mount Email routes: {e}")

# Mount the mail routes (EMAIL_CLIENT.md)
try:
    from vaf.api.mail_routes import router as mail_v2_router
    app.include_router(mail_v2_router)
    log("WebServer", "Mail engine v2 routes mounted at /api/mail")
except Exception as e:
    log("WebServer", f"Failed to mount Mail v2 routes: {e}")

# Mount Cloud Storage routes (OAuth2 PKCE + sync + accounts CRUD)
try:
    from vaf.api.cloud_routes import router as cloud_router
    app.include_router(cloud_router)
    log("WebServer", "Cloud storage routes mounted at /api/cloud")
except Exception as e:
    log("WebServer", f"Failed to mount Cloud routes: {e}")

# Mount GitHub connection routes (OAuth2 + accounts CRUD)
try:
    from vaf.api.github_routes import router as github_router
    app.include_router(github_router)
    log("WebServer", "GitHub integration routes mounted at /api/github")
except Exception as e:
    log("WebServer", f"Failed to mount GitHub routes: {e}")

# Mount Calendar routes (status + events; uses same OAuth as email)
try:
    from vaf.api.calendar_routes import router as calendar_router
    app.include_router(calendar_router)
    log("WebServer", "Calendar routes mounted at /api/calendar")
except Exception as e:
    log("WebServer", f"Failed to mount Calendar routes: {e}")

# Mount Logs routes (admin-only debug log reader)
try:
    from vaf.api.logs_routes import router as logs_router
    app.include_router(logs_router)
    log("WebServer", "Logs routes mounted at /api/logs")
except Exception as e:
    log("WebServer", f"Failed to mount Logs routes: {e}")

# Mount Supervisor/Watchdog routes (live sub-agent units + kill)
try:
    from vaf.api.supervisor_routes import router as supervisor_router
    app.include_router(supervisor_router)
    log("WebServer", "Supervisor routes mounted at /api/supervisor")
except Exception as e:
    log("WebServer", f"Failed to mount Supervisor routes: {e}")

# Mount Security overview routes (admin-only protection-module status)
try:
    from vaf.api.security_routes import router as security_router
    app.include_router(security_router)
    log("WebServer", "Security routes mounted at /api/security")
except Exception as e:
    log("WebServer", f"Failed to mount Security routes: {e}")

# Mount System routes (admin-only service status/repair + update)
try:
    from vaf.api.system_routes import router as system_router
    app.include_router(system_router)
    log("WebServer", "System routes mounted at /api/system")
except Exception as e:
    log("WebServer", f"Failed to mount System routes: {e}")

# Mount Background-agent status routes (admin-only, Logs > Overview panel)
try:
    from vaf.api.thinking_routes import router as thinking_router
    app.include_router(thinking_router)
    log("WebServer", "Thinking status routes mounted at /api/thinking")
except Exception as e:
    log("WebServer", f"Failed to mount Thinking status routes: {e}")

# Mount Agent Brain routes (working memory, plan, intent, team state)
try:
    from vaf.api.brain_routes import router as brain_router
    app.include_router(brain_router)
    log("WebServer", "Agent Brain routes mounted at /api/agent")
except Exception as e:
    log("WebServer", f"Failed to mount Agent Brain routes: {e}")

# Add authentication middleware if local network is enabled
if Config.get("local_network_enabled", False):
    try:
        from vaf.auth.middleware import AuthMiddleware, IPValidationMiddleware
        from vaf.auth.rate_limit import RateLimitMiddleware

        # Add rate limiting first (outermost)
        app.add_middleware(RateLimitMiddleware)
        # Add IP validation
        app.add_middleware(IPValidationMiddleware)
        # Add auth middleware (innermost - closest to route handlers)
        app.add_middleware(AuthMiddleware)

        log("WebServer", "Authentication middleware enabled for local network mode")
    except ImportError as e:
        log("WebServer", f"WARNING: Auth middleware import failed - network mode is INSECURE: {e}")
    except Exception as e:
        log("WebServer", f"WARNING: Failed to add auth middleware - network mode is INSECURE: {e}")

    # Auto-generate SSL certificates if TLS enabled but no certs configured
    if Config.get("local_network_tls_enabled", False):
        try:
            from vaf.network.ssl_utils import ensure_ssl_certificates
            _ssl_cert, _ssl_key = ensure_ssl_certificates()
            if _ssl_cert and _ssl_key:
                log("WebServer", f"SSL certificates ready: {_ssl_cert}")
            else:
                log("WebServer", "TLS enabled but no certificates available")
        except Exception as e:
            log("WebServer", f"SSL certificate setup failed: {e}")

# Security headers middleware (always active, stronger in network mode)
from starlette.middleware.base import BaseHTTPMiddleware as _BaseHM
from starlette.requests import Request as _Req
from starlette.responses import Response as _Resp

class _SecurityHeadersMiddleware(_BaseHM):
    """Add security headers to all responses."""

    # The one lane that is MEANT to be framed, and only by us. The interactive
    # browser's viewer is a page served by this server and displayed inside the
    # web UI's own browser window; a blanket DENY makes the browser fetch it
    # (200 in the access log) and then refuse to paint it, which reads as a
    # window that loads and dies. SAMEORIGIN rather than dropping the header:
    # the viewer travels the same front door as the UI framing it, so an
    # embedder on any other origin is still refused. Everything else keeps DENY.
    _FRAMEABLE_PREFIX = "/api/browser-vnc/"

    async def dispatch(self, request: _Req, call_next):
        response: _Resp = await call_next(request)
        response.headers["X-Content-Type-Options"] = "nosniff"
        _path = request.url.path or ""
        response.headers["X-Frame-Options"] = (
            "SAMEORIGIN" if _path.startswith(self._FRAMEABLE_PREFIX) else "DENY"
        )
        response.headers["X-XSS-Protection"] = "1; mode=block"
        response.headers["Referrer-Policy"] = "strict-origin-when-cross-origin"
        response.headers["Permissions-Policy"] = "camera=(), microphone=(), geolocation=()"
        # HSTS only when TLS is active
        if Config.get("local_network_tls_enabled", False):
            response.headers["Strict-Transport-Security"] = "max-age=31536000; includeSubDomains"
        return response

app.add_middleware(_SecurityHeadersMiddleware)

log("WebServer", "Module initialization complete")

def _mcp_servers_payload(agent) -> list:
    """
    Build the MCP-server list shown in the Settings UI: each configured server from the manifest,
    merged with its live connection status (connected / tool_count / error) from the agent's last
    discovery. Used by `get_mcp_servers` and returned directly in create/update/delete replies so
    the UI list updates without a separate refetch round-trip.
    """
    from vaf.core.mcp_registry import load_mcp_manifest
    servers = (load_mcp_manifest() or {}).get("servers", {}) or {}
    status = dict(getattr(agent, "_mcp_server_status", {}) or {})
    out = []
    for name, cfg in servers.items():
        if not isinstance(cfg, dict):
            continue
        st = status.get(name, {})
        out.append({
            "name": name,
            "command": cfg.get("command", ""),
            "transport": cfg.get("transport", "stdio"),
            "url": cfg.get("url", ""),
            "enabled": bool(cfg.get("enabled", True)),
            "permission_level": cfg.get("permission_level", "write"),
            "env": cfg.get("env") if isinstance(cfg.get("env"), dict) else {},
            "connected": bool(st.get("connected", False)),
            "tool_count": int(st.get("tool_count", 0) or 0),
            "error": st.get("error"),
        })
    return out


def _attach_learned_states(tools_list: list) -> list:
    """Annotate each tool entry with its Whare Wananga learned_state (best-effort).

    learned_state in {"unlearned","learning","learned","stale"} (derived from the
    tool_knowledge store). Surfaced in the Settings tool list. Never raises.
    """
    try:
        from vaf.whare_wananga import learned_states
        from vaf.whare_wananga.preconditions import tool_precondition
        names = [e.get("name", "") for e in tools_list]
        states = learned_states(names)
        for e in tools_list:
            nm = e.get("name", "")
            e["learned_state"] = states.get(nm, "unlearned")
            pc = tool_precondition(nm)
            e["requires_config"] = pc["requires_config"]
            e["configured"] = pc["configured"]
    except Exception:
        pass
    return tools_list


async def _broadcast_tools_update(manager) -> None:
    """
    Push a refreshed tools_list to every connected WebSocket client.
    Called after any custom-tool mutation (create / update / delete / permissions)
    so all open browser tabs see the change without a manual refresh.

    Each client gets a filtered list based on their own scope/role, which means
    we must send individual responses rather than one broadcast payload.
    """
    if not manager:
        return
    try:
        from vaf.core.config import get_local_admin_scope_id
        from vaf.core.custom_tools_registry import (
            get_all_custom_tool_names,
            get_visible_tool_names_for_user,
            get_tool_manifest_entry,
        )

        from vaf.core.tool_contract import tool_category

        agent          = manager.agent_instance
        local_admin    = get_local_admin_scope_id()
        all_custom     = set(get_all_custom_tool_names())

        # Iterate over all currently connected websockets
        for ws in list(manager.active_connections):
            try:
                _scope = manager.get_connection_user(ws)
                _role  = manager.get_connection_user_role(ws)
                _is_admin = (
                    _role == "admin"
                    or (_scope is not None and str(_scope) == str(local_admin))
                )
                _filter_scope = None if _is_admin else _scope
                visible_custom = set(get_visible_tool_names_for_user(_filter_scope))

                if agent and hasattr(agent, "tools"):
                    tools_list = []
                    for name, tool in agent.tools.items():
                        is_custom = name in all_custom
                        if is_custom and not _is_admin and name not in visible_custom:
                            continue
                        # Same account filter as the get_tools handler - a refreshed
                        # list must not quietly hand back what the allowlist removed.
                        try:
                            from vaf.core.tool_dispatch import account_allows_tool
                            if not account_allows_tool(name, _scope, _role):
                                continue
                        except Exception:
                            continue          # fail closed, like the funnel
                        entry = {
                            "name":        name,
                            "description": getattr(tool, "description", ""),
                            "category":    tool_category(name, tool),
                            "is_custom":   is_custom,
                            "can_manage":  _is_admin,
                        }
                        if is_custom:
                            meta = get_tool_manifest_entry(name)
                            if meta:
                                entry["shared_with"] = meta.get("shared_with", ["*"])
                                entry["created_by"]  = meta.get("created_by", "")
                                entry["updated_at"]  = meta.get("updated_at", "")
                        tools_list.append(entry)
                    _attach_learned_states(tools_list)
                    await ws.send_json({"type": "tools_list", "tools": tools_list})
            except Exception:
                pass  # Ignore disconnected / errored sockets
    except Exception as exc:
        import logging
        logging.getLogger(__name__).warning(
            "_broadcast_tools_update failed: %s", exc
        )


async def _broadcast_skills_update(manager) -> None:
    """
    Push a refreshed skills_list to every connected WebSocket client after any
    skill mutation (create / update / delete / permissions / upload). Like
    _broadcast_tools_update, each client gets a list filtered by their own
    scope/role, so we send individual responses rather than one payload.
    """
    if not manager:
        return
    try:
        from vaf.core.config import get_local_admin_scope_id
        from vaf.skills.templates import list_skills

        local_admin = get_local_admin_scope_id()
        for ws in list(manager.active_connections):
            try:
                _scope = manager.get_connection_user(ws)
                _role  = manager.get_connection_user_role(ws)
                _is_admin = (
                    _role == "admin"
                    or (_scope is not None and str(_scope) == str(local_admin))
                )
                _filter_scope = None if _is_admin else _scope
                # Admins manage all and also see invalid skills (so they can fix them).
                skills = list_skills(user_scope_id=_filter_scope, include_invalid=_is_admin)
                if _is_admin:
                    from vaf.core import skills_registry as _skr_b
                    for s in skills:
                        s["can_manage"] = True
                        s["source"] = _skr_b.get_skill_md_source(s["id"]) or ""
                else:
                    for s in skills:
                        s["can_manage"] = False
                await ws.send_json({"type": "skills_list", "skills": skills})
            except Exception:
                pass  # Ignore disconnected / errored sockets
    except Exception as exc:
        import logging
        logging.getLogger(__name__).warning(
            "_broadcast_skills_update failed: %s", exc
        )


def _scan_tool_modules() -> List[dict]:
    """
    Fallback tool list based on available Python modules.
    Tries to find Tool classes without full Agent initialization.
    """
    try:
        import pkgutil
        import importlib
        import inspect
        from vaf.tools.base import BaseTool
        from vaf.core.tool_contract import tool_category
        import vaf.tools
        
        tools = []
        package_path = os.path.dirname(vaf.tools.__file__)
        excluded_mods = {"__init__", "base"}
        
        for _, name, _ in pkgutil.iter_modules([package_path]):
            if name in excluded_mods or name.startswith("_"):
                continue
            try:
                module = importlib.import_module(f"vaf.tools.{name}")
                for _, obj in inspect.getmembers(module):
                    if inspect.isclass(obj) and issubclass(obj, BaseTool) and obj is not BaseTool:
                        tools.append({
                            "name": getattr(obj, "name", name),
                            "description": getattr(obj, "description", "Python tool"),
                            "category": tool_category(getattr(obj, "name", name), obj)
                        })
            except Exception:
                # Fallback to module name if import fails
                tools.append({
                    "name": name,
                    "description": "Python tool module (load failed)",
                    # Modul liess sich nicht laden - es gibt keine Klasse,
                    # an der eine Kategorie haengen koennte.
                    "category": "general"
                })
        return sorted(tools, key=lambda t: t["name"])
    except Exception:
        return []

def get_autosuggest():
    global autosuggest
    if autosuggest is None:
        log("WebServer", "Lazy loading SmartAutoSuggest...")
        autosuggest = SmartAutoSuggest()
    return autosuggest


def _get_trusted_sources_for_ui():
    """
    Build trusted sources list for Web UI from vaf/sources/*.json and config overrides.
    Order: custom categories first (right after creation box), then predefined from JSON.
    Returns: { "categories": [ { "id", "name", "description", "is_custom", "sources": [...] } ] }
    """
    import json as _json
    try:
        from vaf.core import sources as _sources_mod
        sources_dir = Path(_sources_mod.__file__).resolve().parent.parent / "sources"
    except Exception:
        sources_dir = Path(__file__).resolve().parents[1] / "sources"
    disabled = set(Config.get("trusted_sources_disabled") or [])
    custom = Config.get("trusted_sources_custom") or {}
    custom_only = []
    predefined = []
    seen_category_ids = set()
    # 1) Custom-only categories first (so they appear right after the creation box)
    for cid, cust_list in custom.items():
        custom_only.append({
            "id": cid,
            "name": cid if isinstance(cid, str) else cid.replace("_", " ").title(),
            "description": "Custom category",
            "is_custom": True,
            "sources": [
                {"name": c.get("name", ""), "url": c.get("url", ""), "domains": c.get("domains") or [], "trust_score": c.get("trust_score", 6), "is_custom": True}
                for c in (cust_list or [])
            ],
        })
        seen_category_ids.add(cid)
    # 2) Predefined categories from JSON
    if sources_dir.exists():
        for jpath in sources_dir.glob("*.json"):
            try:
                with open(jpath, "r", encoding="utf-8") as f:
                    data = _json.load(f)
                for cid, cdata in (data.get("categories") or {}).items():
                    if cid in seen_category_ids:
                        continue
                    seen_category_ids.add(cid)
                    sources_list = []
                    for s in (cdata.get("sources") or []):
                        doms = s.get("domains") or []
                        if any(d.lower() in disabled for d in doms):
                            continue
                        sources_list.append({
                            "name": s.get("name", ""),
                            "url": s.get("url", ""),
                            "domains": doms,
                            "trust_score": s.get("trust_score", 5),
                            "is_custom": False,
                        })
                    for cust in (custom.get(cid) or []):
                        doms = cust.get("domains") or []
                        sources_list.append({
                            "name": cust.get("name", ""),
                            "url": cust.get("url", ""),
                            "domains": doms,
                            "trust_score": cust.get("trust_score", 6),
                            "is_custom": True,
                        })
                    predefined.append({
                        "id": cid,
                        "name": cdata.get("name", cid),
                        "description": cdata.get("description", ""),
                        "is_custom": False,
                        "sources": sources_list,
                    })
            except Exception:
                continue
    categories_out = custom_only + predefined
    return {"categories": categories_out}

_auth_db_init_task = None  # module-level ref so the background retry task is not garbage collected
_auth_db_init_gate = threading.Lock()  # in TLS mode the same app runs on 8001 AND 8005 -> two lifespans
_legacy_claim_done = False  # the unscoped-session claim runs once, not per lifespan


@app.on_event("startup")
async def start_skills_rescan():
    # Periodic skill re-scan (post-install tamper detection for the security
    # dashboard). start_periodic_rescan is idempotent - required because this
    # startup hook runs twice in TLS mode (8001 + 8005 lifespans).
    try:
        from vaf.skills.rescan import start_periodic_rescan
        if start_periodic_rescan():
            log("WebServer", "Skills periodic re-scan armed")
    except Exception as e:
        log("WebServer", f"Skills re-scan not armed: {e}")


@app.on_event("startup")
async def startup_event():
    # Initialize auth database tables (creates if not exist). The Docker stack starts in a
    # thread parallel to this server, so PostgreSQL may still be booting here - a failed
    # first attempt must self-heal in the background instead of leaving the auth tables
    # uncreated until the next restart (a fresh install then shows login instead of setup).
    try:
        from vaf.auth.database import init_auth_db
        await init_auth_db()
        log("WebServer", "Auth database tables initialized")
    except Exception as e:
        log("WebServer", f"Auth database not ready yet ({e}); retrying in background")

        async def _bg_auth_init():
            from vaf.auth.database import init_auth_db_with_retry
            await init_auth_db_with_retry()
            log("WebServer", "Auth database tables initialized (background retry)")

        # This startup event runs once per uvicorn server sharing the app (8001 + the
        # internal 8005 channel in TLS mode, on different threads/loops) - spawn only ONE
        # retry lane; a task stranded on a stopped loop (backend restart) counts as stale.
        global _auth_db_init_task
        with _auth_db_init_gate:
            t = _auth_db_init_task
            if t is None or t.done() or not t.get_loop().is_running():
                _auth_db_init_task = asyncio.create_task(_bg_auth_init())
            else:
                log("WebServer", "Auth DB background retry already running (other server lifespan) - not spawning a duplicate")
    
    # Set the event loop for thread-safe broadcasting
    loop = asyncio.get_running_loop()
    manager.set_server_loop(loop)
    log("WebServer", "VAF Web Interface: Event loop registered")

    # One-time repair: pre-scoping sessions have no owner and the list's
    # legacy rule shows them to EVERY user; claim them for the machine owner
    # (multi-user arrived WITH scoping, so they cannot be anyone else's).
    # Gated like the auth init: this startup event runs once per uvicorn
    # server sharing the app (8001 + 8005 in TLS mode).
    global _legacy_claim_done
    with _auth_db_init_gate:
        run_claim = not _legacy_claim_done
        _legacy_claim_done = True
    if run_claim:
        try:
            from vaf.core.config import get_local_admin_scope_id
            claimed = await asyncio.to_thread(
                session_mgr.claim_unscoped, str(get_local_admin_scope_id()))
            if claimed:
                log("WebServer", f"Claimed {claimed} legacy session(s) for the machine owner")
        except Exception as e:
            log("WebServer", f"Legacy session claim skipped: {e}")

    if run_claim:
        # At-rest migration + retention, under the SAME once-per-app gate. The CLI
        # start path calls this too: a repair wired into one lane only is how a
        # cleanup once ran for the terminal and never for the web app.
        try:
            from vaf.core.at_rest_migration import run_once
            report = await asyncio.to_thread(run_once)
            if report and not report.get("skipped"):
                log("WebServer", f"At-rest migration: {report}")
        except Exception as e:
            log("WebServer", f"At-rest migration skipped: {e}")
        # Embedding-model reconcile (returns immediately; worker + waiting run
        # on its own thread). Wired into every start lane like run_once above.
        try:
            from vaf.memory.reembed import ensure_embedding_model_current
            ensure_embedding_model_current()
        except Exception as e:
            log("WebServer", f"Embedding-model reconcile skipped: {e}")
        # Progress lane for machine-level maintenance jobs (re-embed today):
        # polls the shared state, frames go out only while a job is active.
        try:
            from vaf.core.web_interface import start_maintenance_broadcast
            start_maintenance_broadcast()
        except Exception as e:
            log("WebServer", f"Maintenance broadcast not started: {e}")

    # Whare Wananga EAGER: opt-in background scanner that proactively trains safe, configured,
    # not-yet-learned tools (off by default; tolerates a not-yet-built agent). Guarded.
    try:
        from vaf.whare_wananga import eager
        eager.start(lambda: getattr(manager, "agent_instance", None))
        log("WebServer", "Whare Wananga eager scanner started (opt-in)")
    except Exception as e:
        log("WebServer", f"Whare Wananga eager scanner not started: {e}")
    
    # Register TTS Callbacks for UI sync
    from vaf.core.speech import SpeechManager
    sm = SpeechManager.get_instance()
    
    def on_tts_loading():
        if manager._server_loop:
            asyncio.run_coroutine_threadsafe(
                manager.broadcast({"type": "tts_state", "status": "loading"}),
                manager._server_loop
            )

    def on_tts_start(text):
        if manager._server_loop:
            asyncio.run_coroutine_threadsafe(
                manager.broadcast({"type": "tts_state", "status": "playing", "text": text}),
                manager._server_loop
            )
            
    def on_tts_end():
        if manager._server_loop:
            asyncio.run_coroutine_threadsafe(
                manager.broadcast({"type": "tts_state", "status": "stopped"}),
                manager._server_loop
            )
            
    sm.on_speech_loading = on_tts_loading
    sm.on_speech_start = on_tts_start
    sm.on_speech_end = on_tts_end
    
    # Setup firewall rules if local network and firewall are enabled.
    # IMPORTANT: run OFF the startup critical path in a daemon thread. Firewall setup shells out to sudo;
    # a blocking sudo password prompt on the controlling TTY would otherwise stall this startup handler
    # and delay port 8001 readiness (the "Backend not ready after 30s" symptom). firewall.py now passes
    # 'sudo -n' so it fails fast instead of prompting, and this thread guarantees startup is never
    # blocked regardless. (Correct LAN port + firewalld/pkexec elevation are handled separately.)
    if Config.get("local_network_enabled", False) and Config.get("local_network_firewall_enabled", True):
        def _setup_firewall_bg():
            try:
                import time as _time
                from vaf.network.firewall import setup_firewall, register_cleanup_on_exit
                tls_on = bool(Config.get("local_network_tls_enabled", False))
                if tls_on:
                    # Open the port the integrated proxy ACTUALLY bound (e.g. 8443 after the 443→8443
                    # fallback), NOT the configured 443. Wait briefly for the proxy thread to report it.
                    from vaf.network import runtime_status
                    access_port = None
                    for _ in range(20):  # up to ~10s for the proxy to bind + report
                        st = runtime_status.get_proxy_status()
                        if st.get("bound") and st.get("effective_https_port"):
                            access_port = int(st["effective_https_port"])
                            break
                        _time.sleep(0.5)
                    if access_port is None:
                        configured = int(Config.get("local_network_https_port", 443) or 443)
                        # Proxy hasn't reported yet → assume its standard 443→8443 fallback.
                        access_port = 8443 if configured == 443 else configured
                    port = access_port
                    port_frontend = int(Config.get("local_network_port", 8001) or 8001)
                else:
                    port = int(Config.get("local_network_port", 8001) or 8001)
                    port_frontend = int(Config.get("local_network_port_frontend", 3000) or 3000)

                success = setup_firewall(port, port_frontend)
                if success == "present":
                    # No elevation ran: the rule was already active, or the twin
                    # lifespans of TLS mode already attempted this start. The old
                    # single "created" message hid this difference, which made a
                    # nightly password dialog unattributable from the log.
                    log("WebServer", f"Firewall rule already in place for port {port} - no password dialog needed")
                elif success:
                    register_cleanup_on_exit()
                    log("WebServer", f"Firewall rules created for ports {port}, {port_frontend}")
                else:
                    log("WebServer", f"Firewall setup skipped for ports {port}, {port_frontend} - needs elevated privileges (no passwordless sudo). Open the port manually or use the in-app firewall step.")
            except Exception as e:
                log("WebServer", f"Firewall setup error: {e}")
        import threading as _threading
        _threading.Thread(target=_setup_firewall_bg, daemon=True, name="vaf-firewall-setup").start()

    # Start garbage collector (logs, temp, cache, old thinking sessions) when running without tray
    try:
        from vaf.core.garbage_collector import GarbageCollector
        GarbageCollector.get_instance().start()
        log("WebServer", "Garbage collector started")
    except Exception as e:
        log("WebServer", f"Garbage collector start warning: {e}")
    
    # Mail sync (EMAIL_CLIENT.md): the engine supervisor is the ONLY mail sync
    # lane. The 30-minute legacy auto-sync loop that used to run beside it was
    # removed with the legacy stack.
    try:
        from vaf.mail.supervisor import MailSyncSupervisor
        asyncio.create_task(MailSyncSupervisor().run())
        log("WebServer", "Mail v2 sync supervisor task started (flag-gated)")
    except Exception as e:
        log("WebServer", f"Mail v2 supervisor start warning: {e}")

    # Cloud storage background sync
    if Config.get("cloud_sync_enabled", False):
        try:
            from vaf.cloud.sync_worker import cloud_sync_loop
            asyncio.create_task(cloud_sync_loop())
            log("WebServer", "Cloud storage background sync started")
        except Exception as e:
            log("WebServer", f"Cloud sync loop start error: {e}")

    # Thinking mode: background reflection when user idle
    if Config.get("thinking_enabled", True):
        try:
            from vaf.core.thinking_mode import start_thinking_mode_background
            start_thinking_mode_background()
            log("WebServer", "Thinking mode background loop started")
        except Exception as e:
            log("WebServer", f"Thinking mode start error: {e}")
    log("WebServer", "Email auto-sync background task started (every 30 min)")

    # Start the process-wide automation scheduler for existing timed automations.
    try:
        from vaf.core.automation import ensure_scheduler_started

        _scheduler_mgr, started_now = ensure_scheduler_started(origin="web_server_startup")
        if started_now:
            log("WebServer", "Automation scheduler started")
        else:
            log("WebServer", "Automation scheduler already running")
    except Exception as e:
        log("WebServer", f"Automation scheduler start error: {e}")

    # In Docker mode, the tray app doesn't run, so we must start the headless runner here.
    if Config.is_docker_mode():
        log("WebServer", "Docker mode detected - starting headless agent runner...")
        try:
            _ensure_headless_agent_runner(origin="docker_startup")
        except Exception as e:
            log("WebServer", f"Failed to start headless agent: {e}")
            import traceback
            log("WebServer", traceback.format_exc())

    # Start Auto-Capture Queue Worker - DISABLED: causes 13GB+ memory spikes
    # The memory leak occurs during pipeline.ingest() ~13 seconds after embedding
    # TODO: Investigate asyncpg/SQLAlchemy memory leak in auto_capture_memory()
    # if Config.get("memory_enabled", True) and Config.get("memory_auto_capture", True):
    #     asyncio.create_task(_auto_capture_worker())
    #     log("WebServer", "Auto-capture queue worker started")
    log("WebServer", "Auto-capture worker DISABLED (memory leak investigation)")

    # Auto-start Telegram bridge when configured and enabled (so Web UI shows "Connected" after restart)
    try:
        telegram_config = Config.get("telegram_config") or {}
        if isinstance(telegram_config, dict) and telegram_config.get("verified") and telegram_config.get("bot_token") and telegram_config.get("enabled"):
            from vaf.api.telegram_bridge import start_bridge, is_bridge_running
            if not is_bridge_running() and start_bridge():
                log("WebServer", "Telegram bridge auto-started (configured and enabled)")
            elif is_bridge_running():
                log("WebServer", "Telegram bridge already running")
    except Exception as e:
        log("WebServer", f"Telegram bridge auto-start skipped or failed: {e}")

    # Auto-start Discord bridge when configured and enabled
    try:
        discord_config = Config.get("discord_config") or {}
        if isinstance(discord_config, dict) and discord_config.get("verified") and discord_config.get("bot_token") and discord_config.get("admin_user_id") and discord_config.get("enabled"):
            from vaf.api.discord_bridge import start_bridge, is_bridge_running
            if not is_bridge_running() and start_bridge():
                log("WebServer", "Discord bridge auto-started (configured and enabled)")
            elif is_bridge_running():
                log("WebServer", "Discord bridge already running")
    except Exception as e:
        log("WebServer", f"Discord bridge auto-start skipped or failed: {e}")

    # Auto-start WhatsApp bridge when configured and enabled
    try:
        whatsapp_config = Config.get("whatsapp_config") or {}
        if isinstance(whatsapp_config, dict) and whatsapp_config.get("enabled"):
            whitelist = whatsapp_config.get("whitelist") or []
            if any(isinstance(e, dict) and e.get("phone_number") for e in whitelist):
                from vaf.api.whatsapp_bridge import start_bridge, is_bridge_running
                if not is_bridge_running() and start_bridge():
                    log("WebServer", "WhatsApp bridge auto-started (configured and enabled)")
                elif is_bridge_running():
                    log("WebServer", "WhatsApp bridge already running")
                asyncio.create_task(_whatsapp_reconnect_worker())
                log("WebServer", "WhatsApp auto-reconnect worker started")
    except Exception as e:
        log("WebServer", f"WhatsApp bridge auto-start skipped or failed: {e}")


async def _whatsapp_reconnect_worker():
    """
    Periodically check WhatsApp connection; if bridge is running but disconnected,
    restart the bridge so it reconnects with stored credentials (no user action needed).
    """
    from vaf.api.whatsapp_bridge import is_bridge_running, get_connection_status, restart_bridge
    loop = asyncio.get_running_loop()
    last_restart_at = 0.0
    check_interval = 120.0   # check every 2 minutes
    disconnected_since = None  # time when we first saw disconnected
    disconnect_grace = 90.0   # restart only if disconnected for this long (seconds)
    cooldown = 180.0          # after a restart, wait this long before next check

    while True:
        await asyncio.sleep(check_interval)
        try:
            whatsapp_config = Config.get("whatsapp_config") or {}
            if not isinstance(whatsapp_config, dict) or not whatsapp_config.get("enabled"):
                disconnected_since = None
                continue
            whitelist = whatsapp_config.get("whitelist") or []
            if not any(isinstance(e, dict) and e.get("phone_number") for e in whitelist):
                disconnected_since = None
                continue
            if not is_bridge_running():
                disconnected_since = None
                continue
            # Cooldown after a restart
            now = loop.time()
            if now - last_restart_at < cooldown:
                continue
            # Sync call in executor to avoid blocking
            connected = await loop.run_in_executor(
                None, lambda: get_connection_status("admin", wait_timeout=2.0)
            )
            if connected:
                disconnected_since = None
                continue
            if disconnected_since is None:
                disconnected_since = now
            if now - disconnected_since < disconnect_grace:
                continue
            # Restart bridge to reconnect
            log("WebServer", "WhatsApp disconnected; auto-restarting bridge to reconnect")
            await loop.run_in_executor(None, restart_bridge)
            last_restart_at = loop.time()
            disconnected_since = None
        except Exception as e:
            log("WebServer", f"WhatsApp reconnect worker error: {e}")
            disconnected_since = None


async def _auto_capture_worker():
    """
    Background worker that processes auto-capture queue in the main event loop.

    This is MEMORY-LEAK SAFE because:
    - Runs in main event loop (no daemon threads with asyncio.run())
    - ONNX model and DB connections are reused from main thread
    - Processes max 2 tasks per cycle to avoid blocking
    """
    from vaf.memory.rag import process_auto_capture_queue, get_auto_capture_queue_size

    while True:
        try:
            # Process pending captures (max 2 per cycle)
            processed = await process_auto_capture_queue(max_tasks=2)
            if processed > 0:
                log("WebServer", f"Auto-capture: processed {processed} tasks")
        except Exception as e:
            log("WebServer", f"Auto-capture worker error: {e}")

        # Check queue every 5 seconds (balance between responsiveness and CPU)
        await asyncio.sleep(5)

def _detect_language_simple(text: str) -> str:
    """Detect language for TTS: use langid when available, else simple heuristic (de/en)."""
    if not (text and text.strip()):
        return "en"
    # Prefer langid for reliable detection (e.g. German text -> de)
    try:
        from vaf.vendor import langid
        import re
        # Strip thinking/code so we classify on actual spoken content
        clean = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL).strip()
        if not clean:
            clean = text.strip()
        # Use first ~2k chars to avoid slow classify on huge text
        sample = clean[:2000] if len(clean) > 2000 else clean
        if sample:
            code, _ = langid.classify(sample)
            if code and len(code) >= 2:
                return code[:2].lower()
    except ImportError:
        pass
    except Exception:
        pass
    # Fallback: simple heuristic (de/en)
    t = text.lower()
    de_words = [" das ", " und ", " der ", " die ", " ist ", " nicht ", " ich ", " sie ", " es ", " wie ", " was ", " eine ", " ein ", " mit ", " von ", " für ", " auf ", " sind ", " kann ", " auch ", " dann ", " haben ", " wird "]
    if any(w in t for w in de_words):
        return "de"
    if any(ch in t for ch in ["ä", "ö", "ü", "ß"]):
        return "de"
    # "das " / "der " / "die " at start
    if t.startswith(("das ", "der ", "die ", "den ", "dem ", "ein ", "eine ", "und ", "ist ")):
        return "de"
    return "en"


def _tts_lang_for(text: str, base_lang: str, sm) -> str:
    """Pick the TTS language for an LLM-GENERATED spoken reply. If the reply text's
    detected language differs from the call/base language AND the speech lane can
    actually speak it (`SpeechManager.call_lane_speaks`: a Docker voice is installed,
    or a cloud TTS provider is multilingual), follow the reply's language - so a
    Turkish answer is spoken by a Turkish voice, not the German one; the Docker lane
    switches voice by the `language` it is sent, the API lane gets the language. Falls
    back to base_lang on an unclear detection or an unavailable voice (never a mid-call
    download). For model-written text only; vocab lines are already in base_lang."""
    try:
        if not (text or "").strip():
            return base_lang  # nothing to speak -> never switch away from the call language
        rl = _detect_language_simple(text)
        base = (base_lang or "")[:2].lower()
        if rl and rl != base and sm.call_lane_speaks(rl):
            return rl
    except Exception:
        pass
    return base_lang


def _build_artifact_payload(session, session_id: str = None):
    if not session:
        return None
    try:
        state = session.get_provider_state("artifact")
    except Exception:
        state = None
    if not state:
        return None
    payload = {
        "type": "artifact_update",
        "file": state.get("file", ""),
        "code": state.get("code", ""),
        "updatedAt": state.get("updatedAt"),
        "source": state.get("source", "backend")
    }
    if session_id:
        payload["sessionId"] = session_id
    return payload

@app.get("/")
async def root():
    return {"status": "VAF Backend Online", "version": __version__}

from pydantic import BaseModel, ConfigDict

class WorkflowUpdate(BaseModel):
    """Everything a SUBPROCESS sends to /api/workflow/update, on its way to the browser.

    This model is a bottleneck the whole subprocess lane passes through, and an undeclared
    field is dropped here without a trace. That is the CLAUDE.md Rule-2 field-forwarding trap
    (it has already cost `diffs` and `activity` in page.tsx), and it bit again on 2026-07-20:
    `success` was never declared, so every `workflow_done` from the separate workflow runner
    arrived at the browser without it. The panel reads `data.success ? 'completed' : 'failed'`,
    so a run whose steps were all green and whose document was written showed FAILED.

    Two defences, deliberately both:
      - every field a producer sends is DECLARED below, so the payload is documented and typed;
      - extra="allow" keeps an undeclared field from ever being silently dropped again.
    tests/test_workflow_update_payload.py cross-checks the producers against this model.
    """
    model_config = ConfigDict(extra="allow")

    type: str
    sessionId: Optional[str] = None
    workflowId: Optional[str] = None
    name: Optional[str] = None
    steps: Optional[List] = None
    stepId: Optional[str] = None
    status: Optional[str] = None
    progress: Optional[int] = None
    result: Optional[str] = None
    # workflow_done
    success: Optional[bool] = None
    error: Optional[str] = None
    # workflow_output_stream
    line: Optional[str] = None
    # document_ready payload (from notify_document_created)
    filePath: Optional[str] = None
    title: Optional[str] = None
    openMode: Optional[str] = None  # "viewer" routes research reports to the Document Viewer; else editor

class Heartbeat(BaseModel):
    client_id: str
    timestamp: float = 0.0

@app.post("/api/heartbeat")
async def receive_heartbeat(hb: Heartbeat):
    """Receive heartbeat from CLI clients to keep server active."""
    tray_context.register_activity()
    return {"status": "ok", "active": True}

@app.get("/api/heartbeat")
async def healthcheck():
    """Health check endpoint for Docker/monitoring."""
    return {"status": "ok", "healthy": True}

async def _open_report_in_viewer(session_id: str, file_path: str, title: Optional[str]) -> None:
    """Open a freshly-created report (e.g. from research_agent) in the Document Viewer (sidebar) rather
    than the editor: render + RAG-index it like an attachment and broadcast sidebar_documents_set so the
    read-only viewer shows it. Mirrors the set_sidebar_documents WS handler. Existing sidebar attachments
    are preserved (merge + de-dup by name); falls back silently if anything fails."""
    try:
        from pathlib import Path as _Path
        from vaf.core.session import Session as _Session
        import base64 as _b64, mimetypes as _mt
        p = _Path(file_path)
        if not p.exists():
            return
        name = title or p.name
        data_b64 = _b64.b64encode(p.read_bytes()).decode("ascii")
        mime = _mt.guess_type(p.name)[0] or "application/octet-stream"
        contents = await process_files_to_sidebar_list([{"name": name, "data": data_b64, "mimeType": mime}])
        if not contents:
            return
        try:
            loaded = session_mgr.load(session_id)
        except Exception:
            loaded = _Session(id=session_id, name=f"Session {session_id}")
        if not getattr(loaded, "runtime_state", None):
            loaded.runtime_state = {}
        user_scope_id = (getattr(loaded, "metadata", None) or {}).get("user_scope_id") or None
        existing = [d for d in (loaded.runtime_state.get("sidebar_documents") or []) if d.get("name") != name]
        # Persist a SLIM copy (drop the big base64 'data') to keep the session JSON small, but BROADCAST the
        # full contents WITH 'data' (the Gotenberg-rendered PDF). Without 'data' the viewer falls back to the
        # raw extracted text; with it, the docx renders natively as A4 pages -- exactly like a manual attach.
        slim_new = [{k: v for k, v in c.items() if k != "data"} for c in contents]
        loaded.runtime_state["sidebar_documents"] = existing + slim_new
        session_mgr.save(loaded, sync_state=False)
        if bool(Config.get("attachment_rag_enabled", False)):
            try:
                await _notify_attachment_index(manager, session_id, "attachment_indexing", count=len(contents))
                _spawn_attachment_index(manager, session_id, user_scope_id, contents)
            except Exception as e:
                log("WebServer", f"report viewer RAG index failed: {e}")
        await manager.broadcast_to_session(session_id, {
            "type": "sidebar_documents_set",
            "contents": existing + contents,
            "sessionId": session_id,
        })
    except Exception as e:
        append_domain_log("webui", f"[ERROR] open report in viewer failed: {e}")


@app.post("/api/workflow/update")
async def receive_workflow_update(update: WorkflowUpdate):
    """Receive workflow updates from external processes (like separate terminals)."""
    data = update.dict(exclude_none=True)
    try:
        if update.type == "document_ready" and update.openMode == "viewer" and update.sessionId and update.filePath:
            # research_agent reports open read-only in the Document Viewer (sidebar + RAG), not the editor.
            await _open_report_in_viewer(update.sessionId, update.filePath, update.title)
        elif update.sessionId:
            await manager.broadcast_to_session(update.sessionId, data)
        else:
            await manager.broadcast(data)
    except Exception as e:
        append_domain_log("webui", f"[ERROR] broadcast failed in /api/workflow/update: {e}")
    # When a file is created, store its project directory so the agent can edit it later.
    # Shared setter (vaf/core/session.py record_created_file): this endpoint only sees
    # SUBPROCESS notifications - in-process writes anchor via notify_file_created directly.
    if data.get("type") == "file_created" and data.get("filePath") and data.get("sessionId"):
        try:
            from vaf.core.session import record_created_file
            record_created_file(data["sessionId"], data["filePath"])
        except Exception:
            pass
    return {"status": "ok"}


class SubAgentStreamUpdate(BaseModel):
    """Model for subagent output stream updates from separate processes."""
    type: str
    sessionId: Optional[str] = None
    taskId: Optional[str] = None
    agentType: Optional[str] = None
    agentName: Optional[str] = None
    line: Optional[str] = None
    status: Optional[str] = None
    presence: Optional[str] = None
    provider: Optional[str] = None
    model: Optional[str] = None
    file: Optional[str] = None
    code: Optional[str] = None
    steps: Optional[List] = None
    # Workflow-specific fields
    workflowId: Optional[str] = None
    stepId: Optional[str] = None
    progress: Optional[int] = None
    name: Optional[str] = None

    model_config = {"extra": "allow"}


@app.get("/api/whare_wananga/tool_knowledge/{name}")
async def whare_wananga_tool_knowledge(name: str):
    """Return the stored tool_knowledge record + learned state for a tool (or null)."""
    try:
        from vaf.whare_wananga import load as _ww_load, learned_state as _ww_state
        return {"ok": True, "tool": name, "state": _ww_state(name), "record": _ww_load(name)}
    except Exception as e:
        return {"ok": False, "tool": name, "state": "unlearned", "record": None, "error": str(e)}


@app.post("/api/whare_wananga/train/{name}")
async def whare_wananga_train(name: str):
    """Start a Whare Wananga predict-then-verify training pass for a tool (background job)."""
    try:
        from vaf.core.log_helper import append_domain_log
        append_domain_log("backend", f"[WHARE-WANANGA] train requested: {name}")
    except Exception:
        pass
    agent = manager.agent_instance if manager else None
    if agent is None or not hasattr(agent, "tools"):
        return {"ok": False, "tool": name, "state": "error", "message": "Agent not available"}
    if name not in getattr(agent, "tools", {}):
        return {"ok": False, "tool": name, "state": "error", "message": "Unknown tool"}
    # Precondition: never train a tool whose connection is not configured.
    try:
        from vaf.whare_wananga.preconditions import tool_precondition
        pc = tool_precondition(name)
        if pc.get("requires_config") and not pc.get("configured"):
            return {"ok": False, "tool": name, "state": "skipped", "message": "Connection not configured"}
    except Exception:
        pass
    try:
        from vaf.whare_wananga import jobs
        st = jobs.start_training(agent, name)
        print(f"[WHARE-WANANGA] training started: {name} (validate/refine loop)")
        return {"ok": True, "tool": name, **st}
    except Exception as e:
        return {"ok": False, "tool": name, "state": "error", "message": str(e)}


@app.get("/api/whare_wananga/training_status/{name}")
async def whare_wananga_training_status(name: str):
    """Live status of a Whare Wananga training job (polled by the dashboard)."""
    try:
        from vaf.whare_wananga import jobs
        return {"ok": True, "tool": name, "status": jobs.get_status(name)}
    except Exception as e:
        return {"ok": False, "tool": name, "status": None, "error": str(e)}


@app.get("/api/whare_wananga/active_runs")
async def whare_wananga_active_runs():
    """Every Whare Wananga training run in flight, for the tools modal header.

    The sibling route above answers for ONE named tool, which is the wrong shape
    for a header that does not know the name: it would have to poll once per
    installed tool. This one consumes the framework reader instead of filtering
    the job table here, so the two can never disagree.

    Process-local, like the job table it reads: a `vaf ww train` run started in a
    separate shell is invisible to this route by construction.
    """
    try:
        from vaf.whare_wananga import active_runs
        return {"ok": True, "runs": active_runs()}
    except Exception as e:
        return {"ok": False, "runs": [], "error": str(e)}


@app.post("/api/subagent/stream")
async def receive_subagent_stream(update: SubAgentStreamUpdate):
    """
    Receive subagent output stream updates from external processes.
    This endpoint bridges subprocess output to WebSocket clients.
    """
    data = update.dict(exclude_none=True)
    # <ipc-notification>: a sub-agent subprocess signalling that its result is READY.
    # Wake the in-process headless runner so it consumes the result immediately instead
    # of waiting for its ~1s idle poll. Internal control signal — not broadcast to the UI.
    if data.get("type") == "ipc_notification":
        try:
            from vaf.core.subagent_ipc import notify_result_ready
            notify_result_ready()
        except Exception as e:
            append_domain_log("webui", f"[ERROR] ipc_notification handling failed: {e}")
        return {"status": "ok"}
    # We're in an async handler, so we can directly await without checking the loop
    try:
        # ROOM FIRST: a room turn may run with NO session at all (the runner's
        # room frame binds no chat), so the room stamp the producers put on the
        # event is the primary anchor - same trust as sessionId, it is the same
        # loopback-only internal lane. The stamp routes per USER only when the
        # room's tenant can be PROVEN from the manifest; an unresolvable room
        # loses the stamp and falls through, so a forged or dead room id never
        # widens delivery.
        room_hint = str(data.get("roomId") or "").strip()
        room_scope = manager.room_audience(room_hint) if room_hint else None
        if room_hint and room_scope is None:
            data.pop("roomId", None)
            room_hint = ""
        if room_hint:
            await manager.broadcast_to_user(room_scope, data)
        elif update.sessionId:
            # A room-ordered task that DID run under a session still routes to
            # the room watcher: the IPC task record carries the ordering room.
            # No room task: the session lane stays exactly what it was.
            route = manager.room_route_for_session(update.sessionId)
            if route:
                data["roomId"] = route[0]
                await manager.broadcast_to_user(route[1], data)
            else:
                await manager.broadcast_to_session(update.sessionId, data)
        else:
            await manager.broadcast(data)
    except Exception as e:
        append_domain_log("webui", f"[ERROR] broadcast failed in /api/subagent/stream: {e}")
    return {"status": "ok"}


# ── interactive browser stream (KasmVNC proxy) ──────────────────────────────
#
# The browser window's interactive mode loads the KasmVNC client from the
# vaf-browser container THROUGH these two routes, never directly: the container
# port is loopback-only, and a LAN user reaches it exclusively here. The lease
# TICKET in the path is the credential (auth-middleware-exempt, the room-seat
# pattern): the KasmVNC client builds every asset and websocket URL relative to
# its own directory, so the ticket rides along on all of them by construction.
# Validated per request against the CURRENT lease - a superseded or stopped
# lease invalidates its ticket immediately.

@app.get("/api/browser-vnc/t/{ticket}/{path:path}")
async def browser_vnc_asset(ticket: str, path: str):
    """Static KasmVNC client files, fetched from the container and passed on."""
    from fastapi.responses import JSONResponse, Response
    # The ticket names the browser: with the per-user pool on there is more
    # than one instance, and the manager that issued the ticket knows which
    # vnc endpoint its stream lives on.
    from vaf.core.browser_interactive import get_manager_by_ticket
    mgr = get_manager_by_ticket(ticket)
    if mgr is None:
        return JSONResponse(status_code=403, content={"detail": "Invalid or expired stream ticket"})
    if ".." in path:
        return JSONResponse(status_code=400, content={"detail": "Invalid path"})
    import httpx
    url = mgr.vnc_base().rstrip("/") + "/" + (path or "index.html")
    try:
        # follow_redirects: only content-type and cache-control are passed on below,
        # so a 3xx would reach the browser WITHOUT its Location header - a redirect
        # to nowhere, which presents as an empty response. Resolving it here keeps
        # the answer a real body.
        async with httpx.AsyncClient(timeout=10, follow_redirects=True) as client:
            upstream = await client.get(url)
    except Exception as e:
        append_domain_log("webui", f"[browser-vnc] asset fetch failed: {e}")
        return JSONResponse(status_code=502, content={"detail": "Browser stream backend unreachable"})
    headers = {}
    for h in ("content-type", "cache-control"):
        if h in upstream.headers:
            headers[h] = upstream.headers[h]
    return Response(content=upstream.content, status_code=upstream.status_code, headers=headers)


@app.websocket("/api/browser-vnc/t/{ticket}/websockify")
async def browser_vnc_stream(websocket: WebSocket, ticket: str):
    """Byte pipe between the viewer and the container's KasmVNC socket.

    Every open pipe counts as viewer presence for the lease janitor; the last
    disconnect starts the grace timer that eventually releases the lease.
    """
    from vaf.core.browser_interactive import get_manager_by_ticket, AgentStream
    mgr = get_manager_by_ticket(ticket)
    if mgr is None:
        # Accept-then-close: a close before accept surfaces as a handshake
        # error in the browser console instead of a readable code.
        await websocket.accept()
        await websocket.close(code=4403)
        return
    # WATCH-ONLY IS ENFORCED HERE, not in the viewer. A run's grant carries
    # view_only in the client's URL settings and the window puts pointer-events
    # off on top of it, but both live in the page: anything that speaks this
    # socket could still send RFB PointerEvent and KeyEvent frames and drive the
    # browser the agent is working in. The relay therefore drops everything
    # travelling client-to-container on an agent grant - the protocol needs no
    # client frames to deliver a picture (the server pushes framebuffer updates
    # on its own), so a silent drop costs the viewer nothing.
    watch_only = isinstance(mgr.validate_ticket(ticket), AgentStream)

    # Mirror the client's offered subprotocol (the KasmVNC client offers
    # "binary"); inventing or dropping it makes some browsers abort the socket.
    offered = websocket.scope.get("subprotocols") or []
    chosen = offered[0] if offered else None

    from websockets.asyncio.client import connect as _ws_connect
    target = (mgr.vnc_base().rstrip("/")
              .replace("http://", "ws://").replace("https://", "wss://") + "/websockify")
    counted = False
    try:
        # KasmVNC refuses a handshake without an Origin header ("missing
        # Sec-WebSocket-Origin", measured against 1.5.0); browsers always send
        # one, a server-side client must say so explicitly.
        # ping_interval=None: the client library pings every 20s and closes after a
        # 20s pong timeout, and the VNC server does not answer protocol pings - so
        # every stream died after EXACTLY 40 seconds of the person reading a page
        # (measured against the container log: connect 07:14:12, "Clean
        # disconnection" 07:14:52). The VNC protocol has its own liveness, and the
        # lease janitor already notices a dead container, so nothing is lost by
        # switching this off.
        async with _ws_connect(target, subprotocols=[chosen] if chosen else None,
                               origin=mgr.vnc_base(), max_size=None,
                               ping_interval=None) as upstream:
            await websocket.accept(subprotocol=chosen)
            if not mgr.stream_connected(ticket):
                await websocket.close(code=4403)
                return
            counted = True

            async def _to_upstream():
                while True:
                    msg = await websocket.receive()
                    if msg.get("type") == "websocket.disconnect":
                        return
                    if watch_only:
                        # Keep reading (the socket must stay drained and the
                        # disconnect must still be seen), forward nothing.
                        continue
                    if msg.get("bytes") is not None:
                        await upstream.send(msg["bytes"])
                    elif msg.get("text") is not None:
                        await upstream.send(msg["text"])

            # Report relayed bytes until the picture is up: an open socket is not a
            # picture, and the window's connecting cover waits for THIS, not for the
            # accept above. Stops reporting on the first True, so a running stream
            # costs nothing.
            _pixels_up = False

            async def _to_client():
                nonlocal _pixels_up
                async for frame in upstream:
                    if isinstance(frame, (bytes, bytearray)):
                        raw = bytes(frame)
                        if not _pixels_up:
                            _pixels_up = mgr.stream_bytes(ticket, len(raw))
                        await websocket.send_bytes(raw)
                    else:
                        await websocket.send_text(frame)

            tasks = [asyncio.create_task(_to_upstream()), asyncio.create_task(_to_client())]
            try:
                done, pending = await asyncio.wait(tasks, return_when=asyncio.FIRST_COMPLETED)
                for t in pending:
                    t.cancel()
                for t in done:
                    # Surface a pump crash (not a normal disconnect) in the log.
                    exc = t.exception()
                    if exc is not None and not isinstance(exc, (WebSocketDisconnect,)):
                        append_domain_log("webui", f"[browser-vnc] stream pump ended: {exc}")
            finally:
                for t in tasks:
                    t.cancel()
    except WebSocketDisconnect:
        pass
    except Exception as e:
        append_domain_log("webui", f"[browser-vnc] stream proxy failed: {e}")
    finally:
        if counted:
            mgr.stream_disconnected(ticket)
        try:
            await websocket.close()
        except Exception:
            pass


@app.get("/api/tools/{name}/source")
async def get_tool_source(name: str):
    """Get the source code of a tool."""
    print(f"[DEBUG] Fetching source for tool: {name}")
    try:
        manager = get_web_interface()
        if not manager.agent_instance:
            print("[DEBUG] Agent instance not found in manager")
            return {"error": "Agent not initialized"}
            
        if not hasattr(manager.agent_instance, 'tools'):
            print("[DEBUG] Agent has no tools attribute")
            return {"error": "Agent has no tools"}
            
        tool = manager.agent_instance.tools.get(name)
        if tool:
             try:
                 # Try to get source of the class
                 src = inspect.getsource(tool.__class__)
                 print(f"[DEBUG] Found source, length: {len(src)}")
                 return {"name": name, "code": src}
             except Exception as e:
                 print(f"[DEBUG] inspect.getsource failed: {e}")
                 # Fallback: Try to read from __file__ if available
                 try:
                     import os
                     file_path = inspect.getfile(tool.__class__)
                     if os.path.exists(file_path):
                         with open(file_path, 'r', encoding='utf-8') as f:
                             src = f.read()
                             return {"name": name, "code": src}
                 except Exception as ex:
                     print(f"[DEBUG] File read fallback failed: {ex}")
                 
                 return {"error": f"Failed to read source: {e}"}
        else:
             print(f"[DEBUG] Tool '{name}' not found in agent.tools")
             print(f"[DEBUG] Available tools: {list(manager.agent_instance.tools.keys())}")
             
    except Exception as e:
        print(f"[DEBUG] General error in get_tool_source: {e}")
        return {"error": str(e)}
    return {"error": "Tool not found"}

# Sounds directory for Web UI completion/notification (vaf/media/sounds)
SOUNDS_DIR = Path(__file__).resolve().parents[1] / "media" / "sounds"
ALLOWED_SOUND_FILES = {"tts01.mp3", "sst.mp3"}  # Only serve known files

def _resolve_room_workspace(room_id: str, request: Request, create: bool = False) -> str:
    """Workspace dir of an a2a room ('' if none, or the requester is not in it).

    Deliberately a fallback INSIDE _resolve_session_workspace rather than its own
    endpoint: the browser, the upload lane and the delete lane all resolve through
    that one function, so a room id resolving here is what makes all three work for
    rooms with no second code path. Membership is the same question the room WS
    commands ask (_room_rows): the requester must be IN the room through one of
    their lanes - merely being logged in on the host is not enough, because a room
    holds other people's words and its folder holds their files.
    """
    try:
        from vaf.api.config_routes import get_current_user_or_local_admin
        from vaf.core.a2a.room import Room
        from vaf.core.session import _room_rows
        user = get_current_user_or_local_admin(request) or {}
        user_scope_id = user.get("user_scope_id")
        if room_id not in {row["room_id"] for row in _room_rows(user_scope_id)}:
            return ""
        path = Room.open(room_id).workspace_dir(create=create)
        if path is None or not path.is_dir():
            return ""
        return str(path)
    except Exception:
        return ""


def _requester_name(request) -> str:
    """Who is making this request, for an audit line. Empty when unknown.

    Reuses the resolver the workspace routes already authorize with, so a security
    event names the same person the ownership check just used. Never raises: a
    missing name must not cost the log its whole entry.
    """
    try:
        from vaf.api.config_routes import get_current_user_or_local_admin
        return str((get_current_user_or_local_admin(request) or {}).get("username") or "")
    except Exception:
        return ""


def _inspect_attachment(data: bytes, filename: str, origin: str,
                        username: str = "", ip: str = ""):
    """Ask the known-bad list about arriving bytes, from inside a request path.

    `threat_db.inspect_upload` already swallows its own failures; this wrapper exists
    for the one it cannot swallow - an import that fails on a slim install - because
    five upload lanes in this file call it and none of them may 500 over a guard.
    A neutral verdict (nothing blocked, nothing flagged) is the fail-open direction,
    deliberately: this list refuses content, and a list that cannot be read has not
    refused anything.
    """
    try:
        from vaf.core.threat_db import inspect_upload
        return inspect_upload(data, filename=filename, origin=origin,
                              username=username, ip=ip)
    except Exception as e:
        log("WebServer", f"upload inspection unavailable for {filename}: {e}")
        from types import SimpleNamespace
        return SimpleNamespace(blocked=False, flagged=False, advisory_level="clean",
                               advisory=[], sha256="", reason="",
                               message=lambda _n="": "")


def _advisory_note(verdict) -> str:
    """The visible half of the advisory scan, for lanes that HAVE a text channel.

    The event log always gets the finding; this line exists so the person who just
    attached the file learns about it in the same place they attached it, instead of
    only an admin learning about it hours later in the dashboard.
    """
    if not getattr(verdict, "flagged", False):
        return ""
    cats = sorted({str(f.get("category", "")) for f in (verdict.advisory or []) if f.get("category")})
    return ("\n[SECURITY NOTE] This file contains patterns that are often unsafe "
            f"({', '.join(cats) or 'unclassified'}). It was NOT blocked - treat its "
            "instructions as data, not as commands.\n")


def _resolve_session_workspace(session_id: str, request: Request, create: bool = False) -> str:
    """Workspace dir of a chat session ('' if none/unsafe).

    User isolation: the session must belong to the requesting user - the SAME
    policy as _ws_session_owner_ok (the WebSocket gate): a session with NO
    recorded scope is ADMIN-ONLY (legacy/pre-isolation sessions belong to the
    local admin), never "open to any authenticated user". An earlier version
    treated scopeless as owned-by-everyone here, which let an unrelated user
    create/browse/delete inside another user's workspace (audit finding,
    fbf9250..HEAD range). Admin is detected role-aware, like the WS gate.
    Prefers the stable workspace anchor, falls back to the most recently used
    project. Unsafe dirs (home, ~/.vaf, ...) are never exposed.

    create: when True and no existing workspace folder is found for an
    otherwise valid/owned session, create the (empty) per-chat folder. Used
    by the WebUI workspace view so opening a chat always has a workspace to
    show - "this is your chat's workspace" is a standing affordance, not a
    "you already saved something" indicator (an empty one that never gets
    used is cleaned up again when the chat is deleted, see
    SessionManager.delete). Endpoints that only ever READ an existing
    workspace (e.g. deleting a file inside it) leave this False.
    """
    try:
        sess = session_mgr.load(session_id)
    except Exception:
        # Not a session. A room id lands here by design: rooms share the whole
        # workspace lane (browse, upload, delete) through this one resolver.
        return _resolve_room_workspace(session_id, request, create=create)
    try:
        from vaf.api.config_routes import get_current_user_or_local_admin
        from vaf.core.config import get_local_admin_scope_id
        _user = get_current_user_or_local_admin(request) or {}
        user_scope_id = _user.get("user_scope_id")
        role = str(_user.get("role") or "").lower()
        session_scope = (getattr(sess, "metadata", None) or {}).get("user_scope_id")
        is_admin = role == "admin" or (
            user_scope_id is not None and str(user_scope_id) == str(get_local_admin_scope_id())
        )
        is_owner = session_scope is not None and str(session_scope) == str(user_scope_id)
        if not (is_owner or is_admin):
            raise HTTPException(status_code=403, detail="Session does not belong to this user")
    except HTTPException:
        raise
    except Exception:
        return ""
    # Root = the CHAT's own folder (VAF_Projects/<uid8>/<session_id>), which can
    # hold several project folders. Fall back to the per-project path for
    # legacy sessions created before per-chat folders existed.
    import re as _re_ws
    _sid_folder = _re_ws.sub(r'[^a-zA-Z0-9_-]', '', str(session_id))[:32]
    path = ""
    try:
        from vaf.core.session import get_session_workspace_dir
        _ws = get_session_workspace_dir(session_id)
        if _ws:
            path = str(_ws)
    except Exception:
        path = ""
    if not path:
        path = getattr(sess, "project_path", "") or ""
        if path and os.path.isdir(path) and _sid_folder and os.path.basename(os.path.dirname(path)) == _sid_folder:
            path = os.path.dirname(path)
    if not path or not os.path.isdir(path):
        path = (getattr(sess, "runtime_state", None) or {}).get("last_project_path", "") or ""
    if (not path or not os.path.isdir(path)) and create:
        try:
            from vaf.core.session import get_session_workspace_dir as _gwd_create
            _created = _gwd_create(session_id, create=True)
            if _created:
                path = str(_created)
        except Exception:
            pass
    if not path or not os.path.isdir(path):
        return ""
    try:
        from vaf.tools.coder import is_unsafe_project_dir
        if is_unsafe_project_dir(path):
            return ""
    except Exception:
        return ""
    return path


def _resolve_workspace_subdir(root: str, subpath: str) -> str:
    """Join a browse subpath onto the workspace root, refusing escapes.

    Containment is `contained_path`, so it holds against a symlink inside the
    workspace as well: this used to compare normalized strings, and a link
    planted in a chat folder carried every caller of this - browse, upload,
    delete, mkdir - to the link's target.
    """
    try:
        return str(contained_path(root, subpath, must_exist=True))
    except PathEscape as e:
        if "not found" in str(e):
            raise HTTPException(status_code=404, detail="Folder not found")
        raise HTTPException(status_code=400, detail="Invalid subpath")


@app.get("/api/session/workspace")
async def get_session_workspace(
    request: Request,
    sessionId: str = Query(..., description="Chat session id"),
    subpath: str = Query("", description="Folder inside the workspace to list"),
):
    """Browse the chat's workspace folder (feeds the WebUI workspace window).

    create=True: every chat opened in the WebUI gets a workspace folder to
    point at immediately, even before anything has ever been saved into it
    (see _resolve_session_workspace); an unused one is cleaned up again when
    the chat is deleted.
    """
    root = _resolve_session_workspace(sessionId, request, create=True)
    if not root:
        # Orphan drill-in (central explorer): the session JSON is gone but the workspace folder may still
        # exist under the caller's OWN uid8 root. Strictly scoped to the requester's own VAF_Projects/<uid8>.
        try:
            from vaf.api.config_routes import get_current_user_or_local_admin
            from vaf.core.session import get_user_projects_root
            import re as _re_orph
            _scope = str((get_current_user_or_local_admin(request) or {}).get("user_scope_id") or "")
            _uroot = get_user_projects_root(_scope)
            _sid = _re_orph.sub(r'[^a-zA-Z0-9_-]', '', str(sessionId))[:32]
            if _uroot and _sid and (_uroot / _sid).is_dir():
                root = str(_uroot / _sid)
        except Exception:
            root = ""
    if not root:
        return {"path": "", "name": "", "subpath": "", "dirs": [], "files": []}
    target = _resolve_workspace_subdir(root, subpath)
    dirs, files = [], []
    try:
        from datetime import datetime as _dt2
        for entry in sorted(os.listdir(target)):
            if entry.startswith('.'):
                continue
            fp = os.path.join(target, entry)
            if os.path.isdir(fp):
                try:
                    items = len([e for e in os.listdir(fp) if not e.startswith('.')])
                except Exception:
                    items = 0
                dirs.append({"name": entry, "items": items})
            elif os.path.isfile(fp):
                st = os.stat(fp)
                files.append({
                    "name": entry,
                    "size": st.st_size,
                    "modified": _dt2.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d %H:%M"),
                })
    except Exception:
        pass
    rel = os.path.relpath(target, root)
    _live_title = None
    try:
        _live_title = getattr(session_mgr.load(sessionId), "name", None)
    except Exception:
        _live_title = None
    from vaf.core.session import read_workspace_label, resolve_workspace_display_name
    return {
        "path": root,
        "name": os.path.basename(root),
        "displayName": resolve_workspace_display_name(Path(root), os.path.basename(root), _live_title),
        "label": read_workspace_label(Path(root)),
        "subpath": "" if rel == "." else rel.replace(os.sep, "/"),
        "dirs": dirs,
        "files": files,
    }


class WorkspaceUploadRequest(BaseModel):
    sessionId: str
    filename: str
    content_base64: str
    subpath: str = ""


@app.post("/api/session/workspace/upload")
async def upload_session_workspace_file(req: WorkspaceUploadRequest, request: Request):
    """Upload a file into the chat's workspace (gives the agent direct access)."""
    # create=True: uploading IS "saving something" - create the folder on
    # demand rather than 404ing on a chat whose workspace was never touched.
    root = _resolve_session_workspace(req.sessionId, request, create=True)
    if not root:
        raise HTTPException(status_code=404, detail="Session has no workspace")
    path = _resolve_workspace_subdir(root, req.subpath)
    try:
        name = safe_entry_name(req.filename)
    except PathEscape:
        raise HTTPException(status_code=400, detail="Invalid filename")
    import base64 as _b64
    try:
        data = _b64.b64decode(req.content_base64 or "", validate=True)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid base64 content")
    if len(data) > 25 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File too large (max 25 MB)")
    # A workspace upload is the shortest path from a browser to a folder the agent
    # reads and writes with real file tools, so the gate sits before the write, not
    # after it. 403 rather than 400: this is a refusal, not a malformed request.
    _verdict = _inspect_attachment(data, name, "workspace_upload",
                                   username=_requester_name(request),
                                   ip=(request.client.host if request.client else ""))
    if _verdict.blocked:
        raise HTTPException(status_code=403, detail=_verdict.message(name))
    target = os.path.join(path, name)
    try:
        with open(target, "wb") as f:
            f.write(data)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Write failed: {e}")
    return {"ok": True, "name": name, "size": len(data)}


class WorkspaceMkdirRequest(BaseModel):
    sessionId: str
    subpath: str = ""
    name: str


@app.post("/api/session/workspace/mkdir")
async def create_session_workspace_folder(req: WorkspaceMkdirRequest, request: Request):
    """Create one folder inside the chat's workspace (the idle explorer's
    context menu). Same jail as every workspace route: the session must be the
    requester's own (_resolve_session_workspace) and the subpath cannot escape
    the root (_resolve_workspace_subdir). The name is validated the way the
    upload validates filenames - basename only, no dotfiles - because a name
    is the only thing the client contributes here."""
    root = _resolve_session_workspace(req.sessionId, request, create=True)
    if not root:
        raise HTTPException(status_code=404, detail="Session has no workspace")
    path = _resolve_workspace_subdir(root, req.subpath)
    try:
        name = safe_entry_name(req.name)
    except PathEscape:
        raise HTTPException(status_code=400, detail="Invalid folder name")
    target = os.path.join(path, name)
    if os.path.exists(target):
        raise HTTPException(status_code=409, detail="A file or folder with this name exists")
    try:
        os.mkdir(target)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Create failed: {e}")
    return {"ok": True, "name": name}


class WorkspaceDeleteRequest(BaseModel):
    sessionId: str
    name: str
    subpath: str = ""


@app.post("/api/session/workspace/delete")
async def delete_session_workspace_entry(req: WorkspaceDeleteRequest, request: Request):
    """Delete a file or folder inside the chat's workspace.

    The confirmation dialog lives in the UI; this endpoint only enforces the
    boundaries: same ownership rules as browsing, the target must stay inside
    the workspace root, and the root itself cannot be deleted.
    """
    root = _resolve_session_workspace(req.sessionId, request)
    if not root:
        raise HTTPException(status_code=404, detail="Session has no workspace")
    folder = _resolve_workspace_subdir(root, req.subpath)
    try:
        # `folder` already came back contained in the workspace, so containing
        # the entry within `folder` is the whole remaining question - and it is
        # asked against the folder rather than the root because a relative path
        # rebuilt between the two would compare a resolved folder with an
        # unresolved root, which produces an upward walk on any host whose
        # home or documents directory is itself a link.
        name = safe_entry_name(req.name)
        resolved = contained_path(folder, name)
    except PathEscape:
        raise HTTPException(status_code=400, detail="Invalid target")
    target = str(resolved)
    if resolved == Path(root).resolve():
        raise HTTPException(status_code=400, detail="Invalid target")
    if not os.path.exists(target):
        raise HTTPException(status_code=404, detail="Not found")
    try:
        if os.path.isdir(target):
            import shutil as _sh
            _sh.rmtree(target)
        else:
            os.remove(target)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Delete failed: {e}")
    return {"ok": True, "name": name}


# ── Central per-user Data Explorer (all of a user's workspaces, incl. orphans) ───────────────────────
# These endpoints operate on the per-user root VAF_Projects/<uid8>/ — derived ONLY from the AUTHENTICATED
# user's scope, never a client-supplied value — so a user can only ever see/rename/delete their OWN
# workspaces. Opaque session-id handles are returned, never absolute paths (keeps the /api/file surface
# closed). Orphans = workspace folders whose chat session JSON was deleted (left on disk by design).

@app.post("/api/speaker/test")
async def speaker_recognition_test(request: Request):
    """Live recognition test (Settings > Voice): score a short recording
    against the requesting user's voice DB. Strictly read-only."""
    from vaf.api.config_routes import get_current_user_or_local_admin
    from vaf.core.config import get_local_admin_scope_id
    from vaf.core import speaker_id as _sid
    import base64 as _b64t
    user = get_current_user_or_local_admin(request) or {}
    scope = str(user.get("user_scope_id") or get_local_admin_scope_id())
    body = await request.json()
    try:
        wav = _b64t.b64decode(body.get("audio") or "")
    except Exception:
        wav = b""
    if not wav:
        return {"ok": False, "error": "no_audio"}
    if not _sid.is_enabled() or _sid.load_profile(scope) is None:
        return {"ok": False, "error": "no_profile"}
    loop = asyncio.get_running_loop()
    res = await loop.run_in_executor(None, lambda: _sid.score_wav(wav, scope))
    if res is None:
        return {"ok": False, "error": "no_speech"}
    segs = await loop.run_in_executor(None, lambda: _sid.analyze_segments(wav, scope))
    prof = _sid.load_profile(scope)
    return {
        "ok": True, **res,
        "display_name": ((prof or {}).get("meta") or {}).get("display_name", "Ich"),
        "segments": segs or [],
        "threshold": _sid._threshold(), "band": _sid._band(),
    }


@app.post("/api/speaker/feedback")
async def speaker_recognition_feedback(request: Request):
    """Verdict on a recognition test ('correct'|'wrong'). Feeds the per-user
    threshold-calibration store. A WRONG verdict may carry a name: the voice
    is stored as a NAMED third-party profile. OWNER-CONFIRMED clips
    ('correct' on a self label, or the "it was me" false-reject path) also
    train the owner profile adaptively (add_owner_sample guardrails;
    speaker_id_adaptive_enabled) - the verdict comes from the authenticated
    owner session, never from the audio itself."""
    from vaf.api.config_routes import get_current_user_or_local_admin
    from vaf.core.config import get_local_admin_scope_id
    from vaf.core import speaker_id as _sid
    import base64 as _b64f
    user = get_current_user_or_local_admin(request) or {}
    scope = str(user.get("user_scope_id") or get_local_admin_scope_id())
    body = await request.json()
    verdict = body.get("verdict")
    if verdict not in ("correct", "wrong"):
        return {"ok": False, "error": "bad_verdict"}
    try:
        score = float(body.get("score"))
    except Exception:
        return {"ok": False, "error": "bad_score"}
    label = str(body.get("label") or "")
    saved_profile = None
    owner_claim = False
    _name = (body.get("name") or "").strip()
    _who = (body.get("who") or "").strip().lower()  # "owner" | "other" | ""
    _was = None
    if verdict == "wrong":
        if _who == "owner":
            # "It was me" button: a false reject - owner-side calibration
            # data AND (below) the most valuable adaptive sample there is:
            # exactly the border cases the profile misses today.
            owner_claim = True
            _was = "owner"
        elif _name:
            prof = _sid.load_profile(scope)
            owner_name = ((prof or {}).get("meta") or {}).get("display_name", "")
            if owner_name and owner_name.strip().lower() == _name.lower():
                # Typed the owner's name (case-insensitive): same as the
                # "me" button - and no owner-double as a named profile.
                owner_claim = True
                _was = "owner"
            else:
                _was = "other"
                if body.get("audio"):
                    try:
                        wav = _b64f.b64decode(body["audio"])
                        loop = asyncio.get_running_loop()
                        got = await loop.run_in_executor(None, lambda: _sid.embed_wav(wav))
                        if got is not None:
                            meta = _sid.save_named_profile(scope, _name, got["embedding"],
                                                           got["net_seconds"])
                            if meta:
                                saved_profile = meta["display_name"]
                    except Exception as _sf_e:
                        log("WebServer", f"speaker feedback named save failed: {_sf_e}")
        elif _who == "other":
            # "Someone else" without a name: still resolves the side.
            _was = "other"
    # Owner-approved adaptive learning (user decision): a clip the
    # authenticated owner confirmed as THEIR OWN voice trains the profile.
    _owner_confirmed = (verdict == "correct" and label == "self") or owner_claim
    if _owner_confirmed and body.get("audio"):
        try:
            from vaf.core.config import Config as _CfgA
            if _CfgA.get("speaker_id_adaptive_enabled", True):
                wav = _b64f.b64decode(body["audio"])
                loop = asyncio.get_running_loop()
                await loop.run_in_executor(
                    None, lambda: _sid.add_owner_sample(scope, wav))
        except Exception as _ad_e:
            log("WebServer", f"speaker feedback adaptive learn failed: {_ad_e}")
    stats = _sid.record_test_feedback(scope, score, label, verdict, was=_was)
    return {"ok": True, "stats": stats, "saved_profile": saved_profile,
            "owner_claim": owner_claim}


@app.get("/api/workspaces")
async def list_my_workspaces(request: Request):
    """List ALL of the requesting user's chat workspaces (live + orphaned) for the central Data Explorer."""
    from vaf.api.config_routes import get_current_user_or_local_admin
    from vaf.core.config import get_local_admin_scope_id
    from vaf.core.session import get_user_projects_root, read_workspace_label, resolve_workspace_display_name
    from datetime import datetime as _dt
    user = get_current_user_or_local_admin(request) or {}
    scope = str(user.get("user_scope_id") or "")
    is_admin = scope == str(get_local_admin_scope_id())
    root = get_user_projects_root(scope)
    if not root or not root.is_dir():
        return {"workspaces": [], "isAdmin": is_admin}
    try:
        live = {s["id"]: s for s in session_mgr.list(limit=100000, user_scope_id=scope)}
    except Exception:
        live = {}
    out = []
    try:
        for child in sorted(root.iterdir()):
            if not child.is_dir() or child.name.startswith('.'):
                continue
            sid = child.name
            live_sess = live.get(sid)
            live_title = (live_sess or {}).get("name") if live_sess else None
            fcount = dcount = 0
            updated = ""
            try:
                latest = 0.0
                for e in os.scandir(child):
                    if e.name.startswith('.'):
                        continue
                    if e.is_dir():
                        dcount += 1
                    elif e.is_file():
                        fcount += 1
                    try:
                        latest = max(latest, e.stat().st_mtime)
                    except Exception:
                        pass
                if latest:
                    updated = _dt.fromtimestamp(latest).strftime("%Y-%m-%d %H:%M")
            except Exception:
                pass
            out.append({
                "sessionId": sid,
                "displayName": resolve_workspace_display_name(child, sid, live_title),
                "label": read_workspace_label(child),
                "liveTitle": live_title,
                "orphan": live_sess is None,
                "fileCount": fcount,
                "folderCount": dcount,
                "updated": updated,
            })
    except Exception:
        pass
    # Live chats first (the ones the user is actually working in), orphaned
    # folders from deleted chats at the end - one flat list, no sections.
    out.sort(key=lambda w: (1 if w.get("orphan") else 0, str(w.get("displayName") or "").lower()))
    return {"workspaces": out, "isAdmin": is_admin}


# Whole-request budgets for /api/workspaces/search, shared across ALL of a user's
# workspaces (per-workspace caps below are the inner bound). A search runs on the
# shared thread pool and the frontend fires it on every debounced keystroke, so a
# crafted tree must never let one request monopolize the pool: the walk stops on the
# FIRST of wall-clock, total files read, or total entries (dirs+files) visited.
_SEARCH_DEADLINE_S = 2.0
_SEARCH_GLOBAL_FILES = 4000
_SEARCH_GLOBAL_ENTRIES = 60000


class _SearchBudget:
    """Mutable whole-request budget shared across every workspace in one search.
    exhausted() is monotonic-clock aware so a superseded walk self-terminates even
    if the frontend already dropped its response."""
    __slots__ = ("deadline", "files", "entries")

    def __init__(self, deadline, files, entries):
        self.deadline = deadline
        self.files = files
        self.entries = entries

    def exhausted(self) -> bool:
        return (self.files <= 0 or self.entries <= 0
                or (self.deadline is not None and time.monotonic() >= self.deadline))


def _search_one_workspace(child, q: str, per_ws_hits: int = 5, max_files: int = 400,
                          max_depth: int = 6, max_bytes: int = 1_000_000,
                          budget: "_SearchBudget" = None, root_real: str = None):
    """Search ONE workspace folder for `q` (case-insensitive substring) in file/folder
    NAMES and text-file CONTENTS. Returns {"files": [...], "truncated": bool} or None
    when nothing matched. Bounded on every axis (hit count, file count, ENTRY count,
    depth, file size, and an optional whole-request wall-clock/file/entry budget) so a
    crafted wide or deep tree can never stall the request; binary files (NUL in the
    first 8KB) match by name only.

    Symlink containment (user-isolation critical): file CONTENTS are read only when the
    file's real path stays inside `root_real` (the resolved workspace root). A symlink
    planted in the workspace that points at another user's VAF_Projects file or any host
    file (agent/coder can create files here) is therefore never opened for content - it
    can still match by NAME (its name genuinely lives in this workspace, no leak). os.walk
    runs with followlinks=False so symlinked directories are never descended."""
    ql = q.lower()
    hits = []
    scanned = 0          # files whose content we considered (per-workspace cap)
    entries = 0          # dirs + files visited (per-workspace breadth cap)
    max_entries = max_files * 40  # generous per-ws breadth bound; global budget is the real guard
    truncated = False
    if root_real is None:
        try:
            root_real = os.path.realpath(str(child))
        except Exception:
            root_real = str(child)
    _root_prefix = root_real.rstrip(os.sep) + os.sep
    try:
        base_depth = len(child.parts)
        for dirpath, dirnames, filenames in os.walk(child):  # followlinks=False (default): symlink dirs not descended
            # Stop before doing any work in this directory if the walk is already
            # over its per-workspace entry cap or the whole-request budget - so a
            # crafted tree cannot accumulate work directory by directory.
            if entries > max_entries or (budget is not None and budget.exhausted()):
                return {"files": hits, "truncated": True} if hits else None
            dirnames[:] = [d for d in dirnames if not d.startswith(".")]
            # Match directory NAMES at this level BEFORE the depth cutoff prunes them,
            # so a folder sitting exactly on the boundary is not silently skipped.
            for d in dirnames:
                if ql in d.lower():
                    rel = os.path.relpath(os.path.join(dirpath, d), str(child))
                    hits.append({"path": rel.replace(os.sep, "/"), "kind": "name"})
                    if len(hits) >= per_ws_hits:
                        return {"files": hits, "truncated": True}
            if len(Path(dirpath).parts) - base_depth >= max_depth:
                dirnames[:] = []
            entries += len(dirnames)
            if entries > max_entries or (budget is not None and budget.exhausted()):
                return {"files": hits, "truncated": True} if hits else None
            for fn in filenames:
                if fn.startswith("."):
                    continue
                entries += 1
                scanned += 1
                if scanned > max_files or entries > max_entries or (budget is not None and budget.exhausted()):
                    return {"files": hits, "truncated": True} if hits else None
                fp = os.path.join(dirpath, fn)
                rel = os.path.relpath(fp, str(child)).replace(os.sep, "/")
                if ql in fn.lower():
                    hits.append({"path": rel, "kind": "name"})
                    if len(hits) >= per_ws_hits:
                        return {"files": hits, "truncated": True}
                    continue  # name hit is enough for this file
                # Content read: only for files that actually stay inside the workspace
                # root (never follow a symlink out to another user's or a host file).
                # Hardlinks are intentionally NOT containment-checked: a hardlink is
                # indistinguishable from a native file by path (its directory entry
                # genuinely lives in the workspace, same inode), and fs.protected_hardlinks
                # (the modern-Linux default) already forbids linking a file the user
                # cannot read - so a hardlink can only surface content the owner could
                # read directly anyway; it is not a cross-user escape.
                try:
                    if os.path.islink(fp):
                        real = os.path.realpath(fp)
                        if real != root_real and not real.startswith(_root_prefix):
                            continue  # escapes the workspace - name-only, no content leak
                    if budget is not None:
                        budget.files -= 1
                    if os.path.getsize(fp) > max_bytes:
                        continue
                    with open(fp, "rb") as f:
                        raw = f.read(max_bytes)
                    if b"\0" in raw[:8192]:
                        continue  # binary: name-only matching
                    text = raw.decode("utf-8", errors="ignore")
                    idx = text.lower().find(ql)
                    if idx < 0:
                        continue
                    start = max(0, idx - 40)
                    snippet = " ".join(text[start:idx + len(q) + 60].split())
                    hits.append({"path": rel, "kind": "content", "snippet": snippet[:140]})
                    if len(hits) >= per_ws_hits:
                        return {"files": hits, "truncated": True}
                except OSError:
                    continue
            if budget is not None:
                budget.entries -= (len(dirnames) + len(filenames))
    except Exception:
        truncated = True
    return {"files": hits, "truncated": truncated} if hits else None


@app.get("/api/workspaces/search")
async def search_my_workspaces(request: Request, q: str = ""):
    """Search the requesting user's workspaces: file/folder names and text-file
    contents (workspace display names are filtered client-side, which already has
    them). Scoped to the caller's own VAF_Projects/<uid[:8]>/ root exactly like
    /api/workspaces; the query is the only client input - never a path."""
    from vaf.api.config_routes import get_current_user_or_local_admin
    from vaf.core.session import get_user_projects_root
    q = str(q or "").strip()
    if len(q) < 2:
        return {"results": {}, "query": q}
    user = get_current_user_or_local_admin(request) or {}
    scope = str(user.get("user_scope_id") or "")
    root = get_user_projects_root(scope)
    if not root or not root.is_dir():
        return {"results": {}, "query": q}

    def _walk_all():
        results = {}
        budget = _SearchBudget(time.monotonic() + _SEARCH_DEADLINE_S,
                               _SEARCH_GLOBAL_FILES, _SEARCH_GLOBAL_ENTRIES)
        try:
            # Only real subdirectories that stay inside the user's root - never follow a
            # top-level symlink pointing at another user's tree (user-isolation critical).
            root_canon = os.path.realpath(str(root))
            root_prefix = root_canon.rstrip(os.sep) + os.sep
            for child in sorted(root.iterdir()):
                if budget.exhausted():
                    break
                if not child.is_dir() or child.name.startswith("."):
                    continue
                child_real = os.path.realpath(str(child))
                if child_real != root_canon and not child_real.startswith(root_prefix):
                    continue  # symlink escaping the user's root
                found = _search_one_workspace(child, q, budget=budget, root_real=child_real)
                if found:
                    results[child.name] = found
        except Exception:
            pass
        return results

    return {"results": await asyncio.to_thread(_walk_all), "query": q}


class WorkspaceLabelRequest(BaseModel):
    sessionId: str
    label: str


@app.post("/api/workspaces/rename")
async def rename_my_workspace(req: WorkspaceLabelRequest, request: Request):
    """Set a workspace's DISPLAY label (rename = display label only; the on-disk folder stays == session_id)."""
    from vaf.api.config_routes import get_current_user_or_local_admin
    from vaf.core.session import get_user_projects_root, write_workspace_label
    import re as _re_wl
    user = get_current_user_or_local_admin(request) or {}
    scope = str(user.get("user_scope_id") or "")
    root = get_user_projects_root(scope)
    if not root:
        raise HTTPException(status_code=404, detail="No workspace store for this user")
    sid = _re_wl.sub(r'[^a-zA-Z0-9_-]', '', str(req.sessionId))[:32]
    folder = root / sid if sid else None
    if not folder or not folder.is_dir():
        raise HTTPException(status_code=404, detail="Workspace not found")
    if not write_workspace_label(folder, req.label):
        raise HTTPException(status_code=500, detail="Could not set label")
    return {"ok": True, "sessionId": sid, "label": (req.label or "").strip()[:200]}


class WorkspaceFolderDeleteRequest(BaseModel):
    sessionId: str


@app.post("/api/workspaces/delete")
async def delete_my_workspace(req: WorkspaceFolderDeleteRequest, request: Request):
    """Delete a WHOLE workspace folder (orphans + manual cleanup). Boundary-guarded to the caller's own root."""
    from vaf.api.config_routes import get_current_user_or_local_admin
    from vaf.core.session import get_user_projects_root
    import re as _re_wd, shutil as _sh
    user = get_current_user_or_local_admin(request) or {}
    scope = str(user.get("user_scope_id") or "")
    root = get_user_projects_root(scope)
    if not root:
        raise HTTPException(status_code=404, detail="No workspace store for this user")
    sid = _re_wd.sub(r'[^a-zA-Z0-9_-]', '', str(req.sessionId))[:32]
    if not sid:
        raise HTTPException(status_code=400, detail="Invalid sessionId")
    root_norm = os.path.normpath(str(root))
    target = os.path.normpath(os.path.join(root_norm, sid))
    if target == root_norm or not target.startswith(root_norm + os.sep):
        raise HTTPException(status_code=400, detail="Invalid target")
    if not os.path.isdir(target):
        raise HTTPException(status_code=404, detail="Workspace not found")
    try:
        _sh.rmtree(target)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Delete failed: {e}")
    return {"ok": True, "sessionId": sid}


@app.get("/api/file")
async def get_file(request: Request, path: str = Query(..., description="Absolute path to local file")):
    """Serve a local file by path (allowed roots: documents, downloads, data dir, VAF output dir).

    FOUR roots, not three - the output dir is easy to miss and this docstring said three until
    2026-07-30. The list matters beyond this endpoint: `document_editor` carries no access check
    of its own precisely because the decision lives here, so anything added to it becomes
    reachable by a tool that never looks at a path. Pinned in tests/test_api_file_allowed_roots.py.
    Used by Web UI."""
    from vaf.core.platform import Platform
    import mimetypes
    try:
        target = Platform.normalize_path(path)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")
    if not target.exists() or not target.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    allowed_roots = [
        Platform.documents_dir().resolve(),
        Platform.downloads_dir().resolve(),
        Platform.data_dir().resolve(),
        Platform.get_vaf_output_dir().resolve(),
    ]
    if not any(target.is_relative_to(root) for root in allowed_roots):
        raise HTTPException(status_code=403, detail="Access denied")
    # User isolation for generated projects: VAF_Projects/<uid[:8]>/... folders
    # belong to one user — only that user (or the local admin) may download
    # from them. Legacy flat projects (no user prefix) stay accessible.
    import re as _re_iso
    _projects_root = (Platform.documents_dir() / "VAF_Projects").resolve()
    if target.is_relative_to(_projects_root):
        _rel = target.relative_to(_projects_root)
        _first_seg = _rel.parts[0] if _rel.parts else ""
        if _re_iso.fullmatch(r"[0-9a-f]{8}", _first_seg):
            # FAIL-CLOSED: if ownership cannot be verified, deny (never serve a per-user file on error).
            try:
                from vaf.api.config_routes import get_current_user_or_local_admin
                from vaf.core.config import is_admin_identity
                _user = get_current_user_or_local_admin(request) or {}
                _scope = str(_user.get("user_scope_id") or "")
                _is_admin = is_admin_identity(_user.get("role"), _scope)
                _allowed = _is_admin or _scope.replace("-", "").lower().startswith(_first_seg)
                if not _allowed:
                    # The one path into another account's tree that is not another
                    # account's business: the shared folder of a room this account was
                    # admitted to. The agent's file tools got the same exception, and
                    # without it here a member could write a file into the room and
                    # then be refused when it clicked its own link.
                    from vaf.tools.filesystem import _shared_room_roots
                    _allowed = any(
                        target == Path(_r).resolve() or target.is_relative_to(Path(_r).resolve())
                        for _r in _shared_room_roots(_scope))
            except HTTPException:
                raise
            except Exception:
                _allowed = False
            if not _allowed:
                raise HTTPException(status_code=403, detail="Access denied")
    mime_type, _ = mimetypes.guess_type(str(target))
    return FileResponse(
        path=str(target),
        media_type=mime_type or "application/octet-stream",
        filename=target.name,
    )


@app.post("/api/image/describe")
async def describe_image(request: Request):
    """One-time vision description of a chat image (for the Image Viewer + agent context).

    Cached per (session, path) in runtime_state so it is generated once; reuses a
    chat-uploaded image's existing base_description when present. Same per-user isolation
    as /api/file: only the owner (or local admin) of VAF_Projects/<uid8> may describe it.
    """
    from pathlib import Path as _Path
    from vaf.core.platform import Platform
    try:
        body = await request.json()
    except Exception:
        body = {}
    session_id = str((body or {}).get("sessionId") or "").strip()
    path = str((body or {}).get("path") or "").strip()
    if not session_id or not path:
        raise HTTPException(status_code=400, detail="sessionId and path required")
    try:
        target = Platform.normalize_path(path)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")
    if not target.exists() or not target.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    # Scope to the user's own VAF_Projects/<uid8> subtree (fail-closed), mirroring /api/file.
    import re as _re_id
    _projects_root = (Platform.documents_dir() / "VAF_Projects").resolve()
    _resolved = target.resolve()
    if not _resolved.is_relative_to(_projects_root):
        raise HTTPException(status_code=403, detail="Access denied")
    _rel_parts = _resolved.relative_to(_projects_root).parts
    _first = _rel_parts[0] if _rel_parts else ""
    if _re_id.fullmatch(r"[0-9a-f]{8}", _first):
        try:
            from vaf.api.config_routes import get_current_user_or_local_admin
            from vaf.core.config import is_admin_identity
            _user = get_current_user_or_local_admin(request) or {}
            _scope = str(_user.get("user_scope_id") or "")
            # Role-aware, exactly like the session check 25 lines below: this function used to
            # answer "is this an admin" two different ways within one request.
            _is_admin = is_admin_identity(_user.get("role"), _scope)
            _allowed = _is_admin or _scope.replace("-", "").lower().startswith(_first)
        except HTTPException:
            raise
        except Exception:
            _allowed = False
        if not _allowed:
            raise HTTPException(status_code=403, detail="Access denied")

    from vaf.core.session import SessionManager
    sm = SessionManager()
    try:
        sess = sm.load(session_id)
    except Exception:
        sess = None
    # User isolation: never read or write another user's session (same ownership rule as
    # _resolve_session_workspace / the WS load_session gate). A session with NO recorded
    # scope is ADMIN-ONLY (legacy/pre-isolation sessions belong to the local admin) -
    # this copy had the same owned-by-everyone hole as _resolve_session_workspace and was
    # tightened with it (audit finding, fbf9250..HEAD). Fail closed if unverifiable.
    if sess is not None:
        try:
            from vaf.api.config_routes import get_current_user_or_local_admin
            from vaf.core.config import get_local_admin_scope_id
            _uinfo = get_current_user_or_local_admin(request) or {}
            _uscope = _uinfo.get("user_scope_id")
            _urole = str(_uinfo.get("role") or "").lower()
            _sscope = (getattr(sess, "metadata", None) or {}).get("user_scope_id")
            _admin = _urole == "admin" or (
                _uscope is not None and str(_uscope) == str(get_local_admin_scope_id())
            )
            _owner = _sscope is not None and str(_sscope) == str(_uscope)
            if not (_owner or _admin):
                raise HTTPException(status_code=403, detail="Session does not belong to this user")
        except HTTPException:
            raise
        except Exception:
            sess = None  # fail closed: don't touch a session we couldn't verify
    _key = str(target.resolve())
    desc = ""
    # 1. Already cached for this session.
    if sess is not None:
        desc = ((getattr(sess, "runtime_state", None) or {}).get("image_descriptions") or {}).get(_key) or ""
    # 2. Reuse a chat-uploaded image's persisted base_description (match by path).
    if not desc and sess is not None:
        for _m in getattr(sess, "messages", None) or []:
            for _img in ((getattr(_m, "metadata", None) or {}).get("images") or []):
                if isinstance(_img, dict) and _img.get("path") and _img.get("base_description"):
                    try:
                        if _Path(_img["path"]).resolve() == target.resolve():
                            desc = _img["base_description"]
                            break
                    except Exception:
                        pass
            if desc:
                break
    # 3. Generate once via the vision backend, then cache.
    if not desc:
        try:
            import asyncio as _asyncio
            import mimetypes as _mt
            from vaf.core.vision_infer import describe_image_cached
            _mime = _mt.guess_type(str(target))[0] or "image/png"
            # Offload the blocking vision call so the async event loop stays responsive.
            # describe_image_cached shares a process-wide cache with the chat-upload base
            # description, so the same image is never described/billed twice (incl. the race
            # where the viewer is opened mid-turn).
            desc = (await _asyncio.to_thread(
                describe_image_cached,
                {"path": str(target), "mime_type": _mime, "name": target.name},
            )) or ""
            if desc and sess is not None:
                if not getattr(sess, "runtime_state", None):
                    sess.runtime_state = {}
                sess.runtime_state.setdefault("image_descriptions", {})[_key] = desc
                sm.save(sess, sync_state=False)
        except Exception as _de:
            log("WebServer", f"describe_image failed: {_de}")
    return {"name": target.name, "description": desc}


def _allowed_file_path(path_str: str):
    """Resolve path and check it is under allowed roots. Returns Path or raises HTTPException."""
    from vaf.core.platform import Platform
    try:
        target = Platform.normalize_path(path_str)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")
    if not target.exists() or not target.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    allowed_roots = [
        Platform.documents_dir().resolve(),
        Platform.downloads_dir().resolve(),
        Platform.data_dir().resolve(),
        Platform.get_vaf_output_dir().resolve(),
    ]
    if not any(target.is_relative_to(root) for root in allowed_roots):
        raise HTTPException(status_code=403, detail="Access denied")
    return target


def _escape_html(s: str) -> str:
    import html
    return html.escape(s, quote=True)


def _office_to_pdf_via_gotenberg(file_path: str, filename: str) -> bytes | None:
    """
    Convert Office docs (DOCX, XLSX, PPTX, ODT, ODS, ODP) to PDF via Gotenberg Docker service.
    Returns PDF bytes or None if unavailable. Gotenberg uses LibreOffice under the hood; MIT license.
    """
    url = (Config.get("document_conversion_docker_url") or "").strip().rstrip("/")
    if not url:
        return None
    try:
        import requests
        api_url = f"{url}/forms/libreoffice/convert"
        with open(file_path, "rb") as f:
            files = {"files": (filename, f, "application/octet-stream")}
            r = requests.post(api_url, files=files, timeout=60)
        if r.status_code == 200 and r.headers.get("content-type", "").lower().startswith("application/pdf"):
            return r.content
    except Exception as e:
        log("WebServer", f"Gotenberg conversion failed for {filename}: {e}")
    return None


def _docx_to_pdf_via_libreoffice(docx_path: str) -> bytes | None:
    """
    Convert DOCX to PDF using LibreOffice headless. Returns PDF bytes or None if unavailable.
    LibreOffice is MPL 2.0 - compatible with MIT projects. Use in Docker: install libreoffice in the image.
    """
    import subprocess
    import shutil
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None
    out_dir = os.path.dirname(docx_path)
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, docx_path],
            capture_output=True,
            timeout=60,
            check=False,
        )
        pdf_path = os.path.splitext(docx_path)[0] + ".pdf"
        if os.path.isfile(pdf_path):
            with open(pdf_path, "rb") as f:
                data = f.read()
            try:
                os.unlink(pdf_path)
            except Exception:
                pass
            return data
    except (subprocess.TimeoutExpired, OSError, Exception):
        pass
    return None


def _docx_to_html(target) -> str:
    from docx import Document
    doc = Document(str(target))
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = (para.style.name or "").lower()
        if "heading 1" in style_name or para.style.name == "Title":
            parts.append(f"<h1>{_escape_html(text)}</h1>")
        elif "heading 2" in style_name:
            parts.append(f"<h2>{_escape_html(text)}</h2>")
        elif "heading 3" in style_name:
            parts.append(f"<h3>{_escape_html(text)}</h3>")
        else:
            parts.append(f"<p>{_escape_html(text)}</p>")
    for table in doc.tables:
        parts.append("<table border=\"1\" cellpadding=\"4\" cellspacing=\"0\">")
        for row in table.rows:
            parts.append("<tr>")
            for cell in row.cells:
                parts.append(f"<td>{_escape_html(cell.text.strip())}</td>")
            parts.append("</tr>")
        parts.append("</table>")
    return "".join(parts) if parts else "<p></p>"


def _xlsx_to_html(target) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(target, read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames[:10]:
        sheet = wb[sheet_name]
        parts.append(f"<h2>Sheet: {_escape_html(sheet_name)}</h2>")
        parts.append("<table border=\"1\" cellpadding=\"4\" cellspacing=\"0\">")
        max_row = min(sheet.max_row, 500)
        max_col = min(sheet.max_column, 30)
        for row in sheet.iter_rows(min_row=1, max_row=max_row, max_col=max_col, values_only=True):
            parts.append("<tr>")
            for cell in row:
                val = "" if cell is None else str(cell)
                parts.append(f"<td>{_escape_html(val)}</td>")
            parts.append("</tr>")
        parts.append("</table>")
    wb.close()
    return "".join(parts) if parts else "<p></p>"


def _pptx_to_html(target) -> str:
    from pptx import Presentation
    prs = Presentation(str(target))
    parts = []
    for i, slide in enumerate(prs.slides[:50], 1):
        parts.append(f"<h2>Slide {i}</h2>")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(f"<p>{_escape_html(shape.text.strip())}</p>")
    return "".join(parts) if parts else "<p></p>"


@app.get("/api/file/as-html")
async def get_file_as_html(path: str = Query(..., description="Path to .docx, .xlsx or .pptx to convert to HTML for editing")):
    """Convert Word (.docx), Excel (.xlsx) or PowerPoint (.pptx) to HTML so the Document Editor can display and edit it."""
    target = _allowed_file_path(path)
    suf = target.suffix.lower()
    if suf not in (".docx", ".xlsx", ".pptx"):
        raise HTTPException(status_code=400, detail="Only .docx, .xlsx and .pptx can be converted to HTML here")
    try:
        if suf == ".docx":
            body = _docx_to_html(target)
        elif suf == ".xlsx":
            body = _xlsx_to_html(target)
        else:
            body = _pptx_to_html(target)
        html = "<!DOCTYPE html><html><head><meta charset=\"utf-8\"/></head><body>" + body + "</body></html>"
        from fastapi.responses import HTMLResponse
        return HTMLResponse(html)
    except ImportError as e:
        raise HTTPException(status_code=503, detail=f"Support not installed: {e}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to convert to HTML: {e}")


@app.get("/api/file/docx-model")
async def get_file_as_docx_model(path: str = Query(..., description="Path to .docx to convert to VAF's native DOCX model")):
    """Convert a DOCX file into the native editor model."""
    target = _allowed_file_path(path)
    if target.suffix.lower() != ".docx":
        raise HTTPException(status_code=400, detail="Only .docx files are supported")
    try:
        from vaf.core.docx_import import import_docx_to_native_model

        model = import_docx_to_native_model(target)
        return model.to_dict()
    except ImportError as e:
        raise HTTPException(status_code=503, detail=f"DOCX support not installed: {e}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to convert docx to native model: {e}")


class FileSaveRequest(BaseModel):
    """Request body for saving a file."""
    path: str
    content: str


class FileSaveDocxRequest(BaseModel):
    """Request body for saving HTML content back as .docx."""
    path: str
    content: str  # HTML from the editor


class FileSaveDocxNativeRequest(BaseModel):
    """Request body for saving native DOCX model content back as .docx."""
    path: str
    document: dict


def _strip_html_to_text(html_fragment: str) -> str:
    import re
    import html
    t = re.sub(r"<[^>]+>", " ", html_fragment)
    t = html.unescape(t)
    return " ".join(t.split()).strip()


@app.post("/api/file/save-docx")
async def save_file_as_docx(request: FileSaveDocxRequest):
    """Save editor content (HTML) back to a Word (.docx) file. Creates/overwrites the file."""
    from vaf.core.platform import Platform
    import re
    try:
        target = Platform.normalize_path(request.path)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")
    allowed_roots = [
        Platform.documents_dir().resolve(),
        Platform.downloads_dir().resolve(),
        Platform.data_dir().resolve(),
        Platform.get_vaf_output_dir().resolve(),
    ]
    if not any(target.is_relative_to(root) for root in allowed_roots):
        raise HTTPException(status_code=403, detail="Access denied")
    if target.suffix.lower() != ".docx":
        raise HTTPException(status_code=400, detail="Path must be a .docx file")
    try:
        from docx import Document
        doc = Document()
        html = request.content or ""
        body_match = re.search(r"<body[^>]*>(.*?)</body>", html, re.DOTALL | re.IGNORECASE)
        if body_match:
            html = body_match.group(1)
        # Find all block elements in order: h1, h2, h3, p, table
        pattern = r"<(h[1-3]|p|table)[^>]*>(.*?)</\1>"
        for m in re.finditer(pattern, html, re.DOTALL | re.IGNORECASE):
            tag, inner = m.group(1).lower(), m.group(2)
            text = _strip_html_to_text(inner)
            if tag.startswith("h") and len(tag) == 2:
                level = int(tag[1])
                if text:
                    doc.add_heading(text, level=level)
            elif tag == "p" and text:
                doc.add_paragraph(text)
            elif tag == "table":
                rows_html = re.findall(r"<tr[^>]*>(.*?)</tr>", inner, re.DOTALL | re.IGNORECASE)
                if rows_html:
                    all_cells = [re.findall(r"<t[dh][^>]*>(.*?)</t[dh]>", r, re.DOTALL | re.IGNORECASE) for r in rows_html]
                    nrows = len(all_cells)
                    ncols = max(len(c) for c in all_cells) if all_cells else 0
                    if ncols > 0:
                        table = doc.add_table(rows=nrows, cols=ncols)
                        for ri, cells in enumerate(all_cells):
                            for ci, cell_html in enumerate(cells):
                                if ci < ncols:
                                    table.rows[ri].cells[ci].text = _strip_html_to_text(cell_html)
        # If no structured blocks found, add whole body as one paragraph
        if len(doc.paragraphs) == 0 and not doc.tables:
            text = _strip_html_to_text(html)
            if text:
                doc.add_paragraph(text)
        doc.save(str(target))
        return {"status": "ok", "path": str(target)}
    except ImportError:
        raise HTTPException(status_code=503, detail="Word support not installed. Run: pip install python-docx")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save docx: {e}")


@app.post("/api/file/save-docx-native")
async def save_file_as_docx_native(request: FileSaveDocxNativeRequest):
    """Save VAF's native DOCX editor model back as a .docx file."""
    target = _allowed_save_path(request.path, ".docx")
    try:
        from vaf.core.docx_export import export_native_docx
        from vaf.core.docx_native_model import NativeDocxDocument

        document = NativeDocxDocument.from_dict(request.document or {})
        saved_path = export_native_docx(document, target)
        return {"status": "ok", "path": str(saved_path)}
    except ImportError as e:
        raise HTTPException(status_code=503, detail=f"DOCX support not installed: {e}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save native docx: {e}")


class FileSaveOfficeRequest(BaseModel):
    """Request body for saving HTML content back as .xlsx or .pptx."""
    path: str
    content: str  # HTML from the editor


def _allowed_save_path(path_str: str, required_suffix: str):
    from vaf.core.platform import Platform
    try:
        target = Platform.normalize_path(path_str)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")
    allowed_roots = [
        Platform.documents_dir().resolve(),
        Platform.downloads_dir().resolve(),
        Platform.data_dir().resolve(),
        Platform.get_vaf_output_dir().resolve(),
    ]
    if not any(target.is_relative_to(root) for root in allowed_roots):
        raise HTTPException(status_code=403, detail="Access denied")
    if target.suffix.lower() != required_suffix:
        raise HTTPException(status_code=400, detail=f"Path must be a {required_suffix} file")
    return target


@app.post("/api/file/save-xlsx")
async def save_file_as_xlsx(request: FileSaveOfficeRequest):
    """Save editor content (HTML tables) back to an Excel (.xlsx) file. First table = first sheet."""
    import re
    target = _allowed_save_path(request.path, ".xlsx")
    html = request.content or ""
    body_match = re.search(r"<body[^>]*>(.*?)</body>", html, re.DOTALL | re.IGNORECASE)
    if body_match:
        html = body_match.group(1)
    try:
        import openpyxl
        from openpyxl import Workbook
        wb = Workbook()
        wb.remove(wb.active)
        tables = re.findall(r"<table[^>]*>(.*?)</table>", html, re.DOTALL | re.IGNORECASE)
        for ti, table_html in enumerate(tables[:20]):
            rows_html = re.findall(r"<tr[^>]*>(.*?)</tr>", table_html, re.DOTALL | re.IGNORECASE)
            if not rows_html:
                continue
            all_cells = [re.findall(r"<t[dh][^>]*>(.*?)</t[dh]>", r, re.DOTALL | re.IGNORECASE) for r in rows_html]
            ncols = max(len(c) for c in all_cells) if all_cells else 0
            if ncols == 0:
                continue
            sheet_name = f"Sheet{ti + 1}" if ti > 0 else "Sheet1"
            ws = wb.create_sheet(sheet_name, ti)
            for ri, cells in enumerate(all_cells):
                for ci, cell_html in enumerate(cells):
                    if ci < ncols:
                        val = _strip_html_to_text(cell_html)
                        ws.cell(row=ri + 1, column=ci + 1, value=val)
        if not wb.sheetnames:
            ws = wb.create_sheet("Sheet1")
            ws.cell(row=1, column=1, value="")
        wb.save(str(target))
        return {"status": "ok", "path": str(target)}
    except ImportError:
        raise HTTPException(status_code=503, detail="Excel support not installed. Run: pip install openpyxl")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save xlsx: {e}")


@app.post("/api/file/save-pptx")
async def save_file_as_pptx(request: FileSaveOfficeRequest):
    """Save editor content (HTML: h2 = slide title, p = body) back to a PowerPoint (.pptx) file."""
    import re
    target = _allowed_save_path(request.path, ".pptx")
    html = request.content or ""
    body_match = re.search(r"<body[^>]*>(.*?)</body>", html, re.DOTALL | re.IGNORECASE)
    if body_match:
        html = body_match.group(1)
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        slide_sections = re.split(r"<h2[^>]*>", html, flags=re.IGNORECASE)
        for i, section in enumerate(slide_sections):
            if i == 0 and not re.search(r"</h2>", section, re.IGNORECASE):
                continue
            parts = re.split(r"</h2>", section, maxsplit=1, flags=re.IGNORECASE)
            title_html = parts[0] if len(parts) > 1 else ""
            body_html = parts[1] if len(parts) > 1 else section
            title = _strip_html_to_text(title_html) or f"Slide {i}"
            body_text = _strip_html_to_text(body_html)
            if not title and not body_text:
                continue
            blank = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank)
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(0.8)
            tb = slide.shapes.add_textbox(left, top, width, height)
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(24)
            if body_text:
                tb2 = slide.shapes.add_textbox(left, Inches(1.5), width, Inches(5.5))
                tf2 = tb2.text_frame
                tf2.text = body_text
        if len(prs.slides) == 0:
            blank = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1)).text_frame.text = "Untitled"
        prs.save(str(target))
        return {"status": "ok", "path": str(target)}
    except ImportError:
        raise HTTPException(status_code=503, detail="PowerPoint support not installed. Run: pip install python-pptx")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save pptx: {e}")


@app.post("/api/file/save")
async def save_file(request: FileSaveRequest):
    """Save content to a local file (allowed roots: documents, downloads, data dir, VAF output dir)."""
    from vaf.core.platform import Platform
    try:
        target = Platform.normalize_path(request.path)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")

    # Security: only allow saving to specific directories
    allowed_roots = [
        Platform.documents_dir().resolve(),
        Platform.downloads_dir().resolve(),
        Platform.data_dir().resolve(),
        Platform.get_vaf_output_dir().resolve(),
    ]
    if not any(target.is_relative_to(root) for root in allowed_roots):
        raise HTTPException(status_code=403, detail="Access denied - file must be in Documents, Downloads, or VAF data directory")

    import re
    content = request.content
    # Markdown files are rendered to HTML for the Document Editor; convert the
    # edited HTML back to Markdown so a .md file never ends up containing HTML.
    if target.suffix.lower() in (".md", ".mdx", ".markdown") and re.search(r"<\s*(?:html|body|div|p|h[1-6]|ul|ol|table)\b", content[:2000], re.IGNORECASE):
        try:
            from markdownify import markdownify as md
            content = re.sub(r"\n{3,}", "\n\n", md(content, heading_style="ATX")).strip() + "\n"
        except Exception:
            pass  # fall back to writing the raw editor content

    try:
        # Ensure parent directory exists
        target.parent.mkdir(parents=True, exist_ok=True)
        # Write content
        target.write_text(content, encoding='utf-8')
        return {"status": "ok", "path": str(target)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save file: {str(e)}")


@app.get("/sounds/{filename}")
async def get_sound(filename: str):
    """Serve sound files from vaf/media/sounds for Web UI (e.g. answer-complete notification)."""
    if filename not in ALLOWED_SOUND_FILES:
        raise HTTPException(status_code=404, detail="Sound not found")
    path = SOUNDS_DIR / filename
    if not path.is_file():
        raise HTTPException(status_code=404, detail="Sound not found")
    import mimetypes
    mime_type, _ = mimetypes.guess_type(filename)
    return FileResponse(
        path=str(path),
        media_type=mime_type or "audio/mpeg",
        filename=filename,
    )


@app.get("/api/download")
async def download_file(request: Request, path: str = Query(..., description="Absolute path to local file")):
    from vaf.core.platform import Platform
    from pathlib import Path
    import mimetypes

    try:
        target = Platform.normalize_path(path)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")

    if not target.exists() or not target.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    allowed_roots = [
        Platform.documents_dir().resolve(),
        Platform.downloads_dir().resolve(),
        Platform.data_dir().resolve(),
        Platform.get_vaf_output_dir().resolve(),
    ]

    if not any(target.is_relative_to(root) for root in allowed_roots):
        raise HTTPException(status_code=403, detail="Access denied")

    mime_type, _ = mimetypes.guess_type(str(target))

    # Local users (127.0.0.1 / ::1) get inline serving — browser opens HTML, PDF, images directly.
    # Remote/LAN users get Content-Disposition: attachment so a download dialog appears.
    client_host = request.client.host if request.client else ""
    is_local = client_host in ("127.0.0.1", "::1", "localhost", "")
    return FileResponse(
        path=str(target),
        media_type=mime_type or "application/octet-stream",
        filename=None if is_local else target.name,
    )


@app.get("/api/notifications")
async def get_notifications_api(
    request: Request,
    limit: int = Query(50, ge=1, le=100),
):
    """Return recent notifications for the current user (thinking, automation, channel replies)."""
    try:
        from vaf.api.config_routes import get_current_user_or_local_admin
        from vaf.core.user_notifications import get_notifications
    except ImportError as e:
        raise HTTPException(status_code=500, detail=f"Notifications not available: {e}")
    user = get_current_user_or_local_admin(request)
    user_scope_id = user.get("user_scope_id")
    notifications = get_notifications(user_scope_id, limit=limit)
    return {"notifications": notifications}


@app.get("/api/workflows/{wf_id}")
async def get_workflow_details(wf_id: str):
    """Get workflow definition and steps."""
    try:
        from vaf.workflows.templates import WORKFLOW_TEMPLATES
        
        # Try to find by ID first
        wf = WORKFLOW_TEMPLATES.get(wf_id)
        
        # If not found, try by name (fallback)
        if not wf:
            for w in WORKFLOW_TEMPLATES.values():
                if w.get("name") == wf_id:
                    wf = w
                    break
        
        if wf:
            # Format steps for visualization
            steps = []
            for idx, step in enumerate(wf.get("steps", [])):
                # Create a pseudo-code representation of the step
                import json
                step_code = json.dumps(step, indent=2)
                steps.append({
                    "id": str(idx),
                    "name": step.get("description", f"Step {idx+1}"),
                    "type": step.get("tool", "unknown"),
                    "code": step_code
                })
                
            return {
                "id": wf_id,
                "name": wf.get("name"),
                "description": wf.get("description"),
                "triggers": wf.get("triggers", []),
                "steps": steps,
            }
        return {"error": "Workflow not found"}
    except Exception as e:
        return {"error": str(e)}

# Auth cookie name (must match auth_routes) so WebSocket can read JWT when frontend doesn't pass ?token=
VAF_TOKEN_COOKIE = "vaf_token"




def _ws_client_ip(websocket) -> str:
    """Real client address for a WS handshake.

    X-Forwarded-For is honoured only when the immediate peer is loopback (i.e. the handshake was
    relayed by the integrated HTTPS proxy, which rebuilds the relayed headers itself). A direct
    non-loopback client therefore cannot forge a loopback identity to skip the network-mode token
    check, the 2FA gate or the local-admin fallback further down.
    """
    peer = websocket.client.host if websocket.client else "unknown"
    try:
        from vaf.network.binding import effective_client_ip
        return effective_client_ip(peer, websocket.headers.get("x-forwarded-for"))
    except Exception:
        return peer


def _emit_sec_ws(detail: str, ip: str = "") -> None:
    """Security-event mirror for rejected NETWORK WebSocket handshakes (lazy, never raises)."""
    try:
        from vaf.core.security_events import log_security_event
        log_security_event("ws_rejected", ip=ip, detail=detail)
    except Exception:
        pass

@app.websocket("/ws")
async def websocket_endpoint(websocket: WebSocket, token: Optional[str] = Query(None)):
    """
    WebSocket endpoint with optional token authentication.
    
    Token from: query param (?token=<jwt>) or cookie (vaf_token). Cookie is used when
    frontend connects without query (e.g. same-host login); ensures user_scope_id is
    set for RAG/memory so "No memories found" does not happen for logged-in users.
    
    When local_network_enabled=True, non-localhost requires valid JWT.
    When local_network_enabled=False, only localhost is allowed.
    """
    # Use cookie if frontend didn't pass token in URL (e.g. after login from same host)
    if not token and hasattr(websocket, "cookies") and websocket.cookies:
        token = websocket.cookies.get(VAF_TOKEN_COOKIE)
    # Prefer X-Forwarded-For (set by the integrated HTTPS proxy) so a LAN device shows its REAL IP in the
    # network map, not 127.0.0.1 (the proxy). Fall back to the direct peer for the desktop/localhost path.
    client_ip = _ws_client_ip(websocket)
    print(f"[WebSocket] Connection attempt from {client_ip}")
    log("API", f"WebSocket connection attempt from {client_ip}")
    
    # First: check if client is localhost or Docker internal network
    try:
        from vaf.network.binding import is_localhost, is_allowed_ip
        import ipaddress
        
        # Check if localhost OR Docker internal network
        is_localhost_client = is_localhost(client_ip)
        
        # In Docker mode, also treat Docker bridge network as "local"
        # This is safe because Docker network is internal to the host
        if not is_localhost_client:
            try:
                ip_obj = ipaddress.ip_address(client_ip)
                docker_network = ipaddress.ip_network("172.16.0.0/12")
                if ip_obj in docker_network:
                    is_localhost_client = True
                    log("API", f"WebSocket: Docker network IP {client_ip} treated as localhost")
            except ValueError:
                pass
                
    except ImportError:
        # Fallback localhost check (including Docker network range 172.16.0.0/12)
        import ipaddress as _ipaddr
        _is_docker = False
        try:
            _is_docker = _ipaddr.ip_address(client_ip) in _ipaddr.ip_network("172.16.0.0/12")
        except ValueError:
            pass
        is_localhost_client = (
            client_ip in ["127.0.0.1", "::1", "localhost"] or _is_docker
        )
    
    # Check local network setting
    local_network_enabled = Config.get("local_network_enabled", False)
    
    # --- CONNECTION TRACKING (Before Auth) ---
    # Register connection immediately to show on map, even if auth fails later
    connection_id = f"ws_{id(websocket)}"
    user_context = None # Will be populated if auth succeeds
    
    try:
        from vaf.network.connection_tracker import (
            get_tracker, ConnectionType, detect_device_type, DeviceType
        )
        tracker = get_tracker()
        
        # Get user agent
        user_agent = None
        if hasattr(websocket, 'headers'):
            user_agent = websocket.headers.get('user-agent', '')
        
        device_type = detect_device_type(user_agent)
        
        tracker.register_connection(
            connection_id=connection_id,
            connection_type=ConnectionType.WEBSOCKET,
            ip=client_ip,
            device_type=device_type,
            user_agent=user_agent,
            username="Guest (Connecting...)", # Temporary status
            metadata={
                "is_localhost": is_localhost_client,
                "status": "handshake"
            }
        )
        log("API", f"Connection tracked (pre-auth): {connection_id}")
    except Exception as e:
        log("API", f"Could not track connection: {e}")

    # If local network is DISABLED, only allow localhost and require auth token.
    if not local_network_enabled:
        if not is_localhost_client:
            log("API", f"WebSocket rejected: Local network disabled, non-localhost IP {client_ip}")
            _emit_sec_ws("local network disabled", ip=client_ip)
            # Update tracker if possible
            try: get_tracker().unregister_connection(connection_id)
            except: pass
            
            await websocket.close(code=4003, reason="Local network feature is disabled")
            return
        if not token:
            log("API", "WebSocket rejected: Localhost connection without token")
            await websocket.close(code=4001, reason="Authentication required")
            return
        try:
            from vaf.auth.crypto import get_jwt_secret
            import jwt
            secret = get_jwt_secret()
            payload = jwt.decode(token, secret, algorithms=["HS256"])
            user_context = {
                "user_id": payload.get("sub"),
                "user_scope_id": payload.get("user_scope_id"),
                "username": payload.get("username"),
                "role": payload.get("role"),
                "session_id": payload.get("session_id"),
            }
            log("API", f"WebSocket (localhost) authenticated: {user_context.get('username')} (scope: {user_context.get('user_scope_id')})")
        except Exception as e:
            log("API", f"WebSocket rejected: Invalid localhost token ({e})")
            await websocket.close(code=4001, reason="Invalid token")
            return
    else:
        # Local network is ENABLED - authenticate network users
        try:
            from vaf.network.binding import is_allowed_ip
            from vaf.auth.crypto import get_jwt_secret
            import jwt
            
            # Validate IP is from local network
            if not is_allowed_ip(client_ip):
                log("API", f"WebSocket rejected: Non-local IP {client_ip}")
                _emit_sec_ws("non-local ip", ip=client_ip)
                try: get_tracker().unregister_connection(connection_id)
                except: pass
                await websocket.close(code=4003, reason="Local network only")
                return
            
            # Require token authentication for all network-mode clients (including localhost).
            if not token:
                log("API", f"WebSocket rejected: No token from {client_ip}")
                _emit_sec_ws("no token", ip=client_ip)
                # Keep tracked as Guest/Unauth for map visibility?
                # Yes, but maybe mark as 'Auth Required'
                try: 
                    tracker = get_tracker()
                    tracker.register_connection(
                        connection_id=connection_id,
                        connection_type=ConnectionType.WEBSOCKET,
                        ip=client_ip,
                        device_type=device_type,
                        username="Unauthenticated Device",
                        metadata={"status": "auth_required"}
                    )
                except: pass
                
                await websocket.close(code=4001, reason="Authentication required")
                return
            
            if token:
                try:
                    secret = get_jwt_secret()
                    payload = jwt.decode(token, secret, algorithms=["HS256"])
                    
                    # Check 2FA status
                    require_2fa = Config.get("local_network_require_2fa", True)
                    if require_2fa and not payload.get("is_2fa_verified", False):
                        if not is_localhost_client or payload.get("role") != "admin":
                            log("API", f"WebSocket rejected: 2FA not verified for {payload.get('username')}")
                            await websocket.close(code=4003, reason="2FA verification required")
                            return
                    
                    user_context = {
                        "user_id": payload.get("sub"),
                        "user_scope_id": payload.get("user_scope_id"),
                        "username": payload.get("username"),
                        "role": payload.get("role"),
                        "session_id": payload.get("session_id"),
                    }
                    log("API", f"WebSocket authenticated: {user_context.get('username')}")
                    
                except jwt.ExpiredSignatureError:
                    log("API", f"WebSocket rejected: Expired token from {client_ip}")
                    _emit_sec_ws("expired token", ip=client_ip)
                    await websocket.close(code=4001, reason="Token expired")
                    return
                except jwt.InvalidTokenError:
                    log("API", f"WebSocket rejected: Invalid token from {client_ip}")
                    _emit_sec_ws("invalid token", ip=client_ip)
                    await websocket.close(code=4001, reason="Invalid token")
                    return
        except ImportError:
            # Auth modules unavailable in network mode: fail CLOSED for any non-localhost
            # client (we cannot verify the token). Localhost stays trusted (desktop flow).
            if not is_localhost_client:
                log("API", f"WebSocket rejected: auth modules unavailable, cannot verify {client_ip}")
                await websocket.close(code=4003, reason="Authentication error")
                return
        except Exception as e:
            log("API", f"WebSocket auth error: {e}")
            # Fail CLOSED: any auth-phase error for a non-localhost client must NOT fall
            # through to an unscoped (None scope == global/admin) connection. Localhost stays trusted.
            if not is_localhost_client:
                await websocket.close(code=4003, reason="Authentication error")
                return

    # Defense-in-depth: a non-localhost connection must NEVER be unscoped — a None scope
    # is treated as global/admin by session.list() and rag search, so this is the final
    # backstop against any future fall-through that left user_context unset.
    if user_context is None and not is_localhost_client:
        log("API", f"WebSocket rejected: no auth context for non-localhost {client_ip}")
        _emit_sec_ws("no auth context", ip=client_ip)
        try:
            get_tracker().unregister_connection(connection_id)
        except Exception:
            pass
        await websocket.close(code=4003, reason="Authentication error")
        return

    # Desktop fallback: a LOCALHOST connection that ended up without a usable user context
    # (e.g. the token was not passed on the WS handshake, or an auth-module hiccup) IS the local
    # admin - the same policy the HTTP routes use (get_current_user_or_local_admin). Without this
    # the desktop gets locked out of its OWN admin-owned sessions ("chat ... not owned by None").
    if is_localhost_client and not (user_context and (user_context.get("user_scope_id") or user_context.get("user_id"))):
        try:
            from vaf.core.config import get_local_admin_scope_id
            _admin = str(get_local_admin_scope_id())
            user_context = {
                "user_id": _admin,
                "user_scope_id": _admin,
                "username": (user_context or {}).get("username") or "admin",
                "role": "admin",
            }
            log("API", "WebSocket localhost without user context -> defaulting to local admin (desktop)")
        except Exception as _e:
            log("API", f"WebSocket local-admin fallback failed: {_e}")

    try:
        await manager.connect(websocket)
        # Store user_scope_id for RAG/task metadata (memory_save, memory search). Fallback to user_id if no scope in token.
        if user_context and (user_context.get("user_scope_id") or user_context.get("user_id")):
            manager.set_connection_user(
                websocket,
                user_context.get("user_scope_id") or user_context.get("user_id"),
                username=user_context.get("username"),
                role=user_context.get("role"),
            )
        os.environ["VAF_WEBUI_ACTIVE"] = "1"
        print(f"[WebSocket] Connected! Active: {len(manager.active_connections)}")
        log("API", f"WebSocket connected! Active: {len(manager.active_connections)}")
        tray_context.set_websocket_count(len(manager.active_connections)) # Update active count
        log("API", f"WebSocket count updated: {tray_context.active_websockets}")
        
        # Update connection tracker with verified info
        try:
            tracker = get_tracker()
            tracker.register_connection(
                connection_id=connection_id,
                connection_type=ConnectionType.WEBSOCKET,
                ip=client_ip,
                device_type=device_type,
                user_agent=user_agent,
                username=user_context.get("username") if user_context else "Guest",
                service_name="WebUI",
                metadata={
                    "role": user_context.get("role") if user_context else "guest",
                    "is_localhost": is_localhost_client,
                    "status": "connected"
                }
            )
            log("API", f"Connection updated: {connection_id}")
        except Exception as e:
            log("API", f"Could not track connection: {e}")

        # Count Web UI open as activity so thinking mode does not start immediately (idle timer starts after connect)
        try:
            from vaf.core.last_interaction import update_last_interaction
            scope_id = (user_context.get("user_scope_id") or user_context.get("user_id")) if user_context else None
            update_last_interaction(user_scope_id=scope_id, source="web", preview="")
        except Exception:
            pass
        
        try:
            provider = Config.get("provider", "local")
            await websocket.send_json({
                "type": "model_state",
                "loaded": tray_context.model_loaded,
                "persistent": tray_context.is_persistent(),
                "provider": provider
            })
            # Send user context if authenticated
            if user_context:
                await websocket.send_json({
                    "type": "auth_state",
                    "authenticated": True,
                    "username": user_context.get("username"),
                    "role": user_context.get("role"),
                })
        except Exception:
            pass
    except Exception as e:
        log("API", f"WebSocket handshake failed: {e}")
        raise e
    try:
        # Get user scope for filtering sessions
        user_scope_id = manager.get_connection_user(websocket)
        
        # Send initial session list (only Web UI chats; channel sessions appear in their dashboards)
        web_sessions = session_mgr.list_ui(limit=SESSION_LIST_LIMIT, user_scope_id=user_scope_id)
        await websocket.send_json({
            "type": "session_list",
            "sessions": session_list_payload(web_sessions)
        })
        # Send cached stats if available, otherwise send defaults
        stats_to_send = manager.last_stats
        if not stats_to_send:
            # Default stats until agent provides real values
            stats_to_send = {"used": 0, "total": 8192, "percent": 0.0, "api": False}
        await websocket.send_json({
            "type": "stats",
            "stats": stats_to_send
        })
        # Send tools list (cached or live) so UI has correct count
        try:
            from vaf.core.tool_contract import tool_category
            agent = manager.agent_instance
            if agent and hasattr(agent, "tools"):
                # The first list a client ever sees. It had NO filter at all, so a
                # restricted account was offered everything until the next refresh
                # replaced it - the shortest-lived leak is still a leak.
                _oc_scope = manager.get_connection_user(websocket)
                _oc_role = manager.get_connection_user_role(websocket)

                def _oc_allows(_name: str) -> bool:
                    try:
                        from vaf.core.tool_dispatch import account_allows_tool
                        return account_allows_tool(_name, _oc_scope, _oc_role)
                    except Exception:
                        return False
                tools_list = [
                    {
                        "name": name,
                        "description": getattr(tool, "description", "No description"),
                        "category": tool_category(name, tool)
                    }
                    for name, tool in agent.tools.items()
                    if _oc_allows(name)
                ]
            elif manager.tools_cache:
                tools_list = manager.tools_cache
            else:
                tools_list = _scan_tool_modules()
            await websocket.send_json({
                "type": "tools_list",
                "tools": tools_list or []
            })
        except Exception:
            pass
        # Auto-load latest Web UI session so WebUI gets a valid sessionId immediately (not a channel chat)
        if web_sessions:
            sid = web_sessions[0]["id"]
        else:
            # Create a new default session for the user if they have none
            new_sess = session_mgr.new(user_scope_id=user_scope_id)
            session_mgr.save(new_sess)
            sid = new_sess.id
            # Refresh the list for the client so they see the new session
            web_sessions = session_mgr.list_ui(limit=SESSION_LIST_LIMIT, user_scope_id=user_scope_id)
            await websocket.send_json({
                "type": "session_list",
                "sessions": session_list_payload(web_sessions)
            })

        try:
            # Subscribe this connection to the session for scoped updates
            manager.subscribe_to_session(websocket, sid)

            try:
                loaded = session_mgr.load(sid)
            except FileNotFoundError:
                # Should not happen for newly created session, but safety first
                loaded = Session(id=sid, name="New Chat")

            frontend_messages = _history_projection(loaded, sid)
            is_active = _session_is_active(sid)

            await websocket.send_json({
                "type": "history_update",
                "messages": frontend_messages,
                "sessionId": sid,
                "isActive": is_active,
                "currentStatus": _session_status_label(sid) if is_active else "idle"
            })
            artifact_payload = _build_artifact_payload(loaded, sid)
            if artifact_payload:
                await websocket.send_json(artifact_payload)
            # Restore sidebar documents from session on connect (so UI recovers after page refresh)
            try:
                _saved_sidebar = (getattr(loaded, "runtime_state", None) or {}).get("sidebar_documents") or []
                if _saved_sidebar:
                    _slim_sidebar = [{k: v for k, v in d.items() if k != "data"} for d in _saved_sidebar]
                    await websocket.send_json({
                        "type": "sidebar_documents_restored",
                        "contents": _slim_sidebar,
                        "sessionId": sid,
                    })
            except Exception:
                pass
        except Exception as e:
            log("WebServer", f"Auto-load session failed: {e}")

        while True:
            # Listen for client commands
            data_str = await websocket.receive_text()
            tray_context.register_websocket_activity()
            try:
                cmd = json.loads(data_str)
                # Several handlers below read the incoming message via `data`; keep it as an alias of
                # `cmd` so they work on a fresh connection. (Handlers that fetch over HTTP locally
                # rebind `data = resp.json()`, which is unaffected by this alias.)
                data = cmd
                type = cmd.get("type")
                
                # --- SESSION MANAGEMENT ---
                
                if type == "get_sessions":
                    # user_scope_id is required for correct RAG scope and filtered session listing
                    user_scope_id = manager.get_connection_user(websocket)
                    web_sessions = session_mgr.list_ui(limit=SESSION_LIST_LIMIT, user_scope_id=user_scope_id)
                    await websocket.send_json({
                        "type": "session_list",
                        "sessions": session_list_payload(web_sessions)
                    })                
                elif type == "room_task_history":
                    # The room's whole record, ON DEMAND and READ-ONLY.
                    #
                    # Its own branch rather than a sixth entry in the acting-commands
                    # block below: that block JOINS the person into the room when they
                    # are not a member yet, which is right for speaking and wrong for
                    # looking - opening a record must not make anybody a member of
                    # anything. Membership is decided by the same function that decided
                    # the sidebar, exactly as `open_room` does.
                    user_scope_id = manager.get_connection_user(websocket)
                    wanted = str(cmd.get("room_id") or "")
                    try:
                        from vaf.core.a2a.room import Room
                        from vaf.core.session import _room_rows

                        if wanted not in {row["room_id"] for row in _room_rows(user_scope_id)}:
                            await websocket.send_json({"type": "error",
                                                       "message": "Access denied"})
                            continue
                        room = Room.open(wanted)
                        history = [
                            {"id": t["id"], "title": t["title"], "status": t["status"],
                             "assignee": t["assignee_label"],
                             "requester": t["requester_label"],
                             "result": str(t.get("result") or "")[:400],
                             "reports": t["reports"],
                             "started": t["created_ts"], "ts": t["updated_ts"]}
                            for t in room.tasks()
                        ][:400]
                    except Exception as _hist_err:
                        log("API", f"room_task_history failed: {_hist_err}")
                        history = []
                    await websocket.send_json({"type": "room_task_history",
                                               "roomId": wanted, "tasks": history})

                elif type == "open_room":
                    # A room is NOT a session and must never reach the session loader:
                    # Session.save rewrites the whole message list, which with N peers
                    # is exactly the lost update the room store is built to avoid. It
                    # gets its own command, and the reply is read-only.
                    user_scope_id = manager.get_connection_user(websocket)
                    wanted = str(cmd.get("room_id") or "")
                    try:
                        from vaf.core.a2a.room import Room, describe
                        from vaf.core.session import _room_rows

                        # Membership is decided by the SAME function that decided the
                        # sidebar. A second predicate here would be a second answer to
                        # "is this user in that room", and the two would drift.
                        mine = {row["room_id"]: row for row in _room_rows(user_scope_id)}
                        row = mine.get(wanted)
                        if row is None:
                            log("API", f"Access denied: open_room for a room the user has not joined")
                            await websocket.send_json({"type": "error", "message": "Access denied"})
                            continue

                        await _send_room_transcript(websocket, room=Room.open(wanted),
                                                   user_scope_id=user_scope_id)
                    except Exception as e:
                        log("API", f"open_room failed: {e}", "ERROR")
                        await websocket.send_json({"type": "error", "message": "Room could not be opened"})

                elif type == "room_typing":
                    # A keypress in the room's input box. Presence, not
                    # conversation: it lands in an in-memory map the 3s room
                    # poll reads, never in the store and never on the wire (the
                    # protocol deliberately has no typing frame kind). Fire and
                    # forget - no reply, the poll carries the effect to the
                    # other viewers, and the timestamp expires on its own. The
                    # membership check is the same question room_say asks,
                    # answered the cheap way: a room the user is not in takes
                    # no presence from them either.
                    try:
                        user_scope_id = manager.get_connection_user(websocket)
                        wanted = str(cmd.get("room_id") or "")
                        from vaf.core.a2a.room import derive_peer_id, participant_key
                        from vaf.core.session import _room_rows
                        if wanted and any(row["room_id"] == wanted
                                          for row in _room_rows(user_scope_id)):
                            peer = derive_peer_id(
                                participant_key("cli", user_scope_id), wanted)
                            import time as _time
                            _ROOM_KEYS_TYPING.setdefault(wanted, {})[peer] = _time.time()
                    except Exception:
                        pass
                elif type in ("room_say", "close_room", "delete_room", "kick_peer",
                              "rename_room", "set_room_agent_mode",
                              "cast_room_vote"):
                    # The person at the browser acting in a room themselves, rather
                    # than their agent doing it for them. They act on the CLI lane and
                    # not on a lane of their own, deliberately: the lanes separate the
                    # HUMAN from the AGENT, and a browser and a terminal in front of
                    # the same person are the same actor. A "web" lane would derive a
                    # second handle and split one person into two members of the same
                    # room, so that whoever spoke last would look like somebody else.
                    user_scope_id = manager.get_connection_user(websocket)
                    wanted = str(cmd.get("room_id") or "")
                    try:
                        from vaf.core.a2a.room import (Room, RoomError, derive_peer_id,
                                                       participant_key)
                        from vaf.core.session import _room_rows

                        mine = {row["room_id"] for row in _room_rows(user_scope_id)}
                        if wanted not in mine:
                            log("API", f"Access denied: {type} for a room the user is not in")
                            await websocket.send_json({"type": "error", "message": "Access denied"})
                            continue

                        room = Room.open(wanted)
                        key = participant_key("cli", user_scope_id)
                        identity = room.identity_for(key, scope_id=user_scope_id)
                        if identity is None:
                            # The room reached the sidebar because the user's AGENT is
                            # in it. Speaking for themselves makes them a member too,
                            # and Room.join carries the tenant check that decides
                            # whether this account may be in this room at all.
                            display = str(cmd.get("display") or "").strip()
                            if not display:
                                # THIS account's name, not the machine owner's. The
                                # old fallback printed the local admin for whoever
                                # was speaking, so on an installation with several
                                # people the room learned one name for all of them -
                                # and the moment a surface says who belongs to whom,
                                # that becomes the room asserting, with its own
                                # authority, that the admin is somebody else's user.
                                # `resolve_caller_username` is the one definition of
                                # this question; the lookup is one database round
                                # trip, paid once when a person first speaks here.
                                from vaf.core.config import resolve_caller_username
                                display = str(resolve_caller_username(
                                    None, user_scope_id, allow_lookup=True) or "user")
                            identity = room.join(display=display, scope_id=user_scope_id,
                                                 peer_id=derive_peer_id(key, wanted))
                        # A member record still carrying the lane literal as its name is
                        # healed in _send_room_transcript, which every one of these
                        # commands answers with - a second heal here would be a second
                        # copy of the same rule waiting to drift.

                        if type == "delete_room":
                            # What the bin means everywhere else in this product. It
                            # closes first, so a peer reading the room over a wire is
                            # told WHY its access ended rather than finding a
                            # conversation that is simply not there any more.
                            room.delete(identity, reason=Room.TERMINATED_BY_USER)
                        elif type == "close_room":
                            room.close(identity, reason=Room.TERMINATED_BY_USER)
                        elif type == "kick_peer":
                            # The room decides whether this is allowed, including the
                            # rule that its own host handles can never be removed. A
                            # second check here would be a second answer.
                            room.kick(identity, str(cmd.get("peer") or ""),
                                      reason=str(cmd.get("reason") or "").strip())
                        elif type == "set_room_agent_mode":
                            # The mode is THE USER'S standing decision about their own
                            # agent (mode_of's contract), so the person setting it here
                            # is the authority the member file exists to record - and
                            # they may only ever set it for the agent of their OWN
                            # account: the handle is derived from the caller's scope,
                            # so another user's agent is unreachable by construction.
                            # Wirkung: the next room wake reads mode_of live, so
                            # "autonomous" is exactly "work while I sleep", granted
                            # per room and revocable the same way.
                            from vaf.core.a2a.room import ROOM_MODES
                            wanted_mode = str(cmd.get("mode") or "").strip()
                            if wanted_mode not in ROOM_MODES:
                                await websocket.send_json({
                                    "type": "error",
                                    "message": f"Unknown room mode '{wanted_mode}'"})
                                continue
                            agent_identity = room.identity_for(
                                participant_key("agent", user_scope_id))
                            if agent_identity is None:
                                await websocket.send_json({
                                    "type": "error",
                                    "message": "Your agent is not in this room"})
                                continue
                            room.set_mode(agent_identity, wanted_mode)
                        elif type == "cast_room_vote":
                            # The person votes as themselves - the cli lane, the same
                            # handle their messages carry - and voting again replaces
                            # their earlier ballot, exactly as it does on the CLI.
                            vote_id = str(cmd.get("vote_id") or "").strip()
                            choice = str(cmd.get("choice") or "").strip()
                            if not vote_id or not choice:
                                continue
                            try:
                                room.cast(identity, vote_id, choice)
                            except Exception as _vote_err:
                                await websocket.send_json({
                                    "type": "error",
                                    "message": f"Could not cast that vote: {_vote_err}"})
                                continue

                        elif type == "rename_room":
                            title = str(cmd.get("title") or "").strip()
                            if not title:
                                continue
                            # The topic lives in the manifest, not in a frame: it is a
                            # property of the room rather than something somebody said,
                            # and a renamed room must not look like a message nobody
                            # wrote. Only the host may change it.
                            if not room.is_host(identity):
                                await websocket.send_json({
                                    "type": "error",
                                    "message": "Only the room's host can rename it"})
                                continue
                            room.store.update_manifest(topic=title[:200])
                        else:
                            text = str(cmd.get("text") or "").strip()
                            if not text:
                                continue
                            # "@Name ..." is resolved by the ROOM, the one place that
                            # knows who is in it. A lookup here would be a second copy
                            # of the member table.
                            payload = {"kind": "say", "body": {"text": text}}
                            mention = room.address_from_mention(text)
                            if mention:
                                payload["to"] = mention
                            room.ingest(payload, identity=identity)

                        # Answer with the room as it now stands, so the view repaints
                        # from the store rather than from what the browser assumed. A
                        # deleted room has no "as it stands" - the sidebar push below
                        # is the whole answer, and the view closes itself when the row
                        # goes.
                        if type != "delete_room":
                            await _send_room_transcript(websocket, Room.open(wanted),
                                                        user_scope_id)

                        # Anything that changes what the SIDEBAR says has to push the
                        # sidebar. Closing takes the room out of the list and renaming
                        # changes its name, and neither reaches the browser through the
                        # transcript reply - so a closed room kept sitting there with
                        # its own farewell in it, which is precisely what it should
                        # never do. Broadcast rather than reply: the same person may
                        # have the app open twice.
                        if type in ("close_room", "delete_room", "rename_room"):
                            await manager.broadcast_to_user(user_scope_id, {
                                "type": "session_list",
                                "sessions": session_list_payload(
                                    session_mgr.list_ui(limit=SESSION_LIST_LIMIT,
                                                        user_scope_id=user_scope_id)),
                            })
                    except RoomError as e:
                        await websocket.send_json({"type": "error", "message": str(e)})
                    except Exception as e:
                        log("API", f"{type} failed: {e}", "ERROR")
                        await websocket.send_json({"type": "error", "message": "Room action failed"})

                elif type == "load_session":
                    sid = cmd.get("id")
                    user_scope_id = manager.get_connection_user(websocket)
                    try:
                        # 1. Load from disk + verify ownership before subscribing (strict, role-aware gate).
                        allowed, loaded = _ws_session_owner_ok(websocket, sid)
                        if not allowed:
                            log("API", f"Access denied: load_session {sid} does not belong to user {user_scope_id}")
                            await websocket.send_json({"type": "error", "message": "Access denied"})
                            continue
                        if loaded is None:
                            loaded = session_mgr.load(sid)

                        # Subscribe this connection to the session for scoped updates
                        manager.subscribe_to_session(websocket, sid)
                        
                        # Push command to main loop to switch session
                        from vaf.core.task_queue import TaskQueue
                        tq = TaskQueue()
                        tq.add(
                            session_id="system",
                            input_text=f"__CMD__:LOAD_SESSION:{sid}",
                            source="web",
                            metadata={"task_class": "background"},
                        )
                        
                        # Helper to clean historical messages (same logic as run.py)
                        
                        # 3. Send history to Frontend
                        # Serialize messages for JSON (include tool metadata so UI shows tool names, not "Unknown Tool")
                        frontend_messages = _history_projection(loaded, sid)
                        is_active = _session_is_active(sid)

                        await websocket.send_json({
                            "type": "history_update",
                            "messages": frontend_messages,
                            "sessionId": sid,
                            "isActive": is_active,
                            "currentStatus": _session_status_label(sid) if is_active else "idle"
                        })
                        artifact_payload = _build_artifact_payload(loaded, sid)
                        if artifact_payload:
                            await websocket.send_json(artifact_payload)

                        # Restore sidebar documents from session — without base64 data (too large).
                        # The frontend can re-show attached document names from the saved slim entries.
                        try:
                            _saved_sidebar = (getattr(loaded, "runtime_state", None) or {}).get("sidebar_documents") or []
                            if _saved_sidebar:
                                # Strip data field to keep the WS response small
                                _slim_sidebar = [{k: v for k, v in d.items() if k != "data"} for d in _saved_sidebar]
                                await websocket.send_json({
                                    "type": "sidebar_documents_restored",
                                    "contents": _slim_sidebar,
                                    "sessionId": sid,
                                })
                        except Exception:
                            pass

                        # 4. Send initial estimated stats (so UI is not empty)
                        try:
                            # Estimate based on loaded history
                            total_chars = sum(len(str(m.get("content", ""))) for m in frontend_messages)
                            est_tokens = total_chars // 3  # Rough estimate

                            # Get max context from config
                            cfg = Config.load()
                            max_ctx = cfg.get("n_ctx", 8192)
                            is_api = cfg.get("provider", "local") != "local"

                            if is_api and max_ctx <= 16384:
                                max_ctx = 128000

                            stats = {
                                "used": est_tokens,
                                "total": max_ctx,
                                "percent": round((est_tokens / max_ctx) * 100, 1) if max_ctx else 0.0,
                                "api": is_api
                            }
                            await websocket.send_json({
                                "type": "stats",
                                "stats": stats
                            })
                        except Exception as e:
                            print(f"[WebServer] Stats estimation error: {e}")

                        # 5. Send context_status with persisted user_turn_count so the
                        #    compaction progress bar is correct immediately after load/restart.
                        try:
                            from vaf.core.config import Config as _Cfg
                            _compaction_interval = int(_Cfg.get("memory_compaction_interval", 15))
                            _runtime = getattr(loaded, "runtime_state", None) or {}
                            _user_turn_count = _runtime.get("user_turn_count", 0)
                            if _user_turn_count == 0:
                                # Fallback: count from saved messages
                                _user_turn_count = sum(
                                    1 for m in loaded.messages
                                    if (m.get("role") if isinstance(m, dict) else getattr(m, "role", None)) == "user"
                                )
                            await websocket.send_json({
                                "type": "context_status",
                                "sessionId": sid,
                                "stats": {
                                    "user_turn_count": _user_turn_count,
                                    "compaction_interval": _compaction_interval,
                                    "compaction_progress": round((_user_turn_count % _compaction_interval) / _compaction_interval * 100) if _compaction_interval else 0,
                                    "message_count": len(loaded.messages),
                                }
                            })
                        except Exception as e:
                            print(f"[WebServer] context_status on load error: {e}")

                    except Exception as e:
                        import traceback
                        traceback.print_exc()
                        print(f"Load error: {e}")

                elif type == "delete_session":
                    sid = cmd.get("id")
                    if not sid:
                        continue
                    user_scope_id = manager.get_connection_user(websocket)
                    allowed, loaded = _ws_session_owner_ok(websocket, sid)
                    if not allowed:
                        log("API", f"Access denied: delete_session {sid} not owned by {user_scope_id}")
                        await websocket.send_json({"type": "error", "message": "Access denied"})
                        continue
                    # WHETHER TO ASK IS DECIDED HERE, not in the browser.
                    #
                    # The sidebar row is a cached copy and nothing refreshes it
                    # after a turn - no session_list is pushed when a chat gains
                    # messages - so a chat that filled up while the list stood
                    # still went on reporting zero, and the trash icon skipped its
                    # dialog for it. That is the one guess in the product that
                    # destroys something when it is wrong (live incident: a full
                    # conversation deleted with no dialog and no archive).
                    #
                    # `confirmed` is the browser reporting that a person read the
                    # dialog; without it the RECORD has to be untouched, judged by
                    # the framework's one definition. A record that cannot be read
                    # (an admin clearing a corrupt file) counts as not untouched:
                    # being unable to prove a chat is empty must never read as
                    # "it is empty". No title is logged here - a chat name is user
                    # text, and this branch has no encoding-safe printer.
                    if (not bool(cmd.get("confirmed"))
                            and _session_record_exists(sid)
                            and not (loaded is not None and loaded.is_untouched())):
                        await websocket.send_json({
                            "type": "session_delete_result",
                            "id": sid, "deleted": False, "needsConfirm": True,
                        })
                        continue
                    # Archive-then-delete, in that order and never both: archiving
                    # MOVES the file, so a delete afterwards would have nothing to
                    # remove - and if the move failed we must still delete rather
                    # than silently leave a chat the user asked to be gone.
                    gone = False
                    if bool(cmd.get("archive")):
                        if not session_mgr.archive(sid, user_scope_id=user_scope_id):
                            gone = session_mgr.delete(sid)
                        else:
                            gone = True          # the move IS the removal
                    else:
                        gone = session_mgr.delete(sid)
                    # The answer goes out BEFORE the list. The browser decides where
                    # to go next from the row it is about to lose - its title, and
                    # whether it was a thinking run - and the broadcast below takes
                    # that row out of its hands.
                    await websocket.send_json({
                        "type": "session_delete_result",
                        "id": sid, "deleted": bool(gone), "needsConfirm": False,
                    })
                    # Broadcast update ONLY to this user's connections
                    web_sessions = session_mgr.list_ui(limit=SESSION_LIST_LIMIT, user_scope_id=user_scope_id)
                    await manager.broadcast_to_user(user_scope_id, {
                        "type": "session_list",
                        "sessions": session_list_payload(web_sessions)
                    })

                elif type == "hide_session":
                    sid = cmd.get("id")
                    if not sid:
                        continue
                    user_scope_id = manager.get_connection_user(websocket)
                    allowed, loaded = _ws_session_owner_ok(websocket, sid)
                    if not allowed:
                        log("API", f"Access denied: hide_session {sid} not owned by {user_scope_id}")
                        await websocket.send_json({"type": "error", "message": "Access denied"})
                        continue
                    # A thinking run leaves the list the same way a chat leaves the
                    # store, so it is asked about the same way. Hiding keeps the
                    # record (the age-based sweep removes it later), but from the
                    # sidebar it is just as gone, and the row this decision used to
                    # be taken from is just as stale.
                    if (not bool(cmd.get("confirmed"))
                            and _session_record_exists(sid)
                            and not (loaded is not None and loaded.is_untouched())):
                        await websocket.send_json({
                            "type": "session_delete_result",
                            "id": sid, "deleted": False, "needsConfirm": True,
                        })
                        continue
                    hidden = bool(session_mgr.hide(sid))
                    await websocket.send_json({
                        "type": "session_delete_result",
                        "id": sid, "deleted": hidden, "needsConfirm": False,
                    })
                    if hidden:
                        web_sessions = session_mgr.list_ui(limit=SESSION_LIST_LIMIT, user_scope_id=user_scope_id)
                        await manager.broadcast_to_user(user_scope_id, {
                            "type": "session_list",
                            "sessions": session_list_payload(web_sessions)
                        })

                elif type == "new_session":
                    user_scope_id = manager.get_connection_user(websocket)
                    # Push command to main loop to create new session
                    from vaf.core.task_queue import TaskQueue
                    tq = TaskQueue()
                    tq.add(
                        session_id="system",
                        input_text="__CMD__:NEW_SESSION",
                        source="web",
                        metadata={"user_scope_id": user_scope_id, "task_class": "background"},
                    )
                    
                    # Create new session object AND SAVE IT IMMEDIATELY (temp, main loop will take over)
                    new_sess = session_mgr.new(user_scope_id=user_scope_id)
                    session_mgr.save(new_sess)
                    
                    # Subscribe this connection to the new session for scoped updates
                    manager.subscribe_to_session(websocket, new_sess.id)
                    
                    # Refresh list
                    web_sessions = session_mgr.list_ui(limit=SESSION_LIST_LIMIT, user_scope_id=user_scope_id)
                    await websocket.send_json({
                        "type": "session_list",
                        "sessions": session_list_payload(web_sessions)
                    })
                    # Clear frontend chat
                    await websocket.send_json({
                        "type": "history_update",
                        "messages": [],
                        "sessionId": new_sess.id
                    })

                elif type == "rename_session":
                    sid = cmd.get("id")
                    new_name = cmd.get("newName")
                    user_scope_id = manager.get_connection_user(websocket)
                    if sid and new_name:
                        allowed, _ = _ws_session_owner_ok(websocket, sid)
                        if not allowed:
                            log("API", f"Access denied: rename_session {sid} not owned by {user_scope_id}")
                            await websocket.send_json({"type": "error", "message": "Access denied"})
                            continue
                        session_mgr.rename(sid, new_name)
                        # Notify Main Loop to update in-memory object
                        from vaf.core.task_queue import TaskQueue
                        tq = TaskQueue()
                        tq.add(
                            session_id="system",
                            input_text=f"__CMD__:RENAME_SESSION:{sid}:{new_name}",
                            source="web",
                            metadata={"task_class": "background"},
                        )
                        
                        # Broadcast update ONLY to this user's connections
                        web_sessions = session_mgr.list_ui(limit=SESSION_LIST_LIMIT, user_scope_id=user_scope_id)
                        await manager.broadcast_to_user(user_scope_id, {
                            "type": "session_list",
                            "sessions": session_list_payload(web_sessions)
                        })

                elif type == "get_config":
                    # Send config to frontend; non-admins get scoped view (only their own connections)
                    from vaf.core.config import get_local_admin_scope_id
                    user_scope_id = manager.get_connection_user(websocket) if manager else None
                    stored_role = manager.get_connection_user_role(websocket) if manager else None
                    # Admin if stored role says "admin" OR scope matches local admin scope
                    local_admin_scope = get_local_admin_scope_id()
                    is_admin = stored_role == "admin" or (user_scope_id is not None and str(user_scope_id) == str(local_admin_scope))
                    role = "admin" if is_admin else "user"
                    full_cfg = Config.load()
                    cfg = Config.config_for_user(full_cfg, str(user_scope_id) if user_scope_id else None, role)
                    await websocket.send_json({
                        "type": "config_update",
                        "config": cfg
                    })

                elif type == "get_maintenance_status":
                    # Late-joiner snapshot: a client connecting mid-job has
                    # missed every pushed maintenance_progress frame.
                    from vaf.core.web_interface import maintenance_snapshot_frame
                    await websocket.send_json(maintenance_snapshot_frame())

                elif type == "get_models":
                    # Scan models directory for .gguf files
                    # Use absolute path based on project root (parent of 'vaf' package)
                    # core/web_server.py -> core/ -> vaf/ -> VAF/ -> VAF/models
                    project_root = Path(__file__).parent.parent.parent
                    models_dir = project_root / "models"
                    # print(f"[DEBUG] Looking for models in: {models_dir}")
                    models = []
                    if models_dir.exists():
                        from vaf.core.backend import is_selectable_model
                        models = [f.name for f in models_dir.glob("*.gguf")
                                  if is_selectable_model(f.name)]
                        # print(f"[DEBUG] Found models: {models}")
                    else:
                        # print(f"[DEBUG] Models directory not found at {models_dir}")
                        pass
                    
                    await websocket.send_json({
                        "type": "models_list",
                        "models": models
                    })

                elif type == "get_model_preview":
                    repo_id = (cmd.get("repo_id") or "").strip()
                    if not repo_id:
                        await websocket.send_json({
                            "type": "model_preview",
                            "repo_id": "",
                            "error": "repo_id is required (e.g. Nanbeige/Nanbeige4.1-3B)"
                        })
                    else:
                        try:
                            from huggingface_hub import HfApi
                            try:
                                from huggingface_hub import ModelCard
                            except ImportError:
                                ModelCard = None
                            api = HfApi()
                            model_info = api.model_info(repo_id=repo_id, files_metadata=True)
                            siblings = getattr(model_info, "siblings", [])
                            gguf_files = []
                            for f in siblings:
                                rfilename = getattr(f, "rfilename", f) if not isinstance(f, str) else f
                                if isinstance(rfilename, str) and rfilename.endswith(".gguf"):
                                    size_bytes = getattr(f, "size", None) or 0
                                    gguf_files.append({"filename": rfilename, "size_bytes": size_bytes})
                            gguf_files.sort(key=lambda x: x["size_bytes"])
                            card_content = None
                            if ModelCard is not None:
                                try:
                                    card = ModelCard.load(repo_id)
                                    card_content = getattr(card, "content", None) or getattr(card, "text", None) or ""
                                except Exception:
                                    pass
                            if not gguf_files:
                                try:
                                    all_files = api.list_repo_files(repo_id=repo_id)
                                    gguf_from_list = [p for p in all_files if isinstance(p, str) and p.endswith(".gguf")]
                                    if gguf_from_list:
                                        gguf_files = [{"filename": p, "size_bytes": 0} for p in gguf_from_list]
                                except Exception:
                                    pass
                                if not gguf_files:
                                    await websocket.send_json({
                                        "type": "model_preview",
                                        "repo_id": repo_id,
                                        "error": f"No GGUF files found in {repo_id}. This repo may be a base model (safetensors only). Try a GGUF repo instead, e.g. search on Hugging Face for \"<model name> GGUF\".",
                                        "gguf_files": []
                                    })
                                else:
                                    await websocket.send_json({
                                        "type": "model_preview",
                                        "repo_id": repo_id,
                                        "card_content": card_content,
                                        "gguf_files": gguf_files
                                    })
                            else:
                                await websocket.send_json({
                                    "type": "model_preview",
                                    "repo_id": repo_id,
                                    "card_content": card_content,
                                    "gguf_files": gguf_files
                                })
                        except Exception as e:
                            await websocket.send_json({
                                "type": "model_preview",
                                "repo_id": repo_id,
                                "error": str(e),
                                "gguf_files": []
                            })

                elif type == "download_model":
                    repo_id = (cmd.get("repo_id") or "").strip()
                    filename = (cmd.get("filename") or "").strip() or None
                    if not repo_id:
                        await websocket.send_json({
                            "type": "model_download_done",
                            "success": False,
                            "error": "repo_id is required (e.g. Nanbeige/Nanbeige4.1-3B)",
                            "models": []
                        })
                    else:
                        project_root = Path(__file__).parent.parent.parent
                        models_dir = project_root / "models"
                        progress_queue: queue.Queue = queue.Queue()
                        cancel_event = threading.Event()
                        ws_id = id(websocket)
                        _active_model_download_cancels[ws_id] = cancel_event

                        def make_progress_tqdm(pq: queue.Queue, ce: threading.Event):
                            class ProgressTqdm:
                                def __init__(self, total=None, **kwargs):
                                    self.n = 0
                                    self.total = total or 0
                                    self._pq = pq
                                    self._ce = ce
                                    self._start = time.time()
                                def update(self, n=1):
                                    self.n += n
                                    if self._ce.is_set():
                                        raise InterruptedError("Download cancelled")
                                    if self.total and self._pq is not None:
                                        pct = 100.0 * self.n / self.total
                                        elapsed = time.time() - self._start
                                        speed = self.n / elapsed if elapsed > 0 else 0
                                        if speed >= 1024 * 1024:
                                            speed_str = f"{speed / (1024 * 1024):.2f} MB/s"
                                        else:
                                            speed_str = f"{speed / 1024:.2f} KB/s"
                                        try:
                                            self._pq.put_nowait({
                                                "bytes_done": self.n, "bytes_total": self.total,
                                                "progress_pct": round(pct, 2), "speed_str": speed_str
                                            })
                                        except queue.Full:
                                            pass
                                def close(self): pass
                                def __enter__(self): return self
                                def __exit__(self, *a): return None
                            return ProgressTqdm

                        def _do_download() -> tuple[bool, str | None, list]:
                            try:
                                from huggingface_hub import HfApi, hf_hub_download
                                models_dir.mkdir(parents=True, exist_ok=True)
                                api = HfApi()
                                tqdm_class = make_progress_tqdm(progress_queue, cancel_event)
                                if not filename:
                                    try:
                                        model_info = api.model_info(repo_id=repo_id, files_metadata=True)
                                        siblings = getattr(model_info, "siblings", [])
                                    except Exception:
                                        siblings = []
                                        for f in api.list_repo_files(repo_id=repo_id):
                                            class _F:
                                                rfilename = f
                                                size = 0
                                            siblings.append(_F())
                                    gguf = [f for f in siblings if (getattr(f, "rfilename", f) if not isinstance(f, str) else f).endswith(".gguf")]
                                    if not gguf:
                                        return False, f"No GGUF files in {repo_id}", []
                                    gguf.sort(key=lambda x: getattr(x, "size", 0) or 0)
                                    chosen = gguf[0]
                                    fname = getattr(chosen, "rfilename", chosen) if not isinstance(chosen, str) else chosen
                                    hf_hub_download(repo_id=repo_id, filename=fname, local_dir=str(models_dir), tqdm_class=tqdm_class)
                                else:
                                    hf_hub_download(repo_id=repo_id, filename=filename, local_dir=str(models_dir), tqdm_class=tqdm_class)
                                from vaf.core.backend import is_selectable_model
                                new_models = [f.name for f in models_dir.glob("*.gguf")
                                              if is_selectable_model(f.name)]
                                return True, None, new_models
                            except InterruptedError:
                                return False, "Download cancelled", []
                            except Exception as e:
                                return False, str(e), []

                        async def run_download():
                            nonlocal progress_queue, ws_id
                            download_task = asyncio.create_task(asyncio.to_thread(_do_download))

                            async def drain_progress():
                                while not download_task.done():
                                    await asyncio.sleep(0.2)
                                    while True:
                                        try:
                                            msg = progress_queue.get_nowait()
                                            await websocket.send_json({
                                                "type": "model_download_progress",
                                                "progress_pct": msg.get("progress_pct"),
                                                "bytes_done": msg.get("bytes_done"),
                                                "bytes_total": msg.get("bytes_total"),
                                                "speed_str": msg.get("speed_str")
                                            })
                                        except queue.Empty:
                                            break
                                # Final drain
                                while True:
                                    try:
                                        msg = progress_queue.get_nowait()
                                        await websocket.send_json({
                                            "type": "model_download_progress",
                                            "progress_pct": msg.get("progress_pct"),
                                            "bytes_done": msg.get("bytes_done"),
                                            "bytes_total": msg.get("bytes_total"),
                                            "speed_str": msg.get("speed_str")
                                        })
                                    except queue.Empty:
                                        break

                            drain_task = asyncio.create_task(drain_progress())
                            try:
                                success, err, new_models = await download_task
                            except Exception as e:
                                success, err, new_models = False, str(e), []
                            finally:
                                _active_model_download_cancels.pop(ws_id, None)
                                drain_task.cancel()
                                try:
                                    await drain_task
                                except asyncio.CancelledError:
                                    pass
                            try:
                                await websocket.send_json({
                                    "type": "model_download_done",
                                    "success": success,
                                    "error": err,
                                    "models": new_models
                                })
                            except Exception:
                                pass

                        asyncio.create_task(run_download())

                elif type == "cancel_model_download":
                    ws_id = id(websocket)
                    if ws_id in _active_model_download_cancels:
                        _active_model_download_cancels[ws_id].set()

                elif type == "get_api_models":
                    # Fetch available models from API providers. The discovery
                    # URL + auth kind come from the provider registry (single
                    # source of truth, includes the config-overridable Veyllo
                    # base); the response FILTERING below stays provider-specific.
                    provider = cmd.get("provider", "openai")
                    api_key = cmd.get("api_key", "")
                    # The stored key, when the client sends none. The client USED to send
                    # one because the form field carried it - but since keys moved into the
                    # encrypted store the field is empty by design (the value never travels
                    # to the browser), so a payload key only exists while the user is
                    # typing a new one. Without this fallback, "fetch models" refused for
                    # every already-configured provider.
                    if not str(api_key or "").strip():
                        try:
                            from vaf.core.api_keys import resolve_api_key
                            api_key = resolve_api_key(provider)
                        except Exception:                    # noqa: BLE001 - unreadable store = no key here
                            api_key = ""
                    models = []

                    try:
                        from vaf.core.provider_registry import models_discovery

                        disc = models_discovery(provider)
                        # None = no remote listing for this provider (e.g. local):
                        # keep today's empty-result behavior.
                        if disc is not None and api_key:
                            url, auth = disc
                            headers = {}
                            params = {}
                            if auth == "bearer":
                                headers["Authorization"] = f"Bearer {api_key}"
                            elif auth == "x-api-key":
                                headers["X-Api-Key"] = api_key
                                headers["anthropic-version"] = "2023-06-01"
                            elif auth == "query-key":
                                params["key"] = api_key
                            if provider == "google":
                                params["pageSize"] = 1000
                            import httpx
                            async with httpx.AsyncClient() as client:
                                resp = await client.get(
                                    url,
                                    headers=headers,
                                    params=params,
                                    timeout=10.0
                                )
                                if resp.status_code == 200:
                                    data = resp.json()
                                    if provider == "openai":
                                        # Filter to chat models only
                                        models = sorted([
                                            m["id"] for m in data.get("data", [])
                                            if "gpt" in m["id"] or "o1" in m["id"] or "o3" in m["id"]
                                        ])
                                    elif provider == "google":
                                        raw_models = data.get("models", [])
                                        unique_models = []
                                        for model in raw_models:
                                            methods = model.get("supportedGenerationMethods", [])
                                            if "generateContent" not in methods:
                                                continue
                                            model_id = model.get("baseModelId") or model.get("name", "")
                                            if model_id.startswith("models/"):
                                                model_id = model_id.split("/", 1)[1]
                                            if model_id and model_id not in unique_models:
                                                unique_models.append(model_id)
                                        models = sorted(unique_models)
                                    elif provider == "deepseek":
                                        models = [m["id"] for m in data.get("data", [])]
                                    elif provider == "openrouter":
                                        models = [m["id"] for m in data.get("data", [])][:50]  # Limit
                                    else:
                                        # anthropic, veyllo and any future
                                        # OpenAI-compatible provider: plain id list.
                                        models = [m["id"] for m in data.get("data", []) if m.get("id")]
                                        if provider == "veyllo":
                                            # /v1/models also lists veyllo-transcribe (STT); keep chat only.
                                            from vaf.core.provider_registry import is_veyllo_chat_model
                                            models = [mid for mid in models if is_veyllo_chat_model(mid)]
                    except Exception as e:
                        log("WebServer", f"Failed to fetch models for {provider}: {e}")
                    
                    await websocket.send_json({
                        "type": "api_models_list",
                        "provider": provider,
                        "models": models
                    })

                elif type == "save_config":
                    new_config = cmd.get("config")
                    if new_config:
                        from vaf.core.config import get_local_admin_scope_id
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        stored_role = manager.get_connection_user_role(websocket) if manager else None
                        local_admin_scope = get_local_admin_scope_id()
                        is_admin = stored_role == "admin" or (user_scope_id is not None and str(user_scope_id) == str(local_admin_scope))
                        existing = Config.load()
                        if not is_admin:
                            new_filtered, scope_toggles = Config.extract_connection_toggles_for_scope(new_config, str(user_scope_id) if user_scope_id else None)
                            new_config = Config.filter_for_non_admin(new_filtered)
                            if scope_toggles:
                                by_scope = existing.get("connection_enabled_by_scope") or {}
                                if not isinstance(by_scope, dict):
                                    by_scope = {}
                                for scope_id, toggles in scope_toggles.items():
                                    by_scope[scope_id] = {**(by_scope.get(scope_id) or {}), **toggles}
                                existing["connection_enabled_by_scope"] = by_scope
                        # Keys leave the payload for the encrypted store; see
                        # api_keys.absorb_config_keys for why the write side had to move in the
                        # same change as the read side - otherwise a Settings save keeps writing
                        # to a file nobody asks any more.
                        from vaf.core.api_keys import absorb_config_keys
                        merged = Config.merge_preserving_nonempty_sensitive(
                            existing, absorb_config_keys(new_config))
                        Config.save(merged)
                        provider_changed = existing.get("provider") != merged.get("provider")

                        # Prefetch the dedicated voice GGUF when the user (admin) picks the
                        # local voice lane, so the recommended Gemma default is fetched at
                        # selection instead of only lazily at the first call. Download-only
                        # (no server swap - the live server may hold the MAIN model here);
                        # the helper no-ops when the file is already present. Non-admin saves
                        # cannot change voice_agent_* keys, so voice_changed stays False.
                        try:
                            voice_local = merged.get("voice_agent_provider") == "local"
                            voice_changed = (
                                existing.get("voice_agent_provider") != merged.get("voice_agent_provider")
                                or existing.get("voice_agent_model") != merged.get("voice_agent_model")
                            )
                            if voice_local and voice_changed:
                                from vaf.core import voice_model as _vvm
                                _vvm.ensure_voice_model_downloaded_async()
                        except Exception as e:
                            log("WebServer", f"Voice model prefetch failed: {e}")

                        try:
                            if "tray_autostart" in merged:
                                from vaf.core.platform import Platform
                                Platform.set_tray_autostart(bool(merged.get("tray_autostart")))
                        except Exception as e:
                            log("WebServer", f"Tray autostart update failed: {e}")

                        # Use TaskQueue for commands (headless_runner only reads from TaskQueue)
                        # Priority 1 so RELOAD_CONFIG is processed before any pending chat (priority 10)
                        from vaf.core.task_queue import TaskQueue
                        tq = TaskQueue()
                        tq.add(
                            session_id="system",
                            input_text="__CMD__:RELOAD_CONFIG",
                            source="web",
                            priority=1,
                            metadata={"task_class": "background"},
                        )
                        await websocket.send_json({
                            "type": "config_saved",
                            "status": "success",
                            "requires_refresh": provider_changed
                        })

                elif type == "get_autosuggest":
                    text = cmd.get("text", "")
                    if text:
                        # Use the internal _get_best_suggestion method
                        suggestion = get_autosuggest()._get_best_suggestion(text)
                        await websocket.send_json({
                            "type": "autosuggest_result",
                            "suggestion": suggestion
                        })
                
                elif type == "artifact_edit":
                    session_id = cmd.get("sessionId") or manager.get_session_for_connection(websocket)
                    user_scope_id = manager.get_connection_user(websocket)
                    if not session_id:
                        # Use user-scoped default to prevent crosstalk
                        safe_scope = str(user_scope_id or "default").replace("-", "")[:8]
                        session_id = f"web-default-{safe_scope}"
                    allowed, _ = _ws_session_owner_ok(websocket, session_id, allow_missing=True)
                    if not allowed:
                        log("API", f"Access denied: artifact_edit {session_id} not owned by {user_scope_id}")
                        await websocket.send_json({"type": "error", "message": "Access denied"})
                        continue
                    file = cmd.get("file", "")
                    code = cmd.get("code", "")
                    source = cmd.get("source", "web")
                    try:
                        from datetime import datetime
                        updated_at = datetime.now().isoformat()
                        loaded = session_mgr.load(session_id)
                        loaded.update_runtime_state("artifact", {
                            "file": file,
                            "code": code,
                            "updatedAt": updated_at,
                            "source": source
                        })
                        session_mgr.save(loaded, sync_state=False)
                        await manager.broadcast_to_session(session_id, {
                            "type": "artifact_update",
                            "file": file,
                            "code": code,
                            "updatedAt": updated_at,
                            "source": source
                        })
                    except Exception as e:
                        log("WebServer", f"Artifact update failed: {e}")

                elif type == "contact_reply_decision":
                    reply_id = cmd.get("replyId")
                    decision = cmd.get("decision")
                    if not reply_id or decision not in ("approve", "reject"):
                        await websocket.send_json({"type": "contact_reply_result", "ok": False, "error": "missing replyId or decision", "replyId": reply_id or ""})
                    else:
                        try:
                            from vaf.core.contact_reply_pending import get_and_remove
                            payload = get_and_remove(reply_id)
                            if not payload:
                                await websocket.send_json({"type": "contact_reply_result", "ok": False, "error": "expired", "replyId": reply_id})
                            elif decision == "approve":
                                src = payload.get("source")
                                text = payload.get("text") or "[No reply text]"
                                if src == "telegram":
                                    from vaf.core.telegram_reply import send_telegram_reply
                                    send_telegram_reply(str(payload["chat_id_or_jid"]), text)
                                    await websocket.send_json({"type": "contact_reply_result", "ok": True, "decision": "approve", "replyId": reply_id})
                                elif src == "whatsapp":
                                    from vaf.core.whatsapp_reply import send_whatsapp_reply
                                    send_whatsapp_reply(payload.get("username") or "admin", str(payload["chat_id_or_jid"]), text)
                                    await websocket.send_json({"type": "contact_reply_result", "ok": True, "decision": "approve", "replyId": reply_id})
                                else:
                                    await websocket.send_json({"type": "contact_reply_result", "ok": False, "error": "unknown source", "replyId": reply_id})
                            else:
                                await websocket.send_json({"type": "contact_reply_result", "ok": True, "decision": "reject", "replyId": reply_id})
                        except Exception as e:
                            log("WebServer", f"contact_reply_decision failed: {e}")
                            await websocket.send_json({"type": "contact_reply_result", "ok": False, "error": str(e)[:200], "replyId": reply_id})

                elif type == "speaker_confirm_reply":
                    # Web-card answer to a speaker confirmation (yes / no /
                    # no+name). Ownership: the reply only applies to the
                    # PENDING RECORD OF THIS CONNECTION'S OWN SCOPE - a name
                    # answer writes into the voice DB, so this is fail-closed
                    # on the connection identity, never on client-sent ids.
                    try:
                        from vaf.core import speaker_confirm as _scr
                        from vaf.core.config import get_local_admin_scope_id as _glas
                        _scope_r = str(manager.get_connection_user(websocket) or _glas())
                        _answer = cmd.get("answer")
                        _name_r = (cmd.get("name") or "").strip() or None
                        if _answer not in ("yes", "no"):
                            await websocket.send_json({
                                "type": "speaker_confirm_result", "ok": False,
                                "error": "bad answer", "confirmId": cmd.get("confirmId") or "",
                            })
                        else:
                            loop = asyncio.get_running_loop()
                            _res = await loop.run_in_executor(
                                None, lambda: _scr.resolve(
                                    _scope_r, _answer,
                                    _name_r if _answer == "no" else None,
                                    confirm_id=cmd.get("confirmId") or None))
                            _res_payload = {
                                "type": "speaker_confirm_result",
                                "ok": _res.get("ok", False),
                                "outcome": _res.get("outcome", ""),
                                "ack": _res.get("ack", ""),
                                "confirmId": cmd.get("confirmId") or "",
                            }
                            await websocket.send_json(_res_payload)
                    except Exception as e:
                        log("WebServer", f"speaker_confirm_reply failed: {e}")
                        await websocket.send_json({
                            "type": "speaker_confirm_result", "ok": False,
                            "error": str(e)[:200], "confirmId": cmd.get("confirmId") or "",
                        })

                elif type == "gate_response":
                    # User clicked Allow Once / Allow Always / Cancel in the trust gate dialog.
                    decision = cmd.get("decision")  # "allow_once" | "allow_always" | "cancel"
                    if decision in ("allow_once", "allow_always", "cancel"):
                        _gate_session = manager.get_session_for_connection(websocket)
                        from vaf.core.web_interface import get_web_interface as _gwi
                        _gwi().resolve_gate(_gate_session or "", decision)

                elif type == "chat":
                    content = cmd.get("content")
                    files = cmd.get("files", [])  # List of file objects with {name, data, mimeType}
                    sidebar_docs_payload = cmd.get("sidebarDocuments") or []  # Document Viewer docs to inject into this turn

                    if content or files:
                        tray_context.register_activity()
                        # Learn from user input
                        if content:
                            get_autosuggest().learn(content)

                        # Separate image files from document/text files
                        image_files = [f for f in files if (f.get("mimeType") or "").startswith("image/")]
                        text_files  = [f for f in files if not (f.get("mimeType") or "").startswith("image/")]

                        # Process text/document files (extract content as before)
                        if text_files:
                            print(f"[WebUI] Processing {len(text_files)} attached file(s)...")
                            file_contents = await process_uploaded_files(text_files)
                            if file_contents:
                                content = content + "\n\n" + file_contents if content else file_contents

                        # Build raw image list for vision providers (base64, no data-URI prefix)
                        attached_images = []
                        if image_files:
                            import base64 as _b64
                            print(f"[WebUI] Passing {len(image_files)} image(s) to vision pipeline...")
                            for img_f in image_files:
                                raw = img_f.get("data", "")
                                if raw.startswith("data:"):
                                    raw = raw.split(",", 1)[1] if "," in raw else raw
                                attached_images.append({
                                    "data": raw,
                                    "mime_type": img_f.get("mimeType", "image/jpeg"),
                                    "name": img_f.get("name", "image"),
                                })
                        
                        # Get session ID: prefer explicit message field, then safe connection session, then fallback.
                        requested_session_id = cmd.get("sessionId")
                        connection_session_id = manager.get_session_for_connection(websocket)
                        session_id = requested_session_id or connection_session_id
                        user_scope_id = manager.get_connection_user(websocket)
                        if (
                            not requested_session_id
                            and isinstance(connection_session_id, str)
                            and connection_session_id.startswith(("telegram_", "discord_", "whatsapp_"))
                        ):
                            # Defensive guard: never route WebUI chat into channel sessions implicitly.
                            session_id = None
                        if not session_id:
                            # Use user-scoped default to prevent crosstalk
                            safe_scope = str(user_scope_id or "default").replace("-", "")[:8]
                            session_id = f"web-default-{safe_scope}"

                        # Ownership gate: never subscribe to / enqueue into a session owned by another user
                        # (subscribe itself leaks the live stream). allow_missing lets the first message into
                        # a brand-new session through.
                        _chat_allowed, _ = _ws_session_owner_ok(websocket, session_id, allow_missing=True)
                        if not _chat_allowed:
                            log("API", f"Access denied: chat into {session_id} not owned by {user_scope_id}")
                            await websocket.send_json({"type": "error", "message": "Access denied"})
                            continue

                        # Ensure this connection is subscribed to the session we're queueing for,
                        # so streaming (agent_message_update) reaches this client (fixes non-admin/LAN).
                        manager.subscribe_to_session(websocket, session_id)

                        # Editor document: prepend to this turn so agent has current editor content (like Document Viewer)
                        editor_doc = cmd.get("editorDocument")
                        if editor_doc and isinstance(editor_doc, dict):
                            name = editor_doc.get("name") or "Document"
                            ed_content = editor_doc.get("content") or ""
                            if ed_content:
                                block = f"--- CURRENT DOCUMENT (Editor): {name} ---\n{ed_content}\n----------------\n\n"
                                content = (block + content) if content else block
                            try:
                                loaded = session_mgr.load(session_id)
                            except FileNotFoundError:
                                loaded = Session(
                                    id=session_id,
                                    name=f"Session {session_id}",
                                    runtime_state={"editor_document": editor_doc},
                                )
                                session_mgr.save(loaded, sync_state=False)
                            else:
                                if not getattr(loaded, "runtime_state", None):
                                    loaded.runtime_state = {}
                                loaded.runtime_state["editor_document"] = editor_doc
                                session_mgr.save(loaded, sync_state=False)
                        else:
                            try:
                                loaded = session_mgr.load(session_id)
                                if getattr(loaded, "runtime_state", None) and "editor_document" in loaded.runtime_state:
                                    loaded.runtime_state["editor_document"] = {}
                                    session_mgr.save(loaded, sync_state=False)
                            except Exception:
                                pass

                        # Code viewer file: store in runtime_state so headless_runner injects it per-turn.
                        # Do NOT prepend to content — that would store the full file in message history.
                        code_viewer_file = cmd.get("codeViewerFile")
                        if code_viewer_file and isinstance(code_viewer_file, dict) and code_viewer_file.get("content"):
                            try:
                                loaded = session_mgr.load(session_id)
                            except FileNotFoundError:
                                loaded = Session(id=session_id, name=f"Session {session_id}", runtime_state={})
                                session_mgr.save(loaded, sync_state=False)
                            if not getattr(loaded, "runtime_state", None):
                                loaded.runtime_state = {}
                            loaded.runtime_state["code_viewer_file"] = code_viewer_file
                            session_mgr.save(loaded, sync_state=False)
                        else:
                            try:
                                loaded = session_mgr.load(session_id)
                                if getattr(loaded, "runtime_state", None) and "code_viewer_file" in loaded.runtime_state:
                                    del loaded.runtime_state["code_viewer_file"]
                                    session_mgr.save(loaded, sync_state=False)
                            except Exception:
                                pass

                        # Which sub-agent window this chat has open, if any. The browser's
                        # equivalent below is read from its LEASE and needs nothing from the
                        # client; these windows exist only in the tab, so the client states
                        # the fact and this scopes it to the message's own session (already
                        # ownership-checked on this path). Validated against the known kinds
                        # rather than stored raw: it ends up in a prompt block, and an
                        # unchecked client string in a prompt is an injection lane.
                        _saw_kind = cmd.get("subAgentWindow")
                        _known_kinds = ("coder", "research", "document", "librarian")
                        try:
                            loaded = session_mgr.load(session_id)
                        except FileNotFoundError:
                            loaded = None
                        except Exception:
                            loaded = None
                        if loaded is not None:
                            if not getattr(loaded, "runtime_state", None):
                                loaded.runtime_state = {}
                            if isinstance(_saw_kind, str) and _saw_kind in _known_kinds:
                                _sw_dirty = loaded.runtime_state.get("subagent_window") != _saw_kind
                                loaded.runtime_state["subagent_window"] = _saw_kind
                                # What the person is doing inside the window (the idle
                                # librarian's folder). Client-sent prose that ends up in a
                                # prompt, so it is flattened and capped rather than trusted:
                                # newlines/control chars would let it impersonate the
                                # block's own structure.
                                _saw_detail = cmd.get("subAgentWindowDetail")
                                if isinstance(_saw_detail, str) and _saw_detail.strip():
                                    _clean = " ".join(_saw_detail.split())[:200]
                                    if loaded.runtime_state.get("subagent_window_detail") != _clean:
                                        loaded.runtime_state["subagent_window_detail"] = _clean
                                        _sw_dirty = True
                                elif "subagent_window_detail" in loaded.runtime_state:
                                    del loaded.runtime_state["subagent_window_detail"]
                                    _sw_dirty = True
                                if _sw_dirty:
                                    session_mgr.save(loaded, sync_state=False)
                            elif "subagent_window" in loaded.runtime_state:
                                # The delete branch IS the feature: close the window and the
                                # block stops appearing, without any close event of its own.
                                del loaded.runtime_state["subagent_window"]
                                loaded.runtime_state.pop("subagent_window_detail", None)
                                session_mgr.save(loaded, sync_state=False)

                        # Interactive browser context: while THIS chat's window drives the
                        # sandbox browser, the current page rides along with the message -
                        # URL + selected text into runtime_state (the code-viewer pattern,
                        # injected per turn by headless_runner and cleared there), and a
                        # screenshot through the EXISTING attached-images lane, so it shows
                        # under the user's message like any upload and reaches the model as
                        # vision input. Captured server-side at send time: fresher than
                        # anything the client could claim, and nothing to trust.
                        try:
                            # Whichever browser (shared or pooled) this chat's
                            # lease lives on answers; no lease, no context.
                            from vaf.core.browser_interactive import manager_for_session as _bi_for
                            _bi = _bi_for(str(session_id))
                            if _bi is not None:
                                _snap = await asyncio.get_running_loop().run_in_executor(
                                    None, lambda: _bi.snapshot_context(str(session_id)))
                                if _snap and _snap.get("url"):
                                    try:
                                        loaded = session_mgr.load(session_id)
                                    except FileNotFoundError:
                                        loaded = Session(id=session_id, name=f"Session {session_id}", runtime_state={})
                                    if not getattr(loaded, "runtime_state", None):
                                        loaded.runtime_state = {}
                                    loaded.runtime_state["browser_context"] = {
                                        "url": _snap["url"], "selection": _snap.get("selection", "")}
                                    session_mgr.save(loaded, sync_state=False)
                                    if _snap.get("screenshot_b64"):
                                        attached_images.append({
                                            "data": _snap["screenshot_b64"],
                                            "mime_type": "image/jpeg",
                                            "name": "browser-view.jpg",
                                        })
                            else:
                                try:
                                    loaded = session_mgr.load(session_id)
                                    if getattr(loaded, "runtime_state", None) and "browser_context" in loaded.runtime_state:
                                        del loaded.runtime_state["browser_context"]
                                        session_mgr.save(loaded, sync_state=False)
                                except Exception:
                                    pass
                        except Exception:
                            pass

                        # Image viewer context: while the Image Viewer is open, the frontend sends the
                        # open image's vision description each turn. Store it in runtime_state so
                        # headless_runner injects it per-turn (kept in context only while the viewer is
                        # open), exactly like code_viewer_file.
                        image_viewer_ctx = cmd.get("imageViewerContext")
                        if image_viewer_ctx and isinstance(image_viewer_ctx, dict) and image_viewer_ctx.get("description"):
                            try:
                                loaded = session_mgr.load(session_id)
                            except FileNotFoundError:
                                loaded = Session(id=session_id, name=f"Session {session_id}", runtime_state={})
                                session_mgr.save(loaded, sync_state=False)
                            if not getattr(loaded, "runtime_state", None):
                                loaded.runtime_state = {}
                            loaded.runtime_state["image_viewer_context"] = image_viewer_ctx
                            session_mgr.save(loaded, sync_state=False)
                        else:
                            try:
                                loaded = session_mgr.load(session_id)
                                if getattr(loaded, "runtime_state", None) and "image_viewer_context" in loaded.runtime_state:
                                    del loaded.runtime_state["image_viewer_context"]
                                    session_mgr.save(loaded, sync_state=False)
                            except Exception:
                                pass

                        # Marked region: the user drew a yellow box on an image and is asking about it.
                        # Run vision ONCE on the annotated full image + zoomed crop with the user's
                        # question, and stash the answer text so headless injects it this turn (one-shot,
                        # like image_viewer_context). Vision is offloaded so the event loop stays free.
                        marked_region = cmd.get("markedRegion")
                        if marked_region and isinstance(marked_region, dict) and (marked_region.get("annotated") or marked_region.get("crop")):
                            try:
                                import asyncio as _asyncio
                                from vaf.core.vision_infer import vision_infer as _vinfer
                                _rname = marked_region.get("name") or "image"
                                _mr_imgs = []
                                for _k in ("annotated", "crop"):
                                    _d = marked_region.get(_k) or ""
                                    if not _d:
                                        continue
                                    _mime = "image/png"
                                    if _d.startswith("data:"):
                                        if ";" in _d:
                                            _mime = _d[5:].split(";", 1)[0] or "image/png"
                                        _d = _d.split(",", 1)[1] if "," in _d else _d
                                    _mr_imgs.append({"data": _d, "mime_type": _mime, "name": _rname})
                                _q = (content or "").strip() or "What is in the marked region?"
                                _prompt = (
                                    f"The user marked a region of the image '{_rname}' with a YELLOW rectangle. "
                                    f"The first image is the full image with the box; the second is a zoomed crop of it. "
                                    f"Answer specifically about what is INSIDE the marked region:\n{_q}"
                                )
                                _ans = await _asyncio.to_thread(_vinfer, _mr_imgs, _prompt)
                                try:
                                    loaded = session_mgr.load(session_id)
                                except FileNotFoundError:
                                    # First message into a fresh session — create it so the
                                    # (already-paid) vision answer isn't silently dropped.
                                    loaded = Session(id=session_id, name=f"Session {session_id}", runtime_state={})
                                if not getattr(loaded, "runtime_state", None):
                                    loaded.runtime_state = {}
                                if _ans:
                                    loaded.runtime_state["marked_region_answer"] = {"name": _rname, "text": _ans}
                                else:
                                    loaded.runtime_state.pop("marked_region_answer", None)
                                session_mgr.save(loaded, sync_state=False)
                            except Exception as _mre:
                                log("WebServer", f"markedRegion vision failed: {_mre}")
                                # Don't leave a stale answer from a previous turn if this one failed.
                                try:
                                    _l = session_mgr.load(session_id)
                                    if getattr(_l, "runtime_state", None) and "marked_region_answer" in _l.runtime_state:
                                        del _l.runtime_state["marked_region_answer"]
                                        session_mgr.save(_l, sync_state=False)
                                except Exception:
                                    pass
                        else:
                            try:
                                loaded = session_mgr.load(session_id)
                                if getattr(loaded, "runtime_state", None) and "marked_region_answer" in loaded.runtime_state:
                                    del loaded.runtime_state["marked_region_answer"]
                                    session_mgr.save(loaded, sync_state=False)
                            except Exception:
                                pass

                        # Editor selections: store for replace_editor_selection tool (start/end per index)
                        editor_selections = cmd.get("editorSelections")
                        if isinstance(editor_selections, list) and editor_selections:
                            try:
                                loaded = session_mgr.load(session_id)
                            except FileNotFoundError:
                                loaded = Session(
                                    id=session_id,
                                    name=f"Session {session_id}",
                                    runtime_state={"editor_selections": editor_selections},
                                )
                                session_mgr.save(loaded, sync_state=False)
                            else:
                                if not getattr(loaded, "runtime_state", None):
                                    loaded.runtime_state = {}
                                loaded.runtime_state["editor_selections"] = editor_selections
                                session_mgr.save(loaded, sync_state=False)
                        else:
                            try:
                                loaded = session_mgr.load(session_id)
                                if getattr(loaded, "runtime_state", None) and "editor_selections" in loaded.runtime_state:
                                    loaded.runtime_state["editor_selections"] = []
                                    session_mgr.save(loaded, sync_state=False)
                            except Exception:
                                pass

                        # Ensure sidebar documents are in session before queueing (so headless has context).
                        # NOTE: the frontend now sends only name+mimeType (no base64) so we skip
                        # process_files_to_sidebar_list unless a payload entry actually has file data.
                        # Documents are already in the session from the earlier set_sidebar_documents WS call.
                        if sidebar_docs_payload:
                            try:
                                has_data = any(doc.get("data") for doc in sidebar_docs_payload)
                                if has_data:
                                    contents = await process_files_to_sidebar_list(
                                        sidebar_docs_payload, session_id=session_id,
                                        user_scope_id=user_scope_id)
                                    if contents:
                                        _slim = [{k: v for k, v in doc.items() if k != "data"} for doc in contents]
                                        try:
                                            loaded = session_mgr.load(session_id)
                                        except FileNotFoundError:
                                            loaded = Session(
                                                id=session_id,
                                                name=f"Session {session_id}",
                                                runtime_state={"sidebar_documents": _slim},
                                            )
                                            session_mgr.save(loaded, sync_state=False)
                                        else:
                                            if not getattr(loaded, "runtime_state", None):
                                                loaded.runtime_state = {}
                                            loaded.runtime_state["sidebar_documents"] = _slim
                                            session_mgr.save(loaded, sync_state=False)
                                        log("WebServer", f"Injected {len(contents)} sidebar doc(s) for session {session_id} before chat")
                                        if bool(Config.get("attachment_rag_enabled", False)):
                                            try:
                                                from vaf.memory.attachment_rag import index_session_attachments_async
                                                await index_session_attachments_async(
                                                    session_id=session_id,
                                                    user_scope_id=user_scope_id,
                                                    documents=contents,
                                                )
                                            except Exception as e:
                                                log("WebServer", f"attachment index (chat inject) failed: {e}")
                                # else: no file data — session already has documents from set_sidebar_documents
                            except Exception as e:
                                log("WebServer", f"sidebar_documents on chat failed: {e}")

                        # Use TaskQueue for serialized execution
                        from vaf.core.task_queue import TaskQueue
                        tq = TaskQueue()
                        # user_scope_id is required for correct RAG scope (Auto-Recall and memory_save)
                        user_scope_id = manager.get_connection_user(websocket)
                        username = manager.get_connection_username(websocket)
                        user_role = manager.get_connection_user_role(websocket)
                        metadata = {}
                        if user_scope_id:
                            metadata["user_scope_id"] = user_scope_id
                        if username:
                            metadata["username"] = username
                        if user_role:
                            metadata["role"] = user_role
                        metadata["origin_channel"] = "web"
                        metadata["enqueue_session_id"] = str(session_id)
                        if user_scope_id:
                            metadata["enqueue_user_scope_id"] = str(user_scope_id)
                        metadata["task_class"] = "interactive"
                        log("WebServer", f"Chat message from user_scope_id={user_scope_id}, username={username}")
                        # Mark user activity for thinking mode (idle detection)
                        try:
                            from vaf.core.last_interaction import update_last_interaction
                            update_last_interaction(
                                user_scope_id=user_scope_id,
                                source="web",
                                preview=(content or "")[:80],
                            )
                        except Exception:
                            pass
                        # NOTE: do NOT clear waiting_for_reply here. If the background run asked a
                        # question and the user is now replying, the MAIN agent's chat_step needs the
                        # waiting state to reconstruct the question context (_thinking_reply_context) and
                        # to advance the tracked request (asked -> confirmed/done) + mark the source
                        # note/todo handled. It reads the context and then clears the waiting state itself
                        # (see Agent.chat_step). Clearing it here first destroyed that handoff — the agent
                        # found nothing and answered without the question's context (observed 2026-06-22).
                        if attached_images:
                            # Persist images as files in the user-siloed chat folder; entries
                            # then reference them by path (lean session.json, real file location).
                            attached_images = _persist_attached_images_to_files(
                                attached_images, session_id, user_scope_id
                            )
                            metadata["images"] = attached_images
                        # The task is queued BEFORE the upload is announced, because the
                        # announcement needs the turn's identity to say which message the
                        # image belongs to. Announcing first is what made an uploaded
                        # image appear under the PREVIOUS answer: at that moment the turn
                        # it belongs to did not exist yet, and the browser could only
                        # attach it to the newest answer it already had.
                        _task = tq.add(session_id=session_id, input_text=content, source="web",
                                       metadata=metadata)
                        if attached_images:
                            # Surface the stored file(s) to the Web UI so the chat's workspace box
                            # appears (same notify the agent's file-creating tools use). Uploading an
                            # image now creates a real file in the chat folder, so it should show up
                            # there just like agent-created files.
                            try:
                                from vaf.core.web_interface import notify_file_created
                                for _ai in attached_images:
                                    if _ai.get("path"):
                                        notify_file_created(session_id, _ai["path"], title=_ai.get("name"),
                                                            turn_id=getattr(_task, "turn_id", None))
                            except Exception as _nfe:
                                log("WebServer", f"notify_file_created (image upload) failed: {_nfe}")
                        try:
                            if is_debug_logging_enabled():
                                from datetime import datetime as _dt
                                qpath = get_dated_log_path("queue", "log")
                                qpath.parent.mkdir(parents=True, exist_ok=True)
                                qsize = tq.get_queue_size()
                                with open(qpath, "a", encoding="utf-8") as f:
                                    f.write(f"{_dt.now().isoformat()} QUEUE_ADD session_id={session_id} preview={repr((content or '')[:60])} queue_size_after={qsize}\n")
                        except Exception:
                            pass
                        # Ack to console only; do not push to chat UI (avoids duplicating user message as system "Queued input...")
                        file_info = f" [{len(files)} file(s)]" if files else ""
                        print(f"[WebUI] Queued input{file_info} for session {session_id}: {content[:50]}...")

                elif type == "learn_document_start":
                    # The learn button: start the batched learn job for ONE
                    # persisted chat attachment. The click IS the confirmation.
                    session_id = cmd.get("sessionId") or manager.get_session_for_connection(websocket)
                    user_scope_id = manager.get_connection_user(websocket)
                    doc_name = str(cmd.get("name") or "").strip()

                    def _deny(reason: str):
                        return websocket.send_json({
                            "type": "learn_denied", "sessionId": session_id,
                            "name": doc_name, "reason": reason,
                        })

                    if not session_id or not doc_name:
                        await _deny("missing_session_or_name")
                        continue
                    allowed, _loaded_sess = _ws_session_owner_ok(websocket, session_id)
                    if not allowed:
                        log("API", f"Access denied: learn_document_start {session_id} not owned by {user_scope_id}")
                        await _deny("not_owned")
                        continue
                    # The document must carry a REAL persisted path inside the
                    # session's attachments dir - never a client-supplied path.
                    _doc_path = ""
                    try:
                        _docs = ((getattr(_loaded_sess, "runtime_state", None) or {})
                                 .get("sidebar_documents") or [])
                        for _d in _docs:
                            if str((_d or {}).get("name") or "") == doc_name:
                                _doc_path = str((_d or {}).get("path") or "")
                                break
                    except Exception:
                        _doc_path = ""
                    import os as _os
                    from vaf.core.session import get_session_attachments_dir as _gsad
                    try:
                        _adir = str(_gsad(session_id, user_scope_id, create=False))
                    except Exception:
                        _adir = ""
                    if (not _doc_path or not _os.path.isfile(_doc_path)
                            or not _adir
                            or not _os.path.realpath(_doc_path).startswith(_os.path.realpath(_adir))):
                        await _deny("no_persisted_file")
                        continue
                    from vaf.core.subagent_ipc import get_ipc as _get_ipc
                    try:
                        if _get_ipc().has_live_task("learn_agent", session_id):
                            await _deny("already_learning")
                            continue
                    except Exception:
                        pass
                    from vaf.core.subagent_spawn import spawn_subagent
                    from vaf.tools.learn_document import _clean_title, _normalize_doc_tag
                    from vaf.tools.learn_job import LearnJobSpec
                    _title = _clean_title(doc_name)
                    # Already learned? Ask BEFORE spawning so the answer is a
                    # friendly frame instead of a spawned child's refusal. The
                    # child re-checks authoritatively via the same helper.
                    try:
                        from uuid import UUID as _UUID

                        from vaf.core.learn_ledger import file_sha256 as _fsha
                        from vaf.tools.learn_job import find_completed_learn as _fcl
                        _scope_uuid = _UUID(str(user_scope_id)) if user_scope_id else None
                        _done = await _fcl(_fsha(Path(_doc_path)), _scope_uuid)
                    except Exception:
                        _done = None  # best-effort; the job is the authority
                    if _done:
                        await websocket.send_json({
                            "type": "learn_denied", "sessionId": session_id,
                            "name": doc_name, "reason": "already_learned",
                            "message": (
                                f'"{_done["doc_title"] or _title}" is already learned '
                                f'({_done["total_pages"]} page(s), '
                                f'{_done["sections"]} section(s)).'),
                        })
                        continue
                    _spec = LearnJobSpec(path=_doc_path, document_title=_title,
                                         doc_tag=_normalize_doc_tag(_title))
                    _sub_env = {}
                    if user_scope_id:
                        _sub_env["VAF_USER_SCOPE_ID"] = str(user_scope_id)
                    _role = manager.connection_roles.get(websocket)
                    if _role:
                        _sub_env["VAF_USER_ROLE"] = str(_role)
                    spawned = spawn_subagent(
                        "learn_agent", f"Learn document: {_title}",
                        include_task_arg=False, payload=_spec.to_json(),
                        extra_env=_sub_env, session_id=session_id,
                        marker_note=f"Learning \"{_title}\" batch by batch.",
                    )
                    if spawned:
                        await websocket.send_json({
                            "type": "learn_started", "sessionId": session_id,
                            "name": doc_name, "taskId": spawned.task_id,
                        })
                    else:
                        await _deny("spawn_failed")

                elif type == "learn_document_cancel":
                    # Graceful cancel: touch the flag file the job polls at batch
                    # boundaries. One writer (this handler), one reader (the job) -
                    # deliberately NOT a field on the IPC record, which the
                    # heartbeat rewrites every 3 seconds.
                    session_id = cmd.get("sessionId") or manager.get_session_for_connection(websocket)
                    task_id = str(cmd.get("taskId") or "").strip()
                    allowed, _ = _ws_session_owner_ok(websocket, session_id)
                    if not allowed or not task_id:
                        continue
                    # The task must belong to THIS session (never cancel another
                    # tenant's job by guessing ids).
                    try:
                        from vaf.core.subagent_ipc import get_ipc as _get_ipc
                        _mine = any(getattr(t, "task_id", "") == task_id
                                    for t in _get_ipc().get_active_tasks(session_id=session_id))
                        if not _mine:
                            continue
                        from vaf.tools.learn_job import cancel_flag_path
                        cancel_flag_path(task_id).touch()
                    except Exception:
                        pass

                elif type == "set_sidebar_documents":
                    session_id = cmd.get("sessionId") or manager.get_session_for_connection(websocket)
                    user_scope_id = manager.get_connection_user(websocket)
                    if not session_id:
                        # Use user-scoped default to prevent crosstalk
                        safe_scope = str(user_scope_id or "default").replace("-", "")[:8]
                        session_id = f"web-default-{safe_scope}"
                    # OWNERSHIP. Without it this handler loads the named session, replaces its
                    # sidebar_documents and saves the file - on BOTH paths, with no scope filter
                    # anywhere: an empty payload wipes another user's sidebar, a full one plants
                    # documents into it. The victim finds them on next load (nothing is pushed),
                    # confirms, and learn_attached_knowledge ingests foreign content under THEIR
                    # scope.
                    # NOT via the attachment RAG, although that was the first reading: both RAG
                    # calls below carry the CALLER's scope, and _scope_filters never yields an
                    # empty filter (a None scope becomes `user_scope_id IS NULL`), so
                    # caller-scope plus foreign-session matches zero rows. That filter is a real
                    # backstop and is named here so nobody removes it as redundant - but it is
                    # not what was broken. The session FILE was.
                    allowed, _ = _ws_session_owner_ok(websocket, session_id, allow_missing=True)
                    if not allowed:
                        log("API", f"Access denied: set_sidebar_documents {session_id} not owned by {user_scope_id}")
                        # Say so. A silent `continue` is indistinguishable from "attachments are
                        # broken", and one legitimate caller lands here: the owner of a session
                        # whose file went corrupt or 0-byte, who is refused because ownership of
                        # an unreadable file cannot be established. They deserve to know why.
                        await websocket.send_json({
                            "type": "sidebar_documents_denied",
                            "sessionId": session_id,
                            "reason": "not_owned_or_unreadable",
                        })
                        continue
                    documents = cmd.get("documents") or []
                    from vaf.core.log_helper import log_attachment
                    log_attachment("WS_RECEIVED", session=session_id, doc_count=len(documents),
                        doc_names=[d.get("name","?") for d in documents[:5]])
                    contents = []
                    try:
                        if not documents:
                            log_attachment("CLEAR_PATH", session=session_id)
                            loaded = session_mgr.load(session_id)
                            if not getattr(loaded, "runtime_state", None):
                                loaded.runtime_state = {}
                            loaded.runtime_state["sidebar_documents"] = []
                            session_mgr.save(loaded, sync_state=False)
                            if bool(Config.get("attachment_rag_enabled", False)):
                                try:
                                    from vaf.memory.attachment_rag import clear_session_attachments_async
                                    await clear_session_attachments_async(session_id=session_id, user_scope_id=user_scope_id)
                                except Exception as e:
                                    log("WebServer", f"attachment clear failed: {e}")
                            await websocket.send_json({
                                "type": "sidebar_documents_set",
                                "contents": [],
                                "sessionId": session_id
                            })
                        else:
                            log_attachment("PROCESS_START", session=session_id, doc_count=len(documents))
                            contents = await process_files_to_sidebar_list(
                                documents, session_id=session_id, user_scope_id=user_scope_id)
                            log_attachment("PROCESS_DONE", session=session_id, results=len(contents),
                                content_lens=[len(c.get("content","")) for c in contents],
                                content_types=[
                                    "ERROR" if "[ERROR]" in (c.get("content",""))[:80]
                                    else "SCANNED" if "[Scanned PDF" in (c.get("content",""))[:200]
                                    else "OK"
                                    for c in contents
                                ])
                            for _dc in contents:
                                _dc_name = _dc.get("name", "?")
                                _dc_content = _dc.get("content", "")
                                log("WebServer", f"set_sidebar_documents: doc={_dc_name!r} content_len={len(_dc_content)} preview={repr(_dc_content[:120])}")
                            try:
                                loaded = session_mgr.load(session_id)
                            except Exception as _load_err:
                                # Corrupted or 0-byte session file (pre-atomic-write crash).
                                # Create a minimal in-memory session so sidebar_documents can be saved.
                                log_attachment("SESSION_LOAD_FAILED", session=session_id, error=str(_load_err))
                                log("WebServer", f"set_sidebar_documents: session load failed ({_load_err!r}), creating minimal session for {session_id}")
                                # MEASURED REMAINING GAP, named rather than left to be found again.
                                # WHAT GETS THROUGH: the ownership gate above runs with
                                #   allow_missing=True, and _ws_session_owner_ok answers a LOAD
                                #   FAILURE with (bool(allow_missing) or is_admin, None) - so a
                                #   session that cannot be read passes the gate.
                                # UNDER WHICH CONDITION: the session file is missing, corrupt or
                                #   0 bytes - which is exactly the case this fallback exists for.
                                # WHAT IT CAUSES: the minimal session created here is stamped with
                                #   the CALLER's scope, so a stranger acquires the id. Takeover
                                #   plus lock-out of the owner, delivered by the very gate added
                                #   to prevent cross-user writes. Milder sibling: a not-yet-used
                                #   id can be claimed in advance. Not a read leak - the file was
                                #   already unreadable.
                                # UNTESTED today, deliberately recorded as such.
                                # Belongs with _open_report_in_viewer, which is the same family:
                                # what is the authority when the usual one is absent - there
                                # because no authenticated user exists (tokenless loopback IPC),
                                # here because the session cannot be loaded. Solved apart they
                                # become two special rules; measured together, probably one.
                                loaded = Session(id=session_id, name=f"Session {session_id}")
                                loaded.metadata["user_scope_id"] = str(user_scope_id or "")
                            if not getattr(loaded, "runtime_state", None):
                                loaded.runtime_state = {}
                            # Strip base64 data before persisting to session runtime_state.
                            # The 'data' field can be 10-25 MB for large PDFs, bloating the
                            # session JSON file and slowing every subsequent session read/write.
                            # The frontend already received 'data' in the WS response below;
                            # the headless_runner only needs 'content' (extracted text).
                            def _sanitize_str(s):
                                """Strip lone Unicode surrogates (lone surrogates from PDF emoji)  
                                that break UTF-8 JSON serialization in session.save()."""
                                if not isinstance(s, str):
                                    return s
                                return s.encode("utf-8", errors="replace").decode("utf-8")
                            _slim = [
                                {k: (_sanitize_str(v) if k == "content" else v)
                                 for k, v in doc.items() if k != "data"}
                                for doc in contents
                            ]
                            loaded.runtime_state["sidebar_documents"] = _slim
                            session_mgr.save(loaded, sync_state=False)
                            log_attachment("SAVE_OK", session=session_id, docs=len(_slim),
                                names=[d.get("name","?") for d in _slim])
                            if bool(Config.get("attachment_rag_enabled", False)):
                                # Show the banner immediately, then index in a cancellable
                                # background task so the WS loop stays free for the stop button.
                                await _notify_attachment_index(manager, session_id, "attachment_indexing", count=len(contents))
                                _spawn_attachment_index(manager, session_id, user_scope_id, contents)
                            await websocket.send_json({
                                "type": "sidebar_documents_set",
                                "contents": contents,
                                "sessionId": session_id
                            })
                    except FileNotFoundError:
                        # Session file missing — still do RAG ops but do NOT write a new
                        # empty session to disk. Creating Session(id=session_id, messages=[])
                        # and saving it would overwrite a valid session that was being written
                        # concurrently, destroying the user's chat history.
                        # The session is always created by the WS connection handler; sidebar
                        # state is transient and will be re-synced on the next attachment.
                        if not documents:
                            if bool(Config.get("attachment_rag_enabled", False)):
                                try:
                                    from vaf.memory.attachment_rag import clear_session_attachments_async
                                    await clear_session_attachments_async(session_id=session_id, user_scope_id=user_scope_id)
                                except Exception as e:
                                    log("WebServer", f"attachment clear failed (new session): {e}")
                        else:
                            if bool(Config.get("attachment_rag_enabled", False)):
                                await _notify_attachment_index(manager, session_id, "attachment_indexing", count=len(contents))
                                _spawn_attachment_index(manager, session_id, user_scope_id, contents)
                        await websocket.send_json({
                            "type": "sidebar_documents_set",
                            "contents": contents,
                            "sessionId": session_id
                        })
                    except Exception as e:
                        import traceback as _tb
                        _tb_str = _tb.format_exc()
                        log("WebServer", f"set_sidebar_documents FAILED: {e}\n{_tb_str}")
                        log_attachment("SAVE_FAILED", session=session_id, error=str(e),
                            tb=_tb_str[-400:])
                        # Use whatever contents were computed before the failure.
                        # Sending [] here would wipe the user's attached documents from
                        # the viewer even when file extraction succeeded but session save
                        # raised (e.g. stale corrupt JSON from before the atomic-write fix).
                        await websocket.send_json({
                            "type": "sidebar_documents_set",
                            "contents": contents,
                            "sessionId": session_id,
                            "error": str(e)
                        })

                elif type == "get_tools":
                    # Return list of available tools, filtered by the requesting
                    # user's role and custom-tool visibility permissions.
                    try:
                        from vaf.core.config import get_local_admin_scope_id
                        from vaf.core.custom_tools_registry import (
                            get_all_custom_tool_names,
                            get_visible_tool_names_for_user,
                            get_tool_manifest_entry,
                        )

                        _gt_scope = manager.get_connection_user(websocket) if manager else None
                        _gt_role  = manager.get_connection_user_role(websocket) if manager else None
                        _gt_local_admin = get_local_admin_scope_id()
                        _gt_is_admin = (
                            _gt_role == "admin"
                            or (
                                _gt_scope is not None
                                and str(_gt_scope) == str(_gt_local_admin)
                            )
                        )

                        def _gt_account_allows(_name: str) -> bool:
                            """The account allowlist, as a LISTING question.

                            Fails CLOSED per entry, the same polarity the funnel uses: a
                            resolver that crashed must not turn a restricted account into
                            an unrestricted list. One tool vanishing from a menu is a far
                            smaller harm than one that should not be there at all."""
                            try:
                                from vaf.core.tool_dispatch import account_allows_tool
                                return account_allows_tool(_name, _gt_scope, _gt_role)
                            except Exception:
                                return False

                        # Admins pass None so get_visible_tool_names_for_user returns ALL
                        _gt_filter_scope = None if _gt_is_admin else _gt_scope

                        all_custom_names   = set(get_all_custom_tool_names())
                        visible_custom     = set(get_visible_tool_names_for_user(_gt_filter_scope))
                        from vaf.core.tool_contract import tool_category

                        agent = manager.agent_instance
                        if agent and hasattr(agent, "tools"):
                            tools_list = []
                            for name, tool in agent.tools.items():
                                is_custom = name in all_custom_names
                                # Non-admins: skip custom tools they can't see
                                if is_custom and not _gt_is_admin and name not in visible_custom:
                                    continue
                                # ...and skip what the admin took away from this ACCOUNT.
                                # This list feeds the '/' suggestions and the sub-agent
                                # hotbar; without the filter a user whose admin disabled
                                # coding_agent was offered it and only learned at dispatch
                                # time that it was refused. The exemptions (no scope,
                                # admin) live inside account_allows_tool, so this cannot
                                # drift from the funnel's own answer.
                                if not _gt_account_allows(name):
                                    continue
                                entry = {
                                    "name":        name,
                                    "description": getattr(tool, "description", "No description"),
                                    "category":    tool_category(name, tool),
                                    # Frontend uses these two flags to render management controls
                                    "is_custom":   is_custom,
                                    "can_manage":  _gt_is_admin,
                                }
                                if is_custom:
                                    meta = get_tool_manifest_entry(name)
                                    if meta:
                                        entry["shared_with"]  = meta.get("shared_with", ["*"])
                                        entry["created_by"]   = meta.get("created_by", "")
                                        entry["updated_at"]   = meta.get("updated_at", "")
                                tools_list.append(entry)

                            _attach_learned_states(tools_list)
                            manager.tools_cache = tools_list
                            await websocket.send_json({
                                "type": "tools_list",
                                "tools": tools_list,
                            })
                        elif manager.tools_cache:
                            await websocket.send_json({
                                "type": "tools_list",
                                "tools": manager.tools_cache,
                            })
                        else:
                            await websocket.send_json({
                                "type": "tools_list",
                                "tools": _scan_tool_modules(),
                            })
                    except Exception as e:
                        await websocket.send_json({
                            "type": "tools_list",
                            "tools": [],
                            "error": str(e),
                        })

                # ── Custom Tool Management (admin-only) ───────────────────────────
                # All four handlers share the same admin-check pattern used
                # throughout the WS loop (see line ~2192 for the reference pattern).

                elif type == "create_custom_tool":
                    # Upload or write a new custom tool from the WebUI editor.
                    # Payload: { name: str, code: str, shared_with: list[str] }
                    try:
                        from vaf.core.config import get_local_admin_scope_id
                        from vaf.core import custom_tools_registry as _ctr

                        _ct_scope = manager.get_connection_user(websocket) if manager else None
                        _ct_role  = manager.get_connection_user_role(websocket) if manager else None
                        _ct_local_admin = get_local_admin_scope_id()
                        _ct_is_admin = (
                            _ct_role == "admin"
                            or (_ct_scope is not None and str(_ct_scope) == str(_ct_local_admin))
                        )
                        if not _ct_is_admin:
                            await websocket.send_json({
                                "type": "custom_tool_error",
                                "error": "Admin permission required to create tools.",
                            })
                        else:
                            _ct_name       = (data.get("name") or "").strip()
                            _ct_code       = data.get("code", "")
                            _ct_shared     = data.get("shared_with", ["*"])
                            _ct_username   = manager.get_connection_username(websocket) or "admin"

                            # Basic name validation: snake_case identifiers only
                            import re as _re
                            if not _re.match(r'^[a-z][a-z0-9_]*$', _ct_name):
                                raise ValueError(
                                    "Tool name must be lowercase snake_case (e.g. my_tool)."
                                )

                            # Write file + validate BaseTool subclass, then register
                            _ct_filename = f"{_ct_name}.py"
                            _ctr.save_tool_file(_ct_filename, _ct_code)
                            # register_tool also validates via load_custom_tool_class
                            _ctr.register_tool(
                                tool_name=_ct_name,
                                filename=_ct_filename,
                                created_by=_ct_username,
                                shared_with=_ct_shared,
                            )

                            # Hot-reload so the live agent immediately has the new tool
                            agent = manager.agent_instance
                            if agent and hasattr(agent, "reload_custom_tools"):
                                agent.reload_custom_tools()

                            await websocket.send_json({
                                "type": "custom_tool_created",
                                "name": _ct_name,
                            })
                            # Broadcast updated tool list to all connected clients
                            # so other open tabs / admin panels refresh automatically.
                            await _broadcast_tools_update(manager)

                    except Exception as e:
                        await websocket.send_json({
                            "type": "custom_tool_error",
                            "error": str(e),
                        })

                elif type == "update_custom_tool":
                    # Edit the source code of an existing custom tool.
                    # Payload: { name: str, code: str }
                    try:
                        from vaf.core.config import get_local_admin_scope_id
                        from vaf.core import custom_tools_registry as _ctr

                        _ut_scope = manager.get_connection_user(websocket) if manager else None
                        _ut_role  = manager.get_connection_user_role(websocket) if manager else None
                        _ut_is_admin = (
                            _ut_role == "admin"
                            or (_ut_scope is not None and str(_ut_scope) == str(get_local_admin_scope_id()))
                        )
                        if not _ut_is_admin:
                            await websocket.send_json({
                                "type": "custom_tool_error",
                                "error": "Admin permission required to edit tools.",
                            })
                        else:
                            _ut_name     = (data.get("name") or "").strip()
                            _ut_code     = data.get("code", "")
                            _ut_username = manager.get_connection_username(websocket) or "admin"

                            # update_tool_source validates BaseTool presence before overwriting
                            _ctr.update_tool_source(_ut_name, _ut_code, _ut_username)

                            agent = manager.agent_instance
                            if agent and hasattr(agent, "reload_custom_tools"):
                                agent.reload_custom_tools()

                            await websocket.send_json({
                                "type": "custom_tool_updated",
                                "name": _ut_name,
                            })
                            await _broadcast_tools_update(manager)

                    except Exception as e:
                        await websocket.send_json({
                            "type": "custom_tool_error",
                            "error": str(e),
                        })

                elif type == "delete_custom_tool":
                    # Remove a custom tool permanently.
                    # Payload: { name: str }
                    try:
                        from vaf.core.config import get_local_admin_scope_id
                        from vaf.core import custom_tools_registry as _ctr

                        _dt_scope = manager.get_connection_user(websocket) if manager else None
                        _dt_role  = manager.get_connection_user_role(websocket) if manager else None
                        _dt_is_admin = (
                            _dt_role == "admin"
                            or (_dt_scope is not None and str(_dt_scope) == str(get_local_admin_scope_id()))
                        )
                        if not _dt_is_admin:
                            await websocket.send_json({
                                "type": "custom_tool_error",
                                "error": "Admin permission required to delete tools.",
                            })
                        else:
                            _dt_name = (data.get("name") or "").strip()
                            _ctr.delete_tool(_dt_name)

                            # Remove from the live agent immediately
                            agent = manager.agent_instance
                            if agent and hasattr(agent, "reload_custom_tools"):
                                agent.reload_custom_tools()

                            await websocket.send_json({
                                "type": "custom_tool_deleted",
                                "name": _dt_name,
                            })
                            await _broadcast_tools_update(manager)

                    except Exception as e:
                        await websocket.send_json({
                            "type": "custom_tool_error",
                            "error": str(e),
                        })

                elif type == "get_mcp_servers":
                    # List configured MCP servers + live connection status for the Settings UI.
                    try:
                        agent = manager.agent_instance if manager else None
                        await websocket.send_json({"type": "mcp_servers", "servers": _mcp_servers_payload(agent)})
                    except Exception as e:
                        await websocket.send_json({"type": "mcp_server_error", "error": str(e)})

                elif type in ("create_mcp_server", "update_mcp_server"):
                    # Add or edit an MCP server in mcp_servers.json, then hot-reload.
                    # Payload: { name, command, transport?, url?, enabled?, permission_level? }
                    try:
                        from vaf.core.config import get_local_admin_scope_id
                        from vaf.core.mcp_registry import load_mcp_manifest, save_mcp_manifest

                        _ms_scope = manager.get_connection_user(websocket) if manager else None
                        _ms_role  = manager.get_connection_user_role(websocket) if manager else None
                        _ms_admin = (
                            _ms_role == "admin"
                            or (_ms_scope is not None and str(_ms_scope) == str(get_local_admin_scope_id()))
                        )
                        if not _ms_admin:
                            await websocket.send_json({"type": "mcp_server_error", "error": "Admin permission required to manage MCP servers."})
                        else:
                            import re as _re
                            _ms_name = (data.get("name") or "").strip()
                            if not _re.match(r'^[A-Za-z][A-Za-z0-9_-]*$', _ms_name):
                                raise ValueError("Server name must start with a letter and use only letters, digits, _ or -.")
                            _ms_transport = (data.get("transport") or "stdio").strip()
                            _ms_cmd = (data.get("command") or "").strip()
                            if _ms_transport == "stdio" and not _ms_cmd:
                                raise ValueError("A command is required for stdio transport.")
                            _ms_perm = (data.get("permission_level") or "write").strip().lower()
                            if _ms_perm not in ("read", "write", "dangerous"):
                                _ms_perm = "write"
                            _manifest = load_mcp_manifest() or {}
                            _srv = _manifest.get("servers")
                            if not isinstance(_srv, dict):
                                _srv = {}
                            _ms_env = data.get("env")
                            if not isinstance(_ms_env, dict):
                                _ms_env = {}
                            _srv[_ms_name] = {
                                "command": _ms_cmd,
                                "transport": _ms_transport,
                                "url": (data.get("url") or "").strip(),
                                "enabled": bool(data.get("enabled", True)),
                                "permission_level": _ms_perm,
                                "env": {str(k): str(v) for k, v in _ms_env.items()},
                            }
                            _manifest["servers"] = _srv
                            save_mcp_manifest(_manifest)
                            agent = manager.agent_instance if manager else None
                            if agent and hasattr(agent, "reload_mcp_tools"):
                                agent.reload_mcp_tools()
                            # Return the refreshed list with the reply so the UI updates without a refetch.
                            await websocket.send_json({"type": "mcp_server_saved", "name": _ms_name, "servers": _mcp_servers_payload(agent)})
                            await _broadcast_tools_update(manager)
                    except Exception as e:
                        await websocket.send_json({"type": "mcp_server_error", "error": str(e)})

                elif type == "delete_mcp_server":
                    # Remove an MCP server from mcp_servers.json, then hot-reload. Payload: { name }
                    try:
                        from vaf.core.config import get_local_admin_scope_id
                        from vaf.core.mcp_registry import load_mcp_manifest, save_mcp_manifest

                        _md_scope = manager.get_connection_user(websocket) if manager else None
                        _md_role  = manager.get_connection_user_role(websocket) if manager else None
                        _md_admin = (
                            _md_role == "admin"
                            or (_md_scope is not None and str(_md_scope) == str(get_local_admin_scope_id()))
                        )
                        if not _md_admin:
                            await websocket.send_json({"type": "mcp_server_error", "error": "Admin permission required to manage MCP servers."})
                        else:
                            _md_name = (data.get("name") or "").strip()
                            _manifest = load_mcp_manifest() or {}
                            _srv = _manifest.get("servers")
                            if isinstance(_srv, dict) and _md_name in _srv:
                                del _srv[_md_name]
                                _manifest["servers"] = _srv
                                save_mcp_manifest(_manifest)
                            agent = manager.agent_instance if manager else None
                            if agent and hasattr(agent, "reload_mcp_tools"):
                                agent.reload_mcp_tools()
                            # Return the refreshed list with the reply so the UI updates without a refetch.
                            await websocket.send_json({"type": "mcp_server_deleted", "name": _md_name, "servers": _mcp_servers_payload(agent)})
                            await _broadcast_tools_update(manager)
                    except Exception as e:
                        await websocket.send_json({"type": "mcp_server_error", "error": str(e)})

                elif type == "test_mcp_server":
                    # Probe a server config (without saving) so the admin can validate it in the
                    # editor. Payload: { command, transport?, url? }. Reply: mcp_server_test_result.
                    try:
                        from vaf.core.config import get_local_admin_scope_id, Config as _CfgMcp
                        from vaf.core.mcp_registry import probe_mcp_server

                        _tm_scope = manager.get_connection_user(websocket) if manager else None
                        _tm_role  = manager.get_connection_user_role(websocket) if manager else None
                        _tm_admin = (
                            _tm_role == "admin"
                            or (_tm_scope is not None and str(_tm_scope) == str(get_local_admin_scope_id()))
                        )
                        if not _tm_admin:
                            await websocket.send_json({"type": "mcp_server_test_result", "connected": False, "tool_count": 0, "tools": [], "error": "Admin permission required."})
                        else:
                            _tm_env = data.get("env") if isinstance(data.get("env"), dict) else {}
                            _tm_cfg = {
                                "command": (data.get("command") or "").strip(),
                                "transport": (data.get("transport") or "stdio").strip(),
                                "url": (data.get("url") or "").strip(),
                                "env": {str(k): str(v) for k, v in _tm_env.items()},
                            }
                            _tm_timeout = float(_CfgMcp.get("mcp_discovery_timeout_seconds", 5) or 5)
                            _tm_res = probe_mcp_server(_tm_cfg, _tm_timeout)
                            await websocket.send_json({"type": "mcp_server_test_result", **_tm_res})
                    except Exception as e:
                        await websocket.send_json({"type": "mcp_server_test_result", "connected": False, "tool_count": 0, "tools": [], "error": str(e)})

                elif type == "update_custom_tool_permissions":
                    # Change which users can see a custom tool.
                    # Payload: { name: str, shared_with: list[str] }
                    #   shared_with: ["*"] → all users
                    #   shared_with: []    → admin only
                    #   shared_with: ["<scope_id>", ...] → specific users
                    try:
                        from vaf.core.config import get_local_admin_scope_id
                        from vaf.core import custom_tools_registry as _ctr

                        _pp_scope = manager.get_connection_user(websocket) if manager else None
                        _pp_role  = manager.get_connection_user_role(websocket) if manager else None
                        _pp_is_admin = (
                            _pp_role == "admin"
                            or (_pp_scope is not None and str(_pp_scope) == str(get_local_admin_scope_id()))
                        )
                        if not _pp_is_admin:
                            await websocket.send_json({
                                "type": "custom_tool_error",
                                "error": "Admin permission required to change tool permissions.",
                            })
                        else:
                            _pp_name       = (data.get("name") or "").strip()
                            _pp_shared     = data.get("shared_with", ["*"])
                            _ctr.update_tool_permissions(_pp_name, _pp_shared)

                            await websocket.send_json({
                                "type": "custom_tool_permissions_updated",
                                "name":        _pp_name,
                                "shared_with": _pp_shared,
                            })
                            # Broadcast so other clients refresh their tool lists
                            await _broadcast_tools_update(manager)

                    except Exception as e:
                        await websocket.send_json({
                            "type": "custom_tool_error",
                            "error": str(e),
                        })

                elif type == "get_custom_tool_source":
                    # Return the Python source code of a custom tool for the editor.
                    # Payload: { name: str }
                    # Admin-only — non-admins must not be able to exfiltrate source.
                    try:
                        from vaf.core.config import get_local_admin_scope_id
                        from vaf.core import custom_tools_registry as _ctr

                        _gs_scope = manager.get_connection_user(websocket) if manager else None
                        _gs_role  = manager.get_connection_user_role(websocket) if manager else None
                        _gs_is_admin = (
                            _gs_role == "admin"
                            or (_gs_scope is not None and str(_gs_scope) == str(get_local_admin_scope_id()))
                        )
                        if not _gs_is_admin:
                            await websocket.send_json({
                                "type": "custom_tool_error",
                                "error": "Admin permission required to view tool source.",
                            })
                        else:
                            _gs_name   = (data.get("name") or "").strip()
                            _gs_source = _ctr.get_tool_source(_gs_name)
                            await websocket.send_json({
                                "type": "custom_tool_source",
                                "name":   _gs_name,
                                "source": _gs_source or "",
                            })

                    except Exception as e:
                        await websocket.send_json({
                            "type": "custom_tool_error",
                            "error": str(e),
                        })
                
                elif type == "get_custom_tool_users":
                    # Return non-admin users for the share picker in CustomToolEditor.
                    # Admin-only: non-admins have no reason to query the user list.
                    try:
                        from vaf.core.config import get_local_admin_scope_id

                        _gu_scope = manager.get_connection_user(websocket) if manager else None
                        _gu_role  = manager.get_connection_user_role(websocket) if manager else None
                        _gu_is_admin = (
                            _gu_role == "admin"
                            or (_gu_scope is not None and str(_gu_scope) == str(get_local_admin_scope_id()))
                        )
                        if not _gu_is_admin:
                            await websocket.send_json({
                                "type": "custom_tool_error",
                                "error": "Admin permission required.",
                            })
                        else:
                            # Reuse the existing /api/users logic by querying the DB directly
                            try:
                                from vaf.auth.database import get_auth_db
                                from vaf.auth.models import LocalUser
                                from sqlalchemy import select as _sa_select
                                async with get_auth_db() as _db:
                                    _result = await _db.execute(
                                        _sa_select(LocalUser).where(LocalUser.is_active == True)
                                    )
                                    _users = _result.scalars().all()
                                    _user_list = [
                                        {
                                            "id":            str(u.id),
                                            "username":      u.username,
                                            "user_scope_id": str(u.user_scope_id),
                                            "role":          u.role,
                                        }
                                        for u in _users
                                        if u.role != "admin"  # admins always see everything; no point including them
                                    ]
                            except Exception:
                                _user_list = []
                            await websocket.send_json({
                                "type":  "custom_tool_users",
                                "users": _user_list,
                            })
                    except Exception as e:
                        await websocket.send_json({
                            "type": "custom_tool_error",
                            "error": str(e),
                        })

                elif type == "get_workflows":
                    # Return list of available workflow templates
                    try:
                        from vaf.workflows.templates import list_templates
                        workflows = list_templates()
                        await websocket.send_json({
                            "type": "workflows_list",
                            "workflows": workflows
                        })
                    except Exception as e:
                        await websocket.send_json({
                            "type": "workflows_list",
                            "workflows": [],
                            "error": str(e)
                        })

                elif type == "get_workflow_run_state":
                    # "Is the run my panel is showing actually still going?"
                    #
                    # Workflow panel events are fire and forget (no queue, no replay, no
                    # sequence numbers), so a socket that dies mid-run leaves the panel stuck
                    # on the last state it received, with no way for the client to detect the
                    # gap (live incident 2026-07-20). This is how it asks instead of guessing.
                    #
                    # Scoped to the session the CLIENT names, verified against the connection's
                    # own subscription (Rule 4.4): a fresh connection is auto-subscribed to the
                    # newest session, so answering from the connection alone would report about
                    # a different run than the one on screen.
                    _rs_req = str(data.get("sessionId") or "").strip()
                    _rs_conn = str(manager.get_session_for_connection(websocket) or "").strip()
                    _rs_session = _rs_req or _rs_conn
                    _rs_wf = str(data.get("workflowId") or "").strip()
                    _rs_template = str(data.get("templateId") or "").strip()
                    try:
                        if _rs_req and _rs_conn and _rs_req != _rs_conn:
                            # Never answer about a session this connection is not on.
                            _rs_state = "unknown"
                        else:
                            from vaf.core.subagent_ipc import get_ipc as _rs_get_ipc
                            _rs_state = _rs_get_ipc().workflow_run_state_for_session(
                                _rs_session, _rs_template
                            )
                    except Exception:
                        _rs_state = "unknown"
                    await websocket.send_json({
                        "type": "workflow_run_state",
                        "workflowId": _rs_wf,
                        "sessionId": _rs_session,
                        "state": _rs_state,   # running | paused | ended | unknown
                    })

                elif type == "create_workflow":
                    from vaf.core.config import get_local_admin_scope_id
                    _wf_scope       = manager.get_connection_user(websocket)
                    _wf_role        = manager.get_connection_user_role(websocket)
                    _wf_local_admin = get_local_admin_scope_id()
                    _wf_is_admin    = (
                        _wf_role == "admin"
                        or (_wf_scope is not None and str(_wf_scope) == str(_wf_local_admin))
                    )
                    if not _wf_is_admin:
                        await websocket.send_json({
                            "type": "workflow_error",
                            "error": "Admin permission required to create workflows.",
                        })
                    else:
                        try:
                            import json as _json_wf
                            import os as _os_wf
                            import re as _re_wf
                            from vaf.workflows.templates import reload_workflows, list_templates

                            _wf_id       = str(cmd.get("workflow_id") or "").strip()
                            _wf_name     = str(cmd.get("name") or "").strip()
                            _wf_desc     = str(cmd.get("description") or "").strip()
                            _wf_triggers = [str(t) for t in (cmd.get("triggers") or []) if str(t).strip()]
                            _wf_steps    = cmd.get("steps") or []

                            if not _re_wf.match(r'^[a-z][a-z0-9_]*$', _wf_id):
                                raise ValueError(f"workflow_id must be lowercase snake_case, got '{_wf_id}'")
                            if not _wf_name:
                                raise ValueError("name is required")
                            if not _wf_steps:
                                raise ValueError("at least one step is required")

                            _user_wf_dir = _os_wf.path.expanduser("~/.vaf/workflows")
                            _os_wf.makedirs(_user_wf_dir, exist_ok=True)
                            _wf_path = _os_wf.path.join(_user_wf_dir, f"{_wf_id}.py")

                            if _os_wf.path.exists(_wf_path):
                                raise ValueError(
                                    f"Workflow '{_wf_id}' already exists. "
                                    "Use update_workflow to modify it."
                                )

                            _wf_dict = {
                                "name": _wf_name,
                                "description": _wf_desc,
                                "triggers": _wf_triggers,
                                "steps": _wf_steps,
                            }
                            _wf_content = (
                                f"# User-created workflow: {_wf_name}\n"
                                f"WORKFLOW = {_json_wf.dumps(_wf_dict, indent=4, ensure_ascii=False)}\n"
                            )
                            _wf_tmp = _wf_path + ".tmp"
                            with open(_wf_tmp, "w", encoding="utf-8") as _f:
                                _f.write(_wf_content)
                            _os_wf.replace(_wf_tmp, _wf_path)

                            reload_workflows()
                            await websocket.send_json({"type": "workflow_created", "workflow_id": _wf_id})
                            await websocket.send_json({"type": "workflows_list", "workflows": list_templates()})
                        except Exception as e:
                            await websocket.send_json({"type": "workflow_error", "error": str(e)})

                elif type == "update_workflow":
                    from vaf.core.config import get_local_admin_scope_id
                    _wf_scope       = manager.get_connection_user(websocket)
                    _wf_role        = manager.get_connection_user_role(websocket)
                    _wf_local_admin = get_local_admin_scope_id()
                    _wf_is_admin    = (
                        _wf_role == "admin"
                        or (_wf_scope is not None and str(_wf_scope) == str(_wf_local_admin))
                    )
                    if not _wf_is_admin:
                        await websocket.send_json({
                            "type": "workflow_error",
                            "error": "Admin permission required to update workflows.",
                        })
                    else:
                        try:
                            import json as _json_wf
                            import os as _os_wf
                            from vaf.workflows.templates import reload_workflows, list_templates

                            _wf_id       = str(cmd.get("workflow_id") or "").strip()
                            _wf_name     = str(cmd.get("name") or "").strip()
                            _wf_desc     = str(cmd.get("description") or "").strip()
                            _wf_triggers = [str(t) for t in (cmd.get("triggers") or []) if str(t).strip()]
                            _wf_steps    = cmd.get("steps") or []

                            if not _wf_name:
                                raise ValueError("name is required")
                            if not _wf_steps:
                                raise ValueError("at least one step is required")

                            _user_wf_dir = _os_wf.path.expanduser("~/.vaf/workflows")
                            _wf_path = _os_wf.path.join(_user_wf_dir, f"{_wf_id}.py")

                            if not _os_wf.path.exists(_wf_path):
                                raise ValueError(
                                    f"Workflow '{_wf_id}' is a built-in workflow and cannot be modified."
                                )

                            _wf_dict = {
                                "name": _wf_name,
                                "description": _wf_desc,
                                "triggers": _wf_triggers,
                                "steps": _wf_steps,
                            }
                            _wf_content = (
                                f"# User-created workflow: {_wf_name}\n"
                                f"WORKFLOW = {_json_wf.dumps(_wf_dict, indent=4, ensure_ascii=False)}\n"
                            )
                            _wf_tmp = _wf_path + ".tmp"
                            with open(_wf_tmp, "w", encoding="utf-8") as _f:
                                _f.write(_wf_content)
                            _os_wf.replace(_wf_tmp, _wf_path)

                            reload_workflows()
                            await websocket.send_json({"type": "workflow_updated", "workflow_id": _wf_id})
                            await websocket.send_json({"type": "workflows_list", "workflows": list_templates()})
                        except Exception as e:
                            await websocket.send_json({"type": "workflow_error", "error": str(e)})

                elif type == "delete_workflow":
                    from vaf.core.config import get_local_admin_scope_id
                    _wf_scope       = manager.get_connection_user(websocket)
                    _wf_role        = manager.get_connection_user_role(websocket)
                    _wf_local_admin = get_local_admin_scope_id()
                    _wf_is_admin    = (
                        _wf_role == "admin"
                        or (_wf_scope is not None and str(_wf_scope) == str(_wf_local_admin))
                    )
                    if not _wf_is_admin:
                        await websocket.send_json({
                            "type": "workflow_error",
                            "error": "Admin permission required to delete workflows.",
                        })
                    else:
                        try:
                            import os as _os_wf
                            from vaf.workflows.templates import reload_workflows, list_templates

                            _wf_id = str(cmd.get("workflow_id") or "").strip()
                            _user_wf_dir = _os_wf.path.expanduser("~/.vaf/workflows")
                            _wf_path = _os_wf.path.join(_user_wf_dir, f"{_wf_id}.py")

                            if not _os_wf.path.exists(_wf_path):
                                raise ValueError(
                                    f"Workflow '{_wf_id}' is a built-in workflow and cannot be deleted."
                                )

                            _os_wf.remove(_wf_path)
                            reload_workflows()
                            await websocket.send_json({"type": "workflow_deleted", "workflow_id": _wf_id})
                            await websocket.send_json({"type": "workflows_list", "workflows": list_templates()})
                        except Exception as e:
                            await websocket.send_json({"type": "workflow_error", "error": str(e)})

                # ─────────────────────────────────────────────────────────────
                # SKILLS (Anthropic Agent Skills / SKILL.md). The second routing
                # tier under workflows: file I/O mirrors the workflow handlers;
                # visibility/scoping mirrors the custom-tools handlers.
                # ─────────────────────────────────────────────────────────────
                elif type == "get_skills":
                    try:
                        from vaf.core.config import get_local_admin_scope_id
                        from vaf.skills.templates import list_skills as _list_skills
                        _sk_scope = manager.get_connection_user(websocket) if manager else None
                        _sk_role  = manager.get_connection_user_role(websocket) if manager else None
                        _sk_is_admin = (
                            _sk_role == "admin"
                            or (_sk_scope is not None and str(_sk_scope) == str(get_local_admin_scope_id()))
                        )
                        _filter_scope = None if _sk_is_admin else _sk_scope
                        # Admins also see invalid skills so they can fix them in the editor.
                        _skills = _list_skills(user_scope_id=_filter_scope, include_invalid=_sk_is_admin)
                        if _sk_is_admin:
                            from vaf.core import skills_registry as _skr_g
                            for _s in _skills:
                                _s["can_manage"] = True
                                _s["source"] = _skr_g.get_skill_md_source(_s["id"]) or ""
                        else:
                            for _s in _skills:
                                _s["can_manage"] = False
                        await websocket.send_json({"type": "skills_list", "skills": _skills})
                    except Exception as e:
                        await websocket.send_json({"type": "skills_list", "skills": [], "error": str(e)})

                elif type == "get_skill_source":
                    # Return the raw SKILL.md text for the editor (admin-only).
                    try:
                        from vaf.core.config import get_local_admin_scope_id
                        from vaf.core import skills_registry as _skr
                        _ss_scope = manager.get_connection_user(websocket) if manager else None
                        _ss_role  = manager.get_connection_user_role(websocket) if manager else None
                        _ss_is_admin = (
                            _ss_role == "admin"
                            or (_ss_scope is not None and str(_ss_scope) == str(get_local_admin_scope_id()))
                        )
                        if not _ss_is_admin:
                            await websocket.send_json({"type": "skill_error", "error": "Admin permission required to view skill source."})
                        else:
                            _ss_id = (cmd.get("skill_id") or "").strip()
                            _ss_src = _skr.get_skill_md_source(_ss_id)
                            await websocket.send_json({"type": "skill_source", "skill_id": _ss_id, "source": _ss_src or ""})
                    except Exception as e:
                        await websocket.send_json({"type": "skill_error", "error": str(e)})

                elif type in ("create_skill", "update_skill"):
                    from vaf.core.config import get_local_admin_scope_id
                    _sk_scope = manager.get_connection_user(websocket)
                    _sk_role  = manager.get_connection_user_role(websocket)
                    _sk_is_admin = (
                        _sk_role == "admin"
                        or (_sk_scope is not None and str(_sk_scope) == str(get_local_admin_scope_id()))
                    )
                    if not _sk_is_admin:
                        await websocket.send_json({"type": "skill_error", "error": "Admin permission required to manage skills."})
                    else:
                        try:
                            import yaml as _yaml_sk
                            from vaf.core import skills_registry as _skr
                            from vaf.skills.skill_md import (
                                parse_skill_md_text as _parse_text,
                                parse_skill_meta as _parse_meta,
                            )
                            from vaf.skills.templates import reload_skills as _reload_sk

                            _is_update = (type == "update_skill")
                            _sk_id = _skr.validate_skill_id(cmd.get("skill_id") or "")

                            # The editor may send raw SKILL.md, or name/description/body.
                            _raw = (cmd.get("skill_md") or "").strip()
                            if _raw:
                                _content = _raw if _raw.endswith("\n") else _raw + "\n"
                            else:
                                _name = (cmd.get("name") or "").strip()
                                _desc = (cmd.get("description") or "").strip()
                                _body = cmd.get("body") or ""
                                # MERGE, never replace: the format allows keys this
                                # editor has no field for (metadata, license,
                                # allowed-tools), and rebuilding the frontmatter from
                                # two fields deleted every one of them - so editing an
                                # imported skill once stripped what its author wrote.
                                # Only the two fields the editor owns are overwritten.
                                _keep: dict = {}
                                try:
                                    _existing_folder = _skr.resolve_skill_folder(_sk_id)
                                    if _existing_folder is not None:
                                        _prev = _parse_meta(_existing_folder / "SKILL.md")
                                        if isinstance(_prev, dict):
                                            _keep = dict(_prev.get("frontmatter") or {})
                                except Exception:
                                    _keep = {}
                                _keep["name"] = _name
                                _keep["description"] = _desc
                                _fm = _yaml_sk.safe_dump(
                                    _keep, sort_keys=False, allow_unicode=True,
                                ).strip()
                                _content = f"---\n{_fm}\n---\n\n{_body}\n"

                            # Validate BEFORE writing so a bad edit never clobbers a good skill.
                            _parsed = _parse_text(_content, _sk_id)
                            if not _parsed["valid"]:
                                raise ValueError(f"invalid skill: {_parsed['error']}")

                            # resolve_skill_folder counts the shipped copy too. Updating
                            # a shipped skill IS the supported customisation: the save
                            # lands in the user dir and replaces the shipped copy at
                            # discovery. Creating under a shipped id is refused the same
                            # way any taken id is - it would shadow a working skill.
                            _sk_resolved = _skr.resolve_skill_folder(_sk_id)
                            if _is_update and _sk_resolved is None:
                                raise ValueError(f"Skill '{_sk_id}' does not exist. Use create_skill.")
                            if not _is_update and _sk_resolved is not None:
                                raise ValueError(f"Skill '{_sk_id}' already exists. Use update_skill to modify it.")

                            # Security scan the instruction body. HIGH risk blocks unless the
                            # admin explicitly overrides; the result is recorded either way.
                            from vaf.skills.scanner import scan_skill_md_text as _scan_text_fn, format_findings as _fmt_findings
                            _scan = _scan_text_fn(_content)
                            _blocked = _scan["level"] == "high" and not cmd.get("override")
                            try:
                                from vaf.skills.scanner import emit_skill_security_event as _emit_sk
                                if _blocked:
                                    _emit_sk("skill_blocked", "editor", _sk_id, _scan)
                                elif _scan["level"] == "high":
                                    _emit_sk("skill_override", "editor", _sk_id, _scan)
                            except Exception:
                                pass
                            if _blocked:
                                await websocket.send_json({
                                    "type": "skill_error",
                                    "error": _fmt_findings(_scan),
                                    "scan": {"score": _scan["score"], "level": _scan["level"], "findings": _scan["findings"]},
                                    "can_override": True,
                                })
                            else:
                                _skr.save_skill_md(_sk_id, _content)

                                # shared_with: explicit value wins; on update keep existing if omitted.
                                _shared = cmd.get("shared_with")
                                if _shared is None:
                                    _existing = _skr.get_skill_manifest_entry(_sk_id) or {}
                                    _shared = _existing.get("shared_with", ["*"])
                                _skr.register_skill(
                                    _sk_id,
                                    created_by=(manager.get_connection_username(websocket) or "admin"),
                                    shared_with=_shared,
                                    scan=_scan,
                                )
                                _reload_sk()
                                await websocket.send_json({
                                    "type": "skill_updated" if _is_update else "skill_created",
                                    "skill_id": _sk_id,
                                })
                                await _broadcast_skills_update(manager)
                        except Exception as e:
                            await websocket.send_json({"type": "skill_error", "error": str(e)})

                elif type == "delete_skill":
                    from vaf.core.config import get_local_admin_scope_id
                    _sk_scope = manager.get_connection_user(websocket)
                    _sk_role  = manager.get_connection_user_role(websocket)
                    _sk_is_admin = (
                        _sk_role == "admin"
                        or (_sk_scope is not None and str(_sk_scope) == str(get_local_admin_scope_id()))
                    )
                    if not _sk_is_admin:
                        await websocket.send_json({"type": "skill_error", "error": "Admin permission required to delete skills."})
                    else:
                        try:
                            from vaf.core import skills_registry as _skr
                            from vaf.skills.templates import reload_skills as _reload_sk
                            _sk_id = (cmd.get("skill_id") or "").strip()
                            _skr.delete_skill(_sk_id)
                            _reload_sk()
                            await websocket.send_json({"type": "skill_deleted", "skill_id": _sk_id})
                            await _broadcast_skills_update(manager)
                        except Exception as e:
                            await websocket.send_json({"type": "skill_error", "error": str(e)})

                elif type == "update_skill_permissions":
                    # Change which users can see a skill (mirror update_custom_tool_permissions).
                    try:
                        from vaf.core.config import get_local_admin_scope_id
                        from vaf.core import skills_registry as _skr
                        _sk_scope = manager.get_connection_user(websocket) if manager else None
                        _sk_role  = manager.get_connection_user_role(websocket) if manager else None
                        _sk_is_admin = (
                            _sk_role == "admin"
                            or (_sk_scope is not None and str(_sk_scope) == str(get_local_admin_scope_id()))
                        )
                        if not _sk_is_admin:
                            await websocket.send_json({"type": "skill_error", "error": "Admin permission required to change skill permissions."})
                        else:
                            _sk_id = (cmd.get("skill_id") or "").strip()
                            _shared = cmd.get("shared_with", ["*"])
                            _skr.update_skill_permissions(_sk_id, _shared)
                            await websocket.send_json({"type": "skill_permissions_updated", "skill_id": _sk_id, "shared_with": _shared})
                            await _broadcast_skills_update(manager)
                    except Exception as e:
                        await websocket.send_json({"type": "skill_error", "error": str(e)})

                elif type == "upload_skill":
                    # Import an uploaded skill .zip (folder bundle). Payload:
                    #   { data: base64(zip), shared_with?: list[str] }
                    from vaf.core.config import get_local_admin_scope_id
                    _sk_scope = manager.get_connection_user(websocket)
                    _sk_role  = manager.get_connection_user_role(websocket)
                    _sk_is_admin = (
                        _sk_role == "admin"
                        or (_sk_scope is not None and str(_sk_scope) == str(get_local_admin_scope_id()))
                    )
                    if not _sk_is_admin:
                        await websocket.send_json({"type": "skill_error", "error": "Admin permission required to upload skills."})
                    else:
                        _tmpzip = None
                        try:
                            import base64 as _b64m
                            import os as _os_sk
                            import tempfile as _tf_sk
                            from vaf.core import skills_registry as _skr
                            from vaf.skills.templates import reload_skills as _reload_sk
                            from vaf.skills.scanner import SkillScanBlocked as _SkillScanBlocked

                            _raw_bytes = _b64m.b64decode(cmd.get("data") or "")
                            if not _raw_bytes:
                                raise ValueError("no zip data received")
                            _fd, _tmpzip = _tf_sk.mkstemp(suffix=".zip")
                            with _os_sk.fdopen(_fd, "wb") as _f:
                                _f.write(_raw_bytes)
                            try:
                                _sid = _skr.import_skill_zip(
                                    _tmpzip,
                                    created_by=(manager.get_connection_username(websocket) or "admin"),
                                    shared_with=cmd.get("shared_with") or ["*"],
                                    override=bool(cmd.get("override")),
                                )
                            except _SkillScanBlocked as _blk:
                                await websocket.send_json({
                                    "type": "skill_error",
                                    "error": str(_blk),
                                    "scan": {"score": _blk.scan["score"], "level": _blk.scan["level"], "findings": _blk.scan["findings"]},
                                    "can_override": True,
                                })
                            else:
                                _reload_sk()
                                await websocket.send_json({"type": "skill_created", "skill_id": _sid})
                                await _broadcast_skills_update(manager)
                        except Exception as e:
                            await websocket.send_json({"type": "skill_error", "error": str(e)})
                        finally:
                            if _tmpzip:
                                try:
                                    import os as _os_sk2
                                    _os_sk2.unlink(_tmpzip)
                                except OSError:
                                    pass

                elif type == "get_trusted_sources":
                    try:
                        payload = _get_trusted_sources_for_ui()
                        await websocket.send_json({
                            "type": "trusted_sources_list",
                            **payload
                        })
                    except Exception as e:
                        await websocket.send_json({
                            "type": "trusted_sources_list",
                            "categories": [],
                            "error": str(e)
                        })

                elif type == "add_trusted_source":
                    try:
                        from urllib.parse import urlparse
                        category_id = (cmd.get("category_id") or "").strip() or "custom"
                        name = (cmd.get("name") or "").strip()
                        url = (cmd.get("url") or "").strip()
                        if not name or not url:
                            await websocket.send_json({"type": "trusted_source_updated", "ok": False, "error": "name and url required"})
                        else:
                            try:
                                parsed = urlparse(url if "://" in url else "https://" + url)
                                netloc = parsed.netloc or parsed.path.split("/")[0]
                                domain = netloc.lower().lstrip("www.")
                            except Exception:
                                domain = url.replace("https://", "").replace("http://", "").split("/")[0].lower().lstrip("www.")
                            custom = Config.get("trusted_sources_custom") or {}
                            if category_id not in custom:
                                custom[category_id] = []
                            custom[category_id].append({"name": name, "url": url, "domains": [domain], "trust_score": 6})
                            cfg = Config.load()
                            cfg["trusted_sources_custom"] = custom
                            Config.save(cfg)
                            payload = _get_trusted_sources_for_ui()
                            await websocket.send_json({"type": "trusted_source_updated", "ok": True, "categories": payload.get("categories", [])})
                    except Exception as e:
                        await websocket.send_json({"type": "trusted_source_updated", "ok": False, "error": str(e)})

                elif type == "create_trusted_category":
                    try:
                        name = (cmd.get("name") or "").strip()
                        if not name:
                            await websocket.send_json({"type": "trusted_source_updated", "ok": False, "error": "Category name required"})
                        else:
                            payload = _get_trusted_sources_for_ui()
                            existing_names = { (c.get("name") or c.get("id") or "").strip() for c in payload.get("categories", []) }
                            if name in existing_names:
                                await websocket.send_json({"type": "trusted_source_updated", "ok": False, "error": "Category name already exists"})
                            else:
                                custom = Config.get("trusted_sources_custom") or {}
                                custom[name] = []
                                cfg = Config.load()
                                cfg["trusted_sources_custom"] = custom
                                Config.save(cfg)
                                payload = _get_trusted_sources_for_ui()
                                await websocket.send_json({"type": "trusted_source_updated", "ok": True, "categories": payload.get("categories", [])})
                    except Exception as e:
                        await websocket.send_json({"type": "trusted_source_updated", "ok": False, "error": str(e)})

                elif type == "delete_trusted_category":
                    try:
                        category_id = (cmd.get("category_id") or cmd.get("categoryId") or "").strip()
                        if not category_id:
                            await websocket.send_json({"type": "trusted_source_updated", "ok": False, "error": "category_id required"})
                        else:
                            custom = dict(Config.get("trusted_sources_custom") or {})
                            if category_id not in custom:
                                await websocket.send_json({"type": "trusted_source_updated", "ok": False, "error": "Only custom categories can be deleted"})
                            else:
                                del custom[category_id]
                                cfg = Config.load()
                                cfg["trusted_sources_custom"] = custom
                                Config.save(cfg)
                                payload = _get_trusted_sources_for_ui()
                                await websocket.send_json({"type": "trusted_source_updated", "ok": True, "categories": payload.get("categories", [])})
                    except Exception as e:
                        await websocket.send_json({"type": "trusted_source_updated", "ok": False, "error": str(e)})

                elif type == "remove_trusted_source":
                    try:
                        domain = (cmd.get("domain") or "").strip().lower()
                        is_custom = bool(cmd.get("is_custom"))
                        if not domain:
                            await websocket.send_json({"type": "trusted_source_updated", "ok": False, "error": "domain required"})
                        elif is_custom:
                            custom = Config.get("trusted_sources_custom") or {}
                            for cid in list(custom.keys()):
                                custom[cid] = [s for s in custom[cid] if domain not in (s.get("domains") or [])]
                            cfg = Config.load()
                            cfg["trusted_sources_custom"] = custom
                            Config.save(cfg)
                        else:
                            disabled = list(Config.get("trusted_sources_disabled") or [])
                            if domain not in disabled:
                                disabled.append(domain)
                            cfg = Config.load()
                            cfg["trusted_sources_disabled"] = disabled
                            Config.save(cfg)
                        payload = _get_trusted_sources_for_ui()
                        await websocket.send_json({"type": "trusted_source_updated", "ok": True, "categories": payload.get("categories", [])})
                    except Exception as e:
                        await websocket.send_json({"type": "trusted_source_updated", "ok": False, "error": str(e)})

                elif type == "get_automations":
                    # Return list of saved automations; each user sees only their own (user_scope_id).
                    # Admin (local_admin_scope): also sees root automations (e.g. daily calendar check).
                    try:
                        from vaf.core.automation import AutomationManager
                        from vaf.core.config import get_local_admin_scope_id
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        stored_role = manager.get_connection_user_role(websocket) if manager else None
                        local_admin_scope = get_local_admin_scope_id()
                        is_admin = (stored_role == "admin") or (
                            user_scope_id is not None and str(user_scope_id) == str(local_admin_scope)
                        )
                        if is_admin:
                            # Admin: load root + all user dirs, then filter to root + admin-visible scopes.
                            # Some admin JWTs can have a different scope than local_admin_scope.
                            mgr = AutomationManager()
                            all_tasks = list(mgr.list())
                            tasks = [
                                t for t in all_tasks
                                if (
                                    t.user_scope_id is None
                                    or str(t.user_scope_id) == str(user_scope_id)
                                    or str(t.user_scope_id) == str(local_admin_scope)
                                )
                            ]
                        else:
                            mgr = AutomationManager(user_scope_id=user_scope_id) if user_scope_id else AutomationManager()
                            tasks = list(mgr.list())
                        automations_list = [
                            {
                                "id": task.id,
                                "name": task.name,
                                "description": task.description,
                                "prompt": getattr(task, "prompt", "") or task.description,
                                "frequency": task.frequency,
                                "time": task.time,
                                "weekday": task.weekday,
                                "day": task.day,
                                "enabled": task.enabled,
                                "next_run": task.next_run_iso,
                                "last_run": task.last_run
                            }
                            for task in tasks
                        ]
                        await websocket.send_json({
                            "type": "automations_list",
                            "automations": automations_list
                        })
                    except Exception as e:
                        await websocket.send_json({
                            "type": "automations_list",
                            "automations": [],
                            "error": str(e)
                        })

                elif type == "create_automation":
                    try:
                        from vaf.core.automation import AutomationManager, AutomationTask
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        mgr = AutomationManager(user_scope_id=user_scope_id) if user_scope_id else AutomationManager()
                        prompt = (cmd.get("prompt") or "").strip()
                        if not prompt:
                            await websocket.send_json({
                                "type": "create_automation_result",
                                "ok": False,
                                "error": "prompt is required"
                            })
                            continue
                        frequency = (cmd.get("frequency") or "daily").lower()
                        time_str = (cmd.get("time") or "06:00").strip()
                        if ":" not in time_str:
                            time_str = f"{int(time_str or 0) % 24:02d}:00"
                        else:
                            parts = time_str.split(":", 1)
                            h = max(0, min(23, int(parts[0] or 0)))
                            m = max(0, min(59, int(parts[1] or 0) if len(parts) > 1 else 0))
                            time_str = f"{h:02d}:{m:02d}"
                        name = (cmd.get("name") or "").strip() or (prompt[:50] + ("..." if len(prompt) > 50 else ""))
                        description = (cmd.get("description") or "").strip() or prompt[:200]
                        weekday = (cmd.get("weekday") or "").strip().lower() or None
                        day = cmd.get("day")
                        if frequency == "weekly" and weekday:
                            pass
                        elif frequency == "monthly" and day is not None:
                            day = max(1, min(31, int(day)))
                        else:
                            if frequency == "weekly":
                                weekday = None
                            if frequency == "monthly":
                                day = None
                        task = AutomationTask(
                            name=name,
                            description=description,
                            prompt=prompt,
                            frequency=frequency,
                            time=time_str,
                            weekday=weekday if frequency == "weekly" else None,
                            day=day if frequency == "monthly" else None,
                            enabled=True,
                            user_scope_id=user_scope_id,
                        )
                        can_create, err_msg = mgr.check_can_create_automation(new_time=time_str, new_frequency=frequency)
                        if not can_create and err_msg:
                            await websocket.send_json({
                                "type": "create_automation_result",
                                "ok": False,
                                "error": err_msg[:500]
                            })
                            continue
                        mgr.create(task)
                        await websocket.send_json({
                            "type": "create_automation_result",
                            "ok": True,
                            "automation": {
                                "id": task.id,
                                "name": task.name,
                                "description": task.description,
                                "frequency": task.frequency,
                                "time": task.time,
                                "enabled": task.enabled,
                                "next_run": task.next_run_iso,
                                "last_run": task.last_run,
                            }
                        })
                    except Exception as e:
                        await websocket.send_json({
                            "type": "create_automation_result",
                            "ok": False,
                            "error": str(e)
                        })

                elif type == "delete_automation":
                    try:
                        from vaf.core.automation import AutomationManager
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        mgr = AutomationManager(user_scope_id=user_scope_id) if user_scope_id else AutomationManager()
                        task_id = (cmd.get("task_id") or cmd.get("id") or "").strip()
                        if not task_id:
                            await websocket.send_json({"type": "delete_automation_result", "ok": False, "error": "task_id required"})
                            continue
                        ok = mgr.delete(task_id, permanent=True)
                        if not ok and user_scope_id:
                            root_mgr = AutomationManager()
                            ok = root_mgr.delete(task_id, permanent=True)
                        await websocket.send_json({"type": "delete_automation_result", "ok": ok})
                    except Exception as e:
                        await websocket.send_json({"type": "delete_automation_result", "ok": False, "error": str(e)})

                elif type == "update_automation":
                    try:
                        from vaf.core.automation import AutomationManager
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        mgr = AutomationManager(user_scope_id=user_scope_id) if user_scope_id else AutomationManager()
                        task_id = (cmd.get("task_id") or cmd.get("id") or "").strip()
                        if not task_id:
                            await websocket.send_json({"type": "update_automation_result", "ok": False, "error": "task_id required"})
                            continue
                        task = mgr.get(task_id)
                        if not task and user_scope_id:
                            root_mgr = AutomationManager()
                            task = root_mgr.get(task_id)
                            if task:
                                mgr = root_mgr
                        if not task:
                            await websocket.send_json({"type": "update_automation_result", "ok": False, "error": "Automation not found"})
                            continue
                        update_params = {}
                        for key in ("name", "description", "prompt", "frequency", "time", "weekday", "day", "enabled"):
                            if key in cmd and cmd[key] is not None:
                                if key == "enabled":
                                    update_params[key] = bool(cmd[key])
                                elif key == "day":
                                    try:
                                        update_params[key] = max(1, min(31, int(cmd[key])))
                                    except (TypeError, ValueError):
                                        pass
                                elif key == "weekday" and isinstance(cmd.get(key), str):
                                    update_params[key] = (cmd[key] or "").strip().lower() or None
                                else:
                                    update_params[key] = cmd[key]
                        if not update_params:
                            await websocket.send_json({"type": "update_automation_result", "ok": False, "error": "No fields to update"})
                            continue
                        if "time" in update_params:
                            new_time = update_params["time"]
                            if isinstance(new_time, str) and ":" in new_time:
                                parts = new_time.strip().split(":", 1)
                                h = max(0, min(23, int(parts[0] or 0)))
                                m = max(0, min(59, int(parts[1] or 0) if len(parts) > 1 else 0))
                                update_params["time"] = f"{h:02d}:{m:02d}"
                                new_time = update_params["time"]
                            new_frequency = update_params.get("frequency", task.frequency)
                            can_update, err_msg = mgr.check_can_update_automation(task_id=task_id, new_time=str(new_time), new_frequency=new_frequency)
                            if not can_update and err_msg:
                                await websocket.send_json({
                                    "type": "update_automation_result",
                                    "ok": False,
                                    "error": (err_msg[:500] if err_msg else "Time too close to another automation"),
                                })
                                continue
                        updated = mgr.update(task_id, **update_params)
                        if not updated:
                            await websocket.send_json({"type": "update_automation_result", "ok": False, "error": "Update failed"})
                            continue
                        await websocket.send_json({
                            "type": "update_automation_result",
                            "ok": True,
                            "automation": {
                                "id": updated.id,
                                "name": updated.name,
                                "description": updated.description,
                                "frequency": updated.frequency,
                                "time": updated.time,
                                "enabled": updated.enabled,
                                "next_run": updated.next_run_iso,
                                "last_run": updated.last_run,
                            }
                        })
                    except Exception as e:
                        await websocket.send_json({"type": "update_automation_result", "ok": False, "error": str(e)})

                elif type == "get_automation_notes":
                    try:
                        from vaf.core.automation_planner import list_notes
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        notes = list_notes(user_scope_id)
                        await websocket.send_json({"type": "automation_notes_list", "notes": notes})
                    except Exception as e:
                        await websocket.send_json({"type": "automation_notes_list", "notes": [], "error": str(e)})

                elif type == "get_automation_todos":
                    try:
                        from vaf.core.automation_planner import list_todos
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        todos = list_todos(user_scope_id)
                        await websocket.send_json({"type": "automation_todos_list", "todos": todos})
                    except Exception as e:
                        await websocket.send_json({"type": "automation_todos_list", "todos": [], "error": str(e)})

                elif type == "get_thinking_workspace_handoffs":
                    try:
                        from vaf.core.thinking_workspace import list_pending_handoffs

                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        handoffs = list_pending_handoffs(user_scope_id)
                        await websocket.send_json({"type": "thinking_workspace_handoffs_list", "handoffs": handoffs})
                    except Exception as e:
                        await websocket.send_json({"type": "thinking_workspace_handoffs_list", "handoffs": [], "error": str(e)})

                elif type == "approve_thinking_workspace_handoff":
                    try:
                        from vaf.core.thinking_workspace import approve_handoff, get_handoff
                        from vaf.core.user_notifications import append_notification

                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        task_id = (cmd.get("task_id") or "").strip()
                        handoff_id = (cmd.get("handoff_id") or "").strip()
                        if not task_id or not handoff_id:
                            await websocket.send_json({
                                "type": "thinking_workspace_handoff_result",
                                "ok": False,
                                "action": "approve",
                                "error": "task_id and handoff_id required",
                            })
                            continue
                        ok = approve_handoff(user_scope_id, task_id, handoff_id)
                        handoff = get_handoff(user_scope_id, task_id, handoff_id) if ok else None
                        action_result = handoff.get("automation_action_result") if isinstance(handoff, dict) else None
                        action_failed = isinstance(action_result, dict) and (action_result.get("ok") is False)
                        notif_status = "error" if (not ok or action_failed) else "success"
                        summary_parts = [
                            f"Action: approve",
                            f"Task: {task_id}",
                            f"Handoff: {handoff_id}",
                        ]
                        if isinstance(action_result, dict):
                            op = str(action_result.get("operation") or "").strip()
                            ok_txt = "ok" if action_result.get("ok") else "failed"
                            if op:
                                summary_parts.append(f"Automation action: {op} ({ok_txt})")
                            if action_result.get("task_id"):
                                summary_parts.append(f"Automation id: {action_result.get('task_id')}")
                            if action_result.get("error"):
                                summary_parts.append(f"Error: {action_result.get('error')}")
                        append_notification(
                            user_scope_id,
                            kind="system",
                            title="Workspace handoff approved",
                            status=notif_status,
                            summary="\n".join(summary_parts),
                            task_id=task_id,
                            handoff_id=handoff_id,
                            action="approve",
                            automation_action_result=action_result if isinstance(action_result, dict) else None,
                        )
                        await websocket.send_json(
                            {
                                "type": "new_log",
                                "entry": {
                                    "timestamp": time.time(),
                                    "message": (
                                        f"Thinking Workspace handoff approved ({handoff_id})."
                                        + (
                                            " Automation action failed."
                                            if action_failed
                                            else (
                                                " Automation action applied."
                                                if isinstance(action_result, dict) and action_result.get("ok") is True
                                                else ""
                                            )
                                        )
                                    ),
                                    "source": "System",
                                },
                            }
                        )
                        await websocket.send_json({
                            "type": "thinking_workspace_handoff_result",
                            "ok": ok,
                            "action": "approve",
                            "task_id": task_id,
                            "handoff_id": handoff_id,
                            "automation_action_result": action_result,
                            "error": None if ok else "handoff not found",
                        })
                    except Exception as e:
                        await websocket.send_json({
                            "type": "thinking_workspace_handoff_result",
                            "ok": False,
                            "action": "approve",
                            "error": str(e),
                        })

                elif type == "reject_thinking_workspace_handoff":
                    try:
                        from vaf.core.thinking_workspace import reject_handoff
                        from vaf.core.user_notifications import append_notification

                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        task_id = (cmd.get("task_id") or "").strip()
                        handoff_id = (cmd.get("handoff_id") or "").strip()
                        reason = (cmd.get("reason") or "").strip()
                        if not task_id or not handoff_id:
                            await websocket.send_json({
                                "type": "thinking_workspace_handoff_result",
                                "ok": False,
                                "action": "reject",
                                "error": "task_id and handoff_id required",
                            })
                            continue
                        ok = reject_handoff(user_scope_id, task_id, handoff_id, reason=reason)
                        append_notification(
                            user_scope_id,
                            kind="system",
                            title="Workspace handoff rejected",
                            status="success" if ok else "error",
                            summary="\n".join(
                                [
                                    "Action: reject",
                                    f"Task: {task_id}",
                                    f"Handoff: {handoff_id}",
                                    f"Reason: {(reason or 'n/a')[:500]}",
                                ]
                            ),
                            task_id=task_id,
                            handoff_id=handoff_id,
                            action="reject",
                        )
                        await websocket.send_json(
                            {
                                "type": "new_log",
                                "entry": {
                                    "timestamp": time.time(),
                                    "message": f"Thinking Workspace handoff rejected ({handoff_id}).",
                                    "source": "System",
                                },
                            }
                        )
                        await websocket.send_json({
                            "type": "thinking_workspace_handoff_result",
                            "ok": ok,
                            "action": "reject",
                            "task_id": task_id,
                            "handoff_id": handoff_id,
                            "error": None if ok else "handoff not found",
                        })
                    except Exception as e:
                        await websocket.send_json({
                            "type": "thinking_workspace_handoff_result",
                            "ok": False,
                            "action": "reject",
                            "error": str(e),
                        })

                elif type == "get_notifications":
                    try:
                        from vaf.core.user_notifications import get_notifications
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        limit = min(100, max(1, int(cmd.get("limit") or 50)))
                        notifications = get_notifications(user_scope_id, limit=limit)
                        await websocket.send_json({"type": "notifications_list", "notifications": notifications})
                    except Exception as e:
                        await websocket.send_json({"type": "notifications_list", "notifications": [], "error": str(e)})

                elif type == "create_automation_note":
                    try:
                        from vaf.core.automation_planner import add_note
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        title = cmd.get("title")
                        content = (cmd.get("content") or "").strip()
                        if not content:
                            await websocket.send_json({"type": "create_automation_note_result", "ok": False, "error": "content required"})
                            continue
                        note = add_note(user_scope_id, content, title=title)
                        await websocket.send_json({"type": "create_automation_note_result", "ok": True, "note": note})
                    except Exception as e:
                        await websocket.send_json({"type": "create_automation_note_result", "ok": False, "error": str(e)})

                elif type == "create_automation_todo":
                    try:
                        from vaf.core.automation_planner import add_todo
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        text = (cmd.get("text") or "").strip()
                        if not text:
                            await websocket.send_json({"type": "create_automation_todo_result", "ok": False, "error": "text required"})
                            continue
                        due_at = cmd.get("due_at")
                        todo = add_todo(user_scope_id, text, due_at=due_at)
                        await websocket.send_json({"type": "create_automation_todo_result", "ok": True, "todo": todo})
                    except Exception as e:
                        await websocket.send_json({"type": "create_automation_todo_result", "ok": False, "error": str(e)})

                elif type == "update_automation_todo":
                    try:
                        from vaf.core.automation_planner import update_todo
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        todo_id = (cmd.get("id") or "").strip()
                        if not todo_id:
                            await websocket.send_json({"type": "update_automation_todo_result", "ok": False, "error": "id required"})
                            continue
                        text = cmd.get("text")
                        done = cmd.get("done")
                        due_at = cmd.get("due_at")
                        updated = update_todo(user_scope_id, todo_id, text=text, done=done, due_at=due_at)
                        if not updated:
                            await websocket.send_json({"type": "update_automation_todo_result", "ok": False, "error": "Todo not found"})
                            continue
                        await websocket.send_json({"type": "update_automation_todo_result", "ok": True, "todo": updated})
                    except Exception as e:
                        await websocket.send_json({"type": "update_automation_todo_result", "ok": False, "error": str(e)})

                elif type == "delete_automation_note":
                    try:
                        from vaf.core.automation_planner import delete_note
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        note_id = (cmd.get("id") or "").strip()
                        if not note_id:
                            await websocket.send_json({"type": "delete_automation_note_result", "ok": False, "error": "id required"})
                            continue
                        ok = delete_note(user_scope_id, note_id)
                        await websocket.send_json({"type": "delete_automation_note_result", "ok": ok, "id": note_id if ok else None})
                    except Exception as e:
                        await websocket.send_json({"type": "delete_automation_note_result", "ok": False, "error": str(e)})

                elif type == "delete_automation_todo":
                    try:
                        from vaf.core.automation_planner import delete_todo
                        user_scope_id = manager.get_connection_user(websocket) if manager else None
                        todo_id = (cmd.get("id") or "").strip()
                        if not todo_id:
                            await websocket.send_json({"type": "delete_automation_todo_result", "ok": False, "error": "id required"})
                            continue
                        ok = delete_todo(user_scope_id, todo_id)
                        await websocket.send_json({"type": "delete_automation_todo_result", "ok": ok, "id": todo_id if ok else None})
                    except Exception as e:
                        await websocket.send_json({"type": "delete_automation_todo_result", "ok": False, "error": str(e)})

                elif type == "process_audio":
                    # Process audio for STT: Docker (HTTP) or local (faster-whisper)
                    import base64
                    import tempfile
                    
                    temp_path = None
                    try:
                        audio_b64 = cmd.get("audio")
                        if not audio_b64:
                            await websocket.send_json({
                                "type": "stt_error",
                                "error": "No audio data provided"
                            })
                            continue

                        audio_data = base64.b64decode(audio_b64)
                        # Use format hint from frontend (wav preferred for Whisper compatibility)
                        audio_format = cmd.get("format", "webm")
                        suffix = ".wav" if audio_format == "wav" else ".webm"
                        mime_type = "audio/wav" if audio_format == "wav" else "audio/webm"

                        with tempfile.NamedTemporaryFile(prefix="vaf_", suffix=suffix, delete=False) as temp_audio:
                            temp_audio.write(audio_data)
                            temp_path = temp_audio.name

                        from vaf.core.speech import SpeechManager
                        sm = SpeechManager.get_instance()
                        if not sm.is_stt_enabled():
                            await websocket.send_json({
                                "type": "stt_error",
                                "error": "STT is disabled in settings"
                            })
                            continue

                        text = ""

                        from vaf.core import speech_api as _speech_api
                        if _speech_api.resolve_stt_engine() == "docker":
                            # Shared speech client: cloud provider lane first when
                            # speech_stt_provider is set, then Docker Whisper
                            # (/asr with /transcribe fallback, parsing inside).
                            from vaf.core import speech_client
                            loop = asyncio.get_running_loop()
                            filename = f"audio{suffix}"
                            # Per-speaker language hint keyed on the user's scope (the
                            # language is a trait of the speaker, and the key is
                            # user-isolated); "local" for the no-scope local admin.
                            _stt_key = (manager.get_connection_user(websocket) if manager else None) or "local"
                            stt_text, _stt_lang = await loop.run_in_executor(
                                None,
                                lambda: speech_client.transcribe(temp_path, mime=mime_type, filename=filename, cache_key=_stt_key),
                            )
                            if not stt_text:
                                # Detail is in the server log (speech_client warns).
                                await websocket.send_json({
                                    "type": "stt_error",
                                    "error": "Docker STT failed. Is the STT container running (e.g. docker compose -f docker-compose.memory.yml up -d)?"
                                })
                                continue
                            text = stt_text
                        else:
                            # Local STT: faster-whisper (SINGLETON)
                            try:
                                model = get_whisper_model()
                                segments, info = model.transcribe(temp_path, beam_size=5)
                                text = " ".join([segment.text for segment in segments])
                            except ImportError as import_error:
                                # Name the SETTING that led here and the REAL reason.
                                # "not installed" used to be printed for every
                                # ImportError, including an installed faster-whisper
                                # whose native dependency does not load (a common
                                # arch mismatch on macOS), which sent people looking
                                # in the wrong place. No installer ships the package,
                                # so choosing this engine at all is the usual cause.
                                await websocket.send_json({
                                    "type": "stt_error",
                                    "error": (
                                        "Local STT (faster-whisper) is selected in "
                                        "Settings > Voice, but it is not usable: "
                                        f"{import_error} Switch the engine to Docker "
                                        "in Settings to use the Whisper container instead."
                                    )
                                })
                                continue
                            except Exception as transcribe_error:
                                await websocket.send_json({
                                    "type": "stt_error",
                                    "error": f"Transcription failed: {str(transcribe_error)}"
                                })
                                continue

                        # Speaker identification: label the utterance against the
                        # user's voice profile (fail-closed: any problem -> no label).
                        _speaker = None
                        try:
                            from vaf.core import speaker_id as _sid
                            if suffix == ".wav" and _sid.is_enabled() and text.strip():
                                from vaf.core.config import get_local_admin_scope_id
                                _scope = str(manager.get_connection_user(websocket)
                                             or get_local_admin_scope_id())
                                if _sid.load_profile(_scope) is not None:
                                    _wav_bytes = open(temp_path, "rb").read()
                                    loop = asyncio.get_running_loop()
                                    _speaker = await loop.run_in_executor(
                                        None, lambda: _sid.score_wav(_wav_bytes, _scope))
                        except Exception:
                            _speaker = None

                        _payload = {"type": "stt_result", "text": text.strip()}
                        if _speaker:
                            try:
                                from vaf.core import speaker_id as _sid
                                _prof = _sid.load_profile(_scope)
                                _name = ((_prof or {}).get("meta") or {}).get("display_name", "Ich")
                                _payload["text"] = _sid.label_prefix(_speaker, _name) + _payload["text"]
                                _payload["speaker_label"] = _speaker["label"]
                                _payload["speaker_score"] = _speaker["score"]
                            except Exception:
                                pass
                            if _speaker.get("label") != "self":
                                # Non-owner mic turn: speaker_confirm decides whether to
                                # ask the owner - a spoofing check when the transcript
                                # CLAIMS to be the owner, or a restrained adaptive reclaim
                                # on a plain unsure (one pending per scope, cooldown inside).
                                try:
                                    from vaf.core import speaker_confirm as _sc
                                    _uname = manager.get_connection_username(websocket) or "admin"
                                    await loop.run_in_executor(
                                        None, lambda: _sc.maybe_request_confirmation(
                                            _scope, _uname, _wav_bytes, _speaker,
                                            transcript=text, owner_name=_name))
                                except Exception:
                                    pass
                        await websocket.send_json(_payload)
                    except Exception as e:
                        await websocket.send_json({
                            "type": "stt_error",
                            "error": str(e)
                        })
                    finally:
                        if temp_path and os.path.exists(temp_path):
                            try:
                                os.unlink(temp_path)
                            except Exception:
                                pass
                
                elif type == "speak":
                    # TTS Output: only if Auto TTS is enabled; otherwise skip to avoid loading TTS when disabled.
                    text = cmd.get("text")
                    print(f"[WebSocket] speak request received, text_len={len(text) if text else 0}")
                    if text:
                        tts_enabled = Config.get("speech_tts_enabled", False)
                        print(f"[WebSocket] speech_tts_enabled = {tts_enabled}")
                        if not tts_enabled:
                            print("[WebSocket] TTS disabled, sending stopped state")
                            await websocket.send_json({"type": "tts_state", "status": "stopped"})
                        else:
                            from vaf.core.speech import SpeechManager
                            sm = SpeechManager.get_instance()
                            await websocket.send_json({"type": "tts_state", "status": "loading"})
                            lang = _detect_language_simple(text)
                            print(f"[WebSocket] Detected language: {lang}")
                            import base64
                            # asyncio is already imported at top of file
                            loop = asyncio.get_running_loop()
                            # Local Docker TTS answers within seconds (LAN); cloud
                            # providers legitimately need longer for long texts
                            # (observed: gpt-4o-mini-tts ~45s for a 4-minute answer),
                            # so the budget scales with the configured lane.
                            from vaf.core import speech_api as _sa
                            _TTS_SYNTH_TIMEOUT = 130.0 if _sa.select_tts_backend()[0] else 35.0
                            try:
                                configured_engine = Config.get("speech_tts_engine", "docker")
                                if configured_engine not in ("docker", "chatterbox"):
                                    configured_engine = "docker"
                                print(f"[WebSocket] Calling synthesize_audio with engine={configured_engine}")
                                audio_bytes = await asyncio.wait_for(
                                    loop.run_in_executor(
                                        None,
                                        lambda: sm.synthesize_audio(text, lang, force_engine=configured_engine),
                                    ),
                                    timeout=_TTS_SYNTH_TIMEOUT,
                                )
                                print(f"[WebSocket] synthesize_audio returned {len(audio_bytes) if audio_bytes else 0} bytes")
                                if audio_bytes:
                                    audio_b64 = base64.b64encode(audio_bytes).decode('utf-8')
                                    print(f"[WebSocket] Sending tts_audio, base64 len={len(audio_b64)}")
                                    await websocket.send_json({
                                        "type": "tts_audio",
                                        "audio": audio_b64,
                                        "format": "wav"
                                    })
                                else:
                                    print("[WebSocket] No audio bytes, sending stopped state")
                                    await websocket.send_json({"type": "tts_state", "status": "stopped"})
                            except asyncio.TimeoutError:
                                log("WebServer", f"TTS synthesize timed out after {_TTS_SYNTH_TIMEOUT}s (WebSocket continues)")
                                await websocket.send_json({"type": "tts_state", "status": "stopped"})
                            except Exception as e:
                                print(f"[WebServer] TTS Error: {e}")
                                import traceback
                                traceback.print_exc()
                                await websocket.send_json({"type": "tts_state", "status": "stopped"})
                
                elif type == "stop_speech":
                    # Stop TTS
                    from vaf.core.speech import SpeechManager
                    sm = SpeechManager.get_instance()
                    sm.stop()

                # ---- Speaker identification (voice profile, guided enrollment) ----
                # All replies are per-connection (websocket.send_json) and all
                # state is keyed on the connection's user_scope_id (isolation).
                # ---- Live-Call: voice-agent first layer ----
                # The tuple is the OUTER door and must list every voice_call_*
                # type the branches below handle. A branch whose type is missing
                # here is dead code that looks alive: the streaming endpointer
                # lane shipped that way (chunk/reset handled inside, never
                # admitted), so the browser streamed microphone audio the server
                # silently discarded while the feature read as available.
                # tests/test_voice_call_wiring.py pins door and branches together.
                elif type in ("voice_call_start", "voice_call_turn", "voice_call_end",
                              "voice_call_chunk", "voice_call_reset"):
                    # ONE door into the voice engine: everything engine-shaped in
                    # this block goes through voice_turn (the same discipline as
                    # agent.py consuming the tool pipeline only through
                    # tool_dispatch). Auth, Config, TTS, TaskQueue and tray
                    # concerns are the handler's own half and stay direct.
                    from vaf.core import voice_turn as _vt
                    from vaf.core.config import get_local_admin_scope_id
                    _scope = str(manager.get_connection_user(websocket) or get_local_admin_scope_id())
                    _conn_key = id(websocket)

                    if type == "voice_call_start":
                        # Base language: identity preferred_language first,
                        # then the configured default_language (the user's
                        # chosen default voice language), then the UI locale.
                        # Per-turn language follow can override on top.
                        from vaf.core.config import Config as _CfgL
                        _cfg_default = (_CfgL.get("default_language", "") or "").strip()
                        _lang = (_cfg_default or cmd.get("ui_lang") or "de")[:2].lower()
                        try:
                            from vaf.auth.user_workspace import get_user_workspace
                            _vuser = manager.get_connection_username(websocket)
                            if _vuser:
                                _pref = ((get_user_workspace(_vuser).get_user_identity() or {})
                                         .get("preferred_language") or "")
                                if str(_pref).strip():
                                    _lang = str(_pref).strip()[:2].lower()
                        except Exception:
                            pass
                        # Chat context: a compact structural digest of the OPEN
                        # chat so the voice agent knows what "here" refers to.
                        # Ownership-gated: never build it from another user's
                        # session (fail-closed skip).
                        _chat_ctx = ""
                        try:
                            _vc_sid = (cmd.get("sessionId") or "").strip()
                            if _vc_sid:
                                _ok_sess, _sess = _ws_session_owner_ok(websocket, _vc_sid)
                                if _ok_sess and _sess is not None:
                                    _chat_ctx = _vt.build_chat_digest(
                                        getattr(_sess, "messages", None) or [])
                        except Exception:
                            _chat_ctx = ""
                        # Agent persona name for the wake-word filter (fuzzy
                        # "you were addressed by name" - see voice_agent.
                        # addressed_by_name). Cached per call, fail-open "".
                        _agent_name = ""
                        _agent_soul = ""
                        try:
                            if _vuser:
                                _ws_p = get_user_workspace(_vuser)
                                _agent_name = str((_ws_p.get_identity() or {}).get("name") or "")
                                # Personality for the first layer: capped, the
                                # voice prompt is latency-bound.
                                _agent_soul = str(_ws_p.get_soul() or "")[:500]
                        except Exception:
                            _agent_name, _agent_soul = "", ""
                        _VOICE_CALLS[_conn_key] = {"history": [], "lang": _lang,
                                                   "scope": _scope, "chat_context": _chat_ctx,
                                                   "agent_name": _agent_name,
                                                   "agent_soul": _agent_soul,
                                                   # Reflex system (VOICE_REFLEX.md): a
                                                   # session/scope key for the durable
                                                   # rolling transcript, and a small
                                                   # ring of recent chime-ins for dedup.
                                                   "session": (cmd.get("sessionId") or str(_conn_key)),
                                                   "chime_recent": []}
                        # The turn pipeline lives in the engine (vaf/core/voice_turn.py);
                        # it SHARES this record as its state - the registry keeps its
                        # exact dict shape and id(websocket) key, which the tray's
                        # model-keepalive probe and the teardown pops depend on.
                        try:
                            from vaf.core.voice_turn import VoiceTurnEngine
                            _VOICE_CALLS[_conn_key]["engine"] = VoiceTurnEngine(
                                _VOICE_CALLS[_conn_key],
                                log=lambda msg: log("WebServer", msg))
                        except Exception as _eng_e:
                            log("WebServer", f"voice_call engine init failed: {_eng_e}")
                        # Server-side semantic endpointing (opt-in): one StreamEndpointer
                        # PER CALL (per-stream Silero state + the caller's own audio ring,
                        # Rule 4.4 - never shared). The frontend only streams PCM chunks
                        # when `server_endpoint` in voice_call_started says the lane is on.
                        try:
                            _ep_new = _vt.make_endpointer()
                            if _ep_new is not None:
                                _VOICE_CALLS[_conn_key]["endpointer"] = _ep_new
                        except Exception:
                            pass
                        # Pre-warm the speaker-ID extractor in the background: it lazy-loads
                        # on the first score_wav (~seconds), and during that window the owner
                        # scores as 'unsure' and is treated as a guest (formal replies, a
                        # needless 'did you mean me?', side-talk silence) - live solo call:
                        # the owner was mislabeled for the whole cold-load window. Loading it
                        # now means the very first turn is scored correctly. Best-effort, off
                        # the event loop, only when a profile exists to score against.
                        try:
                            asyncio.get_running_loop().run_in_executor(
                                None, lambda: _vt.prewarm_speaker(_scope))
                        except Exception:
                            pass
                        _lane = _vt.lane_status()
                        _lane_ok = _lane["available"]
                        _lane_reason = None
                        if _lane_ok and _lane["dedicated_model"] and not _vt.subagents_hold_model():
                            # Server up but possibly holding the MAIN model
                            # (e.g. right after a chat): pre-warm the swap so
                            # the first voice turn does not pay it.
                            try:
                                _vt.prepare_dedicated_model_async()
                            except Exception:
                                pass
                        if not _lane_ok:
                            _vm_ref = _lane["dedicated_model"]
                            if _vt.subagents_hold_model():
                                # Never load/swap while a sub-agent computes
                                # on the one model: the call opens in the
                                # busy state and heals after the result.
                                _lane_reason = "model_loading"
                            elif _vm_ref:
                                # Dedicated voice model configured: load THAT
                                # (download on first use), not the main model.
                                # On success push model_state so the frontend
                                # self-heal re-sends voice_call_start.
                                _lane_reason = "model_loading"

                                def _vm_ready(ok: bool) -> None:
                                    if not ok:
                                        return
                                    try:
                                        tray_context.set_model_loaded(True)
                                        get_web_interface().push_update({
                                            "type": "model_state",
                                            "loaded": True,
                                            "persistent": tray_context.is_persistent(),
                                            "provider": "local",
                                        })
                                    except Exception:
                                        pass

                                _vt.prepare_dedicated_model_async(on_ready=_vm_ready)
                            elif _lane["exclusive"]:
                                # Local mode with the llama server down: load
                                # the MAIN model directly in a worker thread
                                # (closes the former headless gap - the tray
                                # watchdog is desktop-only and register_
                                # activity alone did nothing without it).
                                # ensure_local_model is lock-serialized and
                                # start_server reuses/waits, so a tray that
                                # also reacts to the heartbeat cannot spawn a
                                # second server. The frontend shows "loading"
                                # and re-sends voice_call_start once the
                                # model_state push reports loaded.
                                _lane_reason = "model_loading"
                                try:
                                    tray_context.register_activity()
                                except Exception:
                                    pass

                                def _main_load_work():
                                    try:
                                        from vaf.core.backend import (ServerManager,
                                                                      ensure_local_model)
                                        _mgr = ServerManager(skip_cleanup=True)
                                        _mpath = _mgr.ensure_model_present()
                                        if _mpath and os.path.exists(_mpath) and \
                                                ensure_local_model(_mpath, reason="voice call start"):
                                            tray_context.set_model_loaded(True)
                                            get_web_interface().push_update({
                                                "type": "model_state",
                                                "loaded": True,
                                                "persistent": tray_context.is_persistent(),
                                                "provider": "local",
                                            })
                                    except Exception as _ml_e:
                                        log("WebServer", f"voice_call main-model load failed: {_ml_e}")

                                threading.Thread(target=_main_load_work,
                                                 name="voice-main-model-load",
                                                 daemon=True).start()
                            else:
                                _lane_reason = "no_model"
                            log("WebServer",
                                f"voice_call_start: lane not ready (reason={_lane_reason})")
                        await websocket.send_json({
                            "type": "voice_call_started",
                            "ok": _lane_ok,
                            "lang": _lang,
                            # exclusive: ONE local model time-shared with the
                            # main agent - the frontend mutes the voice agent
                            # while a delegated task runs.
                            "exclusive": _lane["exclusive"],
                            "reason": _lane_reason,
                            # server_endpoint: the semantic turn-end lane is armed for
                            # this call - the frontend streams mic PCM chunks and lets
                            # voice_turn_end trigger the stop it would otherwise take
                            # from its own silence timer (which stays as the fallback).
                            "server_endpoint": bool(_VOICE_CALLS.get(_conn_key, {}).get("endpointer")),
                        })
                        # Greeting: the agent opens the call (deterministic
                        # line, no LLM round-trip) - also an audio check for
                        # the user. Skipped when the voice lane is deaf.
                        if _lane_ok:
                            try:
                                _greet = _vt.greeting_for(_VOICE_CALLS[_conn_key])
                                from vaf.core.speech import SpeechManager
                                _gsm = SpeechManager.get_instance()
                                loop = asyncio.get_running_loop()
                                _gaudio = await asyncio.wait_for(
                                    loop.run_in_executor(
                                        None,
                                        lambda: _gsm.synthesize_audio(_greet, _lang, force_engine="docker"),
                                    ),
                                    timeout=30.0,
                                )
                                _call_rec = _VOICE_CALLS.get(_conn_key)
                                if _gaudio and _call_rec is not None:
                                    import base64 as _b64g
                                    # Success-gated on purpose: a greeting whose TTS
                                    # failed was never heard, so it never happened.
                                    _eng_g = _call_rec.get("engine")
                                    if _eng_g is not None:
                                        _eng_g.note_greeting(_greet)
                                    else:
                                        _call_rec["history"].append(
                                            {"role": "assistant", "content": _greet})
                                    await websocket.send_json({
                                        "type": "voice_call_reply",
                                        "user_text": "",
                                        "speaker_label": None,
                                        "reply": _greet,
                                        "audio": _b64g.b64encode(_gaudio).decode("utf-8"),
                                        "delegated": None,
                                        "greeting": True,
                                    })
                            except Exception as _greet_e:
                                log("WebServer", f"voice_call greeting failed: {_greet_e}")

                    elif type == "voice_call_end":
                        _ended = _VOICE_CALLS.pop(_conn_key, None)
                        # Retention: drop the rolling transcript at call end (it is
                        # context, not a record - the age cap would prune it anyway).
                        try:
                            if _ended is not None:
                                _eng_end = _ended.get("engine")
                                if _eng_end is not None:
                                    _eng_end.end()
                                else:
                                    _vt.clear_transcript(_ended.get("scope"), _ended.get("session"))
                        except Exception:
                            pass
                        await websocket.send_json({"type": "voice_call_ended"})

                    elif type == "voice_call_chunk":
                        # Live mic PCM (16 kHz mono int16, base64) for the server-side
                        # endpointer. Pure OBSERVATION: the audio that becomes the turn
                        # still arrives complete via voice_call_turn - this stream only
                        # decides WHEN. The endpointer keys on the connection's call
                        # record, so one user's stream can never touch another's state.
                        _call = _VOICE_CALLS.get(_conn_key)
                        _ep = (_call or {}).get("endpointer")
                        if _ep is None:
                            continue   # lane off or no call: chunks are simply ignored
                        try:
                            import base64 as _b64c
                            _pcm = _b64c.b64decode(cmd.get("pcm") or "")
                            if not _pcm:
                                continue
                            loop = asyncio.get_running_loop()
                            # Silero + (optionally) the judge run off the event loop,
                            # like every other blocking step in the live turn.
                            _verdict = await loop.run_in_executor(None, _ep.feed_pcm, _pcm)
                            if _verdict == "end":
                                await websocket.send_json(
                                    {"type": "voice_turn_end", "reason": "semantic"})
                            elif _verdict == "hold":
                                log("WebServer", "voice_call: semantic endpoint HELD an "
                                                 "unfinished sentence open")
                        except Exception:
                            pass   # fail-open: the browser timer still ends the turn

                    elif type == "voice_call_reset":
                        # Mute toggled / utterance discarded on the client: the stream
                        # state no longer describes reality - drop it entirely.
                        _ep = (_VOICE_CALLS.get(_conn_key) or {}).get("endpointer")
                        if _ep is not None:
                            try:
                                _ep.reset()
                            except Exception:
                                pass

                    elif type == "voice_call_turn":
                        _call = _VOICE_CALLS.get(_conn_key)
                        if _call is None:
                            await websocket.send_json({"type": "voice_call_error", "error": "no_call"})
                            continue
                        # OWNERSHIP. The call itself is keyed on the CONNECTION, so it cannot be
                        # hijacked - but a `sessionId` from the payload is used twice deep in this
                        # turn, and both uses act on whatever it names: a speaker-confirmation
                        # request is routed into it, and a delegation is enqueued into it via
                        # TaskQueue().add(session_id=...) carrying THIS caller's scope in the
                        # metadata. That is an injection into another user's work queue, not
                        # merely a misrouted prompt, which is why this branch is gated even though
                        # it never loads a session.
                        if cmd.get("sessionId"):
                            _ok_vc, _ = _ws_session_owner_ok(websocket, cmd.get("sessionId"),
                                                             allow_missing=True)
                            if not _ok_vc:
                                await websocket.send_json({"type": "voice_call_error",
                                                           "error": "session_not_owned"})
                                continue
                        import base64 as _b64v
                        _audio_b64 = cmd.get("audio") or ""
                        if cmd.get("format") != "wav" or not _audio_b64:
                            await websocket.send_json({"type": "voice_call_error", "error": "bad_format"})
                            continue
                        _wav = _b64v.b64decode(_audio_b64)
                        loop = asyncio.get_running_loop()

                        # The pipeline is the engine (vaf/core/voice_turn.py) - THIS branch is
                        # its consumer: one executor hop in, one wire message out. What stays
                        # here is the wire (decode above, send_json below), TTS with its
                        # per-variant timeouts, and the TaskQueue enqueue - the one external
                        # write, kept visible in the handler on purpose.
                        from vaf.core.voice_turn import VoiceTurnEngine, timings_report
                        _eng = _call.get("engine")
                        if _eng is None:   # engine init failed at start: build it late
                            _eng = VoiceTurnEngine(_call, log=lambda msg: log("WebServer", msg))
                            _call["engine"] = _eng
                        _tm_browser = cmd.get("timings") or {}
                        _uname_t = manager.get_connection_username(websocket) or ""
                        _sid_t = cmd.get("sessionId") or ""
                        _main_busy_t = bool(cmd.get("main_busy"))
                        _outcome = await loop.run_in_executor(
                            None,
                            lambda: _eng.turn(
                                _wav, session_id=_sid_t, main_busy=_main_busy_t,
                                pending_task=(cmd.get("pending_task") or ""),
                                username=_uname_t))

                        if _outcome.error is not None:
                            await websocket.send_json({"type": "voice_call_error",
                                                       "error": _outcome.error})
                            continue

                        if _outcome.kind == "silent":
                            await websocket.send_json({
                                "type": "voice_call_reply", "user_text": _outcome.user_text,
                                "speaker_label": _outcome.speaker_label, "reply": "",
                                "audio": None, "delegated": None, "silent": True,
                            })
                            continue

                        if _outcome.kind in ("clarify", "reask", "chime_in"):
                            # Short spoken lines (30 s) and the chime-in (60 s, follows the
                            # remark's own language) - TTS failure degrades to audio=None,
                            # exactly as before; state was already written by the engine.
                            _v_audio = None
                            try:
                                from vaf.core.speech import SpeechManager as _SMv
                                _v_sm = _SMv.get_instance()
                                _v_lang = (_tts_lang_for(_outcome.reply, _outcome.tts_lang, _v_sm)
                                           if _outcome.tts_follow else _outcome.tts_lang)
                                _va_out = await asyncio.wait_for(
                                    loop.run_in_executor(
                                        None, lambda: _v_sm.synthesize_audio(
                                            _outcome.reply, _v_lang, force_engine="docker")),
                                    timeout=60.0 if _outcome.kind == "chime_in" else 30.0)
                                if _va_out:
                                    _v_audio = _b64v.b64encode(_va_out).decode("utf-8")
                            except Exception:
                                _v_audio = None
                            await websocket.send_json({
                                "type": "voice_call_reply", "user_text": _outcome.user_text,
                                "speaker_label": _outcome.speaker_label,
                                "reply": _outcome.reply, "audio": _v_audio,
                                "delegated": None, **_outcome.flags,
                            })
                            continue

                        # kind == "reply": forensic line first (pre-TTS marks, as always),
                        # then the delegation enqueue, then the reply TTS (130 s, follows
                        # the reply's language), then the payload with full timings.
                        log("WebServer",
                            f"voice_call turn: active={_outcome.active_s:.1f}s "
                            f"label={_outcome.speaker_label or '-'} "
                            f"speaker_ok={_outcome.speaker_ok} text={_outcome.user_text[:80]!r} -> "
                            f"reply_len={len(_outcome.reply)} "
                            f"delegate={'yes' if _outcome.delegate else 'no'} "
                            f"reply={_outcome.reply[:100]!r} "
                            f"timings={timings_report(_outcome.marks, _tm_browser)}")

                        # 4. Delegation -> normal main-agent task in this session. The engine
                        # DECIDED; the enqueue - the one external write - happens here.
                        _delegated = None
                        if _outcome.delegate:
                            try:
                                from vaf.core.task_queue import TaskQueue
                                TaskQueue().add(
                                    session_id=_sid_t,
                                    input_text=_outcome.delegate,
                                    source="web",
                                    metadata={
                                        "user_scope_id": _call["scope"],
                                        "origin_channel": "voice_call",
                                    },
                                )
                                _delegated = _outcome.delegate
                            except Exception as _dele:
                                log("WebServer", f"voice_call delegation failed: {_dele}")

                        # 5. Speak the reply (provider lane first inside synthesize_audio).
                        # Speak it in ITS OWN language, not the input language: a model reply
                        # in another language (e.g. Turkish from a German turn) is voiced by
                        # that language's voice when the lane can speak it, else stays on the
                        # call language.
                        _reply_audio = None
                        try:
                            from vaf.core.speech import SpeechManager
                            _sm = SpeechManager.get_instance()
                            _tts_lang = _tts_lang_for(_outcome.reply, _outcome.tts_lang, _sm)
                            if _tts_lang != _outcome.tts_lang:
                                log("WebServer",
                                    f"voice_call: TTS follows reply language "
                                    f"{_outcome.tts_lang}->{_tts_lang}")
                            _audio_out = await asyncio.wait_for(
                                loop.run_in_executor(
                                    None,
                                    lambda: _sm.synthesize_audio(
                                        _outcome.reply, _tts_lang, force_engine="docker"),
                                ),
                                timeout=130.0,
                            )
                            if _audio_out:
                                _reply_audio = _b64v.b64encode(_audio_out).decode("utf-8")
                        except Exception:
                            _reply_audio = None
                        _outcome.marks["tts"] = __import__("time").perf_counter()

                        await websocket.send_json({
                            "type": "voice_call_reply",
                            "user_text": _outcome.user_text,
                            "speaker_label": _outcome.speaker_label,
                            "reply": _outcome.reply,
                            "audio": _reply_audio,
                            "delegated": _delegated,
                            # Per-stage latency (ms): browser endpoint_wait/encode +
                            # server gate/stt/speaker/policy/llm/tts. Read by the
                            # VoiceCallLayer debug line; absent fields simply not shown.
                            "timings": timings_report(_outcome.marks, _tm_browser),
                        })

                elif type in ("speaker_enroll_start", "speaker_enroll_round",
                              "speaker_enroll_finalize", "speaker_enroll_abort",
                              "speaker_enroll_speak", "voice_call_speak",
                              "speaker_profile_get", "speaker_profile_delete"):
                    from vaf.core import speaker_id as _sid
                    from vaf.core.config import get_local_admin_scope_id
                    _scope = manager.get_connection_user(websocket) or get_local_admin_scope_id()
                    _scope = str(_scope)

                    if type == "speaker_enroll_start":
                        info = _sid.enroll_start(_scope, ui_lang=cmd.get("ui_lang", "en"))
                        await websocket.send_json({"type": "speaker_enroll_started", **info})

                    elif type == "speaker_enroll_round":
                        import base64 as _b64
                        _audio_b64 = cmd.get("audio") or ""
                        _fmt = cmd.get("format", "")
                        if _fmt != "wav" or not _audio_b64:
                            await websocket.send_json({
                                "type": "speaker_enroll_round_result",
                                "ok": False, "quality": "bad_format",
                            })
                        else:
                            _wav = _b64.b64decode(_audio_b64)
                            loop = asyncio.get_running_loop()
                            _res = await loop.run_in_executor(
                                None, lambda: _sid.enroll_round(_scope, _wav))
                            await websocket.send_json(
                                {"type": "speaker_enroll_round_result", **_res})

                    elif type == "speaker_enroll_finalize":
                        _name = (cmd.get("display_name") or "").strip() or "Ich"
                        loop = asyncio.get_running_loop()
                        _meta = await loop.run_in_executor(
                            None, lambda: _sid.enroll_finalize(_scope, _name))
                        await websocket.send_json({
                            "type": "speaker_profile",
                            "profile": _meta,
                            "saved": _meta is not None,
                        })

                    elif type == "speaker_enroll_abort":
                        _sid.enroll_abort(_scope)
                        await websocket.send_json({"type": "speaker_enroll_aborted"})

                    elif type in ("speaker_enroll_speak", "voice_call_speak"):
                        # The agent's voice during enrollment/live calls. Distinct
                        # reply type so the chat TTS handler in page.tsx never
                        # reacts to it. Falls back to captions-only (audio=null).
                        _text = (cmd.get("text") or "").strip()
                        if type == "voice_call_speak" and _text:
                            # Strip any reasoning the main model leaked into the result
                            # content (<think>...</think>, or an UNCLOSED block on a stuck/
                            # truncated stream): reasoning must NEVER reach TTS (VOICE_AGENT.md
                            # invariant 3) and must not pollute the call history either. This
                            # also guards the stuck-local-model case where the streamed
                            # content is thinking-only.
                            from vaf.core.voice_turn import strip_spoken_result as _srx
                            _text = _srx(_text)
                            # A delegated result is being announced: the voice
                            # agent's own call history must know it, or on the
                            # next turn it still believes the task is running
                            # (live incident: "the search is still running").
                            _vc = _VOICE_CALLS.get(id(websocket))
                            if _vc is not None:
                                _eng_s = _vc.get("engine")
                                if _eng_s is not None:
                                    _eng_s.note_spoken(_text)
                                else:
                                    _vc["history"].append(
                                        {"role": "assistant", "content": _text[:800]})
                                    _vc["history"] = _vc["history"][-16:]
                        _audio_b64 = None
                        if _text:
                            try:
                                from vaf.core.speech import SpeechManager
                                _sm = SpeechManager.get_instance()
                                _lang = _detect_language_simple(_text)
                                loop = asyncio.get_running_loop()
                                _audio = await asyncio.wait_for(
                                    loop.run_in_executor(
                                        None,
                                        lambda: _sm.synthesize_audio(_text, _lang, force_engine="docker"),
                                    ),
                                    timeout=60.0,
                                )
                                if _audio:
                                    import base64 as _b64s
                                    _audio_b64 = _b64s.b64encode(_audio).decode("utf-8")
                            except Exception:
                                _audio_b64 = None
                        await websocket.send_json({"type": "speaker_enroll_tts", "audio": _audio_b64})

                    elif type == "speaker_profile_get":
                        _prof = _sid.load_profile(_scope)
                        await websocket.send_json({
                            "type": "speaker_profile",
                            "profile": (_prof or {}).get("meta"),
                            "enabled": _sid.is_enabled(),
                        })

                    elif type == "speaker_profile_delete":
                        _sid.delete_profile(_scope)
                        _sid.enroll_abort(_scope)
                        await websocket.send_json({"type": "speaker_profile", "profile": None})

                elif type == "stop_generation":
                    # Stop the current generation by setting a flag
                    from vaf.core.task_queue import TaskQueue
                    from vaf.core.platform import Platform
                    tq = TaskQueue()
                    session_id = cmd.get("sessionId") or manager.get_session_for_connection(websocket)
                    # OWNERSHIP. The id comes straight from the payload, and below it stops the
                    # generation for that session and cancels its in-flight attachment indexing.
                    # Without the gate that is a denial of service against another user's work -
                    # no data is read or written, which is why this ranks below
                    # set_sidebar_documents, but "only" interrupting someone else's run is still
                    # acting on a session that is not yours.
                    if session_id:
                        _ok_stop, _ = _ws_session_owner_ok(websocket, session_id, allow_missing=True)
                        if not _ok_stop:
                            log("API", f"Access denied: stop_generation {session_id} not owned by caller")
                            continue
                    # Scope (chat-while-subagent-runs): a Stop press while a sub-agent is
                    # active stops ONLY the generation — cancelling a throwaway chat reply
                    # must not destroy a 20-minute coder run. Killing the sub-agent needs
                    # the EXPLICIT scope "all" (the UI's dedicated stop-sub-agent action).
                    # With no active sub-agent, behavior is unchanged.
                    stop_scope = str(cmd.get("scope") or "").strip().lower()
                    if session_id:
                        tq.request_stop(session_id)
                        # Cancel any in-flight attachment indexing for this session — the kill
                        # switch must stop the background RAG indexing too, not just generation.
                        try:
                            for _idx_task in list(_active_index_tasks.get(str(session_id), set())):
                                _idx_task.cancel()
                        except Exception:
                            pass
                        dropped = 0
                        try:
                            # Also drop already queued follow-up tasks for this session.
                            # Otherwise users may need to press stop multiple times.
                            dropped = tq.drop_queued_tasks_for_session(str(session_id))
                        except Exception:
                            dropped = 0
                        _has_active_subagents = False
                        try:
                            from vaf.core.subagent_ipc import get_ipc
                            ipc = get_ipc()
                            _has_active_subagents = bool(ipc.get_active_tasks(session_id=str(session_id)))
                        except Exception:
                            _has_active_subagents = False
                        killed = 0
                        subagents_kept = False
                        if _has_active_subagents and stop_scope != "all":
                            # Scoped stop: generation halted, the sub-agent keeps working.
                            subagents_kept = True
                        else:
                            try:
                                # Hard-stop any running sub-agent processes for this chat session.
                                killed = Platform.stop_webui_subagent_processes(str(session_id))
                            except Exception:
                                killed = 0
                            # Also convert active sub-agent tasks into failed results immediately.
                            try:
                                from vaf.core.subagent_ipc import get_ipc
                                ipc = get_ipc()
                                active = ipc.get_active_tasks(session_id=str(session_id))
                                for t in active:
                                    try:
                                        ipc.fail_task(t.task_id, "[USER_CANCELLED] Stopped/Cancelled by user via stop button.")
                                    except Exception:
                                        pass
                            except Exception:
                                pass
                        log("WebServer", f"Stop requested for session {session_id}; scope={stop_scope or 'default'}; killed_subagents={killed}; subagents_kept={subagents_kept}; dropped_queued={dropped}")
                        await websocket.send_json({
                            "type": "generation_stopped",
                            "sessionId": session_id,
                            "subagentsKept": subagents_kept,
                        })

                elif type == "regenerate_last_reply":
                    # Ask the newest exchange again. The person judged the answer
                    # unwanted and said so with a click; the transcript rewinds to
                    # their message and it is asked once more.
                    #
                    # This is not the bounce the never-erase invariant forbids
                    # (vaf/core/agent.py, the team-await gate): that one fired on a
                    # heuristic, mid-turn, against an answer the person wanted. This
                    # fires on an explicit click, between turns, on a settled bubble.
                    # Which is why the gates below are absolute rather than advisory.
                    from vaf.core.task_queue import TaskQueue
                    from vaf.core.session import truncate_to_last_user_turn
                    session_id = cmd.get("sessionId") or manager.get_session_for_connection(websocket)
                    if not session_id:
                        continue
                    # OWNERSHIP. The id comes from the payload and this DELETES stored
                    # messages, so the gate is the strict one: there is nothing to
                    # regenerate in a session that does not exist.
                    _ok_regen, _ = _ws_session_owner_ok(websocket, session_id)
                    if not _ok_regen:
                        log("API", f"Access denied: regenerate_last_reply {session_id} not owned by caller")
                        await websocket.send_json({"type": "error", "message": "Access denied"})
                        continue
                    _tq_regen = TaskQueue()
                    # Never mid-run. The running turn holds an index into the live
                    # history it captured when it started, so cutting underneath it
                    # misaligns every later slice with no error at all, and the runner
                    # writes the whole session back when it finishes, which would undo
                    # the cut anyway.
                    if _tq_regen.is_busy_for_session(str(session_id)):
                        await websocket.send_json({
                            "type": "regenerate_refused", "sessionId": session_id, "reason": "busy",
                        })
                        continue
                    # Never while a specialist of this chat is still working. Its result
                    # is delivered later and would answer a question no longer in the
                    # transcript, and a sub-agent result is never dropped silently.
                    # PENDING counts as working: a task exists from the moment it is
                    # created, and it only becomes "active" once its process is up, so
                    # checking active alone lets a spawn still in flight slip through.
                    try:
                        from vaf.core.subagent_ipc import get_ipc as _get_ipc_regen
                        _ipc_regen = _get_ipc_regen()
                        if (_ipc_regen.get_active_tasks(session_id=str(session_id))
                                or _ipc_regen.get_pending_tasks(session_id=str(session_id))):
                            await websocket.send_json({
                                "type": "regenerate_refused", "sessionId": session_id, "reason": "subagent",
                            })
                            continue
                    except Exception:
                        pass
                    try:
                        _sm_regen = SessionManager()
                        # repoint=False: reading another session must not change what
                        # "current" means, and must not push its runtime state into the
                        # live registry.
                        _loaded_regen = _sm_regen.load(str(session_id), restore_state=False, repoint=False)
                        _removed_regen = truncate_to_last_user_turn(_loaded_regen)
                        _user_text = str(getattr(_removed_regen, "content", "") or "") if _removed_regen else ""
                        if not _user_text:
                            await websocket.send_json({
                                "type": "regenerate_refused", "sessionId": session_id, "reason": "nothing",
                            })
                            continue
                        _sm_regen.save(_loaded_regen)
                    except Exception as _regen_err:
                        log("WebServer", f"regenerate_last_reply failed for {session_id}: {_regen_err}")
                        await websocket.send_json({
                            "type": "regenerate_refused", "sessionId": session_id, "reason": "error",
                        })
                        continue
                    # The browser keeps its own copy of the transcript and re-injects
                    # anything the server no longer sends as an "orphan". Without this
                    # the discarded answer comes straight back.
                    await websocket.send_json({"type": "context_checkpoint", "session_id": str(session_id)})
                    _regen_msgs = _history_projection(_loaded_regen, str(session_id))
                    # The question goes back on screen right away. It was just removed
                    # from the record and the runner only writes it again when the turn
                    # ends, so without this the new answer would appear under the
                    # previous one with nothing asking it, and a turn that uses a tool
                    # would group the two answers into one bubble.
                    from datetime import datetime as _dt_regen
                    _regen_ask = {"role": "user", "content": _user_text,
                                  "timestamp": _dt_regen.now().isoformat()}
                    _regen_imgs = ((getattr(_removed_regen, "metadata", None) or {}).get("images")) or []
                    if _regen_imgs:
                        from urllib.parse import quote as _q_regen
                        _regen_ask["images"] = [
                            {"url": (f"/api/file?path={_q_regen(str(_i['path']))}" if _i.get("path")
                                     else f"data:{_i.get('mime_type', 'image/jpeg')};base64,{_i.get('data', '')}"),
                             "name": _i.get("name", "image")}
                            for _i in _regen_imgs
                        ]
                    _regen_msgs.append(_regen_ask)
                    # Every view of this chat, not only the tab that pressed: a second
                    # window on the same session would otherwise keep showing the reply
                    # that no longer exists.
                    await manager.broadcast_to_session(str(session_id), {
                        "type": "history_update",
                        "messages": _regen_msgs,
                        "sessionId": str(session_id),
                        "isActive": True,
                        "currentStatus": "thinking",
                    })
                    _regen_scope = manager.get_connection_user(websocket)
                    _regen_username = manager.get_connection_username(websocket)
                    _regen_role = manager.get_connection_user_role(websocket)
                    _regen_meta = {"origin_channel": "web", "task_class": "interactive",
                                   "enqueue_session_id": str(session_id),
                                   # The stored transcript changed under a live agent, whose
                                   # in-memory history is otherwise authoritative: without
                                   # this the worker answers from the reply just discarded.
                                   "force_reload": True}
                    if _regen_scope:
                        _regen_meta["user_scope_id"] = _regen_scope
                        _regen_meta["enqueue_user_scope_id"] = str(_regen_scope)
                    if _regen_username:
                        _regen_meta["username"] = _regen_username
                    if _regen_role:
                        _regen_meta["role"] = _regen_role
                    # The question travels whole, image included. Re-sending the text
                    # alone would re-ask about a picture the model can no longer see,
                    # and the record would lose the attachment with it.
                    if _regen_imgs:
                        _regen_meta["images"] = _regen_imgs
                    manager.subscribe_to_session(websocket, str(session_id))
                    _tq_regen.add(session_id=str(session_id), input_text=_user_text,
                                  source="web", metadata=_regen_meta)
                    log("WebServer", f"Regenerate requested for session {session_id}")

                elif type == "browser_interactive_start":
                    # A person asks to drive the sandbox browser by hand. Identity
                    # comes from the CONNECTION, never the message; ownership of the
                    # chat session whose window will show the stream is the same gate
                    # every session command passes.
                    session_id = cmd.get("sessionId") or manager.get_session_for_connection(websocket)
                    if not session_id:
                        continue
                    _ok_bi, _ = _ws_session_owner_ok(websocket, session_id, allow_missing=True)
                    if not _ok_bi:
                        log("API", f"Access denied: browser_interactive_start {session_id} not owned by caller")
                        continue
                    _bi_scope = manager.get_connection_user(websocket)
                    _bi_role = manager.get_connection_user_role(websocket)
                    try:
                        from vaf.core.config import is_admin_identity as _is_admin
                        _bi_admin = _is_admin(_bi_role, _bi_scope)
                    except Exception:
                        _bi_admin = False
                    from vaf.core.browser_interactive import get_manager_for_scope
                    # Scope-resolved: with the per-user pool on this may START
                    # the caller's own browser instance, which is exactly why
                    # the whole call sits in the executor below.
                    # start() blocks (pool resolve + container probe + cookie
                    # handover): executor, so a slow browser start cannot
                    # starve every other websocket.
                    _bi_result = await asyncio.get_running_loop().run_in_executor(
                        None,
                        lambda: get_manager_for_scope(_bi_scope).start(
                            _bi_scope or "", str(session_id),
                            # Always the persistent mode: the person's browser asks about
                            # saving logins itself; a client-sent flag would be a second
                            # switch for the same question and is deliberately not read.
                            save=True,
                            session_name=str(cmd.get("session") or "default").strip() or "default",
                            is_admin=_bi_admin,
                        ),
                    )
                    # Direct reply on top of the manager's broadcast: the opener needs
                    # the verdict even when it is a refusal that no session-scoped
                    # broadcast would reach (busy for a foreign user's lease).
                    await websocket.send_json({
                        "type": "browser_interactive_state",
                        "sessionId": session_id,
                        **_bi_result,
                    })

                elif type == "browser_interactive_stop":
                    # The ACTION is authorised inside stop() against the lease scope; the
                    # sessionId only addresses the direct reply. Gated anyway: a branch
                    # reading a client sessionId without _ws_session_owner_ok passed the
                    # ownership guard only through a string match in a comment far below,
                    # and an invariant held by accident is not held.
                    session_id = cmd.get("sessionId") or manager.get_session_for_connection(websocket)
                    if session_id:
                        _ok_bi, _ = _ws_session_owner_ok(websocket, session_id, allow_missing=True)
                        if not _ok_bi:
                            log("API", f"Access denied: browser_interactive_stop {session_id} not owned by caller")
                            continue
                    _bi_scope = manager.get_connection_user(websocket)
                    from vaf.core.browser_interactive import get_manager_for_stop
                    # The stop lands on whichever browser actually holds this
                    # session's (or scope's) lease - and deliberately never
                    # starts a container to have something to stop.
                    _bi_mgr = get_manager_for_stop(_bi_scope or "", str(session_id) if session_id else None)
                    # Owner check happens inside stop() against the lease scope; the
                    # cookie export in there blocks, so executor again.
                    _bi_result = await asyncio.get_running_loop().run_in_executor(
                        None,
                        lambda: _bi_mgr.stop("user", requester_scope=_bi_scope or ""),
                    )
                    if session_id:
                        await websocket.send_json({
                            "type": "browser_interactive_state",
                            "sessionId": session_id,
                            **_bi_result,
                        })

            except json.JSONDecodeError:
                pass
                
    except WebSocketDisconnect:
        manager.disconnect(websocket)
        # Drop any live-call state for this socket. voice_call_end is sent by exactly one
        # frontend site (inside endCall), so an abrupt teardown would leave an orphan behind -
        # and the idle watchdog reads this dict as "a call is live, keep the model loaded",
        # which would pin the local model for the life of the process.
        _VOICE_CALLS.pop(id(websocket), None)
        tray_context.set_websocket_count(len(manager.active_connections)) # Update active count
        log("API", f"WebSocket disconnected. Active: {tray_context.active_websockets}")
        
        # Unregister from connection tracker
        try:
            from vaf.network.connection_tracker import get_tracker
            tracker = get_tracker()
            tracker.unregister_connection(connection_id)
        except Exception:
            pass
        
        # VAF_WEBUI_ACTIVE is NOT popped here any more. It used to mean "a browser is
        # attached right now", and the sub-agent spawn decision read it - so a transient
        # socket drop silently changed process-wide spawn behaviour and a host terminal
        # window popped up for a web-launched sub-agent two minutes later (live incident
        # 2026-07-20). It now means "this process serves a web UI", which is a property of
        # the process and is set once at startup.
    except Exception as e:
        print(f"WebSocket error: {e}")
        manager.disconnect(websocket)
        _VOICE_CALLS.pop(id(websocket), None)   # same reason as the disconnect path above
        tray_context.set_websocket_count(len(manager.active_connections)) # Update active count
        log("API", f"WebSocket error; active now {tray_context.active_websockets}")
        
        # Unregister from connection tracker
        try:
            from vaf.network.connection_tracker import get_tracker
            tracker = get_tracker()
            tracker.unregister_connection(connection_id)
        except Exception:
            pass
        
        # VAF_WEBUI_ACTIVE is NOT popped here any more. It used to mean "a browser is
        # attached right now", and the sub-agent spawn decision read it - so a transient
        # socket drop silently changed process-wide spawn behaviour and a host terminal
        # window popped up for a web-launched sub-agent two minutes later (live incident
        # 2026-07-20). It now means "this process serves a web UI", which is a property of
        # the process and is set once at startup.


# ── agent rooms over the wire ────────────────────────────────────────────────
#
# Deliberately NOT built on the WebUI socket's plumbing, and every omission below is a
# decision rather than an oversight:
#
#   manager.connect()        would accept the socket AND append it to active_connections
#                            (so every WebUI broadcast fans out to a room peer), send it
#                            an unasked-for state_full frame, and make the tray count an
#                            agent as a browser. The socket is accepted directly instead.
#   set_connection_user()    would file a room peer in the same identity table every
#                            ownership check reads, making it indistinguishable from a
#                            browser to _ws_session_owner_ok and everything like it.
#   the local-admin fallback would hand the machine owner's seat to anything that
#                            reaches the port with a token that decodes to nothing. The
#                            WebUI has that fallback so the desktop is not locked out of
#                            its own chats; a room has no such problem, so no credential
#                            means no connection.
#
# The identity work itself lives in vaf/core/a2a/wire.py, which is stricter than the
# WebUI handshake in the ways that matter (an access token only, never a refresh one;
# no scope means no connection; every connection lands in the remote participant lane).

# room_id -> Hub, and room_id -> {peer_id: asyncio.Queue}. Process-local by design: the
# files are the record, so a second process serving the same room is not a second truth,
# only a second accelerator.
_A2A_HUBS: dict = {}
_A2A_OUTBOXES: dict = {}


def _a2a_hub(room):
    """The one hub for a room in this process, created on first use."""
    from vaf.core.a2a.hub import Hub

    existing = _A2A_HUBS.get(room.room_id)
    if existing is not None:
        return existing

    boxes = _A2A_OUTBOXES.setdefault(room.room_id, {})

    def _sink(peer_id, message):
        # Never blocks and never raises: the hub treats delivery as best effort, and a
        # queue that is gone means a peer that disconnected between two frames.
        box = boxes.get(peer_id)
        if box is not None:
            box.put_nowait(message)

    hub = Hub(room, sink=_sink)
    _A2A_HUBS[room.room_id] = hub
    return hub


# ── the VAF-less lane: what a stranger's harness downloads before it can join ──
#
# Both endpoints are deliberately unauthenticated (listed in AUTH_EXEMPT_PATHS):
# a guest holds a join ticket, not an account, and the ticket only works on the
# room socket. Neither response carries user data - the CA certificate is public
# by design (the private key never leaves this machine), and the client file is
# public repository content. What makes the unverified fetches safe for the guest
# is the invitation itself: it carries the CA fingerprint and the file's sha256
# by another route, and the client checks both before anything is spoken.

@app.get("/api/a2a/client.py")
async def a2a_guest_client():
    """The single-file wire client named in every invitation's guest section."""
    from fastapi.responses import Response as _Response

    from vaf.core.a2a.invite import GUEST_CLIENT_REPO_URL, guest_client_path
    path = guest_client_path()
    if path is None:
        # A pip-installed host has no examples/ tree to serve; the repository
        # copy is byte-identical and reachable over publicly verifiable TLS.
        raise HTTPException(status_code=404, detail=(
            "this host cannot serve the client file; fetch the same file from "
            f"{GUEST_CLIENT_REPO_URL}"))
    return _Response(content=path.read_bytes(), media_type="text/x-python",
                     headers={"Content-Disposition": 'attachment; filename="a2a_client.py"'})


@app.get("/api/a2a/ca.pem")
async def a2a_guest_ca():
    """The authority a guest pins, checked against the invitation's fingerprint."""
    from fastapi.responses import Response as _Response

    from vaf.network.ssl_utils import ca_certificate_path
    path = ca_certificate_path()
    if path is None or not path.exists():
        raise HTTPException(status_code=404, detail=(
            "this host has no LAN authority; local network TLS is not enabled"))
    return _Response(content=path.read_bytes(), media_type="application/x-pem-file")


# ── the room workspace over the wire: list, fetch, push ─────────────────────
#
# The workspace is a folder on THIS machine; local members reach it through the
# browser's workspace lane. A remote seat holder had no way to it at all - a
# room collaborating on files could talk about them but never exchange them.
# These three endpoints are SEAT-authenticated: the seat travels as a query
# parameter for the same documented reason the socket's credential does (the
# integrated proxy strips Authorization headers), which means the same log
# exposure the ticket already has. Read and write only - DELETING over the
# wire does not exist on purpose: destruction of the shared folder stays with
# the members on the machine that owns it.

_A2A_WORKSPACE_LIST_CAP = 500
_A2A_WORKSPACE_UPLOAD_CAP = 25 * 1024 * 1024


def _a2a_workspace_for_seat(room_id: str, seat: str, request: Request):
    """(room, workspace_dir) for a valid seat, or an HTTPException.

    Only a SEAT opens this door. A ticket is single-use and redeeming it here
    would burn the join a guest still needs; an account token has the whole
    authenticated API. Refusals mirror the socket door: one neutral message,
    no room-existence oracle, one security event.
    """
    from vaf.core.a2a.room import Room
    from vaf.core.a2a.store import StoreError, UnsafeName

    ip = getattr(getattr(request, "client", None), "host", "") or ""

    def _refuse():
        _emit_sec_ws(f"a2a workspace room={room_id}: seat refused", ip=ip)
        return HTTPException(status_code=403,
                             detail="no seat for that room on this credential")

    seat = str(seat or "").strip()
    if not seat.startswith(Room.SEAT_PREFIX):
        raise _refuse()
    try:
        room = Room.open(room_id)
    except (StoreError, UnsafeName):
        raise _refuse() from None
    try:
        room.redeem_seat(seat)
    except Exception:
        raise _refuse() from None
    workspace = room.workspace_dir(create=True)
    if workspace is None:
        raise HTTPException(status_code=404, detail="this room has no workspace")
    return room, workspace


def _a2a_workspace_member_path(workspace, relative: str):
    """The absolute path a relative name may touch, or an HTTPException.

    Relative, no leading slash, no `..`, and the RESOLVED path must stay under
    the resolved workspace - which also catches a symlink pointing out.
    """
    if not str(relative or "").strip():
        raise HTTPException(status_code=400, detail="path must be relative, without ..")
    try:
        return contained_path(workspace, relative)
    except PathEscape as e:
        raise HTTPException(status_code=400, detail=str(e))


@app.get("/api/a2a/rooms/{room_id}/files")
async def a2a_workspace_list(room_id: str, request: Request,
                             seat: str = Query("")):
    """Every file in the room's shared folder, for a remote seat holder."""
    def _list():
        _room, workspace = _a2a_workspace_for_seat(room_id, seat, request)
        rows = []
        for path in sorted(workspace.rglob("*")):
            if not path.is_file():
                continue
            # POSIX on the wire, whatever this host's OS: the listed path is
            # what a remote client hands back to fetch, and a Windows host
            # would otherwise list `sub\a.bin` for a file pushed as `sub/a.bin`.
            rows.append({"path": path.relative_to(workspace).as_posix(),
                         "size": path.stat().st_size,
                         "mtime": path.stat().st_mtime})
            if len(rows) >= _A2A_WORKSPACE_LIST_CAP:
                break
        return {"room": room_id, "files": rows,
                "capped": len(rows) >= _A2A_WORKSPACE_LIST_CAP}

    return await asyncio.to_thread(_list)


@app.get("/api/a2a/rooms/{room_id}/file")
async def a2a_workspace_fetch(room_id: str, request: Request,
                              seat: str = Query(""), path: str = Query("")):
    """One file's bytes, for a remote seat holder."""
    from fastapi.responses import Response as _Response

    def _fetch():
        _room, workspace = _a2a_workspace_for_seat(room_id, seat, request)
        target = _a2a_workspace_member_path(workspace, path)
        if not target.is_file():
            raise HTTPException(status_code=404, detail="no such file in this workspace")
        return target.read_bytes(), target.name

    content, name = await asyncio.to_thread(_fetch)
    return _Response(content=content, media_type="application/octet-stream",
                     headers={"Content-Disposition": f'attachment; filename="{name}"'})


@app.post("/api/a2a/rooms/{room_id}/file")
async def a2a_workspace_push(room_id: str, request: Request,
                             seat: str = Query(""), path: str = Query("")):
    """Put one file into the room's shared folder, for a remote seat holder."""
    body = await request.body()
    if len(body) > _A2A_WORKSPACE_UPLOAD_CAP:
        raise HTTPException(status_code=413, detail=(
            f"the upload cap is {_A2A_WORKSPACE_UPLOAD_CAP} bytes"))

    # A room's shared folder is the one place in this tree where bytes arrive from a
    # machine that is not ours, put there by an agent that may not be VAF. That makes
    # it the lane with the least standing trust, so it is gated like every other - and
    # the refusal travels back over the wire as a 403 the remote peer can read.
    _verdict = _inspect_attachment(body, path or "upload", "a2a_room",
                                   ip=(request.client.host if request.client else ""))
    if _verdict.blocked:
        raise HTTPException(status_code=403, detail=_verdict.message(path or "the file"))

    def _push():
        _room, workspace = _a2a_workspace_for_seat(room_id, seat, request)
        target = _a2a_workspace_member_path(workspace, path)
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_bytes(body)
        return {"room": room_id,
                "path": target.relative_to(workspace.resolve()).as_posix(),
                "size": len(body)}

    return await asyncio.to_thread(_push)


@app.websocket("/ws/a2a/{room_id}")
async def a2a_room_endpoint(websocket: WebSocket, room_id: str,
                            token: Optional[str] = Query(None)):
    """One agent, one room, over a socket.

    The credential is either an access token (an account on this machine) or a join
    ticket (an invitation to exactly this room), and it arrives in the query string
    because that is the only carrier that survives the integrated proxy: Authorization
    headers and subprotocols are stripped on the relayed leg, silently and without an
    error anywhere.
    """
    import asyncio
    import json as _json

    from vaf.core.a2a.wire import HandshakeRefused, admit, open_room

    client_ip = _ws_client_ip(websocket)
    credential = token
    if not credential and getattr(websocket, "cookies", None):
        credential = websocket.cookies.get(VAF_TOKEN_COOKIE)

    try:
        # In a thread, like every store touch on this route: the handshake reads
        # the room and may redeem a ticket (crypto plus store writes), and this
        # loop is the one every socket shares - a remote client's connect storm
        # must not stall the WebUI beside it.
        room = await asyncio.to_thread(open_room, room_id)
        identity, seat = await asyncio.to_thread(admit, room, credential or "")
    except HandshakeRefused as refusal:
        # Mirrored into the security log exactly like a refused WebUI handshake. A room
        # is the one door a stranger can knock on from another machine, so a refusal
        # here is MORE worth seeing on the dashboard than one on the browser socket,
        # not less.
        log("API", f"A2A rejected from {client_ip}: {refusal.reason}")
        _emit_sec_ws(f"a2a room={room_id}: {refusal.reason}", ip=client_ip)
        await websocket.close(code=refusal.code, reason=refusal.reason)
        return
    except Exception as e:                      # pragma: no cover - defensive
        log("API", f"A2A handshake failed from {client_ip}: {e}")
        _emit_sec_ws(f"a2a room={room_id}: handshake failed", ip=client_ip)
        await websocket.close(code=4001, reason="handshake failed")
        return

    hub = _a2a_hub(room)
    outbox: "asyncio.Queue" = asyncio.Queue()
    _A2A_OUTBOXES.setdefault(room.room_id, {})[identity.peer_id] = outbox

    try:
        # The attach scans the store for the backlog - in a thread, so a large
        # room cannot stall the shared event loop while a peer connects.
        lease = await asyncio.to_thread(hub.attach, identity)
    except Exception as e:
        _A2A_OUTBOXES.get(room.room_id, {}).pop(identity.peer_id, None)
        log("API", f"A2A refused {identity.peer_id} in {room.room_id}: {e}")
        await websocket.close(code=4009, reason=str(e))
        return

    # From the moment the lease exists, EVERY exit runs the finally below. The
    # accept and the welcome used to sit outside this protection: a client that
    # vanished between attach and the welcome (a timed-out dialer hanging up)
    # skipped the detach, and its dead lease refused its own reconnects for the
    # full 90 second TTL - each half-successful retry re-arming another dead
    # lease. Measured live as a room gone permanently mute behind "another
    # connection is writing".
    pump = None
    try:
        await websocket.accept()
        log("API", f"A2A joined: {identity.peer_id} in {room.room_id} from {client_ip}")
        welcome = {
            "kind": "welcome", "room": room.room_id, "peer": identity.peer_id,
            "role": identity.role, "protocol": "vaf-a2a", "v": 1,
        }
        # The room's half of the handshake, for a peer that arrived over the wire -
        # the same packet a local join answers with. An ADDED key, which rule 1 makes
        # free: an older client keeps reading the four fields it knows and ignores
        # this one. Without it the remote lane would be the poorer implementation of
        # our own protocol, which is the drift this file has paid for elsewhere.
        try:
            welcome["welcome"] = await asyncio.to_thread(room.welcome, identity)
        except Exception:
            pass
        if seat:
            # Handed over exactly once, on the connection that redeemed the ticket. It
            # never appears again: the server keeps only the hash, so a client that
            # drops it has to be invited again - which is the honest outcome for a
            # bearer credential nobody wrote down.
            welcome["seat"] = seat
        await websocket.send_text(_json.dumps(welcome))

        async def _pump() -> None:
            while True:
                message = await outbox.get()
                await websocket.send_text(_json.dumps(message, ensure_ascii=False))

        pump = asyncio.create_task(_pump())
        while True:
            raw = await websocket.receive_text()
            try:
                payload = _json.loads(raw)
            except Exception:
                await websocket.send_text(_json.dumps(
                    {"kind": "ack", "status": "malformed", "reason": "not json"}))
                continue
            if not isinstance(payload, dict):
                await websocket.send_text(_json.dumps(
                    {"kind": "ack", "status": "malformed", "reason": "not an object"}))
                continue
            if str(payload.get("kind") or "") == "renew":
                # Transport, never a frame - it does not touch the store. This is
                # the server half of protocol contract C9 ("leases are renewed
                # while attached"): the hub renews only on successful submits, so
                # a held line that read and thought for longer than the lease TTL
                # lost its write right while connected and receiving - measured
                # by the first foreign agent to hold a session against this
                # endpoint. A LAPSED lease is not silently re-taken: the client
                # is told not_writer and reconnects, so its own cursor decides
                # the backlog and nothing is skipped in between.
                from vaf.core.a2a.hub import NotWriter as _NotWriter
                try:
                    await asyncio.to_thread(hub.renew, identity, lease)
                    ack = {"kind": "ack", "status": "renewed"}
                except _NotWriter as e:
                    ack = {"kind": "ack", "status": "not_writer", "reason": str(e)}
                await websocket.send_text(_json.dumps(ack, ensure_ascii=False))
                continue
            ack = await asyncio.to_thread(hub.submit, identity, lease, payload)
            await websocket.send_text(_json.dumps(ack, ensure_ascii=False))
    except WebSocketDisconnect:
        pass
    except Exception as e:
        log("API", f"A2A socket error for {identity.peer_id}: {e}")
    finally:
        if pump is not None:
            pump.cancel()
        _A2A_OUTBOXES.get(room.room_id, {}).pop(identity.peer_id, None)
        try:
            hub.detach(identity, lease)
        except Exception:
            pass


async def process_uploaded_files(files: list) -> str:
    """
    Process uploaded files and return their text content.
    Mimics CLI @filename behavior by reading file contents and formatting them.
    
    Args:
        files: List of file objects with {name, data, mimeType}
        
    Returns:
        Formatted string with file contents
    """
    import base64
    import tempfile
    import os
    from pathlib import Path
    
    if not files:
        return ""
    
    results = []
    
    for file_obj in files:
        try:
            filename = file_obj.get("name", "unknown")
            file_data = file_obj.get("data", "")
            mime_type = file_obj.get("mimeType", "")
            
            print(f"[WebUI] Processing file: {filename} ({mime_type})")
            
            # Decode base64 data
            if file_data.startswith("data:"):
                # Remove data URL prefix (e.g., "data:application/pdf;base64,")
                file_data = file_data.split(",", 1)[1] if "," in file_data else file_data
            
            decoded_data = base64.b64decode(file_data)

            # Known-bad gate. Before the bytes reach the Librarian, the extractors or
            # the model's context, they are hashed and looked up. A hit is refused and
            # the refusal takes the file's place in the transcript, so the user and the
            # agent both see WHY there is no content here.
            _verdict = _inspect_attachment(decoded_data, filename, "web_chat")
            if _verdict.blocked:
                results.append(f"\n\n--- FILE: {filename} ---\n[BLOCKED] "
                               f"{_verdict.message(filename)}\n----------------\n")
                continue

            # Save to temporary file
            file_ext = Path(filename).suffix or ".txt"
            with tempfile.NamedTemporaryFile(prefix="vaf_", suffix=file_ext, delete=False) as temp_file:
                temp_file.write(decoded_data)
                temp_path = temp_file.name

            try:
                # Use Librarian to read file contents
                from vaf.tools.librarian import LibrarianTool
                librarian = LibrarianTool()
                
                # Read file using Librarian's _read_file method
                content = librarian._read_file(Path(temp_path), enable_chunking=True)
                
                # Format like CLI does: --- FILE: name ---\nCONTENTS\n----------------
                formatted = (f"\n\n--- FILE: {filename} ---\n{content}"
                             f"{_advisory_note(_verdict)}\n----------------\n")
                results.append(formatted)
                
                print(f"[WebUI] Successfully processed: {filename}")
                
            finally:
                # Clean up temp file
                try:
                    os.unlink(temp_path)
                except:
                    pass
                    
        except Exception as e:
            error_msg = f"\n\n--- FILE: {filename} ---\n[ERROR] Failed to process file: {str(e)}\n----------------\n"
            results.append(error_msg)
            print(f"[WebUI] Error processing {filename}: {e}")
    
    return "".join(results)


def _is_temp_like_filename(name: str) -> bool:
    """True if name looks like our temp file (vaf_xxxxx.ext), so we should not expose it as document name."""
    if not name or not name.startswith("vaf_"):
        return False
    import re
    return bool(re.match(r"^vaf_[a-zA-Z0-9_]+\.[a-zA-Z0-9]+$", name))


def _persist_attached_images_to_files(attached_images: list, session_id, user_scope_id) -> list:
    """Write uploaded chat images to the user-siloed chat attachments folder and return entries
    that reference them by PATH instead of inline base64. Keeps session.json lean and gives the
    agent a real file location (analyze_image / read_file). Per image, on ANY failure the original
    base64 entry is kept so vision never breaks. Only image/* is written; oversized payloads are
    skipped (kept inline)."""
    import base64 as _b64
    import mimetypes as _mt
    import os as _os
    import re as _re
    import time as _time
    import uuid as _uuid
    from vaf.core.session import get_session_attachments_dir
    try:
        attach_dir = get_session_attachments_dir(session_id, user_scope_id, create=True)
    except Exception:
        attach_dir = None
    if not attach_dir:
        return attached_images
    _MAX_BYTES = 25 * 1024 * 1024
    out = []
    for idx, img in enumerate(attached_images):
        data_b64 = img.get("data") or ""
        mime = img.get("mime_type", "image/jpeg") or "image/jpeg"
        name = img.get("name", "image") or "image"
        if not data_b64 or not str(mime).startswith("image/"):
            out.append(img)
            continue
        try:
            raw_bytes = _b64.b64decode(data_b64)
            if not raw_bytes or len(raw_bytes) > _MAX_BYTES:
                out.append(img)
                continue
            # Known-bad gate. Unlike the two document lanes there is no text channel
            # here, so a blocked image is DROPPED from the list entirely rather than
            # replaced: keeping the inline base64 (this loop's usual fallback) would
            # hand the exact bytes to the vision model that were just refused.
            if _inspect_attachment(raw_bytes, name, "web_chat_image").blocked:
                continue
            # On-disk basename: derive the extension from the VALIDATED image mime (never trust
            # the uploaded name's extension), so /api/file always serves an image Content-Type.
            # The display name keeps the original. A uuid token guarantees uniqueness (no overwrite
            # when two same-named images arrive in the same second).
            stem = _re.sub(r'[^A-Za-z0-9._-]', '_', _os.path.splitext(name)[0])[:60] or "image"
            ext = _mt.guess_extension(str(mime).split(";")[0].strip()) or ".img"
            fpath = attach_dir / f"{int(_time.time())}_{_uuid.uuid4().hex[:8]}_{stem}{ext}"
            fpath.write_bytes(raw_bytes)
            out.append({"name": name, "mime_type": mime, "path": str(fpath)})
        except Exception as e:
            log("WebServer", f"attach image to file failed ({name}): {e}")
            out.append(img)  # fall back to inline base64
    return out


async def _notify_attachment_index(manager, session_id: str, kind: str, **extra) -> None:
    """Fail-safe WS broadcast for the attachment indexing-status indicator in the Web UI.
    kind: 'attachment_indexing' (start) | 'attachment_indexed' (done) | 'attachment_index_error'
        | 'attachment_index_cancelled' (stop button)."""
    try:
        await manager.broadcast_to_session(session_id, {"type": kind, "sessionId": session_id, **extra})
    except Exception:
        pass


def _spawn_attachment_index(manager, session_id: str, user_scope_id, contents: list) -> "asyncio.Task":
    """Run attachment RAG indexing as a cancellable background task.

    Indexing used to be awaited inline in the set_sidebar_documents handler, which blocked
    the per-connection WS receive loop for the whole index — so the stop button could not be
    processed mid-index. Running it as a tracked task lets stop_generation cancel it and lets
    the user keep interacting once indexing settles.
    The caller is responsible for sending the initial 'attachment_indexing' notify so the UI
    shows the banner immediately even before this task is scheduled."""
    sid = str(session_id)

    async def _runner():
        try:
            from vaf.memory.attachment_rag import index_session_attachments_async
            await index_session_attachments_async(
                session_id=sid,
                user_scope_id=user_scope_id,
                documents=contents,
            )
            await _notify_attachment_index(manager, sid, "attachment_indexed", count=len(contents))
        except asyncio.CancelledError:
            # Stop button: the frontend already cleared its local indexing state, but emit a
            # cancelled notify so any other connected client unblocks too. Best-effort.
            try:
                await _notify_attachment_index(manager, sid, "attachment_index_cancelled")
            except Exception:
                pass
            raise
        except Exception as e:
            log("WebServer", f"attachment index failed: {e}")
            await _notify_attachment_index(manager, sid, "attachment_index_error")

    task = asyncio.create_task(_runner())
    _active_index_tasks.setdefault(sid, set()).add(task)

    def _cleanup(t: "asyncio.Task") -> None:
        bucket = _active_index_tasks.get(sid)
        if bucket is not None:
            bucket.discard(t)
            if not bucket:
                _active_index_tasks.pop(sid, None)

    task.add_done_callback(_cleanup)
    return task


# Per-file upload ceiling for chat attachments: 100 MB raw leaves the 200 MB
# WebSocket frame (WS_MAX_SIZE_BYTES) headroom after base64's ~1.37x inflation.
# The client gates before encoding with the same number.
_MAX_ATTACH_BYTES = 100 * 1024 * 1024


async def process_files_to_sidebar_list(files: list, session_id: str = None,
                                        user_scope_id: str = None) -> list:
    """
    Process uploaded files and return a list of {name, content, data?, mimeType?, path?} for sidebar_documents.
    Uses same Librarian extraction as process_uploaded_files.
    Passes through data (base64) and mimeType for PDF display.

    With a session, each document is PERSISTED to the user-siloed chat
    attachments folder (the image pattern) and `path` carries the REAL file
    location - which is what makes `learn_document(path=...)` and the learn
    button possible for chat attachments; the runner's own truncation advice
    used to point at a path that literally did not exist. Only when persisting
    fails (or no session is known) does the honest literal
    "Hochgeladen über Web-UI (kein lokaler Pfad)" remain - never a fake path.
    """
    import base64
    import tempfile
    import os
    from pathlib import Path

    if not files:
        return []

    from vaf.core.log_helper import log_attachment

    attach_dir = None
    if session_id:
        try:
            from vaf.core.session import get_session_attachments_dir
            attach_dir = get_session_attachments_dir(session_id, user_scope_id, create=True)
        except Exception:
            attach_dir = None

    results = []
    for file_obj in files:
        try:
            filename = file_obj.get("name", "unknown")
            file_data = file_obj.get("data", "")
            mime_type = file_obj.get("mimeType", "")

            base64_part = file_data
            if base64_part.startswith("data:"):
                base64_part = base64_part.split(",", 1)[1] if "," in base64_part else base64_part

            decoded_data = base64.b64decode(base64_part)
            file_ext = Path(filename).suffix or ".txt"
            if _is_temp_like_filename(filename):
                filename = f"Dokument{file_ext}"

            log_attachment("FILE_RECEIVED",
                name=filename, mime=mime_type, ext=file_ext,
                base64_len=len(base64_part), decoded_bytes=len(decoded_data))

            # Server-side half of the size gate (the client checks before
            # encoding; a non-VAF client does not). Named limit instead of the
            # old behavior - nothing at all, until the 200 MB WS frame ceiling
            # dropped the whole connection with no message.
            if len(decoded_data) > _MAX_ATTACH_BYTES:
                results.append({
                    "name": filename,
                    "content": (f"[ERROR] {filename} exceeds the "
                                f"{_MAX_ATTACH_BYTES // (1024 * 1024)} MB upload limit "
                                f"({len(decoded_data) // (1024 * 1024)} MB). Save it to "
                                f"disk and learn or read it by path instead."),
                    "path": "Hochgeladen über Web-UI (kein lokaler Pfad)",
                })
                continue

            # Known-bad gate, before extraction, before persisting, before the RAG
            # index. Placed AFTER the size check on purpose: hashing is O(size) and
            # runs inline here, so the cheap ceiling turns an oversized upload away
            # first rather than stalling the event loop to fingerprint it.
            #
            # A blocked file still appears in the sidebar as an entry - showing nothing
            # would read as an upload that silently failed - but it carries the refusal
            # instead of content, and no bytes are written anywhere.
            _verdict = _inspect_attachment(decoded_data, filename, "web_chat")
            if _verdict.blocked:
                results.append({
                    "name": filename,
                    "content": f"[BLOCKED] {_verdict.message(filename)}",
                    "path": "Blocked by the security scanner (not saved)",
                })
                continue

            with tempfile.NamedTemporaryFile(prefix="vaf_", suffix=file_ext, delete=False) as temp_file:
                temp_file.write(decoded_data)
                temp_path = temp_file.name

            try:
                from vaf.tools.librarian import LibrarianTool
                librarian = LibrarianTool()
                # Run in a thread — PDF parsing (PyPDF2) is CPU-bound and blocks the
                # event loop, which delays WebSocket streaming for the entire duration.
                content = await asyncio.to_thread(librarian._read_file, Path(temp_path), True)
                # Strip lone Unicode surrogates (e.g. \u{DE16} from PDF emoji).
                # They survive json.dumps but crash UTF-8 encoding in WebSocket sends,
                # session file writes, and API request serialization.
                if content:
                    content = content.encode("utf-8", errors="replace").decode("utf-8")

                # Classify content type for diagnostics
                _ctype = "TEXT"
                if "[ERROR]" in (content or "")[:80]:
                    _ctype = "ERROR"
                elif ("[No text layer" in (content or "")[:400]
                      or "[No readable text" in (content or "")[:400]):
                    # format_pdf_read_result's empty markers (the old hand-rolled
                    # "[Scanned PDF" string is gone).
                    _ctype = "SCANNED_NO_TEXT"
                elif "[INFO]" in (content or "")[:80] and "Auto-Chunk" in (content or "")[:200]:
                    _ctype = "CHUNKED"
                log_attachment("EXTRACT_DONE",
                    name=filename, content_type=_ctype,
                    content_len=len(content or ""),
                    preview=repr((content or "")[:300]))

                # Persist to the session's attachments dir so the entry carries
                # a REAL path (image pattern). Failure keeps the honest literal.
                persisted_path = None
                if attach_dir is not None:
                    try:
                        import re as _re
                        import time as _time
                        import uuid as _uuid
                        _stem = _re.sub(r'[^A-Za-z0-9._-]', '_',
                                        Path(filename).stem)[:60] or "document"
                        fpath = attach_dir / f"{int(_time.time())}_{_uuid.uuid4().hex[:8]}_{_stem}{file_ext}"
                        fpath.write_bytes(decoded_data)
                        persisted_path = str(fpath)
                    except Exception as _pe:
                        log("WebServer", f"attach document to file failed ({filename}): {_pe}")

                entry = {
                    "name": filename,
                    "content": (content or "") + _advisory_note(_verdict),
                    "path": persisted_path or "Hochgeladen über Web-UI (kein lokaler Pfad)",
                }
                if base64_part:
                    entry["data"] = base64_part
                if mime_type:
                    entry["mimeType"] = mime_type
                suf = file_ext.lower()
                if suf in (".docx", ".xlsx", ".pptx", ".odt", ".ods", ".odp"):
                    # Prefer Gotenberg (LibreOffice in Docker): full design fidelity, DOCX/XLSX/PPTX → PDF.
                    pdf_bytes = await asyncio.to_thread(_office_to_pdf_via_gotenberg, temp_path, filename)
                    if not pdf_bytes and suf == ".docx":
                        pdf_bytes = await asyncio.to_thread(_docx_to_pdf_via_libreoffice, temp_path)
                    if pdf_bytes:
                        entry["data"] = base64.b64encode(pdf_bytes).decode("ascii")
                        entry["mimeType"] = "application/pdf"
                    else:
                        try:
                            if suf == ".docx":
                                html_body = _docx_to_html(temp_path)
                            elif suf == ".xlsx":
                                html_body = _xlsx_to_html(temp_path)
                            elif suf == ".pptx":
                                html_body = _pptx_to_html(temp_path)
                            else:
                                html_body = "<p>[ODT/ODS/ODP: Gotenberg nicht erreichbar. Bitte starten Sie den Docker-Stack.]</p>"
                            entry["htmlContent"] = (
                                "<!DOCTYPE html><html><head><meta charset=\"utf-8\"/>"
                                "<style>body{font-family:system-ui,sans-serif;padding:1.5rem;max-width:80ch;}"
                                "table{border-collapse:collapse;}td,th{border:1px solid #ccc;padding:6px;}</style>"
                                "</head><body>" + html_body + "</body></html>"
                            )
                        except Exception as e:
                            log("WebServer", f"Office to HTML failed for {filename}: {e}")
                results.append(entry)
            finally:
                try:
                    os.unlink(temp_path)
                except Exception:
                    pass
        except Exception as e:
            import traceback as _tb2
            err_name = file_obj.get("name", "unknown")
            if _is_temp_like_filename(err_name):
                err_name = "Dokument" + (Path(err_name).suffix or ".txt")
            log_attachment("FILE_ERROR", name=err_name, error=str(e), tb=_tb2.format_exc()[-300:])
            results.append({
                "name": err_name,
                "content": f"[ERROR] Failed to process file: {str(e)}",
                "path": "Hochgeladen über Web-UI (kein lokaler Pfad)",
            })
    return results


def mark_webui_process() -> None:
    """Declare that THIS process serves the web UI, for as long as it lives.

    That is what VAF_WEBUI_ACTIVE means, and it is why sub-agents spawned from here run
    piped into the browser panel instead of opening a host terminal window.

    It used to be set on WebSocket connect and cleared when the last one dropped, i.e. it
    tracked "is a browser attached right now". A socket drop then silently changed
    process-wide spawn behaviour: on 2026-07-20 the browser connection died mid-run and two
    minutes later a sub-agent spawn saw the flag unset and opened a real terminal window on
    the host desktop.

    EVERY entry point that serves this app must call this. The desktop app does NOT go
    through run_server: vaf/tray.py has its own start_uvicorn that imports `app` and drives
    uvicorn itself, so putting the write in run_server alone left the fix inert on exactly
    the path where the incident happened. tests/test_terminal_spawn_lifetime.py pins that
    every uvicorn entry point calls this.
    """
    os.environ["VAF_WEBUI_ACTIVE"] = "1"


def run_server(host="127.0.0.1", port=8001):
    # The MAIN server must not believe it is a sub-agent child: launched from a
    # contaminated shell (e.g. a finished coder's terminal), inherited child
    # markers silently turn every spawn into an invisible in-process run.
    from vaf.core.platform import scrub_inherited_subagent_env
    scrub_inherited_subagent_env()
    """Run the Uvicorn server. Uses TLS (HTTPS/WSS) when config has cert/key paths set."""
    mark_webui_process()

    # Store the loop so the TUI thread can schedule updates
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    manager.set_server_loop(loop)

    tls_enabled = Config.get("local_network_tls_enabled", False)
    ssl_cert = (Config.get("local_network_ssl_cert") or "").strip()
    ssl_key = (Config.get("local_network_ssl_key") or "").strip()

    # Auto-generate certificates if TLS enabled but no valid certs
    if tls_enabled and (not ssl_cert or not ssl_key or not os.path.isfile(ssl_cert) or not os.path.isfile(ssl_key)):
        try:
            from vaf.network.ssl_utils import ensure_ssl_certificates
            ssl_cert, ssl_key = ensure_ssl_certificates()
            if ssl_cert and ssl_key:
                log("WebServer", f"Auto-generated SSL certificate: {ssl_cert}")
        except Exception as e:
            log("WebServer", f"SSL auto-generation failed, falling back to HTTP: {e}")
            ssl_cert, ssl_key = "", ""

    # Allow large PDF/file attachments: base64-encoded image-heavy PDFs can easily exceed
    # uvicorn's default ws_max_size (16 MB), causing the server to close the WebSocket
    # connection mid-upload (frontend shows "Verbindung wird wiederhergestellt").
    # Shared constant: the tray and the HTTPS proxy start their own uvicorns and
    # carried the 16 MB default while this value sat only here.
    from vaf.core.log_helper import WS_MAX_SIZE_BYTES

    if tls_enabled and ssl_cert and ssl_key and os.path.isfile(ssl_cert) and os.path.isfile(ssl_key):
        config = uvicorn.Config(
            app=app, host=host, port=port, loop="asyncio", log_level="error",
            ssl_certfile=ssl_cert, ssl_keyfile=ssl_key,
            ws_max_size=WS_MAX_SIZE_BYTES,
        )
        log("WebServer", f"Starting with TLS (HTTPS/WSS) on {host}:{port}")
    else:
        config = uvicorn.Config(app=app, host=host, port=port, loop="asyncio", log_level="error",
                                ws_max_size=WS_MAX_SIZE_BYTES)
        if tls_enabled:
            log("WebServer", f"WARNING: TLS enabled but no certificates available, running HTTP on {host}:{port}")
    server = uvicorn.Server(config)

    # We run this in the thread provided by the caller
    loop.run_until_complete(server.serve())

def start_background_server(host="127.0.0.1", port=8001):
    """Start server in a daemon thread."""
    t = threading.Thread(target=run_server, args=(host, port), daemon=True)
    t.start()
    return t
