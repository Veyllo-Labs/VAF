# SPDX-FileCopyrightText: 2026 Veyllo GmbH
# SPDX-License-Identifier: AGPL-3.0-or-later
# Additional permissions and terms under AGPL Section 7: see LICENSING.md
"""
VAF Librarian - Smart File & Information Retrieval Agent
Optimized for fast direct execution of simple tasks
Falls back to LLM only for complex queries
"""
import os
import re
import json
import time
import requests
import sys
import subprocess
import platform
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
from rich.live import Live

from vaf.tools.base import BaseTool
from vaf.cli.ui import UI, AnimatedHeader
from vaf.tools.filesystem import (
    ReadFileTool,
    ListFilesTool,
    TreeTool,
    FinderTool,
    WriteFileTool,
    MoveFileTool,
    FolderSizeTool,
    is_safe_path,
)
from vaf.tools.python_sandbox import PythonSandboxTool
from vaf.tools.document_viewer import (
    DocumentViewerTool,
    DocumentEditorTool,
    ReplaceEditorSelectionTool,
    ReplaceEditorTextTool,
)
from vaf.tools.cloud_storage import CloudStorageTool
from vaf.core.fs_map import CachedFilesystemMap

# Try to import psutil for better disk info (optional)
try:
    import psutil
    HAS_PSUTIL = True
except ImportError:
    HAS_PSUTIL = False


class LibrarianTool(BaseTool):
    """
    Smart Librarian that handles file/info retrieval tasks.
    
    OPTIMIZATION: Simple tasks are executed DIRECTLY without LLM calls.
    Only complex queries go through the LLM reasoning loop.
    """
    
    name = "librarian_agent"
    permission_level = "write"
    side_effect_class = "reversible"
    description = """A specialized Sub-Agent for File System, Storage & Information Retrieval.
    **PRIMARY TOOL for:**
    - **Folder Sizes:** "How big is Downloads?", "Check folder size", "Disk usage analysis"
    - **Storage Info:** "Disk space", "Free space", "List drives", "Storage capacity"
    - **File Ops:** "List files", "Find file", "Read file", "Count files"
    - **Cloud Storage:** "What's in my Google Drive?", "Browse cloud", "List Drive folders"
    - **System Info:** "How many drives", "Hardware info"
    
    Use this agent for ANY file system, cloud storage, or storage query. Do not say "I can't access files" - delegate to this agent! Do not use this agent to search or analyze the VAF installation folder; that path is not allowed."""
    
    parameters = {
        "type": "object",
        "properties": {
            "task": {
                "type": "string",
                "description": "The task (e.g., 'Count files in Downloads', 'Find config.json')."
            }
        },
        "required": ["task"]
    }
    
    def __init__(self):
        super().__init__()
        # Initialize tools
        self.tools = {
            "read_file": ReadFileTool(),
            "write_file": WriteFileTool(),
            "list_files": ListFilesTool(),
            "tree": TreeTool(),
            "find_files": FinderTool(),
            "python_sandbox": PythonSandboxTool(),
            "folder_size": FolderSizeTool(),
            "document_viewer": DocumentViewerTool(),
            "document_editor": DocumentEditorTool(),
            "replace_editor_selection": ReplaceEditorSelectionTool(),
            "replace_editor_text": ReplaceEditorTextTool(),
            "move_file": MoveFileTool(),
            "cloud_storage": CloudStorageTool(),
        }
        
        # Cross-platform home directory
        self.home = Path.home()
        
        # Common folder mappings (cross-platform, multilingual)
        # Try common folder names for each language/OS
        self.folder_aliases = {}
        
        # English folder names
        for name in ["Downloads", "Desktop", "Documents", "Pictures", "Videos", "Music"]:
            path = self.home / name
            if path.exists():
                self.folder_aliases[name.lower()] = path
        
        # German folder names (Windows often uses these)
        german_mappings = {
            "downloads": ["Downloads", "Herunterladen"],
            "desktop": ["Desktop", "Arbeitsplatz"],
            "documents": ["Documents", "Dokumente", "Dokumen"],  # Common typo/abbreviation
            "pictures": ["Pictures", "Bilder"],
            "videos": ["Videos"],
            "music": ["Music", "Musik"],
        }
        
        for key, variants in german_mappings.items():
            for variant in variants:
                path = self.home / variant
                if path.exists():
                    self.folder_aliases[variant.lower()] = path
                    self.folder_aliases[key] = path  # Also map English key
        
        # Always add home and ~
        self.folder_aliases["home"] = self.home
        self.folder_aliases["~"] = self.home
        
        # Filesystem Map
        self.fs_map = CachedFilesystemMap()
        self.map_cache = None
        self.last_scan = None
    
    def get_system_prompt_addition(self) -> str:
        """Adds the filesystem map to the system prompt."""
        if self.should_refresh_map():
            UI.event("Librarian", "Scanning filesystem map...", style="dim")
            self.map_cache = self.fs_map.build_map(depth=1)
            self.last_scan = time.time()
        
        return f"""
## FILESYSTEM CONTEXT (Smart Map)

You have access to this filesystem map for fast navigation:

{self.fs_map.format_summary()}

**IMPORTANT RULES:**
- When the user explicitly names a folder (e.g. "in Downloads", "im Downloads Ordner"), use THAT folder. Do NOT default to Documents.
- When asked "how many documents/images/files?", check the map FIRST.
- Default to standard locations (Documents for docs, Pictures for images) only when NO folder is mentioned.
- Only use file_search/find_files tool if the map doesn't have the answer.
- If user asks about a specific file type, check the 'Common' types in the map first.
"""

    def should_refresh_map(self) -> bool:
        """Refresh map every 5 minutes or if not yet scanned."""
        if not self.last_scan:
            return True
        return (time.time() - self.last_scan) > 300  # 5 min

    def run(self, **kwargs) -> str:
        # Per-user filesystem jail: while this librarian run executes, is_safe_path additionally enforces
        # that the agent can read only the caller's OWN data — their VAF_Projects/<uid8> (+ personal folders
        # for the local admin / machine owner), never another user's. Scoped via a contextvar so no other
        # tool is affected. user_scope_id is injected by the agent's tool dispatcher (in-process) or via env.
        from vaf.tools.filesystem import set_librarian_scope, reset_librarian_scope
        scope = kwargs.get("user_scope_id") or os.environ.get("VAF_USER_SCOPE_ID") or None
        token = set_librarian_scope(self._compute_jail(scope))
        try:
            return self._run_impl(**kwargs)
        finally:
            reset_librarian_scope(token)

    def _compute_jail(self, user_scope_id):
        """Librarian jail info. Local admin (no scope OR the local-admin scope) => full access; a remote
        user => jailed to their own VAF_Projects/<uid8> only (no personal folders). Fail-closed on error."""
        try:
            from vaf.core.config import get_local_admin_scope_id
            from vaf.core.session import get_user_projects_root
            scope = str(user_scope_id or "")
            local_admin = str(get_local_admin_scope_id() or "")
            if (not scope) or (scope == local_admin):
                return {"is_admin": True, "uid8": None, "allowed_roots": []}
            own_root = get_user_projects_root(scope)
            return {"is_admin": False, "uid8": scope.replace("-", "").lower()[:8],
                    "allowed_roots": [own_root] if own_root else []}
        except Exception:
            return {"is_admin": False, "uid8": "", "allowed_roots": []}

    def _run_impl(self, **kwargs) -> str:
        task = kwargs.get('task', '').strip()
        if not task:
            return "Error: No task provided."
        
        # ═══════════════════════════════════════════════════════════════════════
        # CHECK IF RUNNING IN SEPARATE TERMINAL MODE
        # ═══════════════════════════════════════════════════════════════════════
        from vaf.core.config import Config
        from vaf.core.platform import Platform
        
        # If already in sub-agent terminal, run normally
        if os.environ.get("VAF_IN_SUBAGENT_TERMINAL", "").strip() in ("1", "true", "yes"):
            # Continue with normal execution below
            pass
        elif Config.get("sub_agents_in_separate_terminals", False):
            # Start in new terminal window with IPC tracking
            import shlex
            from vaf.core.subagent_ipc import get_ipc, get_current_session_id
            
            # Create task in IPC system
            ipc = get_ipc()
            task_id = ipc.create_task("librarian_agent", task_description=task)
            
            # Pass session/task context to the sub-agent via the CHILD env only (not the parent's
            # process-global os.environ), so concurrent workers don't clobber each other's session.
            session_id = get_current_session_id()
            _sub_env = {"VAF_TASK_ID": task_id, "VAF_AGENT_TYPE": "librarian_agent"}
            if session_id:
                _sub_env["VAF_SESSION_ID"] = session_id
            # Carry the user scope into the child so its librarian jail (is_safe_path) applies there too.
            _scope_for_child = kwargs.get("user_scope_id") or os.environ.get("VAF_USER_SCOPE_ID")
            if _scope_for_child:
                _sub_env["VAF_USER_SCOPE_ID"] = str(_scope_for_child)

            # Pass provider configuration to sub-agent (Best Practice: Inherit or override)
            use_separate_provider = Config.get("subagent_use_separate_provider", False)
            if use_separate_provider:
                subagent_provider = Config.get("subagent_provider", "inherit")
                if subagent_provider != "inherit":
                    _sub_env["VAF_PROVIDER"] = subagent_provider
            
            cmd_parts = [sys.executable, '-m', 'vaf.main', 'subagent', 'run', 'librarian_agent', '--task', task, '--task-id', task_id]
            
            if Platform.is_windows():
                # Windows: properly escape for cmd /k
                escaped_parts = []
                for part in cmd_parts:
                    if ' ' in part or '"' in part:
                        escaped = part.replace('"', '\\"')
                        escaped_parts.append(f'"{escaped}"')
                    else:
                        escaped_parts.append(part)
                cmd = ' '.join(escaped_parts)
                title = f"VAF Librarian Agent [{task_id}]"
            else:
                # Unix: use shell quoting
                cmd = ' '.join(shlex.quote(str(part)) for part in cmd_parts)
                title = f"VAF Librarian Agent [{task_id}]"
            
            if Platform.open_new_terminal(cmd, title=title, extra_env=_sub_env):
                # Mark task as running
                ipc.mark_task_running(task_id)
                
                UI.event("Sub-Agent", f"Librarian Agent started in new terminal [Task: {task_id}]", style="bold cyan")
                # Return special marker for main agent to recognize async task
                return f"[SUBAGENT_ASYNC:{task_id}:librarian_agent] Sub-Agent running in separate terminal. Task: {task[:80]}..."
            else:
                # Fallback: run normally if terminal opening fails
                UI.warning("Failed to open new terminal, running in current window")
                ipc.cancel_task(task_id)
        
        # ═══════════════════════════════════════════════════════════════════════
        # FAST PATH: Try to handle simple tasks DIRECTLY (no LLM needed)
        # ═══════════════════════════════════════════════════════════════════════
        
        # The librarian only reads — surface its live state (filesystem map, storage,
        # search) to the read-only explorer window in the WebUI (best-effort, no-op
        # without a session / WebUI; never affects the task result). Clear any folder
        # browsed by a previous task so the overview shows until this task opens a folder.
        self._browse_path = None
        self._emit_librarian_state(task, stage="scanning")

        # Try direct execution first (no animation needed for fast path)
        direct_result = self._try_direct_execution(task)
        if direct_result:
            # Fast path succeeded - show brief static header
            from vaf.cli.tui import _StaticHeader
            header = _StaticHeader("Collaboration Mode Active", "Main Agt", "Librarian")
            live = Live(header, refresh_per_second=12, console=UI.console)
            live.start()
            try:
                time.sleep(0.5)  # Brief display
                UI.event("Sub-Agent", "Librarian activated...", style="bold cyan")
                time.sleep(0.3)
            finally:
                live.stop()
            self._emit_librarian_state(task, stage="done", result=direct_result)
            return direct_result
        
        # ═══════════════════════════════════════════════════════════════════════
        # MAP-REDUCE PATH: Huge content - use Chunking & Summarization
        # ═══════════════════════════════════════════════════════════════════════
        if len(task) > 15000:
            _res = self._summarize_chunks(task)
            self._emit_librarian_state(task, stage="done", result=_res)
            return _res
        
        # ═══════════════════════════════════════════════════════════════════════
        # SLOW PATH: Complex task - use LLM reasoning
        # ═══════════════════════════════════════════════════════════════════════

        # Attach any hint produced by the fast-path read check (path not found)
        hint = getattr(self, "_read_hint_for_llm", None)
        if hint:
            task = task + hint
            self._read_hint_for_llm = None

        # LLM path has its own animation
        _res = self._execute_with_llm(task)
        self._emit_librarian_state(task, stage="done", result=_res)
        return _res
    
    # ═══════════════════════════════════════════════════════════════════════════
    # READ-ONLY LIVE STATE → WebUI explorer window (best-effort, never alters output)
    # ═══════════════════════════════════════════════════════════════════════════

    def _fmt_bytes_short(self, num_bytes) -> str:
        """Human-readable size for activity/log text."""
        try:
            b = float(num_bytes)
        except Exception:
            return "0 B"
        for unit in ("B", "KB", "MB", "GB", "TB"):
            if b < 1024 or unit == "TB":
                return f"{b:.0f} {unit}" if unit in ("B", "KB") else f"{b:.1f} {unit}"
            b /= 1024
        return f"{b:.1f} TB"

    def _ext_type(self, name: str) -> str:
        """Map a file name to a coarse type for the explorer icon (matches the frontend)."""
        ext = (os.path.splitext(name)[1] or "").lower().lstrip(".")
        if ext == "pdf":
            return "pdf"
        if ext in ("png", "jpg", "jpeg", "gif", "bmp", "webp", "svg", "tif", "tiff", "heic", "ico"):
            return "image"
        if ext in ("mp4", "mov", "mkv", "avi", "webm", "wmv", "m4v", "flv", "mpg", "mpeg"):
            return "video"
        if ext in ("doc", "docx", "odt", "rtf", "ppt", "pptx", "xls", "xlsx", "ods", "odp", "md"):
            return "doc"
        if ext in ("csv", "json", "xml", "sql", "db", "sqlite", "parquet", "tsv", "yaml", "yml"):
            return "data"
        if ext in ("zip", "tar", "gz", "rar", "7z", "bz2", "xz", "tgz"):
            return "arch"
        if ext in ("exe", "msi", "appimage", "deb", "rpm", "dmg", "apk", "bin", "app"):
            return "exe"
        return "txt"

    def _build_librarian_state(self, task: str, result: Optional[str] = None) -> dict:
        """Assemble the read-only librarian window state (filesystem map, storage,
        biggest folders, optional search) from the structured filesystem map and disk
        usage. Side-effect-free; used only to feed the WebUI explorer view."""
        entries: List[dict] = []
        total_size = 0
        total_files = 0
        total_folders = 0
        try:
            fmap = self.fs_map.build_map(depth=1) or {}
            for name, data in (fmap.get("locations") or {}).items():
                sz = int(data.get("size_bytes", 0) or 0)
                files = int(data.get("total_files", 0) or 0)
                entries.append({"name": name, "type": "folder", "sizeBytes": sz, "items": files})
                total_size += sz
                total_files += files
                total_folders += int(data.get("folder_count", 0) or 0)
            entries.sort(key=lambda e: e["sizeBytes"], reverse=True)
        except Exception:
            pass

        top_folders = [{"name": e["name"], "sizeBytes": e["sizeBytes"]} for e in entries[:10]]

        # Storage: local disk + home always; Google Drive only if actually connected
        # (we never fake a quota — just surface it as a connected source).
        drives: List[dict] = []
        try:
            import shutil as _shutil
            du = _shutil.disk_usage(str(self.home))
            drives.append({"name": "/", "kind": "disk", "usedBytes": int(du.used), "totalBytes": int(du.total)})
            drives.append({"name": "~/ (Home)", "kind": "home", "usedBytes": int(total_size), "totalBytes": int(du.total)})
        except Exception:
            pass
        try:
            from vaf.tools.cloud_storage import _get_username, _get_cloud_accounts
            for a in (_get_cloud_accounts(_get_username()) or []):
                if (a.get("provider") or "").lower() in ("google_drive", "googledrive", "gdrive", "google"):
                    drives.append({"name": a.get("label") or "Google Drive", "kind": "cloud",
                                   "usedBytes": 0, "totalBytes": 0, "connected": True})
                    break
        except Exception:
            pass

        # Search context: file searches drive the match highlighting. Detect an explicit
        # file type ("alle PDFs", "*.pdf", "pdf files") first, else fall back to a free term.
        search = None
        try:
            tl = (task or "").lower()
            is_search = any(k in tl for k in ("find", "search", "suche", "finde", "locate", "schick", "send"))
            ext_m = re.search(r"\b(pdfs?|jpe?gs?|pngs?|gifs?|txts?|docx?|xlsx?|pptx?|csvs?|jsons?|zips?|mp3s?|mp4s?|movs?)\b", tl)
            query = None
            if ext_m:
                query = "*." + re.sub(r"s$", "", ext_m.group(1))
            elif is_search:
                query = self._extract_search_term(task) or None
            if query:
                hits = None
                if result:
                    hits = sum(1 for ln in result.splitlines()
                               if ("/" in ln or "\\" in ln) and "." in ln) or None
                search = {"query": query, "hits": hits}
        except Exception:
            pass

        activity: List[dict] = []
        try:
            activity.append({"cls": "scan", "text": f"Filesystem-Map: {total_folders} Ordner, {total_files} Dateien"})
            if top_folders:
                activity.append({"cls": "ok", "text": f"Größte: {top_folders[0]['name']} ({self._fmt_bytes_short(top_folders[0]['sizeBytes'])})"})
            if any(d.get("kind") == "cloud" for d in drives):
                activity.append({"cls": "info", "text": "Google Drive verbunden"})
            if search:
                txt = f"Suche: {search['query']}"
                if search.get("hits"):
                    txt += f" -> {search['hits']} Treffer"
                activity.append({"cls": "scan", "text": txt})
            activity.append({"cls": "warn", "text": "read-only - es wird nichts geändert oder erstellt"})
        except Exception:
            pass

        # When the task operates inside a concrete folder, open it for the explorer:
        # scan its direct entries (read-only) so the UI lists the files like a file
        # manager and highlights search matches. Falls back to the overview on any error.
        current_folder = None
        try:
            bp = getattr(self, "_browse_path", None)
            if bp and os.path.isdir(bp):
                import datetime as _dt
                q = (search or {}).get("query") if search else None
                ext_q = name_q = None
                if q:
                    qs = q.strip()
                    if qs.startswith("*.") or qs.startswith("."):
                        ext_q = qs.lstrip("*").lstrip(".").lower()
                    else:
                        name_q = qs.lower()

                def _is_match(nm: str, is_dir: bool) -> bool:
                    if is_dir:
                        return False
                    low = nm.lower()
                    if ext_q:
                        return low.endswith("." + ext_q)
                    if name_q:
                        return name_q in low
                    return False

                items: List[dict] = []
                f_count = d_count = tot = hit_n = 0
                type_counts: dict = {}
                with os.scandir(bp) as it:
                    for de in it:
                        try:
                            is_dir = de.is_dir(follow_symlinks=False)
                            st = de.stat(follow_symlinks=False)
                            sz = 0 if is_dir else int(st.st_size)
                            etype = "folder" if is_dir else self._ext_type(de.name)
                            try:
                                mod = _dt.datetime.fromtimestamp(st.st_mtime).strftime("%d.%m.%Y")
                            except Exception:
                                mod = ""
                            m = _is_match(de.name, is_dir)
                            items.append({"name": de.name, "type": etype, "isDir": is_dir,
                                          "sizeBytes": sz, "modified": mod, "match": m})
                            type_counts[etype] = type_counts.get(etype, 0) + 1
                            if is_dir:
                                d_count += 1
                            else:
                                f_count += 1
                                tot += sz
                            if m:
                                hit_n += 1
                        except Exception:
                            continue
                # matches first, then folders, then files by size — so hits stay visible
                items.sort(key=lambda e: (not e["match"], not e["isDir"], -(e["sizeBytes"] or 0), e["name"].lower()))
                current_folder = {
                    "path": bp,
                    "name": os.path.basename(bp.rstrip("/\\")) or bp,
                    "fileCount": int(f_count),
                    "folderCount": int(d_count),
                    "totalSize": int(tot),
                    "types": [{"type": t, "count": c} for t, c in sorted(type_counts.items(), key=lambda kv: -kv[1])],
                    "entries": items[:80],
                }
                if search is not None:
                    search = {"query": search.get("query"), "hits": hit_n or search.get("hits")}
        except Exception:
            current_folder = None

        return {
            "root": "~",
            "readOnly": True,
            "totalSize": int(total_size),
            "totalFiles": int(total_files),
            "totalFolders": int(total_folders),
            "entries": entries,
            "topFolders": top_folders,
            "drives": drives,
            "search": search,
            "activity": activity,
            "currentFolder": current_folder,
        }

    def _emit_librarian_state(self, task: str, stage: str = "scanning", result: Optional[str] = None) -> None:
        """Push the read-only librarian state to the WebUI explorer window. Best-effort:
        resolves the session id like the document agent, never raises, and never changes
        the librarian's return value."""
        try:
            session_id = os.environ.get("VAF_SESSION_ID", "").strip()
            if not session_id:
                try:
                    from vaf.core.subagent_ipc import get_current_session_id
                    session_id = get_current_session_id() or ""
                except Exception:
                    session_id = ""
            if not session_id:
                return
            state = self._build_librarian_state(task, result=result)
            state["stage"] = stage
            from vaf.core.web_interface import get_web_interface
            get_web_interface().emit_librarian_state(state, session_id=session_id)
        except Exception:
            pass

    def _summarize_chunks(self, task: str) -> str:
        """
        Handle excessively large inputs by chunking and summarizing iteratively (Map-Reduce).
        """
        UI.event("Librarian", f"Input text too large ({len(task)} chars). Switching to Map-Reduce strategy...", style="warning")
        
        # 1. Split task into instruction (start) and content (rest)
        # We assume instructions are at the beginning.
        instruction_limit = 2000
        if len(task) > instruction_limit:
            instruction = task[:instruction_limit] + "..."
        else:
            instruction = task
            
        # 2. Chunking configuration
        CHUNK_SIZE = 12000  # Safe size for 8k/16k context
        OVERLAP = 500
        
        chunks = []
        start = 0
        total_len = len(task)
        
        while start < total_len:
            end = min(start + CHUNK_SIZE, total_len)
            chunks.append(task[start:end])
            start = end - OVERLAP
            
        UI.event("Librarian", f"Split content into {len(chunks)} chunks for processing.", style="info")
        
        summaries = []
        
        # 3. Process chunks
        for i, chunk in enumerate(chunks, 1):
            UI.event("Librarian", f"Summarizing chunk {i}/{len(chunks)}...", style="dim")
            
            prompt = f"""You are a helpful researcher.
User Instruction: {instruction}

PARTIAL TEXT CONTENT (Part {i}/{len(chunks)}):
{chunk}

TASK:
Extract key information from this part that is relevant to the user's instruction. 
If it's a list of items, extract the items.
If it's an article, summarize the main points of this section.
Keep it concise but detailed enough to be useful.
Output ONLY the summary/extraction.
"""
            try:
                content = self.query_llm(
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=1024,
                    temperature=0.0
                )
                if content:
                    summaries.append(content)
                else:
                    summaries.append(f"[Error summarizing chunk {i}: No response from LLM]")
            except Exception as e:
                summaries.append(f"[Exception summarizing chunk {i}: {e}]")
        
        # 4. Final Combination
        combined_summaries = "\n\n".join(summaries)
        
        # If we have multiple chunks, do a final pass to synthesize
        if len(chunks) > 1:
            UI.event("Librarian", "Synthesizing final report from chunk summaries...", style="dim")
            final_prompt = f"""You are a helpful researcher.
User Instruction: {instruction}

I have processed the document in parts. Here are the summaries of each part:

{combined_summaries}

TASK:
Combine these partial summaries into a coherent final report that answers the user's instruction.
Remove duplicates and ensure smooth flow.
"""
            try:
                content = self.query_llm(
                    messages=[{"role": "user", "content": final_prompt}],
                    max_tokens=2048,
                    temperature=0.3
                )
                if content:
                    return f"### Librarian Map-Reduce Report\n\n{content}"
            except:
                pass
        
        return f"### Librarian Map-Reduce Report (Raw Combination)\n\n{combined_summaries}"

    # ═══════════════════════════════════════════════════════════════════════════
    # DIRECT EXECUTION (Fast Path)
    # ═══════════════════════════════════════════════════════════════════════════
    
    def _try_direct_execution(self, task: str) -> Optional[str]:
        """
        Try to execute simple tasks directly without LLM.
        Returns result string or None if task is too complex.
        """
        task_lower = task.lower()
        
        # ─────────────────────────────────────────────────────────────────────
        # GUARD: Web search result payloads should NOT trigger filesystem heuristics.
        # Without this, URLs like ".com" can be misread as file extensions (e.g. "*.com"),
        # causing bogus searches like "Searching for '*.com' in \\season".
        # ─────────────────────────────────────────────────────────────────────
        if (
            "### web search results" in task_lower
            or "web search results:" in task_lower
            or "user question:" in task_lower
            or "http://" in task_lower
            or "https://" in task_lower
        ):
            # If this looks like the web_lookup workflow "synthesis" prompt,
            # we MUST NOT return the deterministic summary — the workflow expects an actual answer.
            # We only use the deterministic summary to prevent bogus filesystem heuristics.
            if (
                "answer the user's question" in task_lower
                or "use the web search results" in task_lower
                or "include 2-3 source links" in task_lower
                or "synthesize" in task_lower
            ):
                return None

            summary = self._summarize_web_results(task)
            if summary:
                return summary

            # If we couldn't parse, force LLM slow-path instead of filesystem
            return None
        
        # ─────────────────────────────────────────────────────────────────────
        # SMART FILESYSTEM MAP QUERY (Ultra-Fast Path)
        # ─────────────────────────────────────────────────────────────────────
        # Skip map for find/send/rename - user needs actual file op or path, not folder stats
        if not any(kw in task_lower for kw in ["send", "schick", "find", "search", "suche", "finde", "locate", "rename", "umbenennen", "umbenenn"]):
            if self.should_refresh_map():
                 self.fs_map.build_map(depth=1)
                 self.last_scan = time.time()

            map_answer = self.fs_map.query_fast(task)
            if map_answer:
                 return f"### Filesystem Map Answer\n\n{map_answer}"

        # Normalize common mojibake sequences seen in some Windows consoles (e.g., "groÃŸ" instead of "groß")
        task_norm = (
            task_lower
            .replace("ÃŸ", "ß")
            .replace("ãÿ", "ß")
            .replace("Ã¶", "ö")
            .replace("Ã¤", "ä")
            .replace("Ã¼", "ü")
            .replace("Ã–", "ö")
            .replace("Ã„", "ä")
            .replace("Ãœ", "ü")
        )
        
        # ─────────────────────────────────────────────────────────────────────
        # PATTERN: Storage/Drive information ("how many drives", "datenträger")
        # ─────────────────────────────────────────────────────────────────────
        storage_patterns = [
            r"how many (storage|drive|disk|datenträger|datanräger|datnträger)",
            r"wie viele (speicher|laufwerk|festplatte|datenträger|datanräger|datnträger)",
            r"what (drives|storage|disks|laufwerke)",
            r"welche (laufwerke|festplatten|speicher|datenträger)",
            r"list (drives|storage|disks|laufwerke)",
            r"zeige (laufwerke|festplatten|speicher)",
            r"storage devices",
            r"disk space",
            r"speicherplatz",
        ]
        
        for pattern in storage_patterns:
            if re.search(pattern, task_norm):
                # Return storage information directly
                disk_info = self._get_disk_info()
                
                # Extract device count from the formatted string
                if "**Total:**" in disk_info:
                    # Count is already in the string
                    return f"""### Storage Devices Information

{disk_info}
"""
                else:
                    # Count manually if not in string
                    device_count = disk_info.count("**") // 2
                    if device_count == 0:
                        device_count = max(1, disk_info.count("\n") - 2)
                    
                    return f"""### Storage Devices Information

{disk_info}

**Total:** {device_count} storage device(s) found.
"""
        
        # Extract path from task
        path = self._extract_path(task)

        # Remember the folder this task operates inside so the WebUI explorer can open it
        # and list its files (size/count/list/find all flow through here; storage/drive
        # queries returned earlier and leave this None → overview stays).
        try:
            self._browse_path = str(path)
        except Exception:
            pass

        # ─────────────────────────────────────────────────────────────────────
        # PATTERN: Folder size ("how big", "size of", "wie groß", "ordnergröße")
        # ─────────────────────────────────────────────────────────────────────
        size_patterns = [
            r"how (big|large) is",
            r"size of",
            r"folder size",
            r"directory size",
            # German: allow flexible phrasing like "wie groß mein downloads folder ist"
            r"wie gro(ß|ss)",
            r"ordnergr(ö|oe)(ß|ss)e",
            r"größe.*(downloads|download|ordner|folder)",
            r"speicherverbrauch",
        ]
        for pattern in size_patterns:
            if re.search(pattern, task_lower):
                # Resolve aliases like "downloads"/"desktop" via _extract_path
                # If path doesn't exist, give a precise actionable message
                try:
                    if not path.exists():
                        return (
                            "[ERROR] Folder not found.\n\n"
                            f"Extracted path: {path}\n"
                            "Try one of these:\n"
                            f"- Use a full path (Windows example): {self.home / 'Downloads'}\n"
                            "- Or mention a known folder alias: downloads / desktop / documents\n"
                        )
                except Exception:
                    pass

                # Compute size deterministically (no LLM)
                return self.tools["folder_size"].run(path=str(path), top_n=10)
        
        # ─────────────────────────────────────────────────────────────────────
        # PATTERN: Count files ("how many files", "count files", "wie viele")
        # ─────────────────────────────────────────────────────────────────────
        count_patterns = [
            r"how many (files|items|documents|pdf|pdfs)",
            r"count (files|items|documents|pdf|pdfs)",
            r"wie viele (dateien|files|dokumente|pdf|pdfs)",
            r"anzahl (dateien|files|pdf|pdfs)",
            r"number of (files|items|pdf|pdfs)",
        ]
        
        # Check if specific file type is mentioned (PDF, JPG, etc.)
        file_type = None
        ext_match = re.search(r'\b(pdf|jpg|jpeg|png|txt|doc|docx|zip|mp3|mp4)\b', task_lower)
        if ext_match:
            file_type = ext_match.group(1)
        
        for pattern in count_patterns:
            if re.search(pattern, task_norm):
                return self._count_files(path, file_type=file_type)
        
        # ─────────────────────────────────────────────────────────────────────
        # PATTERN: List files ("list files", "show files", "zeige dateien")
        # ─────────────────────────────────────────────────────────────────────
        list_patterns = [
            r"list (files|items|documents|contents)",
            r"show (files|items|documents|contents)",
            r"zeige (dateien|files|inhalt)",
            r"what('s| is) in",
            r"was ist in",
        ]
        
        for pattern in list_patterns:
            if re.search(pattern, task_norm):
                return self._list_files(path)
        
        # ─────────────────────────────────────────────────────────────────────
        # PATTERN: Find files ("find", "search", "send", "locate", "suche", "schick")
        # ─────────────────────────────────────────────────────────────────────
        find_patterns = [
            r"find (file|files)?",
            r"search (for)?",
            r"send (me )?(the )?file",
            r"schick (mir )?(die )?datei",
            r"locate",
            r"suche",
            r"finde",
            r"where is",
            r"wo ist",
        ]

        # Extract search pattern
        search_term = self._extract_search_term(task)
        
        for pattern in find_patterns:
            if re.search(pattern, task_lower) and search_term:
                return self._find_files(path, search_term)
        
        # ─────────────────────────────────────────────────────────────────────
        # PATTERN: Read file ("read file", "show content", "cat")
        # ─────────────────────────────────────────────────────────────────────
        read_patterns = [
            r"read (file|content|the file)",
            r"show (content|file content|me the file)",
            r"cat ",
            r"lies (mir |die |den |das |mir die |mir den )?",
            r"zeige (inhalt|mir die datei|mir den inhalt)",
            r"gib mir (den )?(inhalt|content)",
            r"datei(inhalt)? (lesen|anzeigen|ausgeben|öffnen)",
            r"open (the )?file",
            r"öffne (die )?datei",
            r"inhalt (von|der|des|aus)",
            r"content of",
        ]

        file_path = self._extract_file_path(task)

        read_pattern_hit = any(re.search(p, task_lower) for p in read_patterns)
        if read_pattern_hit:
            if file_path:
                return self._read_file(file_path, enable_chunking=True)
            # Pattern matched but no path could be extracted — inject hint for LLM path
            self._read_hint_for_llm = (
                f"\n\n[LIBRARIAN_HINT: The task appears to be a file-read request but no "
                f"file path could be extracted. Home dir: {self.home}. "
                f"Use list_files to find the file, or check if the path contains spaces or "
                f"non-ASCII characters. Then call read_file with the correct absolute path.]"
            )
        else:
            self._read_hint_for_llm = None
        
        # ─────────────────────────────────────────────────────────────────────
        # PATTERN: Write file ("write file", "create file", "save to file")
        # ─────────────────────────────────────────────────────────────────────
        write_patterns = [
            r"write (to )?file",
            r"create file",
            r"save (to )?file",
            r"schreibe (in )?datei",
            r"erstelle datei",
            r"speichere (in )?datei",
        ]
        
        # For write operations, we need both path and content
        # If content is not clear, let LLM handle it
        for pattern in write_patterns:
            if re.search(pattern, task_lower) and file_path:
                # Extract content from task (simple cases only)
                content_match = re.search(r'(?:with|content|text|inhalt)[:\s]+(.+?)(?:$|\.)', task_lower, re.IGNORECASE)
                if content_match:
                    content = content_match.group(1).strip()
                    return self._write_file(file_path, content)
                # If no clear content, let LLM handle it
                break
        
        # ─────────────────────────────────────────────────────────────────────
        # PATTERN: Folders on drive ("what folders", "welche ordner", "folders on drive")
        # ─────────────────────────────────────────────────────────────────────
        folder_patterns = [
            r"what (folders|directories|ordner) (on|in|auf)",
            r"welche (ordner|verzeichnisse) (auf|in)",
            r"folders (on|in) (drive|disk|laufwerk)",
            r"ordner (auf|in) (laufwerk|festplatte)",
            r"list (folders|directories|ordner)",
            r"zeige (ordner|verzeichnisse)",
        ]
        
        for pattern in folder_patterns:
            if re.search(pattern, task_lower):
                # Extract drive/path from task
                drive_path = self._extract_path(task)
                return self._list_files(drive_path)
        
        # ─────────────────────────────────────────────────────────────────────
        # PATTERN: Rename file ("rename", "umbenennen")
        # ─────────────────────────────────────────────────────────────────────
        rename_patterns = [
            r"umbenennen",
            r"umbenenn",
            r"rename",
        ]
        for pattern in rename_patterns:
            if re.search(pattern, task_lower):
                # Try full-path extraction first (e.g. "rename file D:\path\a.pdf to D:\path\b.pdf")
                src_path, dst_path = self._extract_rename_full_paths(task)
                if src_path is not None and dst_path is not None:
                    safe_src, res_src = is_safe_path(str(src_path))
                    safe_dst, res_dst = is_safe_path(str(dst_path))
                    if not safe_src:
                        return f"[ERROR] {res_src}"
                    if not safe_dst:
                        return f"[ERROR] {res_dst}"
                    if not Path(res_src).exists():
                        return (
                            f"The file you want to rename does not exist:\n`{res_src}`\n\n"
                            f"It may have been moved, deleted, or is in a different folder (e.g. user Downloads vs. project folder)."
                        )
                    try:
                        return self.tools["move_file"].run(src=str(res_src), dst=str(res_dst))
                    except Exception as e:
                        return f"[ERROR] Rename failed: {e}"
                # Fallback: folder + filename (e.g. "rename file X in Downloads to Y")
                old_name, new_name = self._extract_rename_parts(task)
                if old_name and new_name:
                    src_path = path / old_name
                    if Path(new_name).suffix:
                        dst_path = path / new_name
                    else:
                        ext = Path(old_name).suffix
                        dst_path = path / (new_name + ext if ext else new_name)
                    safe_src, res_src = is_safe_path(str(src_path))
                    safe_dst, res_dst = is_safe_path(str(dst_path))
                    if not safe_src:
                        return f"[ERROR] {res_src}"
                    if not safe_dst:
                        return f"[ERROR] {res_dst}"
                    if not Path(res_src).exists():
                        return (
                            f"The file you want to rename does not exist:\n`{res_src}`\n\n"
                            f"It may have been moved, deleted, or is in a different folder (e.g. user Downloads vs. project folder)."
                        )
                    try:
                        return self.tools["move_file"].run(src=str(res_src), dst=str(res_dst))
                    except Exception as e:
                        return f"[ERROR] Rename failed: {e}"
                break  # Matched rename but couldn't extract - let LLM handle

        # ─────────────────────────────────────────────────────────────────────
        # PATTERN: Tree/Structure ("tree", "structure", "struktur")
        # ─────────────────────────────────────────────────────────────────────
        tree_patterns = [
            r"tree",
            r"structure",
            r"struktur",
            r"directory structure",
            r"folder structure",
        ]

        for pattern in tree_patterns:
            if re.search(pattern, task_lower):
                return self._show_tree(path)

        # No direct pattern matched - needs LLM
        return None
    
    def _extract_path(self, task: str) -> Path:
        """Extract and resolve path from task description."""
        task_lower = task.lower()
        
        # Check for folder aliases (try longest matches first)
        # Sort by length descending to match "documents" before "doc"
        sorted_aliases = sorted(self.folder_aliases.items(), key=lambda x: len(x[0]), reverse=True)
        for alias, path in sorted_aliases:
            # Match whole words or at word boundaries
            if re.search(r'\b' + re.escape(alias) + r'\b', task_lower):
                return path
        
        # Also try fuzzy matching for common typos/abbreviations
        # "dokumen" -> "documents" or "dokumente"
        fuzzy_mappings = {
            "dokumen": ["documents", "dokumente"],
            "dok": ["documents", "dokumente"],
            "down": ["downloads", "herunterladen"],
            "donw": ["downloads", "herunterladen"],
            "donwlodas": ["downloads", "herunterladen"],
            "doemlods": ["downloads", "herunterladen"],
            "dl": ["downloads", "herunterladen"],
            "desk": ["desktop", "arbeitsplatz"],
        }
        
        for fuzzy_key, possible_folders in fuzzy_mappings.items():
            if fuzzy_key in task_lower:
                for folder_name in possible_folders:
                    path = self.home / folder_name
                    if path.exists():
                        return path
        
        # Check for explicit path patterns
        # Windows: C:\Users\... or D:\...
        win_match = re.search(r'([A-Za-z]:\\[^\s"\']+)', task)
        if win_match:
            return Path(win_match.group(1))
        
        # Unix: /home/... or ~/...
        unix_match = re.search(r'((?:/[\w.-]+)+|~[\w/.-]*)', task)
        if unix_match:
            path_str = unix_match.group(1)
            if path_str.startswith('~'):
                return Path(path_str).expanduser()
            return Path(path_str)
        
        # Default to current directory
        return Path.cwd()

    def _summarize_web_results(self, task: str) -> Optional[str]:
        """
        Deterministic summarizer for web_search payloads (no LLM).
        Extracts the query and the top links/snippets so workflows like web_lookup
        don't accidentally trigger filesystem heuristics.
        """
        lower = task.lower()
        if "web search results" not in lower:
            return None

        import re

        # Try to extract "User question:" and "Query:"
        user_q = None
        m = re.search(r"user question:\s*(.+)", task, flags=re.IGNORECASE)
        if m:
            user_q = m.group(1).strip()

        query = None
        m = re.search(r"^query:\s*(.+)$", task, flags=re.IGNORECASE | re.MULTILINE)
        if m:
            query = m.group(1).strip()

        # Extract a few links (supports either "Link: ..." or raw URLs)
        urls = re.findall(r"https?://[^\s\)\]]+", task)
        urls = [u.rstrip(".,") for u in urls]

        # Extract numbered result titles if present ("1. **Title**")
        titles = re.findall(r"^\s*\d+\.\s*\*\*(.+?)\*\*", task, flags=re.MULTILINE)

        lines = []
        lines.append("### Web Lookup (Deterministic Summary)")
        if user_q:
            lines.append(f"**Question:** {user_q}")
        elif query:
            lines.append(f"**Question:** {query}")

        if titles:
            lines.append("\n**Top results:**")
            for t in titles[:5]:
                lines.append(f"- {t.strip()}")

        if urls:
            lines.append("\n**Sources:**")
            for u in urls[:5]:
                lines.append(f"- {u}")
        else:
            lines.append("\n[INFO] No URLs found in the search payload.")

        return "\n".join(lines)
    
    def _extract_search_term(self, task: str) -> Optional[str]:
        """Extract search term/pattern from task."""
        # Look for quoted strings
        quoted = re.search(r'["\']([^"\']+)["\']', task)
        if quoted:
            return quoted.group(1)

        # Look for filename after "file" or "datei" (e.g. "file 26-B001-105272426-97758570.PDF")
        file_match = re.search(r'(?:file|datei)\s+([\w.-]+(?:\.\w+)?)', task, re.IGNORECASE)
        if file_match:
            term = file_match.group(1).strip()
            if '*' not in term:
                return f"*{term}*"
            return term

        # Look for file extensions
        ext_match = re.search(r'\.([\w]+)\b', task)
        if ext_match:
            return f"*.{ext_match.group(1)}"

        # Look for specific file patterns
        patterns = [
            r'named?\s+(\S+)',
            r'called?\s+(\S+)',
            r'for\s+(\S+)',
            r'suche\s+(?:nach\s+)?(\S+)',
        ]

        for pattern in patterns:
            match = re.search(pattern, task.lower())
            if match:
                term = match.group(1)
                # Add wildcards if not present
                if '*' not in term:
                    return f"*{term}*"
                return term

        return None
    
    def _extract_file_path(self, task: str) -> Optional[Path]:
        """Extract specific file path from task.
        Handles quoted paths, spaces, Unicode, no-extension paths, Windows and Unix.
        """
        # 1. Quoted paths — most reliable, handles spaces cleanly
        #    Supports "path", 'path', `path`
        for quote in ('"', "'", '`'):
            m = re.search(rf'{re.escape(quote)}([^{re.escape(quote)}\n]+){re.escape(quote)}', task)
            if m:
                candidate = m.group(1).strip()
                # Must look like an absolute or home-relative path
                if re.match(r'^([A-Za-z]:[/\\]|/|~/)', candidate):
                    return Path(candidate).expanduser()

        # 2. Windows absolute path — C:\... or D:\... (with or without spaces/extension)
        win_m = re.search(r'([A-Za-z]:[/\\][^\n"\'`]{1,300})', task)
        if win_m:
            candidate = win_m.group(1).rstrip(' .,;:!?\'"')
            # Stop before common sentence connectors to avoid grabbing trailing words
            candidate = re.split(r'\s+(?:und|oder|and|or|to|from|nach|von|aus|mit|with)\b', candidate)[0].rstrip()
            return Path(candidate)

        # 3. Tilde path — ~/...
        tilde_m = re.search(r'(~/[^\s"\'`\n]*)', task)
        if tilde_m:
            return Path(tilde_m.group(1).rstrip(' .,;:!?')).expanduser()

        # 4. Unix absolute path — /something/...
        #    Character class covers ASCII word chars + spaces + common European Unicode
        _pc = r'[\w\-. äöüÄÖÜéèêàâîïôùûçßæœ]'
        unix_m = re.search(rf'(/(?:{_pc}+/)* {_pc}+(?:\.\w+)?)', task)
        if not unix_m:
            # Simpler fallback without spaces (original behaviour)
            unix_m = re.search(r'((?:/[\w.-]+)+(?:\.\w+)?)', task)
        if unix_m:
            return Path(unix_m.group(1).rstrip(' .,;:!?'))

        # 5. Known-extension relative path (conservative — avoids grabbing random words)
        _EXTS = r'pdf|docx?|xlsx?|pptx?|txt|md|json|xml|csv|log|py|js|ts|sh|yaml|yml|ini|cfg|toml|env'
        rel_m = re.search(rf'\b([\w./\\-]+\.(?:{_EXTS}))\b', task, re.IGNORECASE)
        if rel_m:
            return Path(rel_m.group(1))

        return None

    def _extract_rename_full_paths(self, task: str) -> Tuple[Optional[Path], Optional[Path]]:
        """Extract (src_path, dst_path) when task contains full Windows/Unix paths.
        Returns (None, None) if extraction fails."""
        # Windows: C:\path\file.ext or D:\path\file.ext
        win_path = r'[A-Za-z]:\\[^\s"\'<>|]+'
        # Unix: /path/file.ext or ~/path/file.ext
        unix_path = r'(?:/[^\s"\'<>|]+|~[^\s"\'<>|]*)'
        path_pat = f'({win_path}|{unix_path})'
        paths = re.findall(path_pat, task)
        if len(paths) >= 2:
            try:
                p1 = Path(paths[0]).expanduser() if paths[0].startswith('~') else Path(paths[0])
                p2 = Path(paths[1]).expanduser() if paths[1].startswith('~') else Path(paths[1])
                return (p1, p2)
            except Exception:
                pass
        return (None, None)

    def _extract_rename_parts(self, task: str) -> Tuple[Optional[str], Optional[str]]:
        """Extract (old_filename, new_name) for rename operations.
        Returns (None, None) if extraction fails."""
        # Old filename: after "Datei"/"file" or a filename pattern with extension
        old_name = None
        file_match = re.search(r'(?:file|datei)\s+([\w.-]+\.\w+)', task, re.IGNORECASE)
        if file_match:
            old_name = file_match.group(1).strip()
        if not old_name:
            # Fallback: any filename-like token (word.word) before "umbenennen"/"rename"
            before_rename = re.split(r'umbenennen|rename', task, flags=re.IGNORECASE)[0]
            ext_match = re.search(r'([\w.-]+\.(?:pdf|docx?|txt|xlsx?|png|jpg|jpeg))\b', before_rename, re.IGNORECASE)
            if ext_match:
                old_name = ext_match.group(1).strip()

        # New name: after "zu"/"to"
        new_name = None
        to_match = re.search(r'(?:zu|to)\s+([\w.-]+)', task, re.IGNORECASE)
        if to_match:
            new_name = to_match.group(1).strip()

        return (old_name, new_name) if (old_name and new_name) else (None, None)
    
    # ═══════════════════════════════════════════════════════════════════════════
    # DIRECT TASK IMPLEMENTATIONS
    # ═══════════════════════════════════════════════════════════════════════════
    
    def _count_files(self, path: Path, file_type: str = None) -> str:
        """Count files in directory, optionally filtered by extension."""
        safe, res = is_safe_path(str(path))
        if not safe:
            return f"[ERROR] {res}"
        path = Path(res)

        filter_text = f" ({file_type.upper()} files)" if file_type else ""
        UI.event("Librarian", f"Counting files{filter_text} in {path}...", style="dim")

        if not path.exists():
            return f"[ERROR] Path does not exist: {path}"
        
        if not path.is_dir():
            return f"[ERROR] Not a directory: {path}"
        
        try:
            # Count files (not directories), optionally filtered by extension
            all_files = [f for f in path.iterdir() if f.is_file()]
            
            if file_type:
                # Filter by extension (case-insensitive)
                ext = file_type.lower().lstrip('.')
                files = [f for f in all_files if f.suffix.lower() == f'.{ext}']
            else:
                files = all_files
            
            dirs = [d for d in path.iterdir() if d.is_dir()]
            
            total_files = len(files)
            total_dirs = len(dirs)
            
            # Get some stats
            extensions = {}
            total_size = 0
            
            for f in files:
                ext = f.suffix.lower() or "(no extension)"
                extensions[ext] = extensions.get(ext, 0) + 1
                try:
                    total_size += f.stat().st_size
                except:
                    pass
            
            # Format size
            if total_size < 1024:
                size_str = f"{total_size} B"
            elif total_size < 1024 * 1024:
                size_str = f"{total_size / 1024:.1f} KB"
            elif total_size < 1024 * 1024 * 1024:
                size_str = f"{total_size / (1024 * 1024):.1f} MB"
            else:
                size_str = f"{total_size / (1024 * 1024 * 1024):.1f} GB"
            
            # Top extensions (only show if not filtering by specific type)
            if file_type:
                ext_str = f".{file_type.lower()}: {total_files}"
            else:
                top_ext = sorted(extensions.items(), key=lambda x: x[1], reverse=True)[:5]
                ext_str = ", ".join([f"{ext}: {count}" for ext, count in top_ext])
            
            return f"""### Folder: {path.name or path}

**Files:** {total_files}
**Folders:** {total_dirs}
**Total Size:** {size_str}

**Top Extensions:** {ext_str}"""
            
        except PermissionError:
            return f"[ERROR] Permission denied: {path}"
        except Exception as e:
            return f"[ERROR] Error counting files: {e}"
    
    def _list_files(self, path: Path, limit: int = 30) -> str:
        """List files in directory."""
        safe, res = is_safe_path(str(path))
        if not safe:
            return f"[ERROR] {res}"
        path = Path(res)

        UI.event("Librarian", f"Listing files in {path}...", style="dim")

        if not path.exists():
            return f"[ERROR] Path does not exist: {path}"
        
        if not path.is_dir():
            return f"[ERROR] Not a directory: {path}"
        
        try:
            items = list(path.iterdir())
            files = sorted([f for f in items if f.is_file()], key=lambda x: x.name.lower())
            dirs = sorted([d for d in items if d.is_dir()], key=lambda x: x.name.lower())
            
            result = [f"### Folder: {path.name or path}\n"]
            
            # Directories first
            for d in dirs[:10]:
                result.append(f"[DIR] {d.name}/")
            
            if len(dirs) > 10:
                result.append(f"   ... and {len(dirs) - 10} more folders")
            
            result.append("")
            
            # Files
            for f in files[:limit]:
                size = f.stat().st_size
                if size < 1024:
                    size_str = f"{size}B"
                elif size < 1024 * 1024:
                    size_str = f"{size // 1024}KB"
                else:
                    size_str = f"{size // (1024 * 1024)}MB"
                result.append(f"[FILE] {f.name} ({size_str})")
            
            if len(files) > limit:
                result.append(f"\n... and {len(files) - limit} more files")
            
            result.append(f"\n**Total:** {len(files)} files, {len(dirs)} folders")
            
            return "\n".join(result)
            
        except PermissionError:
            return f"[ERROR] Permission denied: {path}"
        except Exception as e:
            return f"[ERROR] Error listing files: {e}"
    
    def _find_files(self, path: Path, pattern: str) -> str:
        """Find files matching pattern."""
        safe, res = is_safe_path(str(path))
        if not safe:
            return f"[ERROR] {res}"
        path = Path(res)

        UI.event("Librarian", f"Searching for '{pattern}' in {path}...", style="dim")

        if not path.exists():
            return f"[ERROR] Path does not exist: {path}"
        
        try:
            import os as _os, time as _time, fnmatch as _fnmatch
            # Bounded, pruned search: walk lazily, skip huge/system dirs, stop early.
            # The previous code did list(path.rglob('*')) which materialised the ENTIRE
            # tree (node_modules, venv, /mnt mounts, …) and could hang for minutes.
            _SKIP_DIRS = {
                "node_modules", "venv", ".venv", "env", "__pycache__", ".git",
                "site-packages", "dist", "build", ".next", ".cache", ".npm",
                "Library", "AppData", ".cargo", ".rustup", "proc", "sys", "dev",
            }
            _MAX = 50
            _deadline = _time.monotonic() + 20.0   # soft cap → always responsive
            _glob = ('*' in pattern or '?' in pattern)
            _needle = pattern.lower()
            matches = []
            for _root, _dirs, _files in _os.walk(str(path)):
                # prune heavy/system/hidden dirs in-place so we never descend into them
                _dirs[:] = [d for d in _dirs if d not in _SKIP_DIRS and not d.startswith('.')]
                for _fn in _files:
                    if (_fnmatch.fnmatch(_fn, pattern) if _glob else _needle in _fn.lower()):
                        matches.append(Path(_root) / _fn)
                        if len(matches) >= _MAX:
                            break
                if len(matches) >= _MAX or _time.monotonic() > _deadline:
                    break

            if not matches:
                return f"No files found matching '{pattern}' in {path}"

            result = [f"### Search: '{pattern}'\n"]
            result.append(f"Found {len(matches)} matches:\n")

            file_matches = []
            for f in matches[:30]:
                if f.is_file():
                    full_path = str(f.resolve())
                    file_matches.append(f)
                    result.append(f"[FILE] {f.name}\n  Full path: {full_path}")
                else:
                    result.append(f"[DIR] {f.name}")

            if len(matches) > 30:
                result.append(f"\n... and {len(matches) - 30} more")

            # If exactly one file found, add send hint for main agent
            if len(file_matches) == 1:
                result.append(f"\n**To send via Telegram:** Use send_telegram(message=\"...\", file_path=\"{str(file_matches[0].resolve())}\")")

            return "\n".join(result)
            
        except Exception as e:
            return f"[ERROR] Error searching: {e}"
    
    def _write_file(self, file_path: Path, content: str) -> str:
        """Write content to a file."""
        try:
            # Use the write_file tool
            result = self.tools["write_file"].run(path=str(file_path), content=content)
            return result
        except Exception as e:
            return f"Error writing file: {str(e)}"

    def _pdf_ocr_fallback(self, file_path: Path, max_pages: int) -> str:
        """Thin wrapper -> shared OCR helper (single copy lives in vaf/core/pdf_extract.py)."""
        from vaf.core.pdf_extract import pdf_ocr_fallback
        return pdf_ocr_fallback(file_path, max_pages)
    
    def _read_file(self, file_path: Path, enable_chunking: bool = True) -> str:
        """Read file contents - supports text, PDF, Word, Excel, PowerPoint.

        Args:
            file_path: Path to the file to read
            enable_chunking: If True, automatically chunk large files (default: True)
        """
        safe, res = is_safe_path(str(file_path))
        if not safe:
            return f"[ERROR] {res}"
        file_path = Path(res)

        UI.event("Librarian", f"Reading {file_path}...", style="dim")

        if not file_path.exists():
            parent = file_path.parent
            nearby = ""
            try:
                if parent.exists():
                    entries = [e.name for e in parent.iterdir() if e.is_file()][:8]
                    if entries:
                        nearby = f"\nFiles in {parent}:\n" + "\n".join(f"  - {e}" for e in entries)
            except Exception:
                pass
            return (
                f"[LIBRARIAN_ERROR type=file_not_found]\n"
                f"Path: {file_path}\n"
                f"Home: {self.home}\n"
                f"{nearby}\n"
                f"Suggestions for Main Agent:\n"
                f"- Verify the exact filename and extension\n"
                f"- The path may contain spaces — pass it in quotes\n"
                f"- Use librarian_agent(task='list files in {parent}') to see what is there\n"
                f"- On Windows paths use backslash: C:\\Users\\...\\file.txt"
            )

        if not file_path.is_file():
            return (
                f"[LIBRARIAN_ERROR type=not_a_file]\n"
                f"Path: {file_path}\n"
                f"This path exists but is a directory, not a file.\n"
                f"Suggestions for Main Agent:\n"
                f"- Use librarian_agent(task='list files in {file_path}') to list its contents"
            )
        
        try:
            # Get file extension
            ext = file_path.suffix.lower()
            
            # Check file size with enhanced limits and chunking support
            size = file_path.stat().st_size
            size_kb = size // 1024
            size_mb = size / (1024 * 1024)
            
            # Enhanced size limits based on file type (configurable via config.json)
            from vaf.core.config import Config
            config = Config.load()
            
            size_limits = {
                '.pdf': config.get('librarian_max_pdf_size_mb', 50) * 1024 * 1024,
                '.docx': config.get('librarian_max_doc_size_mb', 20) * 1024 * 1024,
                '.xlsx': config.get('librarian_max_excel_size_mb', 30) * 1024 * 1024,
                '.pptx': config.get('librarian_max_doc_size_mb', 20) * 1024 * 1024,
                'text': config.get('librarian_max_text_size_kb', 500) * 1024,
            }
            
            max_size = size_limits.get(ext, size_limits['text'])
            auto_chunk = config.get('librarian_auto_chunk_large_files', True)
            
            # Check if file is too large
            if size > max_size:
                max_mb = max_size / (1024 * 1024)
                
                # Offer chunking for supported formats (if enabled in config)
                if auto_chunk and enable_chunking and ext in ['.pdf', '.txt', '.md', '.json', '.xml']:
                    return (
                        f"[INFO] File is large ({size_mb:.1f} MB, max {max_mb:.0f} MB direct read)\n\n"
                        f"**Auto-Chunking Enabled:** Reading file in manageable sections...\n\n"
                        f"{self._read_file_chunked(file_path, ext)}\n\n"
                        f"**Configuration:**\n"
                        f"- To change size limits, edit `~/.vaf/config.json`\n"
                        f"- Current limit for {ext} files: {max_mb:.0f} MB\n"
                        f"- Auto-chunking: {'Enabled' if auto_chunk else 'Disabled'}"
                    )
                else:
                    return (
                        f"[ERROR] File too large: {size_mb:.1f} MB\n\n"
                        f"**File:** {file_path.name}\n"
                        f"**Size:** {size_kb:,} KB ({size_mb:.2f} MB)\n"
                        f"**Maximum:** {max_mb:.0f} MB for {ext} files\n\n"
                        f"**Suggestions:**\n"
                        f"- Split the file into smaller parts\n"
                        f"- Extract specific pages/sections you need\n"
                        f"- For PDFs: Use a PDF editor to extract pages\n"
                        f"- For Excel: Export specific sheets as separate files\n"
                        f"- Compress the file if possible"
                    )
            
            # ═══════════════════════════════════════════════════════════
            # PDF Files
            # ═══════════════════════════════════════════════════════════
            if ext == '.pdf':
                try:
                    from vaf.core.config import Config
                    from vaf.core.pdf_extract import extract_pdf_markdown
                    config = Config.load()
                    max_pages = config.get('librarian_pdf_max_pages_preview', 50)
                    use_ocr = config.get("librarian_ocr_fallback_for_pdf", True)

                    # Shared extractor: pdfplumber markdown (headings + tables), PyPDF2 + OCR fallbacks.
                    res = extract_pdf_markdown(file_path, max_pages=max_pages, ocr_fallback=use_ocr)
                    full_text = res["markdown"]
                    num_pages = res["num_pages"]

                    # Truncate if too long
                    if len(full_text) > 15000:
                        full_text = full_text[:15000] + "\n\n... (truncated)"

                    if not full_text.strip():
                        full_text = (
                            "[Scanned PDF: no embedded text detected. For OCR install: "
                            "pip install pdf2image pytesseract, and system tools: poppler (pdf2image), Tesseract (pytesseract). "
                            "Then re-open this file."
                        )

                    return f"### PDF: {file_path.name}\n**Pages:** {num_pages}\n\n{full_text}"

                except ImportError:
                    return f"[ERROR] PDF support not installed. Run: pip install pdfplumber PyPDF2"
                except Exception as e:
                    err_str = str(e)
                    hint = " For AES-encrypted PDFs (e.g. bank statements) run: pip install pycryptodome" if ("PyCryptodome" in err_str or "AES" in err_str) else ""
                    return f"[ERROR] Failed to read PDF: {e}{hint}"
            
            # ═══════════════════════════════════════════════════════════
            # Word Documents (.docx)
            # ═══════════════════════════════════════════════════════════
            elif ext == '.docx':
                try:
                    from docx import Document
                    doc = Document(file_path)
                    
                    content = []
                    
                    # Extract paragraphs
                    for para in doc.paragraphs:
                        text = para.text.strip()
                        if text:
                            content.append(text)
                    
                    # Extract tables
                    if doc.tables:
                        content.append("\n--- Tables ---")
                        for i, table in enumerate(doc.tables[:5], 1):  # Limit to 5 tables
                            content.append(f"\nTable {i}:")
                            for row in table.rows[:10]:  # Limit to 10 rows per table
                                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                                if row_text:
                                    content.append(row_text)
                    
                    full_text = "\n".join(content)
                    
                    # Truncate if too long
                    if len(full_text) > 15000:
                        full_text = full_text[:15000] + "\n\n... (truncated)"
                    
                    return f"### Word Document: {file_path.name}\n**Paragraphs:** {len(doc.paragraphs)}\n**Tables:** {len(doc.tables)}\n\n{full_text}"
                    
                except ImportError:
                    return f"[ERROR] Word document support not installed. Run: pip install python-docx"
                except Exception as e:
                    return f"[ERROR] Failed to read Word document: {e}"
            
            # ═══════════════════════════════════════════════════════════
            # Excel Files (.xlsx, .xls)
            # ═══════════════════════════════════════════════════════════
            elif ext in ['.xlsx', '.xls']:
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                    
                    content = []
                    content.append(f"**Sheets:** {', '.join(wb.sheetnames)}\n")
                    
                    # Read first 3 sheets
                    for sheet_name in wb.sheetnames[:3]:
                        sheet = wb[sheet_name]
                        content.append(f"\n--- Sheet: {sheet_name} ---")
                        
                        # Get dimensions
                        max_row = min(sheet.max_row, 50)  # Limit to 50 rows
                        max_col = min(sheet.max_column, 20)  # Limit to 20 columns
                        
                        # Read data
                        for row in sheet.iter_rows(min_row=1, max_row=max_row, max_col=max_col, values_only=True):
                            row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                            if row_text.strip():
                                content.append(row_text)
                    
                    if len(wb.sheetnames) > 3:
                        content.append(f"\n... ({len(wb.sheetnames) - 3} more sheets not shown)")
                    
                    full_text = "\n".join(content)
                    
                    # Truncate if too long
                    if len(full_text) > 15000:
                        full_text = full_text[:15000] + "\n\n... (truncated)"
                    
                    return f"### Excel File: {file_path.name}\n{full_text}"
                    
                except ImportError:
                    return f"[ERROR] Excel support not installed. Run: pip install openpyxl"
                except Exception as e:
                    return f"[ERROR] Failed to read Excel file: {e}"
            
            # ═══════════════════════════════════════════════════════════
            # PowerPoint Files (.pptx)
            # ═══════════════════════════════════════════════════════════
            elif ext == '.pptx':
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    
                    content = []
                    content.append(f"**Slides:** {len(prs.slides)}\n")
                    
                    # Read first 20 slides
                    for i, slide in enumerate(prs.slides[:20], 1):
                        content.append(f"\n--- Slide {i} ---")
                        
                        # Extract text from shapes
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                content.append(shape.text)
                    
                    if len(prs.slides) > 20:
                        content.append(f"\n... ({len(prs.slides) - 20} more slides not shown)")
                    
                    full_text = "\n".join(content)
                    
                    # Truncate if too long
                    if len(full_text) > 15000:
                        full_text = full_text[:15000] + "\n\n... (truncated)"
                    
                    return f"### PowerPoint: {file_path.name}\n{full_text}"
                    
                except ImportError:
                    return f"[ERROR] PowerPoint support not installed. Run: pip install python-pptx"
                except Exception as e:
                    return f"[ERROR] Failed to read PowerPoint file: {e}"
            
            # ═══════════════════════════════════════════════════════════
            # Text Files (default)
            # ═══════════════════════════════════════════════════════════
            else:
                content = file_path.read_text(encoding='utf-8', errors='replace')
            
            # Truncate if needed. HTML files need higher limit for Document Viewer rendering.
            truncate_limit = 100_000 if ext in ('.html', '.htm') else 5000
            if len(content) > truncate_limit:
                content = content[:truncate_limit] + "\n\n... (truncated)"
            
            return f"### File: {file_path.name}\n\n```\n{content}\n```"
            
        except PermissionError:
            return (
                f"[LIBRARIAN_ERROR type=permission_denied]\n"
                f"Path: {file_path}\n"
                f"The file exists but cannot be read (permission denied).\n"
                f"Suggestions for Main Agent:\n"
                f"- This file is protected by the OS and cannot be opened\n"
                f"- Try a different file or ask the user to check permissions"
            )
        except Exception as e:
            error_type = type(e).__name__
            return (
                f"[LIBRARIAN_ERROR type=read_failed error={error_type}]\n"
                f"Path: {file_path}\n"
                f"Error: {e}\n"
                f"Suggestions for Main Agent:\n"
                f"- The file format may not be supported or the file may be corrupted\n"
                f"- Supported formats: PDF, DOCX, XLSX, PPTX, TXT, MD, JSON, XML, CSV and most text files\n"
                f"- Try: librarian_agent(task='list files in {file_path.parent}') to verify the file exists"
            )

    def _read_file_chunked(self, file_path: Path, ext: str) -> str:
        """Read large files in chunks and provide summary.
        
        This method implements a Map-Reduce strategy for large files:
        1. Split file into manageable chunks
        2. Process each chunk separately
        3. Provide navigation links for detailed sections
        """
        UI.event("Librarian", f"Reading file in chunks: {file_path.name}...", style="dim")
        
        try:
            # PDF chunking
            if ext == '.pdf':
                import PyPDF2
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    num_pages = len(pdf_reader.pages)
                    
                    # Read first 20 pages as preview
                    preview_pages = min(20, num_pages)
                    content = []
                    
                    for page_num in range(preview_pages):
                        page = pdf_reader.pages[page_num]
                        page_text = page.extract_text()
                        if page_text.strip():
                            # Only show first 200 chars per page for overview
                            preview = page_text[:200].strip()
                            content.append(f"**Page {page_num + 1}:** {preview}...")
                    
                    result = f"### PDF Preview: {file_path.name}\n"
                    result += f"**Total Pages:** {num_pages}\n"
                    result += f"**Showing:** First {preview_pages} pages (preview)\n\n"
                    result += "\n".join(content)
                    result += f"\n\n**Note:** This is a preview. The full document has {num_pages} pages."
                    result += f"\n**Tip:** Ask me to 'read pages 10-20 of {file_path.name}' for specific sections."
                    
                    return result
            
            # Text file chunking
            elif ext in ['.txt', '.md', '.json', '.xml', '.csv', '.log']:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    # Read first 10,000 characters
                    content = f.read(10000)
                    
                    result = f"### Text File Preview: {file_path.name}\n"
                    result += f"**Showing:** First 10,000 characters\n\n"
                    result += content
                    result += "\n\n... (file continues)"
                    result += f"\n**Tip:** Ask me to 'read next section' or 'search for [keyword] in {file_path.name}'"
                    
                    return result
            
            else:
                return f"[INFO] Chunked reading not supported for {ext} files yet."
                
        except Exception as e:
            err_str = str(e)
            hint = " For AES-encrypted PDFs run: pip install pycryptodome" if ("PyCryptodome" in err_str or "AES" in err_str) else ""
            return f"[ERROR] Failed to read file in chunks: {e}{hint}"
    
    def _show_tree(self, path: Path, max_depth: int = 3) -> str:
        """Show directory tree."""
        safe, res = is_safe_path(str(path))
        if not safe:
            return f"[ERROR] {res}"
        path = Path(res)

        UI.event("Librarian", f"Building tree for {path}...", style="dim")

        if not path.exists():
            return f"[ERROR] Path does not exist: {path}"
        
        try:
            result = self.tools["tree"].run(path=str(path), depth=max_depth)
            return f"### Directory Tree\n\n```\n{result}\n```"
        except Exception as e:
            return f"[ERROR] Error building tree: {e}"
    
    # ═══════════════════════════════════════════════════════════════════════════
    # DISK/STORAGE INFORMATION (Cross-Platform)
    # ═══════════════════════════════════════════════════════════════════════════
    
    def _get_disk_info(self) -> str:
        """Get information about all storage devices (HDD, SSD, USB) - OS independent."""
        drives_info = []
        
        try:
            if HAS_PSUTIL:
                # Use psutil for cross-platform disk info
                partitions = psutil.disk_partitions()
                for partition in partitions:
                    try:
                        usage = psutil.disk_usage(partition.mountpoint)
                        total_gb = usage.total / (1024**3)
                        free_gb = usage.free / (1024**3)
                        used_gb = usage.used / (1024**3)
                        used_percent = (usage.used / usage.total) * 100
                        
                        # Try to detect device type
                        device_type = "Unknown"
                        if sys.platform == "win32":
                            # Windows: Check if it's a removable drive
                            if partition.fstype == "" or "removable" in partition.opts.lower():
                                device_type = "USB/Removable"
                            elif "fixed" in partition.opts.lower():
                                # Try to detect SSD vs HDD (Windows 10+)
                                device_type = self._detect_windows_drive_type(partition.device)
                            else:
                                device_type = "Fixed Drive"
                        else:
                            # Linux/Mac: Check mount point and device name
                            if "/media/" in partition.mountpoint or "/mnt/" in partition.mountpoint:
                                device_type = "USB/Removable"
                            elif "/dev/sd" in partition.device:
                                # Linux: Check if it's an SSD (look for rotational attribute)
                                device_type = self._detect_linux_drive_type(partition.device)
                            elif "/dev/disk" in partition.device:
                                # macOS
                                device_type = "Fixed Drive"
                        
                        drives_info.append({
                            "device": partition.device,
                            "mountpoint": partition.mountpoint,
                            "fstype": partition.fstype,
                            "type": device_type,
                            "total_gb": total_gb,
                            "free_gb": free_gb,
                            "used_gb": used_gb,
                            "used_percent": used_percent
                        })
                    except PermissionError:
                        # Skip drives we can't access
                        continue
                    except Exception:
                        continue
            else:
                # Fallback: Use OS-specific commands
                if sys.platform == "win32":
                    drives_info = self._get_windows_disk_info()
                elif sys.platform == "darwin":
                    drives_info = self._get_macos_disk_info()
                else:
                    drives_info = self._get_linux_disk_info()
        except Exception as e:
            # If all fails, return basic info
            return f"[WARN] Could not retrieve disk information: {e}"
        
        if not drives_info:
            return "No storage devices found."
        
        # Format the information
        result = "**Storage Devices:**\n\n"
        device_count = len(drives_info)
        
        for i, drive in enumerate(drives_info, 1):
            device = drive.get("device", "Unknown")
            mount = drive.get("mountpoint", "")
            drive_type = drive.get("type", "Unknown")
            total = drive.get("total_gb", 0)
            free = drive.get("free_gb", 0)
            used_pct = drive.get("used_percent", 0)
            
            result += f"{i}. **{device}** ({drive_type})\n"
            result += f"   Mount: {mount}\n"
            result += f"   Total: {total:.1f} GB | Free: {free:.1f} GB | Used: {used_pct:.1f}%\n\n"
        
        result += f"**Total:** {device_count} storage device(s)"
        return result
    
    def _detect_windows_drive_type(self, device: str) -> str:
        """Detect if Windows drive is SSD or HDD."""
        try:
            # Use PowerShell to check if drive is SSD
            # Get drive letter from device path (e.g., "C:" from "C:\\")
            drive_letter = device[0] if device else "C"
            cmd = f'powershell -Command "Get-PhysicalDisk | Get-Disk | Where-Object {{$_.Number -eq (Get-Partition -DriveLetter {drive_letter}).DiskNumber}} | Select-Object -ExpandProperty MediaType"'
            result = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=3)
            if result.returncode == 0:
                media_type = result.stdout.strip()
                if "SSD" in media_type or "ssd" in media_type.lower():
                    return "SSD"
                elif "HDD" in media_type or "HDD" in media_type or "Rotational" in media_type:
                    return "HDD"
        except:
            pass
        # Default: Try to guess based on common patterns
        if device and len(device) > 0:
            # C: is usually the main drive (often SSD on modern systems)
            if device[0].upper() == "C":
                return "Fixed Drive (likely SSD)"
        return "Fixed Drive"
    
    def _detect_linux_drive_type(self, device: str) -> str:
        """Detect if Linux drive is SSD or HDD."""
        try:
            # Check /sys/block for rotational attribute
            block_device = device.replace("/dev/", "").rstrip("0123456789")
            rotational_path = f"/sys/block/{block_device}/queue/rotational"
            if os.path.exists(rotational_path):
                with open(rotational_path, 'r') as f:
                    is_rotational = f.read().strip() == "1"
                return "HDD" if is_rotational else "SSD"
        except:
            pass
        return "Fixed Drive"
    
    def _get_windows_disk_info(self) -> List[Dict]:
        """Get disk info on Windows using wmic."""
        drives = []
        try:
            # Get logical drives
            cmd = 'wmic logicaldisk get deviceid,drivetype,freespace,size,volumename'
            result = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=5)
            if result.returncode == 0:
                lines = result.stdout.strip().split('\n')[1:]  # Skip header
                for line in lines:
                    parts = line.split()
                    if len(parts) >= 2:
                        device = parts[0]
                        drive_type_code = parts[1] if len(parts) > 1 else "3"
                        # Drive types: 2=Removable, 3=Fixed, 4=Network, 5=CD
                        if drive_type_code == "2":
                            drive_type = "USB/Removable"
                        elif drive_type_code == "3":
                            drive_type = "Fixed Drive"
                        else:
                            continue
                        
                        # Get size info
                        free_bytes = int(parts[2]) if len(parts) > 2 and parts[2].isdigit() else 0
                        total_bytes = int(parts[3]) if len(parts) > 3 and parts[3].isdigit() else 0
                        
                        if total_bytes > 0:
                            drives.append({
                                "device": device,
                                "mountpoint": device + "\\",
                                "type": drive_type,
                                "total_gb": total_bytes / (1024**3),
                                "free_gb": free_bytes / (1024**3),
                                "used_gb": (total_bytes - free_bytes) / (1024**3),
                                "used_percent": ((total_bytes - free_bytes) / total_bytes) * 100
                            })
        except:
            pass
        return drives
    
    def _get_linux_disk_info(self) -> List[Dict]:
        """Get disk info on Linux using df."""
        drives = []
        try:
            result = subprocess.run(['df', '-h'], capture_output=True, text=True, timeout=5)
            if result.returncode == 0:
                lines = result.stdout.strip().split('\n')[1:]  # Skip header
                for line in lines:
                    parts = line.split()
                    if len(parts) >= 6:
                        device = parts[0]
                        mount = parts[5]
                        size_str = parts[1]
                        used_str = parts[2]
                        avail_str = parts[3]
                        use_pct = parts[4].rstrip('%')
                        
                        # Determine drive type
                        if "/media/" in mount or "/mnt/" in mount:
                            drive_type = "USB/Removable"
                        else:
                            drive_type = self._detect_linux_drive_type(device)
                        
                        # Parse sizes (convert from human-readable)
                        def parse_size(s):
                            if 'G' in s:
                                return float(s.replace('G', ''))
                            elif 'T' in s:
                                return float(s.replace('T', '')) * 1024
                            elif 'M' in s:
                                return float(s.replace('M', '')) / 1024
                            return 0
                        
                        total_gb = parse_size(size_str)
                        used_gb = parse_size(used_str)
                        free_gb = parse_size(avail_str)
                        
                        drives.append({
                            "device": device,
                            "mountpoint": mount,
                            "type": drive_type,
                            "total_gb": total_gb,
                            "free_gb": free_gb,
                            "used_gb": used_gb,
                            "used_percent": float(use_pct)
                        })
        except:
            pass
        return drives
    
    def _get_macos_disk_info(self) -> List[Dict]:
        """Get disk info on macOS using df and diskutil."""
        drives = []
        try:
            # Use df for basic info
            result = subprocess.run(['df', '-h'], capture_output=True, text=True, timeout=5)
            if result.returncode == 0:
                lines = result.stdout.strip().split('\n')[1:]
                for line in lines:
                    parts = line.split()
                    if len(parts) >= 9:
                        device = parts[0]
                        mount = parts[8]
                        size_str = parts[1]
                        used_str = parts[2]
                        avail_str = parts[3]
                        use_pct = parts[4].rstrip('%')
                        
                        # Determine drive type (check if external)
                        drive_type = "Fixed Drive"
                        try:
                            diskutil_result = subprocess.run(
                                ['diskutil', 'info', device],
                                capture_output=True, text=True, timeout=2
                            )
                            if 'External' in diskutil_result.stdout or 'Removable' in diskutil_result.stdout:
                                drive_type = "USB/Removable"
                        except:
                            pass
                        
                        def parse_size(s):
                            if 'G' in s:
                                return float(s.replace('Gi', '').replace('G', ''))
                            elif 'T' in s:
                                return float(s.replace('Ti', '').replace('T', '')) * 1024
                            return 0
                        
                        total_gb = parse_size(size_str)
                        used_gb = parse_size(used_str)
                        free_gb = parse_size(avail_str)
                        
                        drives.append({
                            "device": device,
                            "mountpoint": mount,
                            "type": drive_type,
                            "total_gb": total_gb,
                            "free_gb": free_gb,
                            "used_gb": used_gb,
                            "used_percent": float(use_pct)
                        })
        except:
            pass
        return drives
    
    # ═══════════════════════════════════════════════════════════════════════════
    # ERROR HANDLING & FORMATTING
    # ═══════════════════════════════════════════════════════════════════════════
    
    def _format_error_response(self, response, task: str, history: List, tools_schema: List) -> str:
        """Format error response with helpful context for Main Agent."""
        status_code = response.status_code
        error_message = ""
        error_type = ""
        
        # Try to extract detailed error from response
        try:
            error_data = response.json()
            if "error" in error_data:
                error_obj = error_data["error"]
                error_message = error_obj.get("message", "")
                error_type = error_obj.get("type", "")
        except:
            error_message = response.text[:500] if response.text else ""
        
        # Analyze the task to provide context
        task_lower = task.lower()
        path_hint = ""
        if "dokumen" in task_lower or "documents" in task_lower:
            # Check if path exists
            possible_paths = [
                self.home / "Documents",
                self.home / "Dokumente",
            ]
            existing_paths = [str(p) for p in possible_paths if p.exists()]
            if existing_paths:
                path_hint = f"\nHint: The Documents folder exists at: {existing_paths[0]}"
            else:
                path_hint = f"\nHint: Documents folder not found. Try checking: {self.home}"
        
        # Analyze what might be wrong
        suggestions = []
        
        if status_code == 400:
            suggestions.append("- **Invalid request format** - The task description might be unclear")
            suggestions.append("- **Missing parameters** - The task might need more specific information")
            suggestions.append("- **Path not found** - The folder path might be incorrect or not exist")
            
            # Check if path extraction failed
            extracted_path = self._extract_path(task)
            if not extracted_path.exists():
                suggestions.append(f"- **Path issue** - Extracted path '{extracted_path}' does not exist")
                # Suggest alternatives
                if "dokumen" in task_lower:
                    suggestions.append(f"- **Try:** Use full path like '{self.home / 'Documents'}' or '{self.home / 'Dokumente'}'")
        
        elif status_code == 404:
            suggestions.append("- **Resource not found** - The requested path or file does not exist")
        
        elif status_code == 500:
            suggestions.append("- **Server error** - Internal server problem, not your fault")
        
        # Detect common server-side format issues for more actionable retries
        extra_retry = ""
        if "Cannot have 2 or more assistant messages" in (error_message or ""):
            extra_retry = (
                "\n\n**Retry Hint (Main Agent):**\n"
                "- This is a server-side request ordering error.\n"
                "- Retry with a fresh tool call (new librarian_agent invocation) and a single, explicit objective.\n"
                "- Prefer deterministic tools for OS queries (folder_size, list_files, tree) instead of LLM reasoning.\n"
            )

        # Build helpful error message
        error_report = f"""[ERROR] **Librarian Agent Error Report**

**Status Code:** {status_code}
**Error Type:** {error_type or 'Unknown'}
**Error Message:** {error_message or 'No detailed error message available'}

**Original Task:** "{task}"

**What went wrong:**
{chr(10).join(suggestions) if suggestions else "- Unable to determine specific issue"}

**Available Tools:**
{chr(10).join([f"- {t.name}: {t.description[:60]}..." for t in self.tools.values()])}

**Suggestions for Main Agent:**
1. **Check the path** - Verify the folder/file path exists
2. **Rephrase the task** - Be more specific about what you want (e.g., "Count PDF files in Documents folder")
3. **Use direct path** - Instead of "Dokumen", try "Documents" or the full path
4. **Check permissions** - Ensure the path is accessible{path_hint}
{extra_retry}

**Context:**
- Home directory: {self.home}
- Available folder aliases: {', '.join(list(self.folder_aliases.keys())[:5])}...
"""
        return error_report
    
    def _format_connection_error(self, error: str, task: str) -> str:
        """Format connection error with helpful context."""
        return f"""[ERROR] **Librarian Agent Connection Error**

**Error:** {error}

**Original Task:** "{task}"

**What happened:**
- Could not connect to the LLM server (http://127.0.0.1:8080)
- The server might be down or not responding

**Suggestions for Main Agent:**
1. **Check server status** - Verify the LLM server is running
2. **Try again later** - The server might be temporarily unavailable
3. **Use alternative approach** - Consider using direct file operations if possible

**Note:** The librarian agent needs the LLM server to process complex queries.
"""
    
    # ═══════════════════════════════════════════════════════════════════════════
    # LLM EXECUTION (Slow Path - for complex queries)
    # ═══════════════════════════════════════════════════════════════════════════
    
    def _execute_with_llm(self, task: str) -> str:
        """Execute complex task using LLM reasoning."""
        
        # Show animation (or static)
        # Disable animation by default for stability (matches Coder/Research agent)
        from vaf.cli.tui import _StaticHeader
        header = _StaticHeader("Collaboration Mode Active", "Main Agt", "Librarian")
        
        # Use Live to show the header, but with minimal updates
        live = Live(header, refresh_per_second=12, console=UI.console)
        live.start()
        
        time.sleep(1.0)
        UI.event("Sub-Agent", "Librarian analyzing complex query...", style="bold cyan")
        
        # Get disk/storage information
        disk_info = self._get_disk_info()
        
        # Get filesystem map
        fs_map_context = self.get_system_prompt_addition()
        
        # System prompt
        system_prompt = f"""<identity>
You are the Librarian, a file/info retrieval specialist.
</identity>

<context>
home: {self.home}
- Downloads: {self.folder_aliases.get('downloads', 'N/A')}
- Desktop: {self.folder_aliases.get('desktop', 'N/A')}
- Documents: {self.folder_aliases.get('documents', 'N/A')}
{fs_map_context}
{disk_info}
</context>

<tools>
- read_file(path): Reads a file's contents. **SUPPORTS: PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx), and text files**
- write_file(path, content): Writes content to a file
- list_files(path, sort_by='name'|'date'|'size', limit=100): Lists files and folders in a directory. Use this to see what folders exist on a drive.
- tree(path, depth): Shows directory tree structure. Use this to explore folder structure on a drive.
- find_files(path, pattern): Finds files by name pattern (glob) recursively
- folder_size(path, top_n=10): Calculates total size of a folder recursively and shows largest files
- cloud_storage(action, query?, file_id?, provider?, account_id?): PREFER action='search_all' + query to search ALL clouds at once (like the UI). search: single provider. Then action='read', 'download', or 'show_in_viewer' with file_id (+ provider or account_id from results). show_in_viewer opens PDF/doc in Document Viewer (Anhänge). action='browse' when user wants folder contents. action='list' for VAF Sync. action='status' for connection.
- python_sandbox(code): Execute Python code safely for mathematical calculations, data processing, and algorithms
</tools>

<document_capabilities>
- PDF files (.pdf): Extracts text from all pages (up to 50 pages shown)
- Word documents (.docx): Reads paragraphs and tables
- Excel files (.xlsx, .xls): Reads all sheets and cells (up to 3 sheets, 50 rows each)
- PowerPoint (.pptx): Extracts text from slides (up to 20 slides)
- Text files: All text-based formats (txt, md, json, xml, csv, etc.)
</document_capabilities>

<tool_selection_guide>
- "Read file content" / "Read file X" / "Read PDF" / "Read Word document" -> Use read_file(path)
- "Write to file" / "Create file" / "Save content to file" -> Use write_file(path, content)
- "What folders on drive X?" / "Welche Ordner auf Laufwerk X?" -> Use list_files(path) to see all folders
- "Show folder structure" -> Use tree(path, depth) for visual tree
- "Find files matching pattern" -> Use find_files(path, pattern)
- "Largest files" / "Biggest files" / "Größte Dateien" -> Use list_files(path, sort_by='size', limit=20) to find largest files
- "Files by size" / "Dateien nach Größe" -> Use list_files(path, sort_by='size') to sort files by size
- "Multiple files analysis" / "Mehrere Dateien analysieren" -> Use list_files with sort_by='size' or find_files, then analyze results
- "Calculate" / "Math" / "Compute" / "Algorithm" -> Use python_sandbox(code) for mathematical calculations and data processing
- "Folder size" / "Ordnergröße" -> Use folder_size(path)
- "Find PDF X" / "Suche Bewilligung in Drive" / "Where is report.pdf?" -> cloud_storage(action='search_all', query='X') FIRST. Then read/download/show_in_viewer with file_id (+ provider or account_id from results).
- "Show PDF in viewer" / "Zeig mir das PDF" / "Open document from cloud" -> cloud_storage(action='search_all', query='...') to find, then action='show_in_viewer', file_id=<id>, provider=... or account_id=<id>. This opens the PDF/doc in the Document Viewer (Anhänge).
- "List cloud folders" / "Drive durchsuchen" -> cloud_storage(action='browse', folder_id='root')
- "Read document from cloud" -> cloud_storage(action='search_all', query='...') to find, then action='read', file_id=<id>
- "Download from cloud" -> cloud_storage(action='download', file_id=<id>)
</tool_selection_guide>

<rules>
1. Call ONE tool, get result, then ANSWER immediately
2. Do NOT think in loops - act decisively
3. Summarize results (don't dump raw data)
4. If unsure about path, use the home directory
</rules>"""

        # Initialize context manager for this librarian agent session
        from vaf.core.context import ContextManager
        max_tokens = 8192  # Same as main agent
        context_manager = ContextManager(max_tokens=max_tokens)

        history = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Task: {task}"}
        ]
        
        # Snapshot history (before processing)
        history_snapshot_len = len(history)

        def _sanitize_history(messages: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
            """
            The local OpenAI-compatible server is strict about message ordering.
            In particular it may reject requests if the message list ends with
            multiple assistant messages (e.g. after retries with empty outputs).
            """
            if not messages:
                return messages

            # Remove trailing assistant messages if there are 2+ at the end
            while len(messages) >= 2 and messages[-1].get("role") == "assistant" and messages[-2].get("role") == "assistant":
                messages.pop(-2)

            return messages
        
        # Reduced max steps
        max_steps = 5
        
        # Descriptive status messages
        # ASCII-only status strings for maximum terminal compatibility
        status_messages = [
            "Thinking...",
            "Analyzing...",
            "Processing...",
            "Refining...",
            "Finalizing...",
        ]
        
        for step in range(max_steps):
            status = status_messages[min(step, len(status_messages) - 1)]
            UI.event("Librarian", status, style="dim")
            
            # Context management - prevent token overflow
            # Proactive compression: Check token usage and compress if > 85% of limit
            estimated_tokens = context_manager.estimate_tokens(history)
            if estimated_tokens > int(context_manager.max_tokens * 0.85):
                UI.event("Librarian", f"Proactive compression: {estimated_tokens}/{context_manager.max_tokens} tokens...", style="dim")
                history = context_manager.compress(history)
            # Also check normal threshold
            elif context_manager.should_compress(history):
                UI.event("Librarian", "Compressing context...", style="dim")
                history = context_manager.compress(history)
            
            try:
                from vaf.core.config import Config
                config = Config.load()
                provider = config.get("provider", "local")
                model_name = config.get("model", "")
                
                # Build tools schema with validation
                tools_schema = []
                for t in self.tools.values():
                    params = getattr(t, 'parameters', {})
                    if not isinstance(params, dict):
                        params = {"type": "object", "properties": {}}
                    if "type" not in params:
                        params["type"] = "object"
                    if "properties" not in params:
                        params["properties"] = {}
                    
                    tools_schema.append({
                        "type": "function",
                        "function": {
                            "name": t.name,
                            "description": t.description or f"Tool: {t.name}",
                            "parameters": params
                        }
                    })
                
                history = _sanitize_history(history)
                response_msg = None
                
                # ---------------------------------------------------------------
                # 1. API Backend Path (OpenAI, Anthropic, DeepSeek, Google, etc.)
                # ---------------------------------------------------------------
                if provider != "local":
                    try:
                        from vaf.core.api_backend import APIBackendManager
                        backend = APIBackendManager(provider)
                        
                        full_content = ""
                        tool_calls = []
                        
                        for chunk in backend.chat_completion(
                            messages=history,
                            temperature=0.1,
                            max_tokens=1024,
                            stream=True,
                            model=model_name,
                            tools=tools_schema
                        ):
                            if chunk.strip().startswith("{") and ("tool_calls" in chunk or "tool_use" in chunk):
                                try:
                                    data = json.loads(chunk)
                                    if "tool_calls" in data:
                                        tool_calls.extend(data["tool_calls"])
                                    elif "tool_use" in data:
                                        # Anthropic conversion
                                        tool_calls.append({
                                            "id": data["tool_use"].get("id"),
                                            "type": "function",
                                            "function": {
                                                "name": data["tool_use"].get("name"),
                                                "arguments": json.dumps(data["tool_use"].get("input", {}))
                                            }
                                        })
                                except:
                                    pass
                            else:
                                full_content += chunk
                        
                        response_msg = {"role": "assistant", "content": full_content}
                        if tool_calls:
                            response_msg["tool_calls"] = tool_calls
                            
                    except Exception as e:
                        live.stop()
                        return self._format_connection_error(f"API Error: {str(e)}", task)

                # ---------------------------------------------------------------
                # 2. Local Server Path (Fallback)
                # ---------------------------------------------------------------
                else:
                    try:
                        res = requests.post(
                            "http://127.0.0.1:8080/v1/chat/completions",
                            json={
                                "model": model_name,
                                "messages": history,
                                "max_tokens": 1024,
                                "temperature": 0.1,
                                "tools": tools_schema,
                                "tool_choice": "auto",
                            },
                            timeout=60,
                        )
                        
                        if res.status_code == 400:
                            # Context Size Error handling
                            error_data = res.json()
                            error_msg = error_data.get("error", {}).get("message", "")
                            if "exceed" in error_msg.lower():
                                history = context_manager.compress(history)
                                continue
                        
                        if res.status_code != 200:
                            live.stop()
                            return self._format_error_response(res, task, history, tools_schema)
                            
                        response_msg = res.json()['choices'][0]['message']
                        
                    except requests.exceptions.RequestException as e:
                        live.stop()
                        return self._format_connection_error(str(e), task)
                
                if not response_msg:
                    continue

                msg = response_msg
                content = msg.get('content', '')
                tool_calls = msg.get('tool_calls', [])
                
                history.append(msg)

                # If model returned neither tool calls nor content, retry with brief prompt
                if not tool_calls and not content:
                    # Remove the empty assistant message
                    if history and history[-1].get('role') == 'assistant':
                        history.pop()
                    
                    # ═══════════════════════════════════════════════════════════════
                    # NEW: Tool-Intent Detection (like main agent)
                    # ═══════════════════════════════════════════════════════════════
                    # Check if agent mentioned a tool name (but didn't actually call it yet)
                    # CRITICAL: First check if response is truly empty (content is already checked above)
                    # In librarian, we already have `not content`, so response is truly empty
                    is_truly_empty = not content
                    
                    # Only check tool intent if response is truly empty
                    if is_truly_empty:
                        # Get available tool names dynamically
                        available_tool_names = list(self.tools.keys()) if hasattr(self, 'tools') and self.tools else []
                        # Also check for common tool names used by librarian
                        common_tool_names = ["read_file", "write_file", "list_files", "find_files", "tree", "folder_size", "cloud_storage", "python_sandbox", "document_viewer", "document_editor", "replace_editor_selection", "replace_editor_text"]
                        all_tool_names = list(set(available_tool_names + common_tool_names))
                        
                        # Check if any tool name appears in the last assistant message (case-insensitive)
                        last_assistant_msg = None
                        for msg in reversed(history):
                            if msg.get('role') == 'assistant' and msg.get('content'):
                                last_assistant_msg = str(msg.get('content', '')).lower()
                                break
                        
                        mentioned_tools = []
                        if last_assistant_msg:
                            # Check if any tool name appears in the response (case-insensitive)
                            mentioned_tools = [tool_name for tool_name in all_tool_names if tool_name.lower() in last_assistant_msg]
                        
                        # CRITICAL: Only reset if BOTH conditions are met:
                        # 1. Response is truly empty (content is None/empty) - checked BEFORE any cleaning
                        # 2. Tool was mentioned but not called
                        # This is language-independent - we only check if response is empty, not what language it's in
                        if mentioned_tools and not tool_calls:
                            tool_hint = mentioned_tools[0]
                            UI.event("Librarian", f"Tool-Intent detected for '{tool_hint}' without action - resetting to snapshot", style="dim")
                            
                            # Check if there's thinking DIRECTLY after user prompt
                            user_prompt_idx = 1  # User message is at index 1 (after system at 0)
                            
                            # Check if there's an assistant message with content directly after user prompt
                            first_assistant_after_user = None
                            first_assistant_idx = None
                            
                            # Look at messages right after user prompt (within next 2 messages)
                            for i in range(user_prompt_idx + 1, min(user_prompt_idx + 3, len(history))):
                                msg = history[i]
                                if msg.get('role') == 'assistant' and msg.get('content'):
                                    content_str = str(msg.get('content', ''))
                                    # Keep if it has substantial content (thinking/reasoning)
                                    if len(content_str.strip()) > 20:
                                        first_assistant_after_user = content_str
                                        first_assistant_idx = i
                                        break  # Only take the FIRST one directly after user prompt
                            
                            if first_assistant_after_user and first_assistant_idx is not None:
                                # Keep user prompt + first thinking - this becomes the new snapshot
                                history = history[:first_assistant_idx + 1]
                                UI.event("Librarian", f"Reset to thinking snapshot (user prompt + {len(first_assistant_after_user)} chars of first thinking)", style="dim")
                            else:
                                # No thinking found - reset to user prompt snapshot (as before)
                                history = history[:history_snapshot_len]
                                UI.event("Librarian", f"Reset to user prompt snapshot", style="dim")
                            
                            # Add a brief system prompt (will work better now because first thinking is preserved)
                            history.append({
                                "role": "system",
                                "content": "You didn't respond. Please answer or continue where you left off."
                            })
                            
                            # Continue the loop - if it fails again, this system message will be removed with the reset
                            continue
                    
                    # ═══════════════════════════════════════════════════════════════
                    # Snapshot System with Thinking Preservation (like main agent)
                    # ═══════════════════════════════════════════════════════════════
                    # Check if there's thinking DIRECTLY after user prompt
                    # The user prompt is at history_snapshot_len (index 1, after system prompt)
                    user_prompt_idx = 1  # User message is at index 1 (after system at 0)
                    
                    # Check if there's an assistant message with content directly after user prompt
                    first_assistant_after_user = None
                    first_assistant_idx = None
                    
                    # Look at messages right after user prompt (within next 2 messages)
                    for i in range(user_prompt_idx + 1, min(user_prompt_idx + 3, len(history))):
                        msg = history[i]
                        if msg.get('role') == 'assistant' and msg.get('content'):
                            content_str = str(msg.get('content', ''))
                            # Keep if it has substantial content (thinking/reasoning)
                            if len(content_str.strip()) > 20:
                                first_assistant_after_user = content_str
                                first_assistant_idx = i
                                break  # Only take the FIRST one directly after user prompt
                    
                    if first_assistant_after_user and first_assistant_idx is not None:
                        # Keep user prompt + first thinking - this becomes the new snapshot
                        history = history[:first_assistant_idx + 1]
                        UI.event("Librarian", f"Reset to thinking snapshot (user prompt + {len(first_assistant_after_user)} chars of first thinking)", style="dim")
                    else:
                        # No thinking found - reset to user prompt snapshot (as before)
                        history = history[:history_snapshot_len]
                        UI.event("Librarian", f"Reset to user prompt snapshot", style="dim")
                    
                    # Add a brief system prompt (will work better now because first thinking is preserved)
                    history.append({
                        "role": "system",
                        "content": "You didn't respond. Please answer or continue where you left off."
                    })
                    
                    # Continue the loop - if it fails again, this system message will be removed with the reset
                    continue
                
                # If no tool calls and has content → done
                if not tool_calls and content:
                    live.stop()
                    return f"### Librarian Report\n\n{content}"
                
                # Execute tools
                for tc in tool_calls:
                    fn_name = tc['function']['name']
                    try:
                        fn_args = json.loads(tc['function']['arguments'])
                    except:
                        fn_args = {}
                    
                    # Show tool action with icon
                    tool_icons = {
                        "list_files": "Listing",
                        "read_file": "Reading",
                        "find_files": "Searching",
                        "tree": "Mapping",
                        "folder_size": "Sizing",
                        "document_viewer": "Opening in viewer",
                        "document_editor": "Opening in editor",
                        "replace_editor_selection": "Replacing marked text in editor",
                        "replace_editor_text": "Replacing unmarked text in editor",
                    }
                    icon = tool_icons.get(fn_name, "Calling")
                    UI.event("Librarian", f"{icon}: {fn_name}", style="bold green")
                    
                    if fn_name in self.tools:
                        try:
                            result = self.tools[fn_name].run(**fn_args)
                        except Exception as e:
                            result = f"Error: {e}"
                    else:
                        result = f"Unknown tool: {fn_name}"
                    
                    history.append({
                        "role": "tool",
                        "tool_call_id": tc['id'],
                        "name": fn_name,
                        "content": str(result)[:2000]  # Limit result size
                    })
                
            except requests.exceptions.Timeout:
                live.stop()
                return "Error: Request timed out"
            except Exception as e:
                live.stop()
                return f"Error: {e}"
        
        live.stop()
        
        # Extract last result if available
        for m in reversed(history):
            if m.get('role') == 'tool':
                return f"### Librarian Report\n\n{m.get('content', '')[:1000]}"
        
        return "Librarian could not complete the task."
